"""Add the project-facing excitement alias while preserving arousal provenance.

The SER model's technical output remains available as audio_arousal. This
migration makes excitement available in data and explanatory artifacts without
breaking existing loaders or historical model comparisons.
"""
from __future__ import annotations

import csv
import json
import re
import shutil
import zipfile
from pathlib import Path

from docx import Document
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
MARKER = "EXCITEMENT_ALIAS_MIGRATION_V1"


def backup(path: Path) -> None:
    target = path.with_name(path.name + ".before_excitement_alias")
    if not target.exists():
        shutil.copy2(path, target)


def migrate_csv(path: Path) -> bool:
    with path.open(newline="", encoding="utf-8-sig") as handle:
        reader = csv.DictReader(handle)
        fields = list(reader.fieldnames or [])
        rows = list(reader)
    additions = []
    for old, new in (("audio_arousal", "audio_excitement"), ("arousal", "excitement")):
        if old in fields and new not in fields:
            additions.append((old, new))
    value_changed = any(
        replace_text(str(cell)) != str(cell)
        for row in rows for cell in row.values()
    )
    if not additions and not value_changed:
        return False
    backup(path)
    for old, new in additions:
        position = fields.index(old) + 1
        fields.insert(position, new)
        for row in rows:
            row[new] = row.get(old, "")
    for row in rows:
        for key, cell in row.items():
            row[key] = replace_text(cell)
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=fields); writer.writeheader(); writer.writerows(rows)
    return True


def add_json_aliases(value):
    if isinstance(value, dict):
        result = {key: add_json_aliases(item) for key, item in value.items()}
        if "audio_arousal" in result and "audio_excitement" not in result:
            result["audio_excitement"] = result["audio_arousal"]
        if "arousal" in result and "excitement" not in result:
            result["excitement"] = result["arousal"]
        return result
    if isinstance(value, list):
        return [add_json_aliases(item) for item in value]
    return replace_text(value) if isinstance(value, str) else value


def migrate_json(path: Path) -> bool:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return False
    migrated = add_json_aliases(data)
    if migrated == data:
        return False
    backup(path)
    path.write_text(json.dumps(migrated, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    return True


def replace_text(text: str) -> str:
    text = re.sub(r"\barousal\b", "excitement", text, flags=re.IGNORECASE)
    return text


def migrate_docx(path: Path) -> bool:
    try:
        doc = Document(path)
    except (OSError, zipfile.BadZipFile, ValueError):
        return False
    changed = False
    for paragraph in list(doc.paragraphs) + [p for table in doc.tables for row in table.rows for cell in row.cells for p in cell.paragraphs]:
        for run in paragraph.runs:
            replacement = replace_text(run.text)
            if replacement != run.text:
                run.text = replacement; changed = True
    if changed:
        backup(path); doc.save(path)
    return changed


def migrate_pptx(path: Path) -> bool:
    try:
        prs = Presentation(path)
    except (OSError, zipfile.BadZipFile, ValueError):
        return False
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    replacement = replace_text(run.text)
                    if replacement != run.text:
                        run.text = replacement; changed = True
    if changed:
        backup(path); prs.save(path)
    return changed


def migrate_text_file(path: Path) -> bool:
    try:
        original = path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError):
        return False
    replacement = replace_text(original)
    if replacement == original:
        return False
    backup(path); path.write_text(replacement, encoding="utf-8")
    return True


def main() -> None:
    changed = {"csv": [], "json": [], "docx": [], "pptx": [], "text": []}
    for path in ROOT.rglob("*.csv"):
        if migrate_csv(path): changed["csv"].append(str(path.relative_to(ROOT)))
    for path in ROOT.rglob("*.json"):
        if migrate_json(path): changed["json"].append(str(path.relative_to(ROOT)))
    for path in ROOT.rglob("*.docx"):
        if ".before_excitement_alias" not in path.name and migrate_docx(path): changed["docx"].append(str(path.relative_to(ROOT)))
    for path in ROOT.rglob("*.pptx"):
        if ".before_excitement_alias" not in path.name and migrate_pptx(path): changed["pptx"].append(str(path.relative_to(ROOT)))
    for suffix in (".md", ".txt", ".mmd"):
        for path in ROOT.rglob(f"*{suffix}"):
            if ".before_excitement_alias" not in path.name and migrate_text_file(path): changed["text"].append(str(path.relative_to(ROOT)))
    print(json.dumps({"marker": MARKER, "changed": changed}, indent=2))


if __name__ == "__main__":
    main()
