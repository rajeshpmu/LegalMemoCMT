"""Document the pilot basis and limitation of the Odyssey arousal threshold."""
from pathlib import Path
import shutil
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt

ROOT=Path(__file__).resolve().parents[2]
MARKER="AROUSAL_THRESHOLD_045_GUIDANCE_V1"
DOCS=[ROOT/"implementation_docments/LegalMemoCMT_Clancy_Corpus_SOP.docx",ROOT/"implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Student_Speaking_Guide.docx",ROOT/"implementation_docments/LegalMemoCMT_Phase2_Corpus_Build_Code_Level_Student_Guide.docx"]
PPTS=[ROOT/"implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Presentation.pptx",ROOT/"implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"]

def update_doc(path):
 doc=Document(path)
 if any(MARKER in p.text for p in doc.paragraphs): return False
 backup=path.with_name(path.name+".before_arousal_threshold_guidance.docx")
 if not backup.exists(): shutil.copy2(path,backup)
 doc.add_heading("Why the Pilot Uses arousal < 0.45", level=1)
 p=doc.add_paragraph(); p.add_run(MARKER).font.size=Pt(8)
 doc.add_paragraph("The 0.45 value is an operational pilot threshold for Odyssey's model-specific arousal output. It is not a clinical threshold, a universal definition of calm speech, or a validated courtroom boundary. It is used to identify relatively low/moderate vocal activation for review.")
 doc.add_paragraph("The choice is data-informed but provisional. In the 200-row pilot, arousal had mean=0.371918, median=0.368324, 75th percentile=0.412108, and maximum=0.577063. Therefore 0.45 lies above the central pilot distribution and provides a transparent separation between typical low/moderate activation and higher observed activation in this small sample.")
 doc.add_paragraph("The threshold is never used alone. The scope-aware rule combines neutral SpeechBrain evidence, arousal below 0.45, compatible non-self emotion scope, and no explicit contrary visual/audio failure. Even when these conditions hold, the output is a NEUTRAL_CANDIDATE requiring human review. CALM_COMPOSED or NEUTRAL_CALM describes delivery and does not prove neutral basic emotion.")
 doc.add_paragraph("A future calibration study must compare thresholds such as 0.40, 0.45, and 0.50 against human-reviewed Clancy examples. It should measure precision, recall, disagreement by source and speaker, and whether the threshold incorrectly converts calm sadness or calm self-expressed distress into neutral. Until that study is complete, 0.45 must be described as a reproducible starting rule, not a discovered truth.")
 doc.add_paragraph("Student explanation: I selected 0.45 because it is slightly above the pilot's normal arousal range, but I do not claim that every value below 0.45 means calm or neutral. It only contributes one piece of evidence to a human-review recommendation.")
 doc.save(path); return True

def add_text(slide,text,x,y,w,h,size=16,color=(31,45,61),bold=False):
 box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); box.text_frame.word_wrap=True
 p=box.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=text; r.font.name="Aptos"; r.font.size=PptPt(size); r.font.bold=bold; r.font.color.rgb=RGBColor(*color)

def update_ppt(path):
 prs=Presentation(path)
 if any(MARKER in sh.text for sl in prs.slides for sh in sl.shapes if hasattr(sh,'text')): return False
 backup=path.with_name(path.name+".before_arousal_threshold_guidance.pptx")
 if not backup.exists(): shutil.copy2(path,backup)
 sl=prs.slides.add_slide(prs.slide_layouts[6]); sl.background.fill.solid(); sl.background.fill.fore_color.rgb=RGBColor(255,255,255)
 bar=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.2)); bar.fill.solid(); bar.fill.fore_color.rgb=RGBColor(31,78,121); bar.line.fill.background()
 add_text(sl,"Why arousal < 0.45?", "0", "0", 0, 0, 1) if False else None
 add_text(sl,"Why arousal < 0.45?",.55,.35,12,.5,23,bold=True)
 add_text(sl,"A transparent pilot heuristic, not a universal calmness threshold",.58,.88,12,.35,11,color=(90,100,110))
 add_text(sl,"AROUSAL_THRESHOLD_045_GUIDANCE_V1\n\n• Pilot statistics: mean 0.371918, median 0.368324, Q75 0.412108, maximum 0.577063.\n• 0.45 is above the central pilot range and separates lower/moderate activation from higher observed activation.\n• It is supporting evidence only, never a standalone neutral rule.\n• Combine neutral SpeechBrain + low arousal + compatible scope + no contrary evidence.\n• Calibration remains necessary using human-reviewed rows and thresholds 0.40 / 0.45 / 0.50.",.8,1.45,11.7,5.3,16)
 prs.save(path); return True

def main():
 print({"updated_docs":[str(p) for p in DOCS if update_doc(p)],"updated_pptx":[str(p) for p in PPTS if update_ppt(p)],"marker":MARKER})

if __name__=="__main__": main()
