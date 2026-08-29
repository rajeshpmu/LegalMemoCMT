"""Append the reviewed utterance example to the Phase 2 PPTX and speaking guide."""
from pathlib import Path
import shutil

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt

ROOT = Path(__file__).resolve().parents[2]
MARKER = "TURN06801_MACHINE_REVIEW_EXAMPLE_V1"
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Student_Speaking_Guide.docx"


def add_text(slide, text, x, y, w, h, size=17, color=(31, 45, 61), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = text
    run.font.name = "Aptos"; run.font.size = PptPt(size); run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def style_slide(slide, title, subtitle):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(31, 78, 121); bar.line.fill.background()
    add_text(slide, title, .55, .35, 12.2, .55, 23, bold=True)
    add_text(slide, subtitle, .58, .92, 12, .35, 11, color=(90, 100, 110))


def update_ppt() -> None:
    if not PPTX.exists(): return
    prs = Presentation(PPTX)
    if any(MARKER in shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")):
        return
    backup = PPTX.with_name(PPTX.name + ".before_turn06801_machine_review.pptx")
    if not backup.exists(): shutil.copy2(PPTX, backup)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide, "Worked Example: DCBWoWhsTpA_turn06801", "Machine-assisted review record, with original model output preserved")
    add_text(slide, "Phase 1: sadness @ 0.544674\nReview candidate: neutral @ 0.75\nCourtroom affect: CALM_COMPOSED @ 0.60\n\nnegative_activation_candidate: YES\ndistress_corroboration_present: NO\n\nspeaker_emotion_evidence_present: NO\nemotion_target_scope: QUOTED_SPEECH", .8, 1.5, 5.7, 4.8, 19)
    add_text(slide, "Interpretation\n\nThe witness reports what another person said. Negative valence and elevated arousal indicate negative activation, not automatically sadness or distress. The witness presents as controlled, so basic emotion and courtroom affect are kept separate.\n\nHuman review remains required; these are machine candidates, not gold labels.", 6.8, 1.5, 5.7, 4.8, 18)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide, "Why the Fields Must Stay Separate", "A courtroom utterance can describe distress without expressing distress")
    add_text(slide, "Phase-1 model\nsadness @ .545\n        ↓\npreserve original prediction\n\nV/A/D evidence\nnegative activation = YES\n        ↓\ncorroboration? NO\n        ↓\nnot DISTRESSED", .8, 1.35, 4.0, 4.9, 18)
    add_text(slide, "Transcript scope\n\"She said ...\"\n        ↓\nQUOTED_SPEECH\n        ↓\nspeaker emotion evidence = NO\n\nIntegrated review\nneutral basic-emotion candidate\nCALM_COMPOSED courtroom affect", 4.75, 1.35, 4.0, 4.9, 18)
    add_text(slide, "Practical conclusion\n\nnegative_activation_candidate=YES\n≠ basic_emotion=sadness\n≠ courtroom_affect=DISTRESSED\n\nThe system records what each signal supports and leaves final annotation to human review.", 8.7, 1.35, 3.9, 4.9, 18)
    prs.save(PPTX)


def update_doc() -> None:
    if not DOCX.exists(): return
    doc = Document(DOCX)
    if any(MARKER in p.text for p in doc.paragraphs): return
    backup = DOCX.with_name(DOCX.name + ".before_turn06801_machine_review.docx")
    if not backup.exists(): shutil.copy2(DOCX, backup)
    doc.add_heading("Worked Example: Clean Machine-Assisted Review Record", level=1)
    p = doc.add_paragraph(); p.add_run(MARKER).font.size = Pt(8)
    doc.add_paragraph(
        "For DCBWoWhsTpA_turn06801, I would accept the following as the final machine-assisted "
        "review record, while still requiring human confirmation. Phase 1 predicts sadness at "
        "0.544674, and that original prediction is preserved as evidence of domain shift. The "
        "integrated basic-emotion review candidate is neutral at 0.75 because the utterance is "
        "quoted/reporting content and does not independently show that the witness is sad."
    )
    doc.add_paragraph(
        "The courtroom-affect candidate is CALM_COMPOSED at 0.60. This describes the witness's "
        "controlled presentation and is intentionally different from the basic-emotion label. "
        "A witness could have basic_emotion=sadness while courtroom_affect=CALM_COMPOSED, so "
        "the two fields must not be merged."
    )
    doc.add_heading("Field-by-field interpretation", level=2)
    for text in [
        "phase1_basic_emotion=sadness and phase1_basic_emotion_confidence=0.544674: the original MELD-trained model output, retained unchanged.",
        "basic_emotion_review_candidate=neutral and confidence=0.75: a machine-assisted review suggestion for the witness's own expressed emotion, not a final label.",
        "proposed_courtroom_affect=CALM_COMPOSED and confidence=0.60: a behavioral presentation candidate, separate from categorical emotion.",
        "negative_activation_candidate=YES: low-level V/A/D evidence of negative acoustic activation; it does not mean sadness or distress.",
        "distress_corroboration_present=NO: the system found no sufficient categorical, linguistic, or visual evidence to promote the row to DISTRESSED.",
        "speaker_emotion_evidence_present=NO: the transcript scope does not provide independent evidence that the current witness expresses the quoted person's emotion.",
        "emotion_target_scope=QUOTED_SPEECH: the witness says 'She said...' and reports another person's words, so emotional content must not automatically be attributed to the witness.",
    ]:
        doc.add_paragraph(text, style="List Bullet")
    doc.add_heading("Reasoning chain to explain", level=2)
    doc.add_paragraph(
        "Phase 1 predicts SADNESS at approximately .545, so I preserve it. Audio V/A/D gives "
        "negative activation, but there is no distress corroboration, so I do not promote the "
        "row to DISTRESSED. The transcript scope is QUOTED_SPEECH and speaker-emotion evidence "
        "is NO, so neutral is the better basic-emotion review candidate. The observed controlled "
        "delivery is represented independently as CALM_COMPOSED. This separation is stronger "
        "than forcing one model output to explain semantic content, speaker emotion, and courtroom "
        "behavior at the same time."
    )
    doc.add_paragraph(
        "This remains a machine-assisted record. A human reviewer must listen to and view the clip "
        "before recording human_basic_emotion=neutral, human_courtroom_affect=CALM_COMPOSED, or "
        "affect_intensity=1. No credibility, truthfulness, deception, or reliability inference is made."
    )
    doc.save(DOCX)


def main():
    update_ppt(); update_doc(); print({"pptx": str(PPTX), "docx": str(DOCX), "marker": MARKER})


if __name__ == "__main__": main()
