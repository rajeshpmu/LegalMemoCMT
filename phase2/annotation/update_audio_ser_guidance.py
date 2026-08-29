"""Append the independent audio-SER stage to Phase 2 guidance artifacts.

This updater is intentionally append-only and idempotent. It creates a backup
before modifying an existing DOCX/PPTX and uses a marker to avoid duplicate
sections on reruns.
"""

from __future__ import annotations

import shutil
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[2]
MARKER = "AUDIO_SER_GUIDANCE_V1"
DOCX_TARGETS = [
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Corpus_SOP.docx",
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Student_Speaking_Guide.docx",
    ROOT / "implementation_docments/LegalMemoCMT_Phase2_Corpus_Build_Code_Level_Student_Guide.docx",
]
PPTX_TARGETS = [
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Presentation.pptx",
    ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx",
]


def backup(path: Path) -> None:
    candidate = path.with_name(path.name + ".before_audio_ser_guidance.docx" if path.suffix == ".docx" else ".before_audio_ser_guidance.pptx")
    if path.exists() and not candidate.exists():
        shutil.copy2(path, candidate)


def add_code(doc: Document, value: str) -> None:
    p = doc.add_paragraph(style="Intense Quote")
    p.add_run(value)


def update_docx(path: Path) -> bool:
    if not path.exists():
        return False
    doc = Document(path)
    if any(MARKER in p.text for p in doc.paragraphs):
        return False
    backup(path)
    doc.add_page_break()
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Independent Audio SER Evidence and Conservative Label Review")
    r.bold = True; r.font.size = Pt(18)
    doc.add_paragraph(MARKER)
    doc.add_heading("Purpose in Phase 2", level=1)
    doc.add_paragraph(
        "The audio speech-emotion-recognition stage is an independent evidence layer. "
        "It is not a replacement for the human-verified witness-role gate, the Phase 1 "
        "MELD checkpoint, or courtroom-affect annotation. The reason for adding it is "
        "to measure whether vocal prosody supports or contradicts a semantic model "
        "prediction, especially when a witness calmly describes another person's distress."
    )
    doc.add_heading("Models and what each contributes", level=1)
    doc.add_paragraph(
        "The primary model is 3loi/SER-Odyssey-Baseline-WavLM-Multi-Attributes. It receives "
        "the utterance WAV and returns continuous valence, arousal, and dominance-like "
        "scores in the model's approximately 0-to-1 output range. These values are model "
        "scores, not physical units and not courtroom-affect labels. A secondary model is "
        "speechbrain/emotion-recognition-wav2vec2-IEMOCAP. It provides a four-class "
        "categorical cross-check (neu, ang, hap, sad) trained on IEMOCAP. Its vocabulary "
        "does not cover the full seven-class MELD vocabulary and therefore must not be "
        "silently remapped into a gold LegalMemoCMT label."
    )
    doc.add_heading("Exact implementation", level=1)
    add_code(doc, "bash phase2/annotation/setup_audio_ser.sh")
    add_code(doc, "PYTHON_BIN=$PWD/.venv-audio-ser/bin/python bash phase2/annotation/run_audio_ser_evidence.sh \\")
    add_code(doc, "  --input-csv data/processed/phase2/clancy/witness_only_v10/clancy_witness_speaking_visual_verified_vit_200.csv \\")
    add_code(doc, "  --output-csv data/processed/phase2/clancy/audio_ser_evidence_200.csv \\")
    add_code(doc, "  --summary-json reports/phase2/clancy_audio_ser_evidence_200.json --max-rows 200 --device cpu")
    doc.add_paragraph(
        "run_audio_ser_evidence.py finds an existing utterance audio path, loads the WAV, "
        "runs both independent models, and writes a derivative CSV. It preserves every "
        "input column and appends audio_ser_* provenance fields. It never overwrites "
        "emotion_label, basic_emotion, courtroom_affect, or human labels. The shell "
        "wrapper selects the repository-local .venv-audio-ser interpreter so the Phase 2 "
        "audio stack is isolated from the incompatible Anaconda installation."
    )
    doc.add_heading("How to interpret the output", level=1)
    doc.add_paragraph(
        "audio_valence, audio_arousal, and audio_dominance are continuous prosody evidence. "
        "audio_emotion_candidate and audio_emotion_confidence are a cross-check only. "
        "The correct workflow is to preserve Phase 1 output as model_basic_emotion, keep "
        "the audio evidence in separate columns, and send disagreement cases to review. "
        "For example, a psychiatrist can speak with low arousal and stable delivery while "
        "the transcript contains words such as hopeless or suicidal. The words describe "
        "another person; they do not establish fear in the speaking witness. This is why "
        "emotion_target_scope should be reviewed separately using SELF_EXPRESSED, "
        "OTHER_PERSON_DESCRIBED, EVENT_DESCRIBED, QUOTED_SPEECH, or UNCLEAR."
    )
    doc.add_heading("Pilot result and limitation", level=1)
    doc.add_paragraph(
        "The first 200-row pilot successfully produced Odyssey outputs for all 200 rows. "
        "The original SpeechBrain call failed because its generic classify_file path "
        "expected compute_features; the adapter was changed to call the checkpoint's "
        "declared wav2vec2, avg_pool, output_mlp, softmax, and label_encoder modules. "
        "The pilot must be accepted only after the rerun reports speechbrain_success_rows "
        "greater than zero. The WavLM/Wav2Vec2 loader also reports unused or newly "
        "initialized positional-convolution weights under the current Transformers stack. "
        "Until that compatibility warning is resolved or documented by a comparison run, "
        "the audio values are provisional evidence and not validated labels."
    )
    doc.add_heading("Student-level decision rule", level=1)
    doc.add_paragraph(
        "I use audio SER to prioritize review, not to manufacture ground truth. Agreement "
        "between the Phase 1 prediction, audio prosody, transcript scope, and visible "
        "speaker can increase review priority; disagreement is scientifically useful and "
        "must be retained. No output may be called deceptive, truthful, credible, or "
        "unreliable, and no audio score alone may assign a courtroom-affect category."
    )
    doc.save(path)
    return True


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 16, color=(31, 45, 61), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.name = "Aptos"; r.font.size = PptPt(size); r.font.bold = bold; r.font.color.rgb = RGBColor(*color)
    return box


def add_slide(prs: Presentation, title: str, subtitle: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(31, 78, 121); bar.line.fill.background()
    add_text(slide, title, .55, .35, 12, .5, 23, bold=True)
    add_text(slide, subtitle, .58, .88, 12, .35, 11, color=(90, 100, 110))
    text = "\n".join("• " + b for b in bullets)
    add_text(slide, text, .8, 1.55, 11.7, 4.9, 17)


def update_pptx(path: Path) -> bool:
    if not path.exists():
        return False
    prs = Presentation(path)
    if any(MARKER in shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")):
        return False
    backup(path)
    add_slide(prs, "Independent Audio SER Evidence", "A separate prosody layer, not a replacement for human labels", [
        "Odyssey WavLM evidence: continuous valence, arousal, and dominance-like scores.",
        "SpeechBrain Wav2Vec2-IEMOCAP: four-class categorical cross-check only.",
        "Both models read the utterance WAV and append provenance-preserving columns.",
        "They do not overwrite Phase 1 predictions or assign courtroom-affect labels.",
    ])
    add_slide(prs, "Conservative Fusion and Review", "Why disagreement is retained instead of hidden", [
        "Compare transcript scope, Phase 1 emotion, audio prosody, role, and visual review.",
        "A witness can calmly describe another person's severe distress; semantic words are not speaker emotion.",
        "Use disagreement to prioritize human review and add emotion_target_scope.",
        "Never infer deception, truthfulness, credibility, or reliability from these models.",
    ])
    prs.save(path)
    return True


def main() -> None:
    changed_docs = [str(p) for p in DOCX_TARGETS if update_docx(p)]
    changed_pptx = [str(p) for p in PPTX_TARGETS if update_pptx(p)]
    print({"updated_docx": changed_docs, "updated_pptx": changed_pptx, "marker": MARKER})


if __name__ == "__main__":
    main()
