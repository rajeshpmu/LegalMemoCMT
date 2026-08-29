"""Append the corrected conservative review policy to guidance DOCX/PPTX files."""
from pathlib import Path
import shutil
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt

ROOT = Path(__file__).resolve().parents[2]
MARKER = "EMOTION_SCOPE_CORRECTION_V1"
DOCS = [
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Corpus_SOP.docx",
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Student_Speaking_Guide.docx",
    ROOT / "implementation_docments/LegalMemoCMT_Phase2_Corpus_Build_Code_Level_Student_Guide.docx",
]
PPTS = [
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Presentation.pptx",
    ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx",
]


def update_doc(path: Path) -> bool:
    doc = Document(path)
    if any(MARKER in p.text for p in doc.paragraphs):
        return False
    backup = path.with_name(path.name + ".before_emotion_scope_correction.docx")
    if not backup.exists(): shutil.copy2(path, backup)
    doc.add_heading("Correction: Conservative Active-Review Outcomes", level=1)
    p = doc.add_paragraph(); p.add_run(MARKER).font.size = Pt(8)
    doc.add_paragraph(
        "The first implementation used low arousal plus SpeechBrain neutral as a general "
        "neutralization rule. That was too aggressive because it changed ordinary Phase 1 "
        "non-neutral predictions without evidence that the transcript described another "
        "person's emotion. The corrected implementation keeps the special "
        "OTHER_PERSON_DESCRIBED rule, but changes ordinary conflicts to UNRESOLVED."
    )
    doc.add_paragraph(
        "The three review outcomes are KEEP_PHASE1_CANDIDATE when comparable predictions "
        "agree or SpeechBrain's class is not comparable; NEUTRAL_CANDIDATE when emotion "
        "language targets another person and neutral, low/moderate-arousal audio supports "
        "a controlled delivery; and UNRESOLVED when Phase 1 and comparable audio evidence "
        "conflict without sufficient scope evidence. UNRESOLVED is a deliberate scientific "
        "decision, not a failed row."
    )
    doc.add_paragraph(
        "SpeechBrain classes are mapped only as neu->neutral, ang->anger, hap->joy, and "
        "sad->sadness. Its lack of fear, disgust, and surprise is not evidence against those "
        "MELD classes. Agreement is not counted as disagreement, and arousal alone is not a "
        "label decision. The original model outputs remain unchanged."
    )
    doc.add_paragraph(
        "After rerunning the 200-row pilot, the corrected distribution was 40 HIGH, 28 "
        "MEDIUM, and 132 LOW priority; proposed outcomes were 133 neutral, 1 sadness, and "
        "66 UNRESOLVED. The queue is now a transparent conflict-review queue. HIGH rows "
        "must be checked against transcript scope, audio delivery, and video before any "
        "human label is entered."
    )
    doc.add_paragraph(
        "For DCBWoWhsTpA_turn06575, the stronger rule remains applicable: the transcript "
        "describes another person's distress, SpeechBrain gives neu, and arousal is below "
        "the conservative threshold. The output is therefore a neutral/CALM_COMPOSED "
        "candidate with HIGH review priority, not an automatic final label."
    )
    doc.save(path)
    return True


def add_text(slide, text, x, y, w, h, size=17, color=(31,45,61), bold=False):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); box.text_frame.word_wrap=True
    p=box.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=text; r.font.name="Aptos"; r.font.size=PptPt(size); r.font.bold=bold; r.font.color.rgb=RGBColor(*color)


def update_ppt(path: Path) -> bool:
    prs=Presentation(path)
    if any(MARKER in shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape,"text")):
        return False
    backup=path.with_name(path.name+".before_emotion_scope_correction.pptx")
    if not backup.exists(): shutil.copy2(path,backup)
    slide=prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=RGBColor(255,255,255)
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.2)); bar.fill.solid(); bar.fill.fore_color.rgb=RGBColor(31,78,121); bar.line.fill.background()
    add_text(slide,"Corrected Active-Review Policy",.55,.35,12,.5,23,bold=True)
    add_text(slide,"Disagreement is a queue for human review, not an automatic relabel",.58,.88,12,.35,11,color=(90,100,110))
    add_text(slide,"EMOTION_SCOPE_CORRECTION_V1\n\n• Keep Phase 1 candidate when comparable outputs agree.\n• Propose NEUTRAL_CANDIDATE only for other-person emotion scope plus neutral low/moderate-arousal audio.\n• Use UNRESOLVED for ordinary Phase 1/audio conflicts.\n• SpeechBrain cannot represent fear, disgust, or surprise; it cannot erase them.\n• Pilot after correction: 40 HIGH, 28 MEDIUM, 132 LOW.\n• Human review is required before final labels.",.8,1.45,11.7,5.3,17)
    prs.save(path); return True


def main():
    print({"updated_docs":[str(p) for p in DOCS if update_doc(p)],"updated_pptx":[str(p) for p in PPTS if update_ppt(p)],"marker":MARKER})


if __name__ == "__main__": main()
