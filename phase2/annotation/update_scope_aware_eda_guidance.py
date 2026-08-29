"""Append scope-aware EDA and scoring explanations to Phase 2 artifacts."""
from pathlib import Path
import shutil
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt

ROOT = Path(__file__).resolve().parents[2]
MARKER = "SCOPE_AWARE_EDA_GUIDANCE_V1"
DOCS = [
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Corpus_SOP.docx",
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Student_Speaking_Guide.docx",
    ROOT / "implementation_docments/LegalMemoCMT_Phase2_Corpus_Build_Code_Level_Student_Guide.docx",
]
PPTS = [
    ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Presentation.pptx",
    ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx",
]


def update_doc(path: Path) -> bool:
    doc = Document(path)
    if any(MARKER in p.text for p in doc.paragraphs): return False
    backup = path.with_name(path.name + ".before_scope_aware_eda_guidance.docx")
    if not backup.exists(): shutil.copy2(path, backup)
    doc.add_heading("Scope-Aware Candidate Labeling and EDA", level=1)
    p=doc.add_paragraph(); p.add_run(MARKER).font.size=Pt(8)
    doc.add_paragraph(
        "The scope-aware layer does not replace Phase 1 or human labels. It creates a "
        "machine_proposed_basic_emotion, review outcome, and review priority by joining "
        "Phase 1 and audio-SER rows on utterance_id. The original phase1_basic_emotion, "
        "audio_emotion_candidate, audio_valence, audio_arousal, and audio_dominance remain "
        "unchanged for auditability."
    )
    doc.add_heading("Rules used by the algorithm", level=2)
    doc.add_paragraph(
        "SpeechBrain labels are mapped only where its four-class vocabulary is comparable: "
        "neu to neutral, ang to anger, hap to joy, and sad to sadness. Fear, disgust, and "
        "surprise are not treated as contradicted merely because SpeechBrain cannot output "
        "them. Neutral audio is defined as mapped audio=neutral and Odyssey arousal below "
        "0.45. CALM_COMPOSED is a behavioral presentation candidate, not proof of neutral "
        "emotion."
    )
    doc.add_paragraph(
        "The stronger rule is: neutral audio + low/moderate arousal + no explicit contrary "
        "visual/audio failure + a non-self scope (OTHER_PERSON_DESCRIBED, EVENT_DESCRIBED, "
        "or QUOTED_SPEECH) produces NEUTRAL_CANDIDATE at 0.72. SELF_EXPRESSED produces "
        "UNRESOLVED. UNCLEAR scope produces a lower-confidence neutral candidate at 0.58 "
        "and should be interpreted as weaker evidence."
    )
    doc.add_heading("Disagreement and priority scoring", level=2)
    doc.add_paragraph(
        "The categorical conflict indicator is: categorical_disagreement = 1 when a Phase 1 "
        "label and a comparable mapped SpeechBrain label differ; otherwise it is 0. The "
        "priority is HIGH when that conflict also has high semantic scope risk, or when both "
        "model confidences are strong (Phase 1 at least 0.70 and SpeechBrain at least 0.80). "
        "It is MEDIUM for a meaningful categorical conflict or ambiguous semantic emotion "
        "scope, and LOW when there is no meaningful conflict. Agreement is not counted as "
        "disagreement. Low arousal alone is not a conflict signal."
    )
    doc.add_paragraph(
        "The review outcome is KEEP_PHASE1_CANDIDATE for compatible agreement or non-comparable "
        "audio classes, NEUTRAL_CANDIDATE for the scoped neutral rule, and UNRESOLVED for a "
        "conflict that lacks enough evidence for a recommendation. human_review_required is "
        "set to YES for HIGH/MEDIUM priority or UNRESOLVED. Therefore, some low-priority "
        "neutral candidates can have human_review_required=NO; this means the row is not in "
        "the first review queue, not that it is gold-labeled. If every proposal must be "
        "human-confirmed, this policy should be tightened before training."
    )
    doc.add_heading("200-row scope-aware EDA interpretation", level=2)
    doc.add_paragraph(
        "The pilot contained 200 rows from one source video and one dev split. Phase 1 had "
        "163 neutral rows. The scope-aware output had 178 neutral, 8 sadness, 5 anger, 5 "
        "disgust, 1 fear, 1 joy, and 2 UNRESOLVED. There were 132 NEUTRAL_CANDIDATE, 66 "
        "KEEP_PHASE1_CANDIDATE, and 2 UNRESOLVED outcomes. The target-scope detector found "
        "151 UNCLEAR, 38 QUOTED_SPEECH, 9 OTHER_PERSON_DESCRIBED, and 2 SELF_EXPRESSED rows."
    )
    doc.add_paragraph(
        "The conclusion is that this output is a useful active-review and leakage-detection "
        "artifact, not a final label manifest. It shows how a courtroom witness may describe "
        "severe distress while speaking with low arousal. It also shows why the Phase 1 "
        "prediction, audio evidence, target scope, proposed label, and human decision must "
        "remain separate fields."
    )
    doc.save(path); return True


def add_text(slide, text, x, y, w, h, size=16, color=(31,45,61), bold=False):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); box.text_frame.word_wrap=True
    p=box.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=text; r.font.name="Aptos"; r.font.size=PptPt(size); r.font.bold=bold; r.font.color.rgb=RGBColor(*color)


def update_ppt(path: Path) -> bool:
    prs=Presentation(path)
    if any(MARKER in shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape,"text")): return False
    backup=path.with_name(path.name+".before_scope_aware_eda_guidance.pptx")
    if not backup.exists(): shutil.copy2(path,backup)
    for title, subtitle, body in [
        ("Scope-Aware Candidate Rules", "Separating speaker emotion from emotion described in testimony", "SCOPE_AWARE_EDA_GUIDANCE_V1\n\n• Map only comparable SpeechBrain classes: neu→neutral, ang→anger, hap→joy, sad→sadness.\n• Neutral candidate requires neutral audio + arousal < 0.45 + compatible target scope.\n• OTHER_PERSON_DESCRIBED / EVENT_DESCRIBED / QUOTED_SPEECH → candidate at 0.72.\n• SELF_EXPRESSED → UNRESOLVED.\n• UNCLEAR → lower-confidence neutral candidate at 0.58, reviewable."),
        ("Agreement, Disagreement, and Confidence", "Machine recommendations are prioritization evidence, not gold labels", "• categorical_disagreement = 1 only when comparable Phase 1 and SpeechBrain labels differ.\n• HIGH: conflict + high scope risk, or strong Phase 1 and audio confidence.\n• MEDIUM: meaningful conflict or ambiguous semantic scope.\n• LOW: agreement or no meaningful conflict.\n• CALM_COMPOSED describes delivery, not basic emotion.\n• human_review_required=NO means low-priority queue exclusion, not human validation."),
        ("Scope-Aware Pilot EDA", "200 rows: one source video, one dev split", "• Phase 1 neutral: 163; scope-aware neutral: 178.\n• Outcomes: 132 NEUTRAL_CANDIDATE, 66 KEEP_PHASE1_CANDIDATE, 2 UNRESOLVED.\n• Scope: 151 UNCLEAR, 38 QUOTED_SPEECH, 9 OTHER_PERSON_DESCRIBED, 2 SELF_EXPRESSED.\n• Priority: 40 HIGH, 28 MEDIUM, 132 LOW.\n• Conclusion: useful review artifact; not a final training-label manifest."),
    ]:
        slide=prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=RGBColor(255,255,255)
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.2)); bar.fill.solid(); bar.fill.fore_color.rgb=RGBColor(31,78,121); bar.line.fill.background()
        add_text(slide,title,.55,.35,12,.5,23,bold=True); add_text(slide,subtitle,.58,.88,12,.35,11,color=(90,100,110)); add_text(slide,body,.8,1.45,11.7,5.6,16)
    prs.save(path); return True


def main():
    print({"updated_docs":[str(p) for p in DOCS if update_doc(p)],"updated_pptx":[str(p) for p in PPTS if update_ppt(p)],"marker":MARKER})


if __name__ == "__main__": main()
