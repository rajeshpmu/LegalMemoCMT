"""Add courtroom-affect EDA slides and detailed speaking notes."""
from pathlib import Path
import shutil
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt

ROOT=Path(__file__).resolve().parents[2]
MARKER="COURTROOM_AFFECT_EDA_GUIDANCE_V1"
DOCS=[ROOT/"implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Student_Speaking_Guide.docx",ROOT/"implementation_docments/LegalMemoCMT_Phase2_Corpus_Build_Code_Level_Student_Guide.docx"]
PPTS=[ROOT/"implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Presentation.pptx",ROOT/"implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"]

def update_doc(path):
 doc=Document(path)
 if any(MARKER in p.text for p in doc.paragraphs): return False
 b=path.with_name(path.name+".before_courtroom_affect_eda_guidance.docx")
 if not b.exists(): shutil.copy2(path,b)
 doc.add_heading("Courtroom-Affect Candidate EDA: Current Position and Next Plan",level=1)
 p=doc.add_paragraph(); p.add_run(MARKER).font.size=Pt(8)
 doc.add_paragraph("The courtroom_affect_candidates_200.csv file is a machine-assisted review artifact. It contains 200 human-verified witness-speaking rows from one Clancy source video and does not contain final human courtroom-affect labels.")
 doc.add_heading("What the EDA shows",level=2)
 doc.add_paragraph("The candidate distribution is NEUTRAL_CALM=148, UNKNOWN=33, ASSERTIVE=11, GUARDED=5, DISTRESSED=2, and TENSE=1. DEFENSIVE and AGITATED have zero candidates. This does not prove that most witnesses are neutral or calm; it shows what the current rules select from this source and feature set.")
 doc.add_paragraph("Audio arousal has mean=0.371918, median=0.368324, and maximum=0.577063. SpeechBrain outputs are neu=151, ang=30, hap=17, and sad=2. Every row is marked courtroom_affect_review_required=YES, so these values remain review candidates.")
 doc.add_heading("How to interpret the candidate rules",level=2)
 doc.add_paragraph("NEUTRAL_CALM is supported mainly by neutral SpeechBrain output, low/moderate arousal, and the existing human visual-speaker-match fields. GUARDED uses qualification or refusal wording with lower dominance. ASSERTIVE uses high dominance and firm-response wording. DISTRESSED requires negative valence and elevated arousal. TENSE requires elevated arousal, negative valence, and hesitation. These rules provide transparent evidence, not validated psychological measurements.")
 doc.add_heading("Important EDA finding",level=2)
 doc.add_paragraph("Fourteen NEUTRAL_CALM rows were assigned without low/moderate-arousal evidence because the scoring threshold could be reached from neutral SpeechBrain plus visual-match evidence. This means the implementation should be tightened so NEUTRAL_CALM requires arousal<0.45 explicitly, or those rows should become UNKNOWN. This is a rule-quality issue, not evidence that the rows are wrong.")
 doc.add_heading("What I infer and what I will do next",level=2)
 doc.add_paragraph("I infer that the pipeline can produce an auditable candidate queue and expose missing evidence, but cannot yet produce final courtroom-affect labels. I will correct the strict NEUTRAL_CALM gate, rerun the 200-row EDA, inspect the changed rows, then obtain human labels for a stratified sample across multiple source videos. I will not use these weak candidates as gold labels or infer credibility, truthfulness, deception, or reliability.")
 doc.add_paragraph("Student explanation: I can say that the EDA tells me how the current rules behave, not what the true emotional distribution is. The next milestone is calibration against human review, followed by evaluation on held-out source videos.")
 doc.save(path); return True

def add_text(slide,text,x,y,w,h,size=16,color=(31,45,61),bold=False):
 box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); box.text_frame.word_wrap=True
 p=box.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=text; r.font.name="Aptos"; r.font.size=PptPt(size); r.font.bold=bold; r.font.color.rgb=RGBColor(*color)

def update_ppt(path):
 prs=Presentation(path)
 if any(MARKER in sh.text for sl in prs.slides for sh in sl.shapes if hasattr(sh,'text')): return False
 b=path.with_name(path.name+".before_courtroom_affect_eda_guidance.pptx")
 if not b.exists(): shutil.copy2(path,b)
 slides=[
  ("Courtroom-Affect Candidate EDA","200-row pilot: candidates, not final labels","COURTROOM_AFFECT_EDA_GUIDANCE_V1\n\n• NEUTRAL_CALM: 148 (74.0%)\n• UNKNOWN: 33 (16.5%)\n• ASSERTIVE: 11 (5.5%)\n• GUARDED: 5 (2.5%)\n• DISTRESSED: 2 (1.0%)\n• TENSE: 1 (0.5%)\n• DEFENSIVE / AGITATED: 0\n\nAll 200 rows are Witness + SPEAKING and require human review."),
  ("Audio Evidence and Rule Coverage","What supports the candidates and where evidence is missing","• Arousal: mean 0.371918, median 0.368324, maximum 0.577063.\n• SpeechBrain: neu 151, ang 30, hap 17, sad 2.\n• Scope: UNCLEAR 151, QUOTED_SPEECH 38, OTHER_PERSON_DESCRIBED 9, SELF_EXPRESSED 2.\n• Candidate rules use transcript wording, audio V/A/D, and visual-review fields.\n• Interaction context, speech-rate change, and verified affect behavior remain limited."),
  ("Current Inference and Next Action","The EDA guides calibration; it does not establish prevalence","• 14 NEUTRAL_CALM rows lacked low/moderate-arousal evidence but passed the score threshold.\n• Tighten the rule: require arousal<0.45 explicitly, or set those rows UNKNOWN.\n• Rerun EDA and inspect changed rows.\n• Human-label a stratified multi-source sample before training.\n• Keep machine evidence and human courtroom-affect labels separate.\n• Never infer credibility, truthfulness, deception, or reliability."),
 ]
 for title,sub,body in slides:
  sl=prs.slides.add_slide(prs.slide_layouts[6]); sl.background.fill.solid(); sl.background.fill.fore_color.rgb=RGBColor(255,255,255)
  bar=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.2)); bar.fill.solid(); bar.fill.fore_color.rgb=RGBColor(31,78,121); bar.line.fill.background()
  add_text(sl,title,.55,.35,12,.5,23,bold=True); add_text(sl,sub,.58,.88,12,.35,11,color=(90,100,110)); add_text(sl,body,.8,1.45,11.7,5.4,16)
 marker=prs.slides[-1].shapes.add_textbox(Inches(0),Inches(0),Inches(.01),Inches(.01)); marker.text=MARKER; marker.text_frame.paragraphs[0].runs[0].font.size=Pt(1)
 prs.save(path); return True

def main(): print({"updated_docs":[str(p) for p in DOCS if update_doc(p)],"updated_pptx":[str(p) for p in PPTS if update_ppt(p)],"marker":MARKER})
if __name__=="__main__": main()
