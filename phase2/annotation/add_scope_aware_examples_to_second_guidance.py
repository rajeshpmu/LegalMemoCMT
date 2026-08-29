"""Add 200-row, data-grounded example callouts to Second Guidance slides 13-31."""
from __future__ import annotations

import csv
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx.before_scope_aware_examples.pptx"
MARKER = "PILOT_EXAMPLE_CALLOUTS_V1"
DATA = ROOT / "data/processed/phase2/clancy/emotion_scope_review_200_scope_aware.csv"


def clean(value: object) -> str:
    return str(value or "").strip()


def load_rows() -> dict[str, dict[str, str]]:
    with DATA.open(newline="", encoding="utf-8-sig") as handle:
        return {row["utterance_id"]: row for row in csv.DictReader(handle)}


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int, color=(31, 45, 61), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.name = "Aptos"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = RGBColor(*color)
    return box


def add_callout(slide, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.55), Inches(6.42), Inches(12.2), Inches(.72))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(232, 242, 250); shape.fill.transparency = 4
    shape.line.color.rgb = RGBColor(144, 180, 205)
    shape.text_frame.clear(); shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Pilot example: " + text; r.font.name = "Aptos"; r.font.size = Pt(10); r.font.color.rgb = RGBColor(31, 45, 61)


def main() -> None:
    rows = load_rows()
    target = rows["DCBWoWhsTpA_turn06575"]
    apology = rows["DCBWoWhsTpA_turn06057"]
    name = rows["DCBWoWhsTpA_turn05916"]
    examples = {
        13: "The 200-row evidence is Clancy-only: DCBWoWhsTpA contains witness rows; no tribunal, Tupac/Keffe D, or Indian-SIM row is claimed by this pilot.",
        14: f"{target['utterance_id']} carries text plus audio/video paths while its behavioral and basic-emotion fields remain separate.",
        15: f"{target['phase1_basic_emotion']} is the Phase 1 label for {target['utterance_id']}; the scope-aware layer preserves it and adds a separate proposed value.",
        16: f"For {target['utterance_id']}, negative words describe another person, while the proposed affect is {target['proposed_courtroom_affect']}; affect is not credibility.",
        17: f"{name['utterance_id']} is traceable to source {name['youtube_id']} and its clip paths; this is one row in the downloaded-source evidence chain.",
        18: f"A concrete row path is {target['utterance_id']}: subtitle-derived text -> timestamped MP4/WAV -> Phase 1/audio evidence -> review fields.",
        19: "The pilot EDA example is 200 Clancy witness-speaking rows; it does not claim the full Clancy corpus or a final emotion distribution.",
        20: "Boundary example: the scope-aware 200-row CSV has no IRMCT/ICTY rows, so tribunal statistics must come from the separate tribunal manifest.",
        21: f"The pilot includes duration values such as {apology['clip_duration_seconds']} seconds; duration filtering and emotion proposals are separate decisions.",
        22: "Boundary example: no Tupac/Keffe D row exists in this 200-row output; it remains a planned branch, not processed evidence.",
        23: "Boundary example: no Indian-SIM row exists in this output; Indian adaptation requires separate acquisition and verification.",
        24: "The 200-row pilot is single-source, so cross-corpus duration comparison cannot be inferred from this file alone.",
        25: f"A guidance-call question can use {target['utterance_id']}: should this be human-reviewed as neutral/CALM_COMPOSED or retained as a Phase 1 conflict case?",
        26: f"{target['utterance_id']} has Odyssey valence={target['audio_valence']}, arousal={target['audio_arousal']} and SpeechBrain={target['audio_emotion_candidate']}; labels are not overwritten.",
        27: f"{target['utterance_id']} shows the key disagreement: Phase 1={target['phase1_basic_emotion']} but SpeechBrain={target['audio_emotion_candidate']}; the transcript describes another person's distress.",
        28: f"{apology['utterance_id']} remains an example of review rather than certainty: Phase 1={apology['phase1_basic_emotion']}, audio={apology['audio_emotion_candidate']}, scope={apology['emotion_target_scope']}.",
        29: f"{target['utterance_id']} satisfies the scope-aware pattern: scope={target['emotion_target_scope']}, audio={target['audio_emotion_candidate']}, arousal={target['audio_arousal']}, proposed={target['proposed_basic_emotion']}.",
        30: f"{name['utterance_id']} illustrates a conflict requiring inspection: Phase 1={name['phase1_basic_emotion']}, SpeechBrain={name['audio_emotion_candidate']}; agreement is not assumed from file existence.",
        31: f"The actual pilot summary is 200 rows, 151 UNCLEAR scopes, 132 NEUTRAL_CANDIDATE outcomes, and 2 UNRESOLVED outcomes; this is a review artifact, not gold annotation.",
    }
    prs = Presentation(PPTX)
    if BACKUP.exists() or any(MARKER in shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")):
        print("Already updated", PPTX); return
    if not BACKUP.exists(): shutil.copy2(PPTX, BACKUP)
    for number, text in examples.items():
        add_callout(prs.slides[number - 1], text)
    # Keep the marker in a tiny hidden text box so the updater is idempotent.
    marker_box = prs.slides[-1].shapes.add_textbox(Inches(0), Inches(0), Inches(.01), Inches(.01))
    marker_box.text = MARKER
    marker_box.text_frame.paragraphs[0].runs[0].font.size = Pt(1)
    prs.save(PPTX)
    print({"updated": str(PPTX), "slides_updated": sorted(examples), "marker": MARKER})


if __name__ == "__main__":
    main()
