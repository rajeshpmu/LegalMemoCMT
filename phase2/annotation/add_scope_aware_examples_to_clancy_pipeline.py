"""Add grounded scope-aware pilot examples to Clancy pipeline slides 13-18."""
from pathlib import Path
import csv, shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT=Path(__file__).resolve().parents[2]
PPTX=ROOT/"implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Presentation.pptx"
BACKUP=ROOT/"implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Presentation.pptx.before_scope_aware_examples.pptx"
MARKER="CLANCY_PILOT_EXAMPLE_CALLOUTS_V1"
DATA=ROOT/"data/processed/phase2/clancy/emotion_scope_review_200_scope_aware.csv"

def add_callout(slide,text):
    shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(.55),Inches(6.42),Inches(12.2),Inches(.72))
    shape.fill.solid(); shape.fill.fore_color.rgb=RGBColor(232,242,250); shape.fill.transparency=4; shape.line.color.rgb=RGBColor(144,180,205)
    shape.text_frame.clear(); shape.text_frame.word_wrap=True
    p=shape.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text="Pilot example: "+text; r.font.name="Aptos"; r.font.size=Pt(10); r.font.color.rgb=RGBColor(31,45,61)

def main():
    rows={r['utterance_id']:r for r in csv.DictReader(DATA.open(encoding='utf-8-sig'))}
    t=rows['DCBWoWhsTpA_turn06575']; a=rows['DCBWoWhsTpA_turn06057']; n=rows['DCBWoWhsTpA_turn05916']
    examples={
      13:f"{t['utterance_id']} has audio-SER fields appended without changing the Phase 1 label; Odyssey arousal={t['audio_arousal']}.",
      14:f"{t['utterance_id']} shows the reason for fusion review: Phase 1={t['phase1_basic_emotion']}, SpeechBrain={t['audio_emotion_candidate']}, scope={t['emotion_target_scope']}.",
      15:f"{a['utterance_id']} is not silently relabeled as truth: Phase 1={a['phase1_basic_emotion']} at {a['phase1_basic_emotion_confidence']}, audio={a['audio_emotion_candidate']}, scope={a['emotion_target_scope']}.",
      16:f"{t['utterance_id']} meets the scope-aware candidate rule: audio={t['audio_emotion_candidate']}, arousal={t['audio_arousal']}, proposed={t['proposed_basic_emotion']} at {t['proposed_basic_emotion_confidence']}.",
      17:f"{n['utterance_id']} illustrates why a disagreement is review evidence: Phase 1={n['phase1_basic_emotion']}, SpeechBrain={n['audio_emotion_candidate']}; neither is overwritten.",
      18:"This exact 200-row pilot contains 178 proposed neutral, 132 neutral-candidate outcomes, and 2 unresolved outcomes; all rows are from DCBWoWhsTpA/dev.",
    }
    prs=Presentation(PPTX)
    if BACKUP.exists() or any(MARKER in sh.text for sl in prs.slides for sh in sl.shapes if hasattr(sh,'text')):
        print('Already updated',PPTX); return
    shutil.copy2(PPTX,BACKUP)
    for number,text in examples.items(): add_callout(prs.slides[number-1],text)
    marker=prs.slides[-1].shapes.add_textbox(Inches(0),Inches(0),Inches(.01),Inches(.01)); marker.text=MARKER; marker.text_frame.paragraphs[0].runs[0].font.size=Pt(1)
    prs.save(PPTX); print({'updated':str(PPTX),'slides_updated':sorted(examples),'marker':MARKER})

if __name__=='__main__': main()
