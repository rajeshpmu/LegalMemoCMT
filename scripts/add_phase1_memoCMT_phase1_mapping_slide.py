from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph
from docx.shared import Pt, RGBColor, Inches
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
READING_DOCX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Reading_Script.docx"


def style_run(run, *, size=10, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def style_ppt_run(run, *, size=10, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color)


def add_box(slide, left, top, width, height, fill, line=None, radius=False):
    shp_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shp_type,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = PptRGBColor.from_string(fill)
    if line:
        shp.line.color.rgb = PptRGBColor.from_string(line)
    else:
        shp.line.color.rgb = PptRGBColor.from_string(fill)
    return shp


def add_text(slide, left, top, width, height, text, *, size=11, bold=False, color="000000", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    style_ppt_run(r, size=size, bold=bold, color=color)
    return tb


def insert_slide_before(prs: Presentation, before_index_0_based: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(before_index_0_based, sldId)
    return slide


def renumber_slide_numbers(prs: Presentation):
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        for shp in slide.shapes:
            if not getattr(shp, "has_text_frame", False):
                continue
            txt = shp.text_frame.text.strip()
            if not re.fullmatch(r"\d+", txt):
                continue
            if float(shp.top) < 400000 and float(shp.left) > 9000000:
                shp.text = str(idx)
                for p in shp.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.RIGHT
                    for run in p.runs:
                        style_ppt_run(run, size=12, color="FFFFFF")


def build_slide():
    prs = Presentation(str(PPTX_PATH))
    if any(
        "MemoCMT Paper vs Phase 1 Implementation Mapping" in (shp.text_frame.text if getattr(shp, "has_text_frame", False) else "")
        for slide in prs.slides
        for shp in slide.shapes
    ):
        print("Comparison slide already exists; skipping slide insertion.")
        renumber_slide_numbers(prs)
        prs.save(str(PPTX_PATH))
        return
    slide = insert_slide_before(prs, 5)  # insert before current slide 6

    add_box(slide, 0, 0, 13.333, 0.58, "122F55")
    add_box(slide, 0, 0.58, 13.333, 0.08, "267377")
    add_text(
        slide,
        0.42,
        0.12,
        11.3,
        0.3,
        "6. MemoCMT Paper vs Phase 1 Implementation Mapping",
        size=17.5,
        bold=True,
        color="FFFFFF",
    )
    add_text(slide, 12.0, 0.12, 0.8, 0.3, "6", size=12, color="FFFFFF", align=PP_ALIGN.RIGHT)
    add_text(
        slide,
        0.5,
        0.68,
        12.3,
        0.28,
        "This slide shows the exact mapping from the paper’s design to the Phase 1 code, and where Phase 1 intentionally extends the paper.",
        size=10.5,
        color="5E5E5E",
    )

    headers = ["Component", "Base paper MemoCMT", "Phase 1 implementation", "Code reference / meaning"]
    rows = [
        [
            "Text encoder",
            "BERT",
            'PretrainedBackboneEncoder("bert-base-uncased") in paper mode',
            "src/models/model.py:249-296, 304-316",
        ],
        [
            "Audio encoder",
            "HuBERT",
            'PretrainedBackboneEncoder("facebook/hubert-base-ls960") in paper mode',
            "src/models/model.py:249-296, 311-316",
        ],
        [
            "CMT fusion",
            "Explicit Wq/Wk/Wv cross-attention equations",
            "BidirectionalCrossAttentionCMT using nn.MultiheadAttention",
            "src/models/model.py:160-246, 347-352",
        ],
        [
            "Pooling",
            "CLS / MEAN / MAX / MIN",
            'Same pooling options; paper-aligned MELD uses min',
            "src/models/model.py:190-208, 345-351; scripts/run_paper_aligned_meld_cv.sh:42-62",
        ],
        [
            "Video modality",
            "N/A",
            "SequenceEncoder over video features, then legacy_fusion or gated_fusion",
            "src/models/model.py:333-445",
        ],
        [
            "Auxiliary video head",
            "N/A",
            "video_aux_classifier attached to the video branch",
            "src/models/model.py:353-360, 429-445",
        ],
        [
            "Evaluation story",
            "MELD case study reported in the paper",
            "5-fold MELD CV + held-out test metrics + raw-mp4 demo path",
            "scripts/run_paper_aligned_meld_cv.sh:4-71; scripts/run_phase1_raw_mp4_demo.sh:9-66",
        ],
    ]

    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        PptInches(0.25),
        PptInches(1.12),
        PptInches(12.82),
        PptInches(5.8),
    ).table
    widths = [1.55, 3.05, 3.7, 4.52]
    for idx, w in enumerate(widths):
        table.columns[idx].width = PptInches(w)

    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptRGBColor.from_string("DDEBF7")
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                style_ppt_run(run, size=9, bold=True, color="122F55")

    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptRGBColor.from_string("F8FBFF")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j != 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    style_ppt_run(run, size=7.6 if j != 3 else 7.2, color="202020")

    add_box(slide, 0.28, 7.05, 12.75, 0.42, "EEF4FA", "C9D7E6", radius=True)
    add_text(
        slide,
        0.45,
        7.12,
        12.35,
        0.16,
        "Key takeaway: the paper-aligned baseline reproduces the paper’s text+audio CMT idea; Phase 1 keeps that path and adds the video extension, gated fusion, and auxiliary video loss.",
        size=9.3,
        bold=True,
        color="122F55",
    )

    renumber_slide_numbers(prs)
    prs.save(str(PPTX_PATH))
    print(f"Updated PPTX: {PPTX_PATH}")


def insert_paragraph_before(anchor_paragraph, text, *, size=11, bold=False, color="000000"):
    new_p = OxmlElement("w:p")
    anchor_paragraph._p.addprevious(new_p)
    paragraph = Paragraph(new_p, anchor_paragraph._parent)
    paragraph.text = text
    paragraph.paragraph_format.space_after = Pt(4)
    paragraph.paragraph_format.space_before = Pt(0)
    paragraph.paragraph_format.left_indent = Pt(0)
    for run in paragraph.runs:
        style_run(run, size=size, bold=bold, color=color)
    return paragraph


def insert_table_after(document, paragraph, headers, rows):
    table = document._body.add_table(rows=1, cols=len(headers), width=Inches(6.5))
    table.style = "Table Grid"
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = value
    paragraph._p.addnext(table._tbl)
    return table


def update_docx():
    doc = Document(str(READING_DOCX_PATH))

    # Shift slide headings from 6 onward by +1 to match the inserted slide.
    for para in doc.paragraphs:
        txt = para.text.strip()
        m = re.match(r"^Slide (\d+): (.+)$", txt)
        if not m:
            continue
        num = int(m.group(1))
        if num >= 6:
            para.text = f"Slide {num + 1}: {m.group(2)}"

    anchor = None
    for para in doc.paragraphs:
        if para.text.strip().startswith("Slide 7: Literature Review 2: MELD"):
            anchor = para
            break
    if anchor is None:
        raise RuntimeError("Could not locate the MELD slide heading in the reading script.")

    new_heading = insert_paragraph_before(anchor, "Slide 6: MemoCMT vs Phase 1 Implementation Mapping", size=14, bold=True, color="122F55")
    insert_paragraph_before(
        new_heading,
        "Use this slide to explain the exact relationship between the MemoCMT paper and the Phase 1 implementation. The safest explanation is that Phase 1 reproduces the paper’s text-plus-audio CMT baseline in paper mode, then extends the system with a video pathway and gated fusion for later experiments.",
        size=11.2,
    )
    insert_paragraph_before(
        new_heading,
        "Say that the paper’s CMT is written mathematically with explicit Wq, Wk, and Wv projections, but Phase 1 implements the same idea using PyTorch’s built-in nn.MultiheadAttention inside BidirectionalCrossAttentionCMT. That is a practical implementation choice, not a change in the underlying fusion principle.",
        size=11.2,
    )
    insert_paragraph_before(
        new_heading,
        "Then explain the important N/A case. The base MemoCMT paper is text-plus-audio, so video is not part of the original design. In Phase 1, video is added separately through SequenceEncoder and then combined with either legacy_fusion or gated_fusion, with an optional video_aux_classifier to support the auxiliary-loss branch.",
        size=11.2,
    )
    insert_paragraph_before(
        new_heading,
        "This means the comparison is not paper versus unrelated model. It is paper baseline versus a faithful Phase 1 reproduction plus a controlled visual extension. That distinction is important because the paper-aligned MELD result should be judged on text and audio, while the video branch should be discussed as a Phase 1 extension.",
        size=11.2,
    )
    insert_paragraph_before(
        new_heading,
        "Code references to mention aloud: src/models/model.py for the paper CMT and video branches, scripts/run_paper_aligned_meld_cv.sh for the paper-aligned text+audio baseline, and scripts/run_phase1_raw_mp4_demo.sh for the raw-mp4 demo route used by the visual experiments.",
        size=11.0,
        bold=True,
        color="122F55",
    )

    # Add a small reference table after the final inserted paragraph.
    tbl_anchor = new_heading
    headers = ["Area", "Paper", "Phase 1", "Exact code refs"]
    rows = [
        ["CMT fusion", "Explicit Wq/Wk/Wv formulas", "BidirectionalCrossAttentionCMT + MultiheadAttention", "src/models/model.py:160-246"],
        ["Text/audio", "BERT + HuBERT", "PretrainedBackboneEncoder in paper mode", "src/models/model.py:249-316"],
        ["Pooling", "CLS / MEAN / MAX / MIN", "Same options; MELD baseline uses min", "src/models/model.py:190-208; scripts/run_paper_aligned_meld_cv.sh:42-62"],
        ["Video", "N/A", "SequenceEncoder + legacy/gated fusion + aux head", "src/models/model.py:333-445"],
    ]
    table_para = insert_paragraph_before(tbl_anchor, "", size=1)
    table = insert_table_after(doc, table_para, headers, rows)
    for c in table.rows[0].cells:
        for p in c.paragraphs:
            for run in p.runs:
                style_run(run, size=9.4, bold=True, color="122F55")
    for row in table.rows[1:]:
        for cell in row.cells:
            for p in cell.paragraphs:
                for run in p.runs:
                    style_run(run, size=8.6, color="202020")

    # remove the empty placeholder paragraph we used for the table insertion
    if table_para.text == "":
        el = table_para._element
        el.getparent().remove(el)

    doc.save(str(READING_DOCX_PATH))
    print(f"Updated reading script: {READING_DOCX_PATH}")


def main():
    build_slide()
    update_docx()


if __name__ == "__main__":
    main()
