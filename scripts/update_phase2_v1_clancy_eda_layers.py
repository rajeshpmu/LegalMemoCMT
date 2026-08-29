"""Add full, diarization, and witness-visible Clancy EDA slides and notes."""
from pathlib import Path
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document

ROOT=Path(__file__).resolve().parents[1]
PPTX=ROOT/'implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Presentation_v1.pptx'
DOCX=ROOT/'implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Speaking_Doc.docx'
MARKER='CLANCY_EDA_LAYERS_V1'
NAVY=RGBColor(20,48,87); TEXT=RGBColor(35,35,35); GREY=RGBColor(95,95,95)

def text(s,v,x,y,w,h,size=16,color=TEXT,bold=False):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); b.text_frame.word_wrap=True; b.text_frame.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p=b.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=v; r.font.name='Aptos'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color

def frame(prs,h,sub):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor(255,255,255)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.22)); bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
    text(s,h,.55,.38,12.2,.55,24,NAVY,True); text(s,sub,.58,.96,12,.35,12,GREY); return s

def table(s,headers,rows,x=.6,y=1.4,w=12.1,h=5.4,size=12):
    t=s.shapes.add_table(len(rows)+1,len(headers),Inches(x),Inches(y),Inches(w),Inches(h)).table
    for c,v in enumerate(headers): t.cell(0,c).text=str(v)
    for ri,row in enumerate(rows,1):
        for ci,v in enumerate(row): t.cell(ri,ci).text=str(v)
    for ri,row in enumerate(t.rows):
        for cell in row.cells:
            cell.text_frame.word_wrap=True
            for p in cell.text_frame.paragraphs:
                for r in p.runs: r.font.name='Aptos'; r.font.size=Pt(size); r.font.color.rgb=NAVY if ri==0 else TEXT; r.font.bold=ri==0

def before(slides,s,index):
    ids=slides._sldIdLst; item=ids[-1]; ids.remove(item); ids.insert(index,item)

def update_ppt():
    prs=Presentation(PPTX)
    if any(MARKER in sh.text for s in prs.slides for sh in s.shapes if hasattr(sh,'text')): return
    backup=PPTX.with_name(PPTX.name+'.before_clancy_eda_layers.pptx')
    if not backup.exists(): shutil.copy2(PPTX,backup)
    index=9
    s=frame(prs,'Clancy EDA Layer 1: Full Source-Shared Utterance Corpus','All subtitle-derived rows before true turn consolidation and witness filtering')
    table(s,['Artifact / measure','Observed result','What it means'],[
        ['clancy_dataset_manifest.csv','77,442 rows; 11 source videos','Source-shared subtitle/utterance rows; many rows point to the same raw MP4/VTT'],
        ['Total represented duration','1,794.23 minutes / 29.904 hours','Sum of row durations, not unique media hours; overlapping/shared-source rows must not be treated as independent video'],
        ['Duration range','0.01–5,106.63 seconds','Contains very long rows and fragments; not yet MELD-style'],
        ['Purpose','Discovery and source coverage','Useful for knowing what was extracted, not sufficient for training selection'],
    ],size=12); before(prs.slides,s,index); index+=1
    s=frame(prs,'Clancy EDA Layer 2: Speaker-Cluster Processing','Diarization connects acoustic speaker intervals to turn rows but does not assign legal roles')
    table(s,['Artifact / measure','Observed result','What it means'],[
        ['clancy_turn_manifest_clipped_diarized.csv','11,926 turn rows; 11 sources','True turn-level processing output enriched with diarization linkage'],
        ['clancy_diarization_segments_all.csv','23,709 segments; 59 local clusters','Pyannote speech intervals and anonymous cluster IDs'],
        ['Segment duration represented','75,261.28 seconds / 20.906 hours','Diarized speech intervals; not automatically witness-only or unique corpus hours'],
        ['Segment-to-turn map','22,729 mapped; 980 without overlap','Mapping evidence is strong but 7,529 segments overlap multiple turns and require cautious interpretation'],
        ['Role limitation','Clusters are anonymous','Manual cluster-role mapping is required for Witness, Prosecutor, Defence, Judge, or Other'],
    ],size=10); before(prs.slides,s,index); index+=1
    s=frame(prs,'Clancy EDA Layer 3: Witness-Visible Speaking Corpus','Rows selected after role mapping, witness-speaking filtering, and persistent rejection rules')
    table(s,['Artifact / measure','Observed result','Training interpretation'],[
        ['clancy_witness_speaking_usable.csv','2,229 rows; 8 sources','Current usable witness-speaking candidate pool'],
        ['Usable duration','284.44 minutes / 4.741 hours','Rows are within approximately 0.8–30 seconds and are the strongest current witness pool'],
        ['Manual-review pool','350 rows; 8 sources','Review/outlier rows; not counted as usable until inspected'],
        ['Review duration','143.35 minutes / 2.389 hours','Includes long/outlier content and must not be added to usable hours automatically'],
        ['Role coverage','2,229 usable rows labelled Witness','Speaker role was manually mapped; visual presence claims still require the defined verification protocol'],
    ],size=11); text(s,'This is the layer relevant to witness-only Phase 2 training. It is different from the first-200 pilot used for model and heuristic validation.',.85,6.45,11.6,.35,13,GREY); before(prs.slides,s,index)
    index+=1
    s=frame(prs,'Clancy Pilot Scope and Final Corpus','Why 200 rows are used first, and what will be processed after pilot validation')
    table(s,['Stage','Rows / scope','Purpose and interpretation'],[
        ['Controlled pilot','200 rows from one controlled input batch','Validate Phase 1 inference, audio-SER evidence, scope-aware rules, affect candidates, gate conflicts, and manual review workflow. Small enough to inspect and rerun.'],
        ['Current witness pool','2,229 usable witness-speaking rows','Candidate pool after role mapping, witness filtering, duration rules, and persistent exclusions; this is larger than the pilot and remains subject to final quality checks.'],
        ['Final Clancy corpus','Full validated eligible pool','Rerun the pipeline after pilot rules are accepted; retain only rows passing role, alignment, audio/video, duration, visual, and annotation gates.'],
        ['Reported final size','Measured after filtering','Final utterance count, minutes, witnesses, sources, and splits must be calculated from the final manifest; they must not be extrapolated from 200 rows.'],
    ],size=11)
    text(s,'The 200-row output is a reproducible quality-control experiment, not the final dataset.',.85,6.45,11.6,.35,14,NAVY,True); before(prs.slides,s,index)
    prs.save(PPTX); print('Updated',PPTX,'slides',len(prs.slides))

def update_doc():
    doc=Document(DOCX)
    if any(MARKER in p.text for p in doc.paragraphs): return
    backup=DOCX.with_name(DOCX.name+'.before_clancy_eda_layers.docx')
    if not backup.exists(): shutil.copy2(DOCX,backup)
    doc.add_heading('Clancy EDA Layers and Pilot-to-Final Corpus Plan',1)
    p=doc.add_paragraph(); p.add_run(MARKER).font.size=Pt(8)
    doc.add_heading('Slide: Full source-shared utterance corpus',2)
    doc.add_paragraph('The full clancy_dataset_manifest.csv contains 77,442 subtitle-derived rows from 11 source videos and represents 1,794.23 summed minutes. This is an extraction/discovery layer, not the final MELD-style corpus. Multiple rows share the same raw video and subtitle files, and some durations are extremely long, so summing these rows does not mean 29.904 unique hours of usable training video. I use this layer to prove source coverage and to identify what needs turn consolidation and filtering.')
    doc.add_heading('Slide: Speaker-cluster processing',2)
    doc.add_paragraph('The diarization layer contains 11,926 turn rows and 23,709 Pyannote segments across 11 sources. The 59 speaker cluster IDs are local anonymous acoustic identities. Diarization answers which voice was active during a time interval; it does not know whether the voice belongs to a witness, lawyer, judge, or news reader. The segment-to-turn map records 22,729 segments with at least one overlapping turn and 980 without overlap. Because 7,529 segments overlap multiple turns, I must not interpret every overlap as a clean one-to-one speaker assignment.')
    doc.add_heading('Slide: Witness-visible speaking corpus',2)
    doc.add_paragraph('After role mapping, witness-speaking filtering, and persistent exclusions, clancy_witness_speaking_usable.csv contains 2,229 candidate witness rows from 8 sources, totalling 284.44 minutes or 4.741 hours. The separate 350-row review file contains 143.35 minutes or 2.389 hours and is not counted as usable until manually inspected. This is the correct layer for witness-only Phase 2 selection, not the 77,442-row source-shared manifest.')
    doc.add_heading('Why the first 200 rows are a pilot',2)
    doc.add_paragraph('The first 200 rows were chosen to validate the end-to-end path: feature availability, Phase 1 inference, audio-SER evidence, scope heuristics, courtroom-affect proposals, and the acceptance gate. A 200-row pilot is small enough to inspect manually and fast enough to rerun after rule changes. It is not a claim that only 200 rows form the final Clancy corpus, and its EDA should not be confused with full-corpus statistics.')
    doc.add_heading('What becomes the final Clancy corpus after the pilot',2)
    doc.add_paragraph('When the pilot rules are accepted, the final corpus is produced by rerunning the same pipeline over the complete validated witness-speaking pool, not by copying the first 200 rows. The expected order is: inspect and revise the 350 review/outlier rows; apply role, source, duration, audio, video, and alignment quality gates; generate feature and scope evidence for the full eligible pool; apply the acceptance policy; create group-aware train/dev/test splits; and retain unresolved rows separately. The final count and usable minutes must be measured from the resulting manifest, because they may decrease after quality controls or increase if reviewed rows are recovered.')
    doc.add_heading('How to explain the three layers',2)
    doc.add_paragraph('I should say: the full source-shared manifest proves extraction coverage; the diarization layer proves that speaker timing and clusters were processed; and the witness-visible speaking layer is the current training-selection pool. These are different denominators. I must not report source-shared minutes as unique usable corpus hours, and I must not treat a diarization cluster as a verified witness without role mapping and visual review.')
    doc.add_heading('Pilot scope versus final Clancy corpus',2)
    doc.add_paragraph('The first 200 rows are a controlled quality-control experiment. They allow me to test Phase 1 inference, audio-SER evidence, emotion-target scope, courtroom-affect heuristics, critical-conflict gates, and the manual-review procedure without spending the time and storage required for the complete pool. The pilot is intentionally not a random claim about the whole corpus and its statistics must be reported as pilot statistics.')
    doc.add_paragraph('After the pilot rules are accepted, I will rerun the same reproducible stages over the complete eligible witness-speaking pool. The final Clancy corpus will contain only rows that pass the configured role, duration, transcript/alignment, audio, video, visual-presence, and annotation-quality checks. Rows that remain unresolved or fail a gate will be retained in separate review or rejection manifests. Therefore the final number of utterances, minutes, witnesses, source videos, and split sizes must be measured from the final output manifest rather than extrapolated from the 200-row pilot.')
    doc.save(DOCX); print('Updated',DOCX)

def main(): update_ppt(); update_doc()
if __name__=='__main__': main()
