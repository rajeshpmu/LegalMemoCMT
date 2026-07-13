from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor
from pptx import Presentation


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
OUT_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Speaker_Taking_Doc.docx"


def style_run(run, *, size=11, bold=False, color="000000", name="Aptos"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_para(doc, text, *, size=11, bold=False, color="000000", indent=0.0):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(3)
    p.paragraph_format.space_before = Pt(0)
    if indent:
        p.paragraph_format.left_indent = Pt(indent)
    r = p.add_run(text)
    style_run(r, size=size, bold=bold, color=color)
    return p


def get_slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                first = text.splitlines()[0].strip()
                if first:
                    return first
    return "Untitled"


def speaker_notes_for_slide(num: int, title: str) -> list[str]:
    notes = {
        1: [
            "Open by naming the project and giving the examiner the one-line purpose: this phase is about proving that the multimodal pipeline is reproducible and explainable before any legal-domain adaptation is claimed.",
            "Do not begin with architecture jargon. Start with why the project exists: courtroom testimony is a real setting where speech, wording, and facial cues all matter, so the framework needs to be built in a way that can later support that setting.",
            "Mention that Phase 1 is the benchmark-validation stage, not the final legal system. That distinction matters because it shows you understand the difference between a controlled research implementation and the longer-term courtroom objective.",
        ],
        2: [
            "Use the agenda to give the story order. Tell the panel that you will first show the problem, then the literature, then the design and architecture, and finally the results and demo path.",
            "The agenda is also your pacing tool: if time is short, you can still keep the presentation coherent by preserving the order from baseline understanding to evaluation and future work.",
        ],
        3: [
            "Explain the problem in practical terms: legal testimony is not just words on a page, because the emotion behind the utterance can influence interpretation in research settings and later in courtroom-adjacent studies.",
            "Point out that Phase 1 uses MELD as a benchmark because it gives you a structured conversational emotion task where you can validate the pipeline before moving to courtroom-specific data.",
            "If asked why MemoCMT matters here, say that it gives a strong multimodal comparison point and a clean text-audio fusion baseline, which makes the Phase 1 implementation defensible.",
        ],
        4: [
            "Break the abstract into two parts: first, the benchmark reproduction; second, the facial-cue support path.",
            "Explain that the paper-aligned baseline is the stable reference, while the ViT branch is the visual extension that lets you study facial information without rebuilding the whole system from scratch.",
            "The scope statement should sound narrow and controlled: reproduce, inspect error patterns, then prepare for later courtroom adaptation.",
        ],
        5: [
            "When speaking about MemoCMT, avoid reading the paper title line by line. Instead say that it is the closest architectural reference because it uses cross-modal fusion rather than simple feature concatenation.",
            "Emphasize why this matters: in a multimodal emotion task, letting one modality influence another is more expressive than appending vectors and hoping the classifier learns the interaction on its own.",
        ],
        6: [
            "MELD is the benchmark because it forces the model to deal with dialogue, not isolated sentences. That is important for your project because courtroom testimony also has conversational and speaker-dependent behavior.",
            "Mention the neutral imbalance and the need for weighted metrics. This shows you understand why accuracy alone is not enough.",
        ],
        7: [
            "Table 10 is useful because it turns the dataset challenge into numbers: shorter turns, heavy neutral skew, and more conversational structure than a simple emotion dataset.",
            "Explain that the statistics justify dialogue-safe fold construction and weighted evaluation. You are not just reporting EDA; you are showing why the pipeline had to be built in a careful way.",
            "A strong viva answer here is that MELD is a realistic benchmark for conversational emotion analysis, not just a convenient dataset.",
        ],
        8: [
            "This slide is about why conversational context matters. Inter-speaker influence means the meaning of an utterance often depends on who spoke just before it.",
            "Emotion shift means the same speaker can change affect across turns, so you need more than a single isolated prediction if you want a meaningful analysis story.",
            "Contextual distance reminds the panel that nearby turns matter most, but longer-range turns can still shift the interpretation.",
            "If challenged, explain that Phase 1 currently handles context indirectly through the data split and fusion model, while explicit speaker-memory modeling is a future extension.",
        ],
        9: [
            "BERT and HuBERT are your pretrained branches. Say clearly that you are reusing strong language and speech representations so Phase 1 can focus on adaptation instead of training encoders from zero.",
            "BERT gives contextual token information from the transcript, while HuBERT contributes acoustic and delivery cues from the waveform.",
            "The technical point is that the baseline is stronger because both branches already contain learned priors from large-scale pretraining.",
        ],
        10: [
            "Use this slide to explain that ViT is not consuming the whole video as a raw motion stream. It consumes sampled RGB frames, turns them into patch tokens, and outputs a compact visual representation.",
            "You can point out that the role of ViT here is facial-cue support: it helps represent visible emotion patterns without disturbing the text-audio baseline.",
        ],
        11: [
            "Weighted cross entropy and focal loss are there because MELD is class-imbalanced. Neutral is common, so a plain loss can hide poor minority-class behavior.",
            "Explain that weighted CE gives minority classes more influence, while focal loss makes the model pay less attention to easy examples and more to difficult ones.",
            "A strong viva point is that the loss function does not solve imbalance by itself; it changes how learning pressure is distributed across classes.",
        ],
        12: [
            "Use this slide to show that you considered review feedback and tightened the presentation around the strongest baseline, fold analysis, and interpretable demo outputs.",
            "It is a good place to say that the project is trying to show honest improvement and honest failure, not just a polished headline number.",
        ],
        13: [
            "The design approach should be explained as a controlled experimental ladder. First the baseline is validated, then the visual branch is attached, then gating or auxiliary changes are tested separately.",
            "That structure matters because it lets you tell the examiner exactly which change caused which effect.",
            "If asked about alternatives, say early fusion, late fusion, and video-only are all valid but they reduce attribution clarity for this project.",
        ],
        14: [
            "This slide is where you defend the engineering assumptions. If manifest rows, cached features, and checkpoint samples do not stay aligned, the comparison is no longer trustworthy.",
            "Also explain that the dependency on pretrained backbones and ffmpeg is not incidental. Those tools are what make the pipeline reproducible and practical.",
        ],
        15: [
            "This slide is the design-details bridge. It is where you can explain the system as a set of layers: data, preprocessing, model, and analysis.",
            "If the panel asks about novelty, say that the novelty is not only in the core architecture but in how the project is being adapted toward courtroom testimony while keeping the implementation transparent.",
        ],
        16: [
            "Discuss the properties of the system as engineering requirements: interoperability, reliability, portability, and maintainability.",
            "It is worth stressing that these are not marketing words. They affect whether the same sample can flow from raw MP4 to manifest to cache to inference without mismatch.",
        ],
        17: [
            "The methodology slide should be treated like a plan-of-work slide. Walk the panel through training, comparison, evaluation, and then the decision about the next experiment.",
            "A useful line here is that the confusion matrix decides whether a change is meaningful, not just the overall accuracy.",
        ],
        18: [
            "Explain the system architecture as the full pipeline from raw inputs to analysis outputs.",
            "The important point is the transition from raw media to cached artifacts: once the preprocessing layer is correct, the same artifact can be reused for training, evaluation, and demo without recomputing everything.",
        ],
        19: [
            "Model architecture is the internal neural network, not the whole pipeline.",
            "Describe the role of each branch: text captures semantics, audio captures delivery and speech cues, and video captures visible facial information.",
            "Then say that fusion decides how much each branch should matter for the final prediction.",
        ],
        20: [
            "Keep this slide short and practical. It is there to show that the project is built on a standard Python + PyTorch stack and can run in both local and GPU environments.",
        ],
        21: [
            "Use the progress slide to tell the examiner that the core benchmark and the facial-cue branch are in place.",
            "The important message is not that the project is perfect; it is that the implementation is traceable and the current behavior is understood through analysis outputs.",
        ],
        22: [
            "When discussing references, do not over-explain. Just state that the citations cover the baseline paper, dataset papers, encoder references, and the imbalance-loss literature.",
        ],
        23: [
            "This is your defense slide for questions about wrong predictions.",
            "Tell the panel that top-3 values are softmax probabilities, not accuracy or F1.",
            "If a prediction is wrong, discuss whether it is a near miss, a confident wrong case, or a neutral-heavy failure, and then connect it back to the confusion matrix.",
        ],
        24: [
            "Use the closing slide to summarize the state of the project: Phase 1 is reproducible, the facial-cue path is explainable, and the error analysis is honest.",
            "Do not sound like the project is finished forever. Sound like the implementation has reached a stable research checkpoint and is ready for the next extension.",
        ],
        25: [
            "Future work should be framed as the next technical step, not as something already claimed in Phase 1.",
            "Explain that explicit speaker modeling, emotion shifts, long-range context, and courtroom-role metadata are useful because the current errors still reflect contextual ambiguity.",
            "Also explain timing: if Phase 1 has time left, only prototype a small ablation; the full context-aware redesign is better treated as the bridge into Phase 2.",
        ],
    }
    return notes.get(num, [f"Explain the slide '{title}' in your own words, focusing on why it matters and how it connects to the project story."])


def build_doc():
    ppt = Presentation(str(PPTX_PATH))
    doc = Document()

    # Title page
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("LegalMemoCMT Phase 1 ESA Speaker-Taking Notes")
    style_run(r, size=18, bold=True, color="122F55")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Ordered by slide, written for live explanation and viva practice")
    style_run(r, size=11, color="5E5E5E")

    intro = [
        "Use this document as a spoken guide, not as a reading script.",
        "The goal is to explain what each slide means, why it exists, and how it connects to the implementation.",
        "When speaking, keep emphasizing the pipeline story: problem -> benchmark -> architecture -> results -> error analysis -> next step.",
    ]
    for item in intro:
        add_para(doc, item)

    doc.add_paragraph("")

    for i, slide in enumerate(ppt.slides, start=1):
        title = get_slide_title(slide)
        h = doc.add_paragraph()
        h.style = "Heading 1"
        r = h.add_run(f"Slide {i}: {title}")
        style_run(r, size=14, bold=True, color="122F55")

        notes = speaker_notes_for_slide(i, title)
        for note in notes:
            add_para(doc, note, indent=12)

        # Add a small bridge line to make the talking style more natural.
        bridge = {
            1: "Start with the big picture, then narrow down to Phase 1.",
            3: "Keep the problem practical and avoid sounding too abstract.",
            7: "Use the numbers to explain the benchmark, not to recite them.",
            13: "This is where you defend your design choices clearly.",
            18: "Move from data flow to model flow without losing the pipeline story.",
            23: "This slide is your answer key for demo questions.",
            25: "Treat this as future work and do not overclaim it as done.",
        }.get(i)
        if bridge:
            add_para(doc, bridge, indent=12, bold=True, color="404040")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    build_doc()
