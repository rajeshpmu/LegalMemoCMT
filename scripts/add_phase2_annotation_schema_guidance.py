from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Corpus_Build_Code_Level_Student_Guide.docx"


def ppt_text(slide, text, x, y, w, h, size=15, bold=False, color=(35, 35, 35)):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.name = "Aptos"; r.font.size = PptPt(size); r.font.bold = bold; r.font.color.rgb = PptRGBColor(*color)
    return box


def add_header(slide, prs, heading, subtitle):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = PptRGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, PptInches(0.22))
    bar.fill.solid(); bar.fill.fore_color.rgb = PptRGBColor(20, 48, 87); bar.line.fill.background()
    ppt_text(slide, heading, 0.55, 0.38, 12.2, 0.55, 24, True, (20, 48, 87))
    ppt_text(slide, subtitle, 0.58, 0.95, 12.1, 0.36, 12, False, (95, 95, 95))


def add_bullets(slide, items, x, y, w, h, size=15):
    text = "\n".join("• " + item for item in items)
    ppt_text(slide, text, x, y, w, h, size)


def add_schema_slides():
    prs = Presentation(PPTX)
    # Slide 1 is the first existing slide. New slides are appended, then moved after slide 6.
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s1, prs, "Annotation Schema: Two Complementary Layers", "The corpus must remain MELD-compatible while adding courtroom-specific interactional information.")
    add_bullets(s1, [
        "Layer 1: basic emotion preserves the ordinary MELD-style emotion classification task.",
        "Layer 2: courtroom affect describes interactional states such as hesitation, guardedness, defensiveness, and distress.",
        "The two layers answer different questions and should not be collapsed into one label.",
        "Text, audio, and video remain aligned at the same utterance or turn row; annotation fields are added to that row.",
        "During migration, UNKNOWN or blank values are permitted for courtroom-specific fields.",
    ], 0.75, 1.45, 5.85, 5.1, 15)
    box = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(7.0), PptInches(1.65), PptInches(5.55), PptInches(1.7))
    box.fill.solid(); box.fill.fore_color.rgb = PptRGBColor(238, 244, 250); box.line.color.rgb = PptRGBColor(150, 175, 200)
    ppt_text(s1, "MELD compatibility", 7.25, 1.9, 5.0, 0.3, 16, True, (20, 48, 87))
    ppt_text(s1, "emotion_label\nemotion_label_source\nemotion_label_confidence", 7.25, 2.3, 5.0, 0.8, 14)
    box2 = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(7.0), PptInches(4.05), PptInches(5.55), PptInches(1.9))
    box2.fill.solid(); box2.fill.fore_color.rgb = PptRGBColor(250, 245, 235); box2.line.color.rgb = PptRGBColor(190, 160, 110)
    ppt_text(s1, "Courtroom-specific layer", 7.25, 4.3, 5.0, 0.3, 16, True, (20, 48, 87))
    ppt_text(s1, "courtroom_affect\naffect_intensity, valence, arousal\naffect confidence", 7.25, 4.7, 5.0, 0.9, 14)

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s2, prs, "Preserve MELD Compatibility During Migration", "Existing training code continues to read emotion_label; canonical aliases make the schema clearer for new code.")
    rows = [
        ["Existing field", "Canonical alias", "Meaning", "Compatibility rule"],
        ["emotion_label", "basic_emotion", "Basic categorical emotion", "Keep both; old scripts use emotion_label"],
        ["emotion_label_source", "basic_emotion_source", "How the basic label was produced", "Copy value, do not overwrite provenance"],
        ["emotion_label_confidence", "basic_emotion_confidence", "Confidence in the basic label", "Copy value, retain uncertainty"],
    ]
    table = s2.shapes.add_table(len(rows), 4, PptInches(0.6), PptInches(1.55), PptInches(12.1), PptInches(3.15)).table
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row): table.cell(ri, ci).text = value
    for ri, row in enumerate(table.rows):
        for cell in row.cells:
            cell.text_frame.word_wrap = True; cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Aptos"; r.font.size = PptPt(13); r.font.color.rgb = PptRGBColor(20, 48, 87) if ri == 0 else PptRGBColor(35, 35, 35)
                    if ri == 0: r.font.bold = True
    add_bullets(s2, [
        "Basic emotion vocabulary: neutral, anger, disgust, fear, joy, sadness, surprise.",
        "If the current data has labels outside this vocabulary, map them explicitly or mark them UNKNOWN; do not silently rename them.",
        "Canonical aliases are a schema migration and documentation step. They are not yet implemented in the current processing scripts.",
        "The training loader can continue using emotion_label until the migration is implemented and tested.",
    ], 0.85, 5.0, 11.7, 1.55, 14)

    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s3, prs, "Courtroom-Specific Affect and Annotation Guardrails", "Affect describes observable interactional behavior; it must never be treated as automatic evidence of truthfulness or deception.")
    add_bullets(s3, [
        "Primary courtroom_affect vocabulary: NEUTRAL_CALM, HESITANT_UNCERTAIN, GUARDED, DEFENSIVE, ASSERTIVE, TENSE, DISTRESSED, AGITATED.",
        "affect_intensity uses an ordinal scale: 0 absent/minimal, 1 mild, 2 moderate, 3 strong.",
        "valence represents the positive-to-negative direction of affect; arousal represents activation or energy. Their annotation scale must be fixed before training.",
        "courtroom_affect_confidence records annotator or model certainty and allows UNKNOWN during migration.",
        "The annotation describes speech, prosody, facial behavior, or interactional stance in context; it does not diagnose mental state.",
    ], 0.72, 1.42, 6.15, 5.3, 14)
    guard = s3.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PptInches(7.2), PptInches(1.75), PptInches(5.35), PptInches(3.85))
    guard.fill.solid(); guard.fill.fore_color.rgb = PptRGBColor(255, 242, 242); guard.line.color.rgb = PptRGBColor(180, 80, 80)
    ppt_text(s3, "Explicit prohibition", 7.5, 2.05, 4.8, 0.35, 17, True, (140, 30, 30))
    ppt_text(s3, "Do not infer or create:\n\n• deceptive\n• truthful\n• lying\n• credible\n• unreliable\n\nEmotion or facial behavior is not automatic evidence of credibility.", 7.5, 2.5, 4.7, 2.7, 15, False, (80, 30, 30))
    ppt_text(s3, "Proposed schema only: no code or manifest changes have been made in this guidance-call update.", 0.75, 6.55, 11.8, 0.35, 11, True, (20, 48, 87))

    # Move the three new slides after the existing courtroom-corpus overview, currently slide 6.
    slide_list = prs.slides._sldIdLst
    new_ids = [slide_list[-3], slide_list[-2], slide_list[-1]]
    for sid in new_ids: slide_list.remove(sid)
    for offset, sid in enumerate(new_ids): slide_list.insert(6 + offset, sid)
    prs.save(PPTX)
    print(f"Added annotation slides; wrote {len(prs.slides)} slides to {PPTX}")


def add_docx_section():
    doc = Document(DOCX)
    doc.add_page_break()
    h = doc.add_heading("Additional Guidance Call Module: MELD Compatibility and Courtroom Affect", level=1)
    doc.add_paragraph("This section describes a proposed annotation schema for the second guidance call. It does not claim that the new fields have already been added to the Phase 2 CSVs or that new labels have already been annotated.")

    doc.add_heading("1. Why Two Annotation Layers Are Needed", level=2)
    doc.add_paragraph("A basic emotion classification task and a courtroom interaction analysis task are related, but they are not identical. Basic emotion labels are useful for compatibility with MELD-style datasets and existing model code. Courtroom affect labels describe interactional states that are more specific to testimony, questioning, uncertainty, and courtroom pressure.")
    doc.add_paragraph("For example, a witness can be classified as basic_emotion=neutral while being courtroom_affect=HESITANT_UNCERTAIN. The witness may not show a conventional emotion such as sadness or anger, but pauses, hedges, corrections, or uncertain responses may still be important interactional behavior. Keeping two fields prevents the project from forcing every courtroom behavior into an ordinary emotion category.")

    doc.add_heading("2. Preserve Existing MELD-Compatible Fields", level=2)
    doc.add_paragraph("The existing emotion_label field must remain because current training and filtering scripts expect that name. The migration introduces canonical aliases without deleting the old fields:")
    table = doc.add_table(rows=1, cols=4); table.alignment = WD_TABLE_ALIGNMENT.CENTER; table.style = "Table Grid"
    for i, v in enumerate(["Existing field", "Canonical alias", "Student explanation", "Implementation rule"]): table.cell(0, i).text = v
    for row in [
        ["emotion_label", "basic_emotion", "The ordinary emotion category used by MELD-style code.", "Keep both names with the same value."],
        ["emotion_label_source", "basic_emotion_source", "How the basic emotion value was created, for example heuristic or manual.", "Copy provenance; never replace it with a generic value."],
        ["emotion_label_confidence", "basic_emotion_confidence", "How certain the annotator or heuristic is about the basic label.", "Copy uncertainty as well as confidence."],
    ]:
        cells = table.add_row().cells
        for i, v in enumerate(row): cells[i].text = v
    doc.add_paragraph("The basic emotion vocabulary is: neutral, anger, disgust, fear, joy, sadness, and surprise. If an existing value does not fit this vocabulary, it should be explicitly mapped or marked UNKNOWN. It should not be silently renamed because silent mapping destroys label provenance.")

    doc.add_heading("3. Proposed Courtroom-Specific Fields", level=2)
    table = doc.add_table(rows=1, cols=3); table.alignment = WD_TABLE_ALIGNMENT.CENTER; table.style = "Table Grid"
    for i, v in enumerate(["Field", "Allowed or expected values", "Meaning"]): table.cell(0, i).text = v
    for row in [
        ["courtroom_affect", "NEUTRAL_CALM, HESITANT_UNCERTAIN, GUARDED, DEFENSIVE, ASSERTIVE, TENSE, DISTRESSED, AGITATED, UNKNOWN", "A context-specific interactional state observed in testimony or questioning."],
        ["courtroom_affect_confidence", "Configured confidence scale or UNKNOWN", "How certain the annotator is about the courtroom-affect decision."],
        ["affect_intensity", "0, 1, 2, 3, or UNKNOWN", "0 absent/minimal; 1 mild; 2 moderate; 3 strong."],
        ["valence", "A documented ordinal or continuous scale, or UNKNOWN", "Direction from more positive to more negative affect."],
        ["arousal", "A documented ordinal or continuous scale, or UNKNOWN", "Activation or energy level, from low to high according to the chosen scale."],
    ]:
        cells = table.add_row().cells
        for i, v in enumerate(row): cells[i].text = v
    doc.add_paragraph("The valence and arousal scales must be formally specified before annotation begins. A field is not reproducible if one annotator uses -1 to +1 and another uses 1 to 5 without a documented conversion.")

    doc.add_heading("4. How These Fields Attach to the Existing Pipeline", level=2)
    for text in [
        "The fields belong in the utterance-level or turn-level metadata row, alongside utterance_text, video_clip, audio_clip, timestamps, and source traceability.",
        "The media-construction scripts create aligned text/audio/video rows; they should not automatically infer courtroom_affect from a transcript keyword or facial expression.",
        "A later annotation stage can read the validated metadata CSV, present the clip and transcript to an annotator, and write the new affect fields to a versioned annotation CSV.",
        "Basic emotion and courtroom affect should be evaluated separately because their class vocabularies and meanings differ.",
        "UNKNOWN is a valid migration value. Missing annotation is preferable to an invented label.",
    ]: doc.add_paragraph(text, style="List Bullet")

    doc.add_heading("5. Annotation Procedure from a Student Perspective", level=2)
    for i, text in enumerate([
        "Open one row using its utterance ID and verify that the transcript, audio clip, and video clip refer to the same time interval.",
        "Read the transcript for context, but do not assign affect from text alone when audio or video contradicts the text.",
        "Listen for prosody, hesitation, intensity, interruptions, pauses, and turn-taking behavior.",
        "Inspect visible behavior only when the speaking person is actually visible; otherwise use UNKNOWN for visual evidence and rely on permitted modalities according to the annotation protocol.",
        "Assign one courtroom_affect value only when the evidence supports it; otherwise assign UNKNOWN and explain the uncertainty.",
        "Assign affect_intensity using the fixed 0–3 definition and record confidence separately.",
        "Never translate the result into deceptive, truthful, lying, credible, or unreliable. Those are prohibited interpretations in this schema.",
    ], 1): doc.add_paragraph(f"{i}. {text}")

    doc.add_heading("6. Example Row Interpretation", level=2)
    doc.add_paragraph("Suppose one aligned row contains a witness answer with a long pause, repeated self-correction, and a visibly tense posture. A cautious annotation could be:")
    code = doc.add_paragraph()
    code.add_run("emotion_label=neutral\n").font.name = "Menlo"
    code.add_run("basic_emotion=neutral\n").font.name = "Menlo"
    code.add_run("courtroom_affect=HESITANT_UNCERTAIN\n").font.name = "Menlo"
    code.add_run("affect_intensity=2\n").font.name = "Menlo"
    code.add_run("valence=negative_or_low\n").font.name = "Menlo"
    code.add_run("arousal=moderate\n").font.name = "Menlo"
    code.add_run("courtroom_affect_confidence=MEDIUM\n").font.name = "Menlo"
    code.add_run("notes=Pause and self-correction observed; no credibility inference.").font.name = "Menlo"
    doc.add_paragraph("This example does not say the witness is lying or unreliable. It only records an interactional presentation under the defined annotation vocabulary. The exact valence/arousal values should be replaced by the project’s finalized scale before annotation.")

    doc.add_heading("7. What Changes Later, and What Does Not Change Now", level=2)
    doc.add_paragraph("No Phase 2 processing code or current manifest was changed for this guidance-call update. The existing emotion_label fields remain the active compatibility fields. The new schema is a design decision for the next annotation/migration stage.")
    for text in [
        "Later: add canonical aliases while retaining emotion_label for old loaders.",
        "Later: add courtroom_affect and its confidence/intensity/valence/arousal fields to a versioned annotation manifest.",
        "Later: define an annotation protocol, examples, adjudication rules, and inter-annotator agreement measurement.",
        "Later: update training loaders only after confirming that old scripts still work with the expanded schema.",
        "Now: explain the two-layer design and its safeguards to the mentor without presenting the proposed fields as completed labels.",
    ]: doc.add_paragraph(text, style="List Bullet")

    doc.add_heading("8. Guidance-Call Explanation", level=2)
    doc.add_paragraph("My dataset remains compatible with MELD because I keep emotion_label and map it to the clearer canonical name basic_emotion. I am adding a separate courtroom-affect layer because courtroom interaction is not fully represented by ordinary emotion categories. The new layer describes observable interactional states such as hesitation, guardedness, defensiveness, tension, and distress, with explicit intensity and confidence. I will not infer deception or credibility from emotion, facial behavior, or affect. At this stage this is the proposed annotation schema; implementation and annotation will be a later controlled step.")

    doc.save(DOCX)
    print(f"Updated {DOCX}")


if __name__ == "__main__":
    add_schema_slides()
    add_docx_section()
