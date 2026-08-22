from pathlib import Path
import shutil

from docx import Document
from docx.enum.section import WD_SECTION
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Active_SemiSupervised_Annotation_Guidance.pptx"
PPTX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Active_SemiSupervised_Annotation_Guidance.before_model_status_update.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Active_SemiSupervised_Annotation_Student_Guide.docx"
DOCX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Active_SemiSupervised_Annotation_Student_Guide.before_model_status_update.docx"
TITLE = "11. Model Status: Implemented, Optional, and Deferred"


def add_box(slide, x, y, w, h, text, color, size=14, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.color.rgb = RGBColor(70, 80, 95)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 1
    for run in p.runs:
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(25, 35, 50)


prs = Presentation(PPTX)
if not PPTX_BACKUP.exists():
    shutil.copy2(PPTX, PPTX_BACKUP)
if not any(any(TITLE in sh.text for sh in slide.shapes if hasattr(sh, "text")) for slide in prs.slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.45))
    title.text_frame.text = TITLE
    title.text_frame.paragraphs[0].runs[0].font.size = PptPt(25)
    title.text_frame.paragraphs[0].runs[0].font.bold = True
    add_box(slide, 0.45, 1.0, 3.85, 1.2, "Operational now\nPyannote diarization\nPhase 1 MELD checkpoint", (208, 224, 220), 17, True)
    add_box(slide, 4.75, 1.0, 3.85, 1.2, "Optional adapters\nDeBERTa text suggestions\nSER-Odyssey audio evidence", (220, 232, 245), 17, True)
    add_box(slide, 9.05, 1.0, 3.85, 1.2, "Not implemented yet\nValidated video-affect model\nAutomatic pseudo-label acceptance", (242, 225, 225), 17, True)
    add_box(slide, 0.75, 3.0, 5.7, 1.25, "Why the optional models help\nThey provide independent text/audio evidence for candidate labels, uncertainty, and disagreement.", (248, 237, 215), 17, True)
    add_box(slide, 6.9, 3.0, 5.7, 1.25, "What remains human\nRole mapping, witness verification, affect labels, seed review, ambiguous cases, and benchmark ground truth.", (232, 224, 245), 17, True)
    note = slide.shapes.add_textbox(Inches(0.75), Inches(4.8), Inches(11.8), Inches(1.35))
    note.text_frame.text = "Current status: models are documented and optionally supported, but DeBERTa and audio-affect inference have not yet been run and validated on Clancy."
    note.text_frame.paragraphs[0].runs[0].font.size = PptPt(19)
    note.text_frame.paragraphs[0].runs[0].font.bold = True
    prs.save(PPTX)


if not DOCX_BACKUP.exists():
    shutil.copy2(DOCX, DOCX_BACKUP)
document = Document(DOCX)
if not any(TITLE in paragraph.text for paragraph in document.paragraphs):
    document.add_section(WD_SECTION.NEW_PAGE)
    document.add_heading(TITLE, level=1)
    document.add_paragraph(
        "The model stack is deliberately separated into operational components, optional suggestion adapters, and "
        "deferred components. This distinction is important for an honest Phase 2 report: a model being named in a "
        "configuration or supported by a script does not mean that it has already generated valid corpus labels."
    )
    document.add_heading("Operational Components", level=2)
    document.add_paragraph(
        "Pyannote speaker diarization is the active model stage. It produces speech intervals and anonymous speaker "
        "clusters. It supports who-spoke-when evidence, cluster-to-role review, and witness-turn extraction, but it "
        "does not predict emotion or courtroom affect. The current execution is a controlled one-source pilot."
    )
    document.add_paragraph(
        "The Phase 1 MELD checkpoint already exists and remains preserved. It can provide a basic-emotion baseline, "
        "provisional pseudo-labels, and a warm-start initialization for later legal-domain adaptation. Its output is "
        "not courtroom-affect ground truth and its training domain is not identical to courtroom testimony."
    )
    document.add_heading("Optional Model-Suggestion Adapters", level=2)
    document.add_paragraph(
        "The new `phase2/annotation/generate_model_suggestions.py` script can optionally call a text zero-shot model "
        "through `--text-model` and an audio classifier through `--audio-model`. These adapters write suggestion and "
        "modality-specific evidence fields. They do not overwrite `basic_emotion`, `courtroom_affect`, or existing "
        "human labels."
    )
    document.add_paragraph(
        "The planned text model is a DeBERTa/NLI zero-shot classifier. It can compare transcript context with phrases "
        "such as `the speaker is guarded and cautious` or `the speaker is defensive or resisting`. Its result is a "
        "candidate based on text semantics, not a courtroom annotation."
    )
    document.add_paragraph(
        "The planned audio model is SER-Odyssey or an equivalent dimensional speech-emotion model. It can provide "
        "continuous valence, arousal, and dominance evidence. Those values are useful for disagreement analysis and "
        "candidate generation, but they are trained on general speech emotion data rather than the LegalMemoCMT "
        "courtroom-affect vocabulary."
    )
    document.add_heading("Components Not Yet Implemented", level=2)
    document.add_paragraph(
        "No validated video-affect model has yet been integrated into this active-annotation stage. Existing Phase 1 "
        "video representations remain available, but they are not automatically treated as courtroom-affect outputs. "
        "Automatic pseudo-label acceptance, iterative active-learning retraining, and contradiction-pair construction "
        "are also deferred until diarization and the human seed set are reliable."
    )
    document.add_heading("Why the Models Help", level=2)
    for item in [
        "They provide independent evidence from text and audio rather than one model making every decision.",
        "They create candidate labels faster than manually annotating the entire corpus.",
        "Their confidence, entropy, margins, and disagreement can prioritize the most informative samples for review.",
        "They support comparison between general MELD emotion and courtroom-specific affect without collapsing the two tasks.",
    ]:
        document.add_paragraph(item, style="List Bullet")
    document.add_heading("What Still Requires Human-in-the-Loop Review", level=2)
    for item in [
        "Mapping anonymous Pyannote clusters to Witness, Counsel, Judge, or Other.",
        "Confirming whether the visible witness is speaking or only being addressed.",
        "Assigning courtroom-affect labels and intensity.",
        "Reviewing text/audio/video disagreement and low-confidence cases.",
        "Creating the 500–1000 utterance seed set.",
        "Protecting the witness-disjoint benchmark from pseudo-label training and active selection.",
    ]:
        document.add_paragraph(item, style="List Bullet")
    document.add_heading("Current Honest Status", level=2)
    document.add_paragraph(
        "The new models are documented and optional adapters are available, but DeBERTa and audio-affect inference "
        "have not yet been run and validated on Clancy. Therefore no claim should be made that these models have already "
        "created trustworthy courtroom-affect labels. The next valid order is: finish diarization, map roles, validate "
        "witness clips, create the human seed, then run model suggestions and inspect disagreement."
    )
    document.add_heading("Suggested Runtime Example", level=2)
    p = document.add_paragraph(style="No Spacing")
    r = p.add_run('''PYTHON_BIN=$PWD/.venv/bin/python \\
python phase2/annotation/generate_model_suggestions.py \\
  --input-csv <schema-extended-manifest.csv> \\
  --output-csv <suggestions-manifest.csv> \\
  --text-model MoritzLaurer/deberta-v3-large-zeroshot-v2.0 \\
  --audio-model 3loi/SER-Odyssey-Baseline-WavLM-Multi-Attributes \\
  --max-rows 200''')
    r.font.name = "Courier New"
    r.font.size = Pt(8)
    document.add_paragraph(
        "This command is an optional future pilot, not a command to run before diarization and role-aware validation. "
        "The generated values remain AUTO_SUGGESTED and must be reviewed before any pseudo-label acceptance."
    )
    document.save(DOCX)

print(f"Updated {PPTX}")
print(f"Updated {DOCX}")
