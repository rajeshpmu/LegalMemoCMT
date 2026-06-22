from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


DECKS = [
    Path("implementation_docments/Zeroth_Review_LegalMemoCMT.pptx"),
    Path("implementation_docments/First_Review_LegalMemoCMT.pptx"),
]
DIAGRAM = Path("artifacts/mermaid/legacy_multimodal_pipeline.png")


def add_bg(slide, prs: Presentation):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None, y: float = 0.35):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(y + 0.52), Inches(12.1), Inches(0.55))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x=0.75, y=1.45, w=6.0, h=5.45, font_size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_diagram(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def add_legacy_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(
        slide,
        "Legacy Multimodal Pipeline",
        "The runnable baseline uses feature-based encoders and transformer fusion; the pretrained/paper path is separate.",
    )
    add_body(
        slide,
        [
            "Legacy mode is the current lightweight path in the codebase and is selected with encoder_mode='legacy'.",
            "Text is represented with simple token IDs and an embedding/Transformer stack instead of a HuggingFace backbone.",
            "Audio is read as fixed-size feature tensors or cached arrays and then encoded as a sequence.",
            "Video stays as compact frame features and can be switched off for ablation-style runs.",
            "The three modality vectors are fused by CrossModalFusion and then passed to the classifier head.",
            "This branch is useful for local runs, sanity checks, and ablations before moving to the pretrained/paper path.",
        ],
        x=0.75,
        y=1.52,
        w=6.0,
        h=5.35,
        font_size=15,
    )
    add_diagram(slide, DIAGRAM, 6.95, 1.72, 5.95)
    tx = slide.shapes.add_textbox(Inches(6.95), Inches(6.55), Inches(5.95), Inches(0.35))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Mermaid diagram of the legacy branch and the separate pretrained/paper path"
    r.font.name = "Aptos"
    r.font.size = Pt(10.5)
    r.font.italic = True
    r.font.color.rgb = RGBColor(95, 95, 95)


def main():
    if not DIAGRAM.exists():
        raise FileNotFoundError(f"Missing diagram: {DIAGRAM}")
    for deck in DECKS:
        prs = Presentation(str(deck))
        add_legacy_slide(prs)
        prs.save(str(deck))
        print(f"Updated {deck}")


if __name__ == "__main__":
    main()
