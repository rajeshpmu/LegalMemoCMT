"""Add separate basic-emotion and courtroom-affect task slides."""
from pathlib import Path
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document

ROOT=Path(__file__).resolve().parents[1]
PPTX=ROOT/'implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Presentation_v1.pptx'
DOCX=ROOT/'implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Speaking_Doc.docx'
FIG=ROOT/'implementation_docments/figures/phase2_separate_basic_affect_prediction.png'
MARK='SEPARATE_BASIC_AFFECT_TASKS_V1'

def text(s,v,x,y,w,h,size=15,color=(35,35,35),bold=False):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); b.text_frame.word_wrap=True
    p=b.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=v
    r.font.name='Aptos'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=RGBColor(*color)

def base(prs,title,sub):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor(255,255,255)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.22)); bar.fill.solid(); bar.fill.fore_color.rgb=RGBColor(20,48,87); bar.line.fill.background()
    text(s,title,.55,.38,12.2,.45,23,(20,48,87),True); text(s,sub,.58,.9,12,.3,12,(95,95,95)); return s

def table(s,headers,rows,x=.55,y=1.35,w=12.2,h=5.5,size=11):
    t=s.shapes.add_table(len(rows)+1,len(headers),Inches(x),Inches(y),Inches(w),Inches(h)).table
    for j,v in enumerate(headers): t.cell(0,j).text=str(v)
    for i,row in enumerate(rows,1):
        for j,v in enumerate(row): t.cell(i,j).text=str(v)
    for i,row in enumerate(t.rows):
        for c in row.cells:
            c.text_frame.word_wrap=True
            for p in c.text_frame.paragraphs:
                for r in p.runs: r.font.name='Aptos'; r.font.size=Pt(size); r.font.bold=i==0; r.font.color.rgb=RGBColor(20,48,87) if i==0 else RGBColor(35,35,35)

def insert_before_final(prs,slide):
    ids=prs.slides._sldIdLst; item=ids[-1]; ids.remove(item); ids.insert(len(ids)-1,item)

def update_ppt():
    prs=Presentation(PPTX)
    if any(MARK in getattr(sh,'text','') for s in prs.slides for sh in s.shapes): return
    backup=PPTX.with_name(PPTX.name+'.before_separate_prediction_tasks.pptx')
    if not backup.exists(): shutil.copy2(PPTX,backup)
    s=base(prs,'Separate Prediction Tasks: Basic Emotion and Courtroom Affect','One shared multimodal representation, two independent targets')
    s.shapes.add_picture(str(FIG),Inches(.45),Inches(1.25),width=Inches(12.45),height=Inches(5.85))
    text(s,'The two outputs can coexist: basic_emotion=neutral while courtroom_affect=CALM_COMPOSED.',.7,7.08,12,.25,11,(20,48,87),True)
    insert_before_final(prs,s)
    s=base(prs,'Training and Evaluation of the Two Heads','Do not turn machine affect candidates into gold labels')
    table(s,['Task','Initial target','Training plan','Evaluation'],[
        ['Basic emotion','final_basic_emotion on approved SILVER rows','Warm-start fold_4 trimodal checkpoint; preserve Phase 1 baseline; group split','Human-validated seven-label test; macro-F1, balanced accuracy, calibration, per-class results'],
        ['Courtroom affect','Human-validated affect seed set','Separate head after reliable labels; optional masked/low-weight weak experiment only','Independent affect test; reviewer agreement, per-class F1, abstention, source/witness breakdown'],
        ['Shared safeguards','Provenance and modality fields','Do not merge labels into combinations such as SADNESS_CALM_COMPOSED','Report text/audio/video ablations and leakage violations'],
    ],size=10)
    text(s,'Expected outcome: adaptation should improve held-out courtroom performance or reveal where MELD transfer fails.',.75,6.55,11.9,.35,12,(20,48,87),True)
    insert_before_final(prs,s); prs.save(PPTX); print('Updated',PPTX,'slides',len(prs.slides))

def update_doc():
    doc=Document(DOCX)
    if any(MARK in p.text for p in doc.paragraphs): return
    backup=DOCX.with_name(DOCX.name+'.before_separate_prediction_tasks.docx')
    if not backup.exists(): shutil.copy2(DOCX,backup)
    doc.add_heading('Separate Prediction Tasks: Basic Emotion and Courtroom Affect',1)
    p=doc.add_paragraph(); p.add_run(MARK).font.size=Pt(8)
    doc.add_paragraph('The proposed architecture uses one shared multimodal representation from aligned text, audio, and video, followed by two independent prediction heads. The first head predicts basic emotion using the seven MELD-compatible labels. The second head predicts courtroom affect, which describes behavioral presentation in the legal interaction. These are deliberately separate dimensions.')
    doc.add_paragraph('For example, a witness can calmly describe another person’s sadness. The correct pair can therefore be basic_emotion=neutral and courtroom_affect=CALM_COMPOSED. Another witness can have basic_emotion=sadness and courtroom_affect=CALM_COMPOSED. Combining the fields into one label such as SADNESS_CALM_COMPOSED would create a sparse and confusing label space, so I will not do that.')
    doc.add_heading('Training targets and warm start',2)
    doc.add_paragraph('The initial supervised experiment uses final_basic_emotion only for approved AUTO_ADJUDICATED/SILVER rows. The seven labels are neutral, anger, disgust, fear, joy, sadness, and surprise. This reuses the existing Phase 1 classification head and keeps compatibility with emotion_label. The executable warm-start checkpoint is results/facial_cues/meld_vit_facecrop_gated_video_aux/fold_4/best_model.pt. I will verify its model configuration, including trimodal use_video=true and video_dim=768, before training.')
    doc.add_paragraph('Courtroom affect should not initially be trained from the current machine candidates as if they were gold labels. Its labels require a human-validated seed set. Until that exists, fields such as CALM_COMPOSED, HESITANT_UNCERTAIN, or DISTRESSED remain machine suggestions and review evidence. A later experiment could use weak affect candidates with a mask or very small loss weight, but that must be compared against a human-labeled result.')
    doc.add_heading('Multi-task loss after human affect labels exist',2)
    doc.add_paragraph('Once both tasks have trustworthy labels, the model can be trained with two classification losses: L_total = L_basic + lambda_affect * L_affect. L_basic is computed only where a basic-emotion target is valid, and L_affect is computed only where a human courtroom-affect target is valid. The coefficient lambda_affect controls how strongly the affect task influences the shared representation. It must be selected on validation data rather than assumed.')
    doc.add_heading('Evaluation plan',2)
    doc.add_paragraph('I will evaluate the two heads independently. For basic emotion, the test set must contain human-validated labels and must be disjoint by source video, hearing, or witness from training. I will report macro-F1, balanced accuracy, per-class precision/recall/F1, confusion matrices, confidence calibration, and abstention or unresolved rates. Accuracy alone can hide collapse into the neutral majority class.')
    doc.add_paragraph('For courtroom affect, I will use a separate human-validated test set with the courtroom-affect vocabulary and intensity values. I will report reviewer agreement, per-class performance, source and witness breakdowns, and the percentage of cases where the model abstains. Machine candidates must not be used as their own test labels.')
    doc.add_paragraph('I will also run text-only, audio-only, video-only, and trimodal ablations. This shows whether the model benefits from all modalities and whether video helps only when the witness is actually visible. I will compare the unchanged Phase 1 checkpoint with the warm-start adapted model and record any split leakage violations.')
    doc.add_heading('How to explain the expected outcome',2)
    doc.add_paragraph('The desired outcome is not simply a higher training score. A useful result would show that the adapted model is better calibrated and more accurate on held-out courtroom data while preserving the distinction between described emotion and the witness’s own presentation. It is also a valid research outcome if the experiment shows that the SILVER set is too small, that the courtroom-affect task needs more human labels, or that a modality does not generalize.')
    doc.add_heading('Speaking lines for the guidance call',2)
    doc.add_paragraph('I will say that Phase 2 has two related but different prediction tasks. I will initially use SILVER rows to adapt the MELD-compatible basic-emotion head, starting from the available trimodal face-crop checkpoint. I will not call the current courtroom-affect suggestions ground truth. After creating a human affect seed set, I can add a second head and evaluate both tasks separately. This design lets the model say both what emotion category is suggested and how the witness presents in court, without forcing those meanings into one label.')
    doc.save(DOCX); print('Updated',DOCX)

if __name__=='__main__': update_ppt(); update_doc()
