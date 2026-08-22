from pathlib import Path
import shutil

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.before_courtroom_exchange.pptx"
MARKER = "Courtroom exchange example"


prs = Presentation(PPTX)
target = None
for slide in prs.slides:
    texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
    if any("Whisper + Pyannote: Building a Witness-Aware Utterance" in text for text in texts):
        target = slide
        break
if target is None:
    raise SystemExit("Architecture slide was not found")
if any(MARKER in shape.text for shape in target.shapes if hasattr(shape, "text")):
    print("Courtroom exchange already present")
else:
    if not BACKUP.exists():
        shutil.copy2(PPTX, BACKUP)
    box = target.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.05), Inches(3.0), Inches(3.45), Inches(2.05),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(242, 245, 250)
    box.line.color.rgb = RGBColor(70, 80, 95)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    lines = [
        (MARKER, True, 13),
        ('Defense: "Did you examine her?"', False, 10),
        ('Witness: "Yes, I did."', False, 10),
        ('Defense: "What did you observe?"', False, 10),
        ('Witness: "She appeared calm."', False, 10),
    ]
    for index, (text, bold, size) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = text
        paragraph.space_after = Pt(2)
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(25, 35, 50)
    prs.save(PPTX)
    print(f"Added courtroom exchange to {PPTX}")
