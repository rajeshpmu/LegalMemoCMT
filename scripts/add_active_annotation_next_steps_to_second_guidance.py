from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.before_active_annotation_next_steps.pptx"
TITLE = "Planned Next Actions: Active Annotation Preparation"


def add_box(slide, x, y, w, h, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.color.rgb = RGBColor(70, 80, 95)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 1
    for run in p.runs:
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = RGBColor(25, 35, 50)


def arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(70, 80, 95)
    line.line.end_arrowhead = True


prs = Presentation(PPTX)
if any(any(TITLE in shape.text for shape in slide.shapes if hasattr(shape, "text")) for slide in prs.slides):
    print("Next-actions slide already exists")
    raise SystemExit(0)
if not BACKUP.exists():
    shutil.copy2(PPTX, BACKUP)

slide = prs.slides.add_slide(prs.slide_layouts[6])
title = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.5))
title.text_frame.text = TITLE
title.text_frame.paragraphs[0].runs[0].font.size = Pt(26)
title.text_frame.paragraphs[0].runs[0].font.bold = True

steps = [
    ("1\nComplete one-source\nPyannote pilot", (248, 237, 215)),
    ("2\nSave segments and\ndiarized manifest", (220, 232, 245)),
    ("3\nMap clusters to\nWitness / Counsel / Judge", (231, 243, 235)),
    ("4\nAlign clusters to turns\nand validate witness clips", (232, 224, 245)),
    ("5\nExtend schema and\nselect diverse seed", (220, 232, 245)),
    ("6\nExport queue for\nhuman review", (208, 224, 220)),
]
for i, (text, color) in enumerate(steps):
    x = 0.35 + (i % 3) * 4.25
    y = 1.15 + (i // 3) * 1.65
    add_box(slide, x, y, 3.45, 0.95, text, color)
    if i % 3 < 2:
        arrow(slide, x + 3.45, y + 0.47, x + 4.05, y + 0.47)

note = slide.shapes.add_textbox(Inches(0.75), Inches(4.75), Inches(11.8), Inches(1.65))
tf = note.text_frame
tf.clear()
tf.word_wrap = True
for i, text in enumerate([
    "Decision gate: do not begin pseudo-labeling until diarization, role mapping, witness extraction, and tri-modal validation pass.",
    "Then create a 500–1000 utterance human-reviewed seed; use model suggestions only as non-canonical evidence.",
    "Keep the benchmark witness-disjoint and isolated from active-learning selection.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(17)

prs.save(PPTX)

ids = prs.slides._sldIdLst
new_id = ids[-1]
ids.remove(new_id)
anchor = "Pool A Readiness Gate Before Adaptation"
insert_at = next(
    (
        i + 1
        for i, sid in enumerate(ids)
        if any(anchor in shape.text for shape in prs.part.related_part(sid.rId).slide.shapes if hasattr(shape, "text"))
    ),
    len(ids),
)
ids.insert(insert_at, new_id)
prs.save(PPTX)
print(f"Added {TITLE} to {PPTX}")
