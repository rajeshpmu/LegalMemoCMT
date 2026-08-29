"""Insert the Clancy corpus-to-200-pilot diagram into the review artifacts."""
from pathlib import Path
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Presentation_v1.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Speaking_Doc.docx"
FIGURE = ROOT / "implementation_docments/figures/phase2_clancy_corpus_to_200_pilot.png"
MARKER = "CLANCY_CORPUS_TO_200_PILOT_DIAGRAM_V1"

def add_text(slide, value, x, y, w, h, size=16, color=(35,35,35), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = value
    run.font.name = "Aptos"; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

def insert_slide():
    prs = Presentation(PPTX)
    if any(MARKER in getattr(sh, "text", "") for s in prs.slides for sh in s.shapes):
        return
    backup = PPTX.with_name(PPTX.name + ".before_clancy_pipeline_diagram.pptx")
    if not backup.exists(): shutil.copy2(PPTX, backup)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255,255,255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.22))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(20,48,87); bar.line.fill.background()
    add_text(slide, "Clancy Corpus Pipeline: From Sources to the 200-Row Pilot", .55, .38, 12.2, .45, 23, (20,48,87), True)
    add_text(slide, "The pilot is selected only after witness-speaking and quality filtering", .58, .9, 12, .3, 12, (95,95,95))
    slide.shapes.add_picture(str(FIGURE), Inches(.65), Inches(1.25), width=Inches(12.05), height=Inches(5.7))
    add_text(slide, "The final corpus is rerun from the complete eligible pool after the pilot rules are accepted.", .8, 7.05, 11.7, .25, 12, (20,48,87), True)
    ids = prs.slides._sldIdLst; item = ids[-1]; ids.remove(item); ids.insert(13, item)
    prs.save(PPTX)
    print("Inserted diagram slide; slides=", len(prs.slides))

def update_doc():
    doc = Document(DOCX)
    if any(MARKER in p.text for p in doc.paragraphs): return
    backup = DOCX.with_name(DOCX.name + ".before_clancy_pipeline_diagram.docx")
    if not backup.exists(): shutil.copy2(DOCX, backup)
    doc.add_heading("Clancy Corpus Pipeline Diagram: Why the First 200 Rows Are Selected", 1)
    p = doc.add_paragraph(); p.add_run(MARKER).font.size = Pt(8)
    doc.add_paragraph(
        "The diagram should be explained from top to bottom. The process begins with official source videos and subtitle files. "
        "These are converted into a source-shared manifest containing 77,442 rows. A source-shared row is not automatically a unique "
        "training example because many rows point to the same raw video and subtitle files."
    )
    doc.add_paragraph(
        "The next stage consolidates subtitle evidence into 11,926 turn rows and applies duration and persistent rejection rules. "
        "Pyannote then produces 23,709 anonymous acoustic speech segments across the source videos. Segment-to-turn mapping tells us "
        "which diarization intervals overlap each turn, but it does not identify the legal role. The cluster IDs are local voice IDs, "
        "so manual review maps them to Witness, Prosecutor, Defence, Judge, or Other."
    )
    doc.add_paragraph(
        "After role mapping, the witness-speaking filter retains rows where the mapped role is Witness and the witness is speaking. "
        "Visual and quality checks then remove known breaks, news mixtures, long unsuitable material, invalid media, and rows that do not "
        "satisfy the configured duration or evidence requirements. This produces the current witness-visible speaking pool of 2,229 "
        "usable rows plus 350 rows held for review."
    )
    doc.add_paragraph(
        "The 200-row pilot is selected from that filtered witness pool. It is deliberately small so I can inspect the Phase 1 model, "
        "audio-SER evidence, semantic scope, courtroom-affect heuristics, critical-conflict checks, and the SILVER/WEAK acceptance gate. "
        "It is not a random estimate of the entire corpus and it is not the final training set."
    )
    doc.add_paragraph(
        "Once the pilot rules are stable, the same pipeline is run over the complete eligible witness pool. The final Clancy corpus is "
        "the rows that pass the final role, alignment, audio, video, duration, visual, and annotation-quality gates. Its final number of "
        "utterances, minutes, witnesses, and source videos must be recomputed from the final manifest; it must not be extrapolated from "
        "the 200-row pilot."
    )
    doc.add_heading("How to explain the diagram in the guidance call", 2)
    doc.add_paragraph(
        "I can say: the first branch proves that the raw sources were extracted; the middle branch converts source text and audio into "
        "time-linked turns and anonymous speaker clusters; the witness branch applies legal-role and visual-selection decisions; and the "
        "200-row branch is a controlled experiment before scaling. This ordering prevents me from treating every subtitle row, diarization "
        "segment, or model prediction as final corpus evidence."
    )
    doc.save(DOCX)
    print("Updated speaking document")

if __name__ == "__main__":
    insert_slide(); update_doc()
