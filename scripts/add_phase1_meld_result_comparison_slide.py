from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
READING_DOCX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Reading_Script.docx"

PAPER_ACC = 64.18
PAPER_WF1 = 62.52

ROWS = [
    {
        "run": "MemoCMT paper CMT+MIN",
        "acc": "64.18%",
        "wf1": "62.52%",
        "macro": "—",
        "uw": "—",
        "takeaway": "Reference point from the published paper.",
    },
    {
        "run": "Paper-aligned MELD CV mean",
        "acc": "62.47%",
        "wf1": "61.95%",
        "macro": "43.95%",
        "uw": "44.17%",
        "takeaway": "Closest overall reproduction in this repository.",
    },
    {
        "run": "Best baseline fold 2",
        "acc": "63.75%",
        "wf1": "62.54%",
        "macro": "44.30%",
        "uw": "43.69%",
        "takeaway": "Very close to the paper on the core MELD numbers.",
    },
    {
        "run": "Gated + aux fold 2",
        "acc": "60.54%",
        "wf1": "60.22%",
        "macro": "43.51%",
        "uw": "44.92%",
        "takeaway": "Selective improvement on some clips, not overall better.",
    },
    {
        "run": "Gated + aux fold 4",
        "acc": "59.92%",
        "wf1": "60.56%",
        "macro": "43.30%",
        "uw": "46.38%",
        "takeaway": "Higher UW-Acc, but lower overall than the baseline.",
    },
]


def style_ppt_run(run, *, size=10, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color)


def style_doc_run(run, *, size=10, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_box(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = PptRGBColor.from_string(fill)
    if line:
        shp.line.color.rgb = PptRGBColor.from_string(line)
    return shp


def add_textbox(slide, left, top, width, height, text, *, size=11, bold=False, color="000000", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    style_ppt_run(r, size=size, bold=bold, color=color)
    return tb


def slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            txt = shape.text_frame.text.strip()
            if txt:
                return txt.splitlines()[0].strip()
    return "Untitled"


def update_pptx():
    prs = Presentation(str(PPTX_PATH))
    existing = any(
        "MELD Result Comparison" in (shape.text_frame.text if getattr(shape, "has_text_frame", False) else "")
        for slide in prs.slides
        for shape in slide.shapes
    )
    if existing:
        print("PPTX already contains the MELD comparison slide; skipping.")
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_no = len(prs.slides)

    add_box(slide, 0.0, 0.0, 13.333, 0.58, "122F55")
    add_box(slide, 0.0, 0.58, 13.333, 0.08, "267377")
    add_textbox(
        slide,
        0.45,
        0.12,
        11.4,
        0.3,
        f"{slide_no}. MELD Result Comparison: MemoCMT vs Baseline vs Gated+Aux",
        size=18,
        bold=True,
        color="FFFFFF",
    )
    add_textbox(
        slide,
        11.95,
        0.12,
        0.95,
        0.3,
        str(slide_no),
        size=12,
        color="FFFFFF",
        align=PP_ALIGN.RIGHT,
    )
    add_textbox(
        slide,
        0.5,
        0.68,
        12.2,
        0.35,
        "Directional comparison only: the paper reports MELD in a different evaluation framing, so compare the trend, not just the raw score.",
        size=10.5,
        color="5E5E5E",
    )

    headers = ["Run", "Acc / W-Acc", "W-F1", "Macro F1", "UW-Acc", "Reading"]
    table = slide.shapes.add_table(
        len(ROWS) + 1,
        len(headers),
        PptInches(0.2),
        PptInches(1.1),
        PptInches(12.9),
        PptInches(5.65),
    ).table

    widths = [2.4, 1.25, 1.0, 1.0, 1.0, 4.25]
    for idx, w in enumerate(widths):
        table.columns[idx].width = PptInches(w)

    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptRGBColor.from_string("DDEBF7")
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                style_ppt_run(r, size=8.5, bold=True, color="122F55")

    for i, row in enumerate(ROWS, start=1):
        values = [row["run"], row["acc"], row["wf1"], row["macro"], row["uw"], row["takeaway"]]
        for j, value in enumerate(values):
            cell = table.cell(i, j)
            cell.text = value
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = PptRGBColor.from_string("F8FBFF")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j in [0, 5] else PP_ALIGN.CENTER
                for r in p.runs:
                    style_ppt_run(r, size=7.6 if j != 5 else 7.2, color="202020")

    add_box(slide, 0.32, 6.82, 12.65, 0.48, "EEF4FA", "C9D7E6")
    add_textbox(
        slide,
        0.45,
        6.91,
        12.35,
        0.25,
        f"Paper gap on the mean baseline is about {PAPER_ACC - 62.47:.2f} points in accuracy and {PAPER_WF1 - 61.95:.2f} points in weighted F1; gated+aux remains selective, not universal.",
        size=9.4,
        bold=True,
        color="122F55",
    )

    prs.save(str(PPTX_PATH))
    print(f"Updated PPTX: {PPTX_PATH}")


def update_reading_doc():
    doc = Document(str(READING_DOCX_PATH))
    if any("Slide 27: MELD Result Comparison" in p.text for p in doc.paragraphs):
        print("Reading script already contains the MELD comparison section; skipping.")
        return

    doc.add_heading("Slide 27: MELD Result Comparison: MemoCMT vs Baseline vs Gated+Aux", level=1)
    paragraphs = [
        "Use this slide to close the MELD results story by comparing the published MemoCMT reference with the implemented paper-aligned baseline and the gated+aux branch. The key is to stay directional and technically honest: the paper reports a different evaluation framing, so the comparison is mainly about whether the current implementation is in the right range and whether the newer branch improves the right error patterns.",
        "Explain the columns before quoting numbers. Accuracy and weighted accuracy are the same sample-level correctness quantity in this codebase, weighted F1 reflects the class distribution, macro F1 shows whether minority classes are being handled better, and unweighted accuracy is the mean per-class accuracy. This is why the table includes both W-F1 and UW-Acc even though the published paper does not report every column in the same form.",
        "Then give the numbers plainly. MemoCMT paper CMT+MIN reports 64.18% accuracy and 62.52% weighted F1. The paper-aligned MELD CV mean in this repository is 62.47% accuracy, 61.95% weighted F1, 43.95% macro F1, and 44.17% unweighted accuracy. The best baseline fold 2 is 63.75% accuracy and 62.54% weighted F1, which is very close to the paper. The gated+aux branch gives 60.54% accuracy and 60.22% weighted F1 on fold 2, and 59.92% accuracy and 60.56% weighted F1 on fold 4.",
        "The main conclusion is that the paper-aligned baseline remains the strongest stable MELD result, while gated+aux is a selective-improvement branch. It fixes some boundary cases, but it is not a universal win on the full MELD benchmark. That makes it useful as a research direction, but not yet as the main claim for Phase 1.",
        "If asked why gated+aux is not always better, explain that the visual branch and gating add expressive capacity, but MELD is still class-imbalanced and heavily neutral-oriented. Extra capacity can help on ambiguous clips where neutral is near a boundary, but it can also sharpen the wrong class on hard or label-ambiguous clips. That is exactly why fold-level confusion matrices and per-clip demos remain necessary.",
        "A concise viva sentence is: the paper-aligned baseline is still the closest overall reproduction of MemoCMT on MELD, and gated+aux is a useful but selective improvement branch that should be presented honestly as an experiment, not as a universal replacement.",
    ]
    for text in paragraphs:
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        r.font.color.rgb = RGBColor.from_string("000000")
        p.paragraph_format.space_after = Pt(4)
        p.paragraph_format.space_before = Pt(0)
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    doc.save(str(READING_DOCX_PATH))
    print(f"Updated reading script: {READING_DOCX_PATH}")


def main():
    update_pptx()
    update_reading_doc()


if __name__ == "__main__":
    main()
