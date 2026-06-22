from __future__ import annotations

import re
import subprocess
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.docx"
MMD_PATH = ROOT / "artifacts" / "mermaid" / "paper_style_pooling_comparison.mmd"
SVG_PATH = ROOT / "artifacts" / "mermaid" / "paper_style_pooling_comparison.svg"
PNG_PATH = ROOT / "artifacts" / "mermaid" / "paper_style_pooling_comparison.png"


MMD_CONTENT = """%%{init: {'themeVariables': {'fontSize': '20px'}}}%%
flowchart TB
    subgraph P["Pretrained / paper-style core"]
        T["Text input ids"] --> TB["BERT backbone"]
        A["Audio waveform"] --> AB["HuBERT backbone"]
        TB --> TS["Text token sequence"]
        AB --> AS["Audio token sequence"]
        TS --> X["Bidirectional cross-attention\n(text ↔ audio)"]
        AS --> X
        X --> F["Fused text+audio sequence"]
        F --> POOL["Masked pooling\nCLS / MEAN / MAX / MIN"]
        POOL --> CLS["Classifier"]
    end

    subgraph H["Video-enabled hybrid branch"]
        V["Video features"] --> VE["Legacy video encoder"]
        VE --> VV["Video vector"]
        TT["Text vector"] -.-> F3["3-token fusion\n(text, audio, video)"]
        AA["Audio vector"] -.-> F3
        VV -.-> F3
        F3 --> P2["Pooling\nCLS / MEAN / MAX / MIN"]
        P2 --> CLS2["Classifier"]
    end

    N1["Pooling acts on token positions in the fused sequence,\nnot on one scalar per modality."]:::note
    POOL -.-> N1
    P2 -.-> N1

    classDef note fill:#fff6d5,stroke:#c9a227,color:#5a4300;
"""


def write_mermaid_assets() -> None:
    MMD_PATH.write_text(MMD_CONTENT, encoding="utf-8")
    cmd = [
        "npx",
        "-y",
        "@mermaid-js/mermaid-cli",
        "-i",
        str(MMD_PATH),
        "-o",
        str(SVG_PATH),
        "-b",
        "transparent",
    ]
    subprocess.run(cmd, check=True, cwd=str(ROOT), stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    cmd = [
        "npx",
        "-y",
        "@mermaid-js/mermaid-cli",
        "-i",
        str(MMD_PATH),
        "-o",
        str(PNG_PATH),
        "-b",
        "transparent",
    ]
    subprocess.run(cmd, check=True, cwd=str(ROOT), stdout=subprocess.PIPE, stderr=subprocess.PIPE)


def add_bg(slide, prs: Presentation):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), prs.slide_width, PptInches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None, y: float = 0.35):
    tx = slide.shapes.add_textbox(PptInches(0.55), PptInches(y), PptInches(12.2), PptInches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = PptPt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)

    if subtitle:
        tx2 = slide.shapes.add_textbox(PptInches(0.58), PptInches(y + 0.52), PptInches(12.1), PptInches(0.55))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = PptPt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x, y, w, h, font_size=15):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = PptPt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_diagram(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), PptInches(left), PptInches(top), width=PptInches(width))


def insert_slide_before(prs: Presentation, slide, before_index_0_based: int):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(before_index_0_based, sldId)


def _title_score(shape) -> tuple[float, float]:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return (0.0, 0.0)
    text = shape.text_frame.text.strip()
    if not text:
        return (0.0, 0.0)
    if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
        return (0.0, 0.0)
    top = float(shape.top)
    max_font = 0.0
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            if run.font.size:
                max_font = max(max_font, float(run.font.size))
    return (max_font, -top)


def renumber_titles(prs: Presentation):
    seq = 1
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        candidates = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            if float(shape.top) > float(PptInches(1.25)):
                continue
            candidates.append(shape)
        if not candidates:
            continue
        title_shape = sorted(candidates, key=_title_score, reverse=True)[0]
        tf = title_shape.text_frame
        old_text = tf.text.strip()
        base = re.sub(r"^\d+\.\s*", "", old_text)
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{seq}. {base}"
        r.font.name = "Aptos Display"
        r.font.size = PptPt(24)
        r.font.bold = True
        r.font.color.rgb = RGBColor(20, 48, 87)
        seq += 1


def update_pptx():
    write_mermaid_assets()
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(
        slide,
        "7. Pretrained Paper-Style Pooling Explained",
        "How cls, mean, max, and min work in the paper-style path and why the masking matters.",
    )
    add_body(
        slide,
        [
            "In the paper-style branch, pooling is applied after the pretrained encoders and after cross-attention fusion.",
            "The pooled object is a sequence of fused token vectors, not a single scalar score for each modality.",
            "CLS takes the first valid token, mean averages all valid tokens, max keeps the strongest per-dimension signal, and min keeps the weakest per-dimension signal.",
            "In the no-video paper-style path, the pooled sequence is the fused text+audio token sequence.",
            "In the video-enabled hybrid path, pooling is applied to the three modality vectors after the legacy-style fusion block.",
        ],
        0.7,
        1.45,
        5.9,
        5.35,
        font_size=14,
    )
    add_diagram(slide, PNG_PATH, 6.9, 1.6, 6.0)
    tx = slide.shapes.add_textbox(PptInches(6.95), PptInches(6.55), PptInches(5.95), PptInches(0.45))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Mermaid diagram exported as SVG and PNG"
    r.font.name = "Aptos"
    r.font.size = PptPt(10.5)
    r.font.italic = True
    r.font.color.rgb = RGBColor(95, 95, 95)

    insert_slide_before(prs, slide, 6)
    renumber_titles(prs)
    prs.save(str(PPTX_PATH))


def update_docx():
    doc = Document(str(DOCX_PATH))
    doc.add_heading("5.1 Pretrained Paper-Style Pooling Explained", level=2)
    doc.add_paragraph(
        "In the pretrained or paper-style branch, pooling is not done on one score per modality. It is done after the encoders have produced token sequences and after the cross-attention block has fused those sequences."
    )
    doc.add_paragraph(
        "The safest way to explain the branch is to say that the model first creates a fused text-plus-audio sequence, then reduces that sequence to a single vector for classification. If video is enabled in the hybrid path, the current implementation falls back to the three-modality fusion block, and pooling then applies across the three modality vectors."
    )
    doc.add_paragraph("The pooling choices mean the following:")
    for text in [
        "CLS: take the first valid token as the summary vector.",
        "Mean: average all valid token vectors dimension by dimension.",
        "Max: keep the strongest activation in each hidden dimension.",
        "Min: keep the lowest activation in each hidden dimension.",
    ]:
        doc.add_paragraph(text, style="List Bullet")
    doc.add_paragraph(
        "For the text-plus-audio paper-style branch, masking matters because the model should ignore padding tokens when it computes mean, max, or min. That is why the implementation uses the attention masks before pooling."
    )
    doc.add_paragraph(
        "For the video-enabled hybrid branch, the pooling is applied over the three modality vectors after the fusion block. That makes the interpretation slightly different: the model is summarizing modality-level evidence instead of token-level evidence."
    )
    doc.add_paragraph(
        "The correct thesis wording is therefore: pooling is a configurable summarization step applied after fusion. It is not a separate model and it does not act on one scalar value per modality."
    )
    doc.add_picture(str(PNG_PATH), width=Inches(5.8))
    cap = doc.add_paragraph("Figure: paper-style pooling and the hybrid video branch")
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.save(str(DOCX_PATH))


def main():
    update_pptx()
    update_docx()
    print("Updated First Review PPTX and DOCX with pretrained pooling explanation.")


if __name__ == "__main__":
    main()
