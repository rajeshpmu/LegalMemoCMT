from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
TITLES = {
    "Hugging Face Models: Why I Use Them",
    "From Diarization Output to Phase 2 Rows",
}

prs = Presentation(PATH)
slide_ids = prs.slides._sldIdLst
selected = []
for slide_id in list(slide_ids):
    slide = prs.part.related_part(slide_id.rId).slide
    title = slide.shapes.title.text if slide.shapes.title is not None else ""
    if title in TITLES:
        selected.append(slide_id)
        slide_ids.remove(slide_id)

# Put the model explanation after slide 4, before corpus expansion and EDA.
insert_at = 4
for offset, slide_id in enumerate(selected):
    slide_ids.insert(insert_at + offset, slide_id)

prs.save(PATH)
print(f"Reordered {len(selected)} Hugging Face slides in {PATH}")
