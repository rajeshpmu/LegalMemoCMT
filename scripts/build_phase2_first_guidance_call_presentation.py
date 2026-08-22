from __future__ import annotations

import csv
import json
import subprocess
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "implementation_docments" / "LegalMemoCMT_Phase2_First_Guidance_Call_Presentation.pptx"
FIG_DIR = ROOT / "implementation_docments" / "phase2_first_guidance_call_figures"
TRIB_MMD = FIG_DIR / "phase2_tribunal_bootstrap_architecture.mmd"
TRIB_SVG = FIG_DIR / "phase2_tribunal_bootstrap_architecture.svg"
TRIB_PNG = FIG_DIR / "phase2_tribunal_bootstrap_architecture.png"
ALIGN_MMD = FIG_DIR / "phase2_alignment_fields_pipeline.mmd"
ALIGN_SVG = FIG_DIR / "phase2_alignment_fields_pipeline.svg"
ALIGN_PNG = FIG_DIR / "phase2_alignment_fields_pipeline.png"
INDIA_MMD = FIG_DIR / "phase2_indian_adaptation_architecture.mmd"
INDIA_SVG = FIG_DIR / "phase2_indian_adaptation_architecture.svg"
INDIA_PNG = FIG_DIR / "phase2_indian_adaptation_architecture.png"

DEEP_DIVE_DOC = ROOT / "implementation_docments" / "LegalMemoCMT_Phase2_Implementation_Student_Deep_Dive.docx"
SUMMARY_JSON = ROOT / "data" / "processed" / "phase2" / "legalmeld_validated" / "dataset_quality_summary.json"
EXPANDED_COMPARE_SUMMARY = ROOT / "reports" / "phase2" / "expanded_planning_vs_verified_inventory_summary.json"
VERIFIED_INVENTORY = ROOT / "data" / "processed" / "phase2" / "verified_case_inventory.csv"
HEARING_MANIFEST = ROOT / "data" / "processed" / "phase2" / "hearing_manifest.csv"
WITNESS_MANIFEST = ROOT / "data" / "phase2" / "source_manifests" / "witness_harvest_manifest_resolved.csv"
LEGALMELD_METADATA = ROOT / "data" / "processed" / "phase2" / "legalmeld_validated" / "legalmeld_metadata_validated.csv"
ALIGNMENT_REVIEW = ROOT / "data" / "processed" / "phase2" / "legalmeld_validated" / "alignment_review_sample.csv"
TRAIN_CSV = ROOT / "data" / "processed" / "phase2" / "legalmeld_validated" / "train.csv"
DEV_CSV = ROOT / "data" / "processed" / "phase2" / "legalmeld_validated" / "dev.csv"
TEST_CSV = ROOT / "data" / "processed" / "phase2" / "legalmeld_validated" / "test.csv"


def load_json(path: Path) -> dict[str, object]:
    if not path.exists():
        return {}
    return json.loads(path.read_text())


def count_rows(path: Path) -> int:
    if not path.exists():
        return 0
    with path.open(newline="", encoding="utf-8-sig") as f:
        return sum(1 for _ in csv.DictReader(f))


def load_rows(path: Path) -> list[dict[str, str]]:
    if not path.exists():
        return []
    with path.open(newline="", encoding="utf-8-sig") as f:
        return list(csv.DictReader(f))


def first_row(rows: list[dict[str, str]], key: str, value: str) -> dict[str, str]:
    for row in rows:
        if (row.get(key) or "").strip().lower() == value.lower():
            return row
    return {}


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_bg(slide, prs: Presentation) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(12.2), Inches(0.65))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(0.95), Inches(12.1), Inches(0.45))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(
    slide,
    bullets: list[str],
    *,
    x: float = 0.75,
    y: float = 1.35,
    w: float = 11.95,
    h: float = 5.65,
    font_size: int = 17,
    indent: float = 0.0,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        p.line_spacing = 1.05
        p.left_margin = Inches(indent)
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)


def add_numbered(slide, bullets: list[str], **kwargs) -> None:
    add_body(slide, [f"{i+1}. {b}" for i, b in enumerate(bullets)], **kwargs)


def add_table_slide(prs: Presentation, title: str, subtitle: str, headers: list[str], rows: list[list[str]], *, x=0.55, y=1.35, w=12.1, h=5.65, font_size=13) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    table.first_row = True
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            table.cell(r_idx, c_idx).text = value
    for row in table.rows:
        for cell in row.cells:
            tf = cell.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor(35, 35, 35)
            if row == table.rows[0]:
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(20, 48, 87)
    return None


def add_diagram(slide, image_path: Path, left: float, top: float, width: float, height: float | None = None) -> None:
    if image_path.exists():
        if height is None:
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
        else:
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def render_mermaid(code: str, mmd: Path, svg: Path, png: Path) -> None:
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    mmd.write_text(code.strip() + "\n")
    for out_path in (svg, png):
        subprocess.run(
            ["npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(mmd), "-o", str(out_path), "-b", "white"],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.PIPE,
        )


def build_diagrams() -> None:
    render_mermaid(
        """
%%{init: {"theme": "base", "themeVariables": {"fontSize": "24px", "nodePadding": 22, "rankSpacing": 80, "nodeSpacing": 80}}}%%
flowchart LR
    P1["Phase 1<br/>baseline"] --> CL["Candidate<br/>ledger"]
    CL --> UCR["UCR<br/>verification"]
    UCR --> VI["Verified case<br/>inventory"]
    VI --> HM["Hearing<br/>manifest"]
    HM --> WM["Witness<br/>manifest"]
    WM --> SEG["Transcript-first<br/>segmentation"]
    SEG --> ALIGN["Forced<br/>alignment"]
    ALIGN --> CLIP["ffmpeg clip<br/>extraction"]
    CLIP --> VAL["Audio / video<br/>validation"]
    VAL --> LM["legalmeld_metadata_<br/>validated.csv"]
    LM --> SPLIT["Leakage-aware<br/>split"]
    SPLIT --> BOOT["Tribunal bootstrap<br/>ready"]
""",
        TRIB_MMD,
        TRIB_SVG,
        TRIB_PNG,
    )

    render_mermaid(
        """
%%{init: {"theme": "base", "themeVariables": {"fontSize": "24px", "nodePadding": 22, "rankSpacing": 65, "nodeSpacing": 55}}}%%
flowchart LR
    subgraph TX["Transcript side"]
        direction TB
        SRC["Transcript source"] --> UTT["utterance_text"]
        UTT --> NORM["transcript_text_normalized"]
    end
    subgraph AS["Audio side"]
        direction TB
        WAV["Audio clip"] --> ASR["asr_text"]
        ASR --> SIM["text_similarity"]
    end
    subgraph AD["Alignment decision"]
        direction TB
        STAT["alignment_status"] --> METH["alignment_method"]
        METH --> CONF["alignment_confidence"]
        CONF --> QT["quality_tier"]
        QT --> SPL["split<br/>train / dev / test / review"]
    end
    subgraph VG["Validation gate"]
        direction TB
        AV["audio_validation_status"] --> GATE["usable row?"]
        VV["video_quality_status"] --> GATE
        GATE --> OUT["model-ready row"]
    end
    NORM --> SIM
    SIM --> STAT
    WAV --> GATE
    UTT --> STAT
    CONF --> GATE
""",
        ALIGN_MMD,
        ALIGN_SVG,
        ALIGN_PNG,
    )

    render_mermaid(
        """
%%{init: {"theme": "base", "themeVariables": {"fontSize": "24px", "nodePadding": 22, "rankSpacing": 80, "nodeSpacing": 80}}}%%
flowchart LR
    A["Indian<br/>acquisition pack"] --> SC["Supreme Court<br/>livestreams"]
    A --> HC["High Court<br/>livestreams"]
    A --> DC["District court<br/>material"]
    A --> MT["Mock trials /<br/>academies"]
    SC --> LEDGER["Indian candidate<br/>ledger"]
    HC --> LEDGER
    DC --> LEDGER
    MT --> LEDGER
    LEDGER --> VERIFY["Indian<br/>verification layer"]
    VERIFY --> INVI["Indian verified<br/>inventory"]
    INVI --> IHM["Indian hearing<br/>manifest"]
    IHM --> IWM["Indian witness<br/>manifest"]
    IWM --> IALIGN["Indian alignment<br/>manifest"]
    IALIGN --> ICORP["Indian courtroom<br/>adaptation corpus"]
""",
        INDIA_MMD,
        INDIA_SVG,
        INDIA_PNG,
    )


def add_content_slide(prs: Presentation, title: str, subtitle: str, bullets: list[str], *, font_size: int = 17) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    add_body(slide, bullets, x=0.75, y=1.48, w=11.85, h=5.55, font_size=font_size)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
    *,
    left_font_size: int = 15,
    right_font_size: int = 15,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    for x, heading, bullets, font_size in [
        (0.55, left_title, left_bullets, left_font_size),
        (6.8, right_title, right_bullets, right_font_size),
    ]:
        box = slide.shapes.add_textbox(Inches(x), Inches(1.47), Inches(5.75), Inches(0.35))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = heading
        r.font.name = "Aptos"
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = RGBColor(20, 48, 87)
        add_body(slide, bullets, x=x, y=1.9, w=5.75, h=4.95, font_size=font_size)


def add_text_diagram_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    bullets: list[str],
    image_path: Path,
    caption: str,
    *,
    text_box=(0.7, 1.45, 6.1, 5.55),
    image_box=(7.0, 1.6, 5.75),
    font_size: int = 16,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    add_body(slide, bullets, x=text_box[0], y=text_box[1], w=text_box[2], h=text_box[3], font_size=font_size)
    add_diagram(slide, image_path, image_box[0], image_box[1], image_box[2])
    cap = slide.shapes.add_textbox(Inches(image_box[0]), Inches(6.55), Inches(image_box[2]), Inches(0.3))
    tf = cap.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = caption
    r.font.name = "Aptos"
    r.font.size = Pt(10.5)
    r.font.italic = True
    r.font.color.rgb = RGBColor(95, 95, 95)


def add_slide_title_only(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)


def add_dual_diagram_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    bullets: list[str],
    tribunal_image: Path,
    india_image: Path,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    add_body(
        slide,
        bullets,
        x=0.65,
        y=1.45,
        w=4.95,
        h=5.6,
        font_size=14,
    )
    add_diagram(slide, tribunal_image, 5.15, 1.48, 7.95, 1.65)
    cap1 = slide.shapes.add_textbox(Inches(5.15), Inches(3.18), Inches(7.95), Inches(0.25))
    cap1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r1 = cap1.text_frame.paragraphs[0].add_run()
    r1.text = "Tribunal bootstrap architecture"
    r1.font.name = "Aptos"
    r1.font.size = Pt(10.5)
    r1.font.italic = True
    r1.font.color.rgb = RGBColor(95, 95, 95)
    add_diagram(slide, india_image, 5.15, 3.55, 7.95, 2.55)
    cap2 = slide.shapes.add_textbox(Inches(5.15), Inches(6.15), Inches(7.95), Inches(0.25))
    cap2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = cap2.text_frame.paragraphs[0].add_run()
    r2.text = "Indian adaptation architecture"
    r2.font.name = "Aptos"
    r2.font.size = Pt(10.5)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor(95, 95, 95)


def build_pptx() -> None:
    build_diagrams()
    summary = load_json(SUMMARY_JSON)
    compare_summary = load_json(EXPANDED_COMPARE_SUMMARY)
    compare_counts = compare_summary.get("category_counts", {}) if isinstance(compare_summary.get("category_counts"), dict) else {}
    inventory_rows = load_rows(VERIFIED_INVENTORY)
    hearing_rows = load_rows(HEARING_MANIFEST)
    witness_rows = load_rows(WITNESS_MANIFEST)
    legalmeld_rows = load_rows(LEGALMELD_METADATA)
    review_rows = load_rows(ALIGNMENT_REVIEW)
    train_rows = count_rows(TRAIN_CSV)
    dev_rows = count_rows(DEV_CSV)
    test_rows = count_rows(TEST_CSV)

    verified_count = len({row.get("resolved_case_number", "") for row in inventory_rows if row.get("resolved_case_number")})
    paired_hearings = sum(1 for row in hearing_rows if (row.get("pairing_status") or "").strip().lower() == "paired")
    eligible_hearings = sum(1 for row in hearing_rows if (row.get("eligible_for_trimodal_dataset") or "").strip().upper() == "YES")
    unresolved_witnesses = sum(1 for row in witness_rows if (row.get("witness_type") or "").strip().lower() == "unresolved_witness")

    prs = Presentation()
    set_slide_size(prs)

    # Slide 1: title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.95))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "LegalMemoCMT"
    r.font.name = "Aptos Display"
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)
    tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(1.25), Inches(12.1), Inches(1.0))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Phase 2 First Guidance Call"
    r2.font.name = "Aptos Display"
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(20, 48, 87)
    tx3 = slide.shapes.add_textbox(Inches(0.58), Inches(2.1), Inches(12.2), Inches(0.95))
    tf3 = tx3.text_frame
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Multilingual Multimodal Emotion Analysis for Indian Courtroom Testimony"
    r3.font.name = "Aptos"
    r3.font.size = Pt(15)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(50, 50, 50)
    tx4 = slide.shapes.add_textbox(Inches(0.58), Inches(3.0), Inches(12.1), Inches(1.0))
    tf4 = tx4.text_frame
    p4 = tf4.paragraphs[0]
    r4 = p4.add_run()
    r4.text = "Student discussion deck for the mentor call: tribunal bootstrap pipeline, Indian adaptation strategy, and what is already implemented in the repo."
    r4.font.name = "Aptos"
    r4.font.size = Pt(12.5)
    r4.font.italic = True
    r4.font.color.rgb = RGBColor(75, 75, 75)
    add_body(
        slide,
        [
            "Phase 2 now proves the method on tribunal data and reuses the same pipeline for the Indian adaptation goal.",
            "The deck stays conservative: it explains what is verified, what is still bootstrap-level, and what still needs to be done.",
        ],
        x=0.85,
        y=4.05,
        w=11.7,
        h=2.5,
        font_size=17,
    )

    # Slide 2 agenda
    add_content_slide(
        prs,
        "Agenda",
        "What I will discuss and how the Phase 2 story is structured.",
        [
            "Phase 1 review and why it leads into Phase 2",
            "Objective, novelty, title, and abstract for Phase 2",
            "Proposed system and the Indian hybrid-corpus plan",
            "Architecture diagram and the literature reading map",
            "What is implemented so far, what remains, and what the Indian sources look like",
        ],
    )

    # Slide 3
    add_content_slide(
        prs,
        "Phase 1 Review",
        "Phase 1 proved that the project can work in a multimodal setting, but it still stayed at the coarse benchmark level.",
        [
            "Phase 1 established the multimodal backbone: text, audio, and video together.",
            "The main limitation was granularity: long hearings are noisy, while the model needs clean utterance-level samples.",
            "Phase 1 also showed why we need stronger validation and traceability before using courtroom data for training.",
            "This is why Phase 2 moves from whole recordings to grounded hearing, witness, and utterance layers.",
        ],
    )

    # Slide 4
    add_content_slide(
        prs,
        "Objective and Novelty of Phase 2",
        "Phase 2 is not just a bigger dataset. It is a grounded data-construction method plus an Indian adaptation plan.",
        [
            "The immediate objective is to build a MELD-style utterance corpus from verified tribunal hearings.",
            "The novelty is the verification-first pipeline: placeholders are rejected, actual case metadata is verified, and only grounded rows move forward.",
            "A second novelty is the reuse strategy: the same manifest logic can later be applied to Indian sources.",
            "So the tribunal corpus is the bootstrap dataset, while the Indian courtroom corpus is the final adaptation target.",
        ],
    )

    # Slide 5
    add_two_column_slide(
        prs,
        "Title and Abstract for Phase 2",
        "The title stays close to the research objective and the abstract explains the bootstrap-and-adapt method.",
        "Proposed title",
        [
            "LegalMemoCMT Phase 2: Multilingual Multimodal Emotion Analysis for Indian Courtroom Testimony",
        ],
        "Proposed abstract",
        [
            "Phase 2 extends LegalMemoCMT from a multimodal baseline into a legal-domain dataset construction and adaptation pipeline.",
            "The tribunal corpus is used as supervised bootstrap data because it provides real witness testimony with synchronized media.",
            "The pipeline verifies case metadata, builds hearing and witness manifests, aligns transcript utterances, and validates the exported clips.",
            "The final goal is Indian courtroom adaptation, not just tribunal data collection.",
        ],
        left_font_size=16,
        right_font_size=14,
    )

    # Slide 6
    add_two_column_slide(
        prs,
        "Proposed System for Phase 2",
        "The system is a pipeline, not a single model. Each stage keeps traceability back to a verified case.",
        "Tribunal bootstrap layer",
        [
            "Curated case ledger",
            "Corrected UCR enrichment",
            "Verified case inventory",
            "Hearing manifest",
            "Witness manifest",
            "Transcript-first alignment",
            "Validated LegalMELD-style export",
        ],
        "Why this matters",
        [
            "Each step filters noise before the next stage begins.",
            "The output is a grounded utterance dataset rather than a long video archive.",
            "The same manifest-and-alignment logic is designed to be reused later for Indian sources.",
        ],
    )

    # Slide 7
    add_two_column_slide(
        prs,
        "Indian Courtroom Adaptation Strategy",
        "The final project objective is an Indian courtroom corpus, so the bootstrap pipeline must be reusable on Indian sources.",
        "Indian hybrid corpus sources",
        [
            "Supreme Court livestreams",
            "High Court livestreams",
            "District court material where public records exist",
            "Mock trials and judicial academies",
        ],
        "Why this is a hybrid corpus",
        [
            "Real Indian testimony at scale is limited in public data.",
            "Court livestreams help with legal language and procedural adaptation.",
            "Mock trials help with speaker roles and witness-style dialogue.",
            "The tribunal corpus remains the bootstrap proof-of-method layer.",
        ],
    )

    # Slide 8
    add_dual_diagram_slide(
        prs,
        "Proposed Mermaid Architecture Diagram",
        "Use the two clearer diagrams to explain the tribunal bootstrap branch and the Indian reuse branch.",
        [
            "The tribunal diagram is the current implemented path.",
            "The Indian diagram shows the reuse path for the final dissertation objective.",
            "The diagrams are split and enlarged so the mentor can read the flow clearly instead of seeing one crowded picture.",
        ],
        TRIB_PNG,
        INDIA_PNG,
    )

    # Slide 9
    add_content_slide(
        prs,
        "Algorithms and Techniques Used",
        "What the pipeline actually does at a technical level, but still in student language.",
        [
            "Placeholder and negative-control rejection stops empty or non-specific case labels before any request is made.",
            "API-based record resolution checks whether the UCR page actually returns case-specific evidence.",
            "Grouping and deduplication keep hearing rows traceable and prevent repeated rows from polluting the corpus.",
            "Witness identity extraction recovers public names or protected codes without deanonymizing protected witnesses.",
            "ASR-based alignment and fuzzy sequence matching turn transcripts into utterance-level timestamped rows. The backend can stay heuristic or switch to WhisperX-style alignment without changing the exported schema.",
            "ffmpeg clip extraction plus audio/video validation make sure the output media is usable.",
            "Group-based splitting prevents train/dev/test leakage across the same hearing or related hearing group.",
        ],
        font_size=15,
    )

    # Slide 10
    add_text_diagram_slide(
        prs,
        "Core Alignment Fields and Dataset Pipeline",
        "This slide explains how the text, ASR, alignment, and split fields work together to produce a usable multimodal row.",
        [
            "The transcript fields define the intended courtroom utterance.",
            "The ASR field lets us compare the audio to the transcript and diagnose mismatch.",
            "The alignment fields decide whether the row is matched, fuzzy, or fallback and how much trust to give it.",
            "The split and quality fields keep the dataset leakage-aware and separate review rows from model-ready rows.",
            "The backend choice is an implementation detail; the row fields keep the same meaning whichever backend is used.",
        ],
        ALIGN_PNG,
        "Mixed horizontal and vertical pipeline view: text side, audio side, alignment decision, and validation gate.",
        text_box=(0.58, 1.42, 4.7, 5.65),
        image_box=(5.45, 1.5, 7.05),
        font_size=14,
    )

    # Slide 11
    add_table_slide(
        prs,
        "Algorithms and Techniques: What To Read",
        "This is the reading map I would use for the literature review and for explaining the implementation choices.",
        ["Technique", "Read first", "Actual link", "Why read it"],
        [
            ["Utterance-level data design", "MELD: A Multimodal Multi-Party Dataset for Emotion Recognition in Conversations", "https://aclanthology.org/P19-1050/", "Best match for the MELD-style row format that Phase 2 is building."],
            ["Speech recognition / timestamps", "Robust Speech Recognition via Large-Scale Weak Supervision (Whisper)", "https://arxiv.org/abs/2212.04356", "Good background for robust ASR on noisy speech and transcript-first alignment."],
            ["Alignment backend option", "WhisperX-style alignment", "https://github.com/m-bain/whisperX", "Explain this as a selectable backend for tighter word timestamps when the dependency is available."],
            ["Forced alignment", "NeMo Forced Aligner docs + NeMo paper + MFA paper", "https://docs.nvidia.com/nemo-framework/user-guide/latest/nemotoolkit/tools/nemo_forced_aligner.html  \nhttps://www.isca-archive.org/interspeech_2023/rastorgueva23_interspeech.pdf  \nhttps://www.isca-archive.org/interspeech_2017/mcauliffe17_interspeech.html", "Explains word- and segment-level timestamp generation for clip extraction."],
            ["Transcript encoder", "BERT: Pre-training of Deep Bidirectional Transformers for Language Understanding", "https://arxiv.org/abs/1810.04805", "Standard reference for contextual text embeddings."],
            ["Visual encoder", "An Image Is Worth 16x16 Words: Transformers for Image Recognition at Scale", "https://arxiv.org/abs/2010.11929", "Useful for the facial-cue branch and frame-level reasoning."],
            ["Imbalance-aware loss", "Focal Loss for Dense Object Detection", "https://arxiv.org/abs/1708.02002", "Helpful when some classes or outcomes dominate the training data."],
            ["Systems tools", "ffmpeg / difflib docs", "https://ffmpeg.org/documentation.html  \nhttps://docs.python.org/3/library/difflib.html", "These are engineering tools, not model papers, but they are necessary for the pipeline."],
        ],
        font_size=10.5,
        y=1.35,
        h=5.65,
    )

    # Slide 12
    add_two_column_slide(
        prs,
        "Source Sites And Download Flow",
        "Show the public tribunal source pages first, then explain how the UCR verification gate prevents generic pages from becoming corpus evidence.",
        "Source sites",
        [
            "ICTY Karadzic Trial: https://www.icty.org/en/content/karadzic-trial-hearing-list",
            "ICTY Mladic Trial: https://www.icty.org/x/cases/mladic/custom11/en/mladic_otp_witness_info.pdf",
            "ICTY Perisic Trial: https://ucr.irmct.org/scasedocs/case/IT-04-81",
            "ICTR Akayesu Trial: https://unictr.irmct.org/en/cases/ictr-96-04",
        ],
        "Download flow",
        [
            "Start from the candidate ledger entry and its public source URL.",
            "Search the UCR inventory by case number, case family, or witness code.",
            "Verify that the record resolves to case-specific metadata, not a generic page shell.",
            "Promote only grounded rows into the verified inventory and manifest builders.",
            "Download transcript or video media only after the row is verified.",
            "Treat generic headings and placeholder manifests as non-evidence.",
        ],
        left_font_size=12,
        right_font_size=12,
    )

    # Slide 13
    add_content_slide(
        prs,
        "What Is Implemented So Far",
        "The bootstrap layer is implemented and reproducible. The Indian acquisition layer is still the planned extension.",
        [
            f"Verified case inventory exists and currently holds {verified_count:,} verified case groups.",
            f"Hearing manifest and witness manifest builders exist; the current hearing manifest has {paired_hearings:,} paired hearings and {eligible_hearings:,} tri-modal-eligible hearings.",
            f"Resolved witness manifest exists; the current file is mostly conservative, with {unresolved_witnesses:,} unresolved witness rows.",
            f"LegalMELD validated outputs exist: {count_rows(LEGALMELD_METADATA):,} utterance rows, {train_rows:,} train rows, {dev_rows:,} dev rows, and {test_rows:,} test rows.",
            f"The current validated run reports {summary.get('high_confidence_alignments', 0)} high-confidence alignments, {summary.get('alignment_confidence_counts', {}).get('MEDIUM', 0) if isinstance(summary.get('alignment_confidence_counts'), dict) else 0} medium-confidence alignments, and {summary.get('split_leakage_violations', 0)} split-leakage violations.",
            f"The current corpus is still small and bootstrap-level, but it already proves the manifest-to-utterance pipeline.",
            "The alignment step is configurable, so the same export format can be produced through the heuristic path or a WhisperX-style path.",
        ],
        font_size=14,
    )

    # Slide 14
    add_table_slide(
        prs,
        "Expanded Planning vs Verified Inventory",
        "The TAP shortlist is exhausted, so this comparison tells us which planning rows are already covered and which rows still broaden the tribunal set.",
        ["Category", "Rows", "Meaning"],
        [
            [
                "Already verified by case name",
                str(compare_counts.get("already_verified_by_case_name", 4)),
                "These rows map to cases that already exist in the verified inventory, so they do not add new corpus coverage.",
            ],
            [
                "Already verified by case number",
                str(compare_counts.get("already_verified_by_case_number", 0)),
                "No row matched this way in the current run, which means the case-name overlap was the useful signal.",
            ],
            [
                "New expansion sources",
                str(compare_counts.get("new_expansion_source", 13)),
                "These rows are still missing from the verified inventory and are the ones that matter if the tribunal set is widened.",
            ],
            [
                "Unresolved placeholders",
                str(compare_counts.get("unresolved_or_placeholder", 0)),
                "None survived the comparison, so the shortlist is grounded even though it is exhausted.",
            ],
        ],
        font_size=10.7,
        y=1.35,
        h=4.55,
    )
    comp_slide = prs.slides[-1]
    comp_note = comp_slide.shapes.add_textbox(Inches(0.72), Inches(6.0), Inches(11.8), Inches(0.72))
    comp_tf = comp_note.text_frame
    comp_p = comp_tf.paragraphs[0]
    comp_p.alignment = PP_ALIGN.LEFT
    comp_r = comp_p.add_run()
    comp_r.text = f"Current summary: {compare_summary.get('expanded_rows', 17)} expanded rows, {compare_summary.get('covered_rows', 4)} covered rows, and {compare_summary.get('missing_rows', 13)} missing-source rows. Next action: inspect the missing-source CSV and broaden the tribunal planning inputs before any new download attempt."
    comp_r.font.name = "Aptos"
    comp_r.font.size = Pt(10.3)
    comp_r.font.italic = True
    comp_r.font.color.rgb = RGBColor(95, 95, 95)

    # Slide 15
    add_table_slide(
        prs,
        "Hearing Manifest EDA",
        "Use this slide while showing hearing_manifest.csv: 2,375 rows, 123 paired rows, 2,240 transcript-only rows, and 12 video-only rows.",
        ["Pattern", "Example hearing", "What to point out", "Why it matters"],
        [
            [
                "Paired / high",
                "hear_992268674c264fe8\nICTR-98-41\n01/06/2007",
                "Both transcript and video URLs are present, and the row is marked paired with high confidence.",
                "This is the strongest hearing-level evidence for the tri-modal branch.",
            ],
            [
                "Paired / high",
                "hear_c149355188328b0b\nIT-95-5/18\n24/03/2016",
                "The same pairing pattern appears in the ICTY branch, so the manifest is not limited to one tribunal.",
                "This shows the bootstrap method works across tribunals.",
            ],
            [
                "Transcript-only",
                "hear_278afb925fbdcba5\nICTR-98-41\n01/03/2006",
                "A transcript exists, but there is no matching video row, so the hearing stays out of the tri-modal set.",
                "This explains why not every grounded hearing becomes a tri-modal training sample.",
            ],
            [
                "Video-only",
                "hear_0a3fd1f78da95a16\nICTR-98-41\n01/04/2011",
                "A video record exists, but the transcript side is missing, so the pairing remains incomplete.",
                "This keeps the corpus honest and prevents false tri-modal counts.",
            ],
        ],
        font_size=11,
        y=1.35,
        h=5.7,
    )

    # Slide 16
    add_table_slide(
        prs,
        "Validated LegalMELD EDA",
        "Use this slide while showing data/processed/phase2/legalmeld_validated/: 130 multimodal rows, 45.2 usable minutes, and all 130 rows valid for text, audio, and video.",
        ["View", "Example / count", "What to point out", "Why it matters"],
        [
            [
                "Overall export",
                "130 utterances\n130 text-valid\n130 audio-valid\n130 video-valid",
                "This is the tribunal-branch utterance dataset, not the hearing manifest. It is already split into train, dev, test, and review rows.",
                "This is the actual MELD-style corpus layer used for later model work, and it is fully multimodal at the row level.",
            ],
            [
                "Usable minutes",
                "45.2 minutes\n(high + medium confidence)",
                "This is the stronger subset if you want to talk about training-quality minutes rather than just raw exported duration.",
                "This is the most honest Phase 2 size figure for early multimodal experimentation.",
            ],
            [
                "High-confidence row",
                "hear_7b2686bf1e5608e9_utt00002\nICTR-98-41\n03/05/2005",
                "A grounded utterance with matched alignment, a real transcript span, audio, and video paths.",
                "This is the clearest example of a usable tribunal utterance row.",
            ],
            [
                "Medium-confidence row",
                "hear_7b2686bf1e5608e9_utt00001\nICTR-98-41\n03/05/2005",
                "The utterance is still grounded, but the alignment is fuzzy rather than exact.",
                "This shows the pipeline keeps usable rows while preserving alignment confidence.",
            ],
            [
                "Review row",
                "hear_7b2686bf1e5608e9_utt00005\nICTR-98-41\n03/05/2005",
                "This row is a fallback alignment and is marked for manual review.",
                "This keeps weak rows visible instead of quietly treating them as high quality, which protects the Phase 2 dataset from noisy supervision.",
            ],
        ],
        font_size=10.5,
        y=1.32,
        h=5.8,
    )

    # Slide 17
    add_table_slide(
        prs,
        "Witness Controlled Validation Comparison",
        "Use this compact slide to show the four anchor hearings versus the six manually promoted hearings in the witness-only path.",
        ["Group", "Hearing ID", "Date", "Witness", "Reason selected"],
        [
            ["Anchor", "hear_2ef83c852251d65c", "21/06/2004", "ANTIPAS NYANJWA", "Already validated in the utterance-level witness subset."],
            ["Anchor", "hear_7b2686bf1e5608e9", "03/05/2005", "DM190 | EMMANUEL NERETSE", "Already validated in the utterance-level witness subset."],
            ["Anchor", "hear_8280237faba9d96f", "16/09/2004", "FILIP REYNTJENS", "Already validated in the utterance-level witness subset."],
            ["Anchor", "hear_f5d485391c1d04cc", "19/01/2004", "ROMÉO DALLAIRE", "Already validated in the utterance-level witness subset."],
            ["Promoted", "hear_dea210cdb4c728e0", "03/02/2004", "BRENT BEARDSLEY | Exanination-in-chief by Mr. White", "Manually promoted from the broader hearing discovery plan for controlled validation."],
            ["Promoted", "hear_8a80539e19e2df44", "05/07/2005", "DK120 | MATHIEU NGIRUMPATSE", "Manually promoted from the broader hearing discovery plan for controlled validation."],
            ["Promoted", "hear_da77e01d6076db4b", "18/06/2003", "OMAR SERUSHAGO", "Manually promoted from the broader hearing discovery plan for controlled validation."],
            ["Promoted", "hear_b3fafce46640756e", "21/09/2006", "ALOYS NTABAKUZE", "Manually promoted from the broader hearing discovery plan for controlled validation."],
            ["Promoted", "hear_23dd766aab93f457", "26/01/2004", "ROMÉO DALLAIRE", "Manually promoted from the broader hearing discovery plan for controlled validation."],
            ["Promoted", "hear_ebd43a70140ec251", "27/01/2004", "ROMÉO DALLAIRE", "Manually promoted from the broader hearing discovery plan for controlled validation."],
        ],
        font_size=10.2,
        y=1.28,
        h=5.95,
    )

    # Slide 18
    add_table_slide(
        prs,
        "Phase 2 Multimodal Summary Box",
        "Use this compact slide to point at the validated artifact and the manifest file in one place.",
        ["File", "Current role", "What to say", "Why it matters"],
        [
            [
                "data/processed/phase2/legalmeld_validated/legalmeld_metadata_validated.csv",
                "Validated multimodal artifact",
                "This is the final row-level export for the tribunal bootstrap branch.",
                "It is the file that actually carries the usable text, audio, and video rows.",
            ],
            [
                "data/processed/phase2/hearing_manifest.csv",
                "Hearing manifest",
                "This is the grounded hearing layer that feeds the utterance-level export.",
                "It shows how verified case evidence becomes hearing rows before utterance alignment.",
            ],
            [
                "45.2 usable minutes",
                "Usable Phase 2 subset",
                "This is the conservative multimodal training size from the high- and medium-confidence rows.",
                "It is the cleanest number to quote when discussing early Phase 2 multimodal readiness.",
            ],
            [
                "130 valid rows",
                "All three modalities present",
                "Every validated row has text, audio, and video available.",
                "This is why the export is genuinely multimodal rather than text-only with placeholders.",
            ],
        ],
        font_size=10.5,
        y=1.32,
        h=5.8,
    )

    # Slide 19
    add_table_slide(
        prs,
        "Validated Multimodal Quality EDA",
        "Use this slide to explain which modality-quality fields are informative and which ones are still uninformative in the current export.",
        ["Field", "Current EDA result", "Why it is useful now"],
        [
            ["audio_present", "YES for all 130 rows", "Confirms every validated row has an audio track."],
            ["audio_rms", "0.0037 to 0.0609, mean 0.0288", "Shows the audio energy range across clips."],
            ["silence_ratio", "0.2928 to 0.9726, mean 0.6231", "Shows how much of each clip is silence versus speech."],
            ["clipping_ratio", "0.0 for all 130 rows", "Confirms there is no clipping problem in the current export."],
            ["sample_rate", "16000 for all 130 rows", "Shows the audio has a consistent model-friendly sample rate."],
            ["audio_validation_status", "valid for all 130 rows", "Confirms the audio is usable for Phase 2."],
            ["video_quality_status", "VALID for all 130 rows", "Confirms the video side is usable for Phase 2."],
            ["split", "train 75, dev 10, test 38, review 7", "Shows the leakage-aware partitioning of the export."],
        ],
        font_size=10.25,
        y=1.28,
        h=5.9,
    )

    # Slide 20
    add_content_slide(
        prs,
        "Clean Implementation Checklist for What Remains",
        "The next work should improve the quality of the current bootstrap corpus before any larger expansion.",
        [
            "Improve alignment quality and reduce fallback rows.",
            "Increase grounded witness diversity and preserve traceability.",
            "Keep group-based splits strict so the same hearing does not leak across train, dev, and test.",
            "Recover more witness-visible clips if the source media allows it.",
            "Add manual review only where the alignment quality is weak, instead of treating every row as equal.",
            "Build the Indian acquisition pack and reuse the same manifest logic on those sources.",
        ],
        font_size=15,
    )

    # Slide 21
    add_table_slide(
        prs,
        "Realistic Indian Sources",
        "These source families can support the Indian adaptation phase, but each one has a different level of usefulness.",
        ["Source family", "What it offers", "Limitation", "Best use"],
        [
            [
                "Supreme Court livestreams",
                "Authentic Indian legal speech from official proceedings.",
                "Mostly arguments by judges and advocates; witness testimony is rare.",
                "Use for legal-language adaptation, judge speech, and procedural phrasing.",
            ],
            [
                "High Court livestreams",
                "Real courtroom discourse with Indian legal context.",
                "Usually argument-heavy and not a large witness-testimony source.",
                "Use for additional domain adaptation and speech-style variety.",
            ],
            [
                "District court material",
                "Closest public source to real witness examination.",
                "Public recordings and transcripts are limited and inconsistent.",
                "Use only when public multimodal testimony is actually available.",
            ],
            [
                "Mock trials and judicial academies",
                "Often include direct and cross-examination with visible speakers.",
                "Simulated, not real proceedings.",
                "Use for speaker-role structure, examination phases, and alignment practice.",
            ],
            [
                "Legal education videos",
                "Useful legal speech and explanation content.",
                "Not witness testimony.",
                "Use only for legal-language adaptation, not for witness emotion modeling.",
            ],
        ],
        font_size=11.5,
        y=1.32,
        h=5.8,
    )

    # Slide 22
    add_content_slide(
        prs,
        "Why a Hybrid Corpus Is the Correct Design",
        "The project needs a bootstrap corpus plus an Indian adaptation corpus, not one perfect source.",
        [
            "The tribunal corpus gives scale, supervision, and real testimony structure.",
            "The Indian corpus gives domain adaptation for accent, procedure, and local legal discourse.",
            "Mock-trial material helps where public real testimony is scarce.",
            "Using a hybrid corpus is more defensible than pretending one source family solves every requirement.",
            "This keeps the dissertation honest: tribunal data proves the method, and Indian data proves the transfer.",
        ],
        font_size=15,
    )

    # Slide 23
    add_content_slide(
        prs,
        "Expected Outcomes",
        "The outcome should be a MELD-style utterance dataset plus a reusable Indian adaptation path.",
        [
            "A grounded utterance-level legal dataset with synchronized text, audio, and video.",
            "A reproducible pipeline that can be rerun on tribunal data and later on Indian sources.",
            "A cleaner unit of training than whole hearings or long videos.",
            "A corpus that supports later emotion, credibility, and speaker-role analysis without inventing labels too early.",
            "A final Indian adaptation direction that is methodologically sound and not overclaimed.",
        ],
        font_size=15,
    )

    # Slide 24
    add_content_slide(
        prs,
        "What I Want to Discuss With You",
        "The guidance call is mainly about whether the scope, evidence, and next step order make sense.",
        [
            "Whether the tribunal bootstrap pipeline is now clear enough as a method proof.",
            "Whether the Indian adaptation scope should stay a formal part of Phase 2.",
            "Whether the next effort should go into better alignment quality or into wider source acquisition.",
            "Whether the literature reading map is the right order for the dissertation write-up.",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build_pptx()
