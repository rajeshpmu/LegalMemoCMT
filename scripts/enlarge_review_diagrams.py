from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_FILES = [
    ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.pptx",
    ROOT / "implementation_docments" / "Zeroth_Review_LegalMemoCMT.pptx",
]


def resize_picture(shape, left: float, top: float, width: float, height: float) -> None:
    shape.left = Inches(left)
    shape.top = Inches(top)
    shape.width = Inches(width)
    shape.height = Inches(height)


def update_first_review(prs: Presentation) -> None:
    # 0-based slide indices.
    targets = {
        4: (6.10, 1.48, 6.95, 5.50),  # Architecture diagram
        11: (7.10, 4.55, 5.95, 1.40),  # Phase 2 pipeline
        13: (6.10, 1.72, 6.95, 2.95),  # Appendices / fusion diagram
        21: (5.45, 3.00, 7.60, 3.10),  # Legacy pipeline
    }
    for idx, (left, top, width, height) in targets.items():
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if shape.shape_type == 13:
                resize_picture(shape, left, top, width, height)
                break


def update_zeroth_review(prs: Presentation) -> None:
    targets = {
        6: (6.00, 1.55, 6.85, 5.15),  # Phase 2 overview
        9: (6.10, 1.72, 6.95, 2.95),   # Operational details / fusion diagram
        16: (6.10, 1.60, 6.95, 2.15),  # Legacy pipeline
    }
    for idx, (left, top, width, height) in targets.items():
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if shape.shape_type == 13:
                resize_picture(shape, left, top, width, height)
                break


def main() -> None:
    for path in PPTX_FILES:
        prs = Presentation(str(path))
        if path.name.startswith("First_Review"):
            update_first_review(prs)
        else:
            update_zeroth_review(prs)
        prs.save(str(path))
        print(f"Updated diagram sizes in {path}")


if __name__ == "__main__":
    main()
