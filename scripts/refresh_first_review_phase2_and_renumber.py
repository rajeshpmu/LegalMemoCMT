from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.pptx"
PHASE2_OVERVIEW = ROOT / "artifacts" / "legalmemocmt_phase2.png"
PHASE2_PIPELINE = ROOT / "artifacts" / "mermaid" / "phase2_methodology_pipeline.png"


def add_bg(slide, prs: Presentation):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def set_textbox_text(shape, text, font_size=15, bold=False, italic=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor(35, 35, 35)


def add_title(slide, title: str, subtitle: str | None = None, y: float = 0.35):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(y + 0.52), Inches(12.1), Inches(0.55))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x, y, w, h, font_size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_diagram(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def clear_slide(slide):
    # Remove all existing shapes so the slide can be rebuilt cleanly.
    for shape in list(slide.shapes):
        el = shape._element
        el.getparent().remove(el)


def replace_slide_title(slide, new_title: str):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip():
            current = shape.text_frame.text.strip()
            if shape.name.startswith("TextBox") and current[:1].isdigit():
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = new_title
                r.font.name = "Aptos Display"
                r.font.size = Pt(24)
                r.font.bold = True
                r.font.color.rgb = RGBColor(20, 48, 87)
                return


def add_table_slide(prs: Presentation):
    slide = prs.slides[10]  # existing Phase 2 methodology slide
    clear_slide(slide)
    add_bg(slide, prs)
    add_title(
        slide,
        "11. Proposed Phase 2 Methodology",
        "Committee-style summary of the proposed courtroom-adaptation plan.",
    )

    # Clean existing body/diagram text boxes by reusing their shapes is harder;
    # instead, use a fresh layout by covering with new shapes.
    add_body(
        slide,
        [
            "Phase 2 adapts the Phase 1 MemoCMT-style pipeline to Indian courtroom testimony.",
            "Observable emotional cues are the target, not legal judgments.",
            "Input modalities are courtroom video testimony, audio testimony, and multilingual transcripts.",
            "Primary outputs are fear, anxiety, anger, stress timeline, and emotional transitions over time.",
        ],
        0.7,
        1.55,
        5.8,
        2.0,
        font_size=14,
    )

    table = slide.shapes.add_table(5, 3, Inches(0.7), Inches(3.0), Inches(11.9), Inches(2.8)).table
    headers = ["Component", "Planned method", "Purpose"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    rows = [
        ("Visual encoder", "ViT-style face feature extractor", "Capture witness facial cues and visible behavior."),
        ("Speech encoder", "XLS-R / IndicWav2Vec", "Capture prosody, pause structure, and stress."),
        ("Text encoder", "IndicBERT v2 / MuRIL", "Handle multilingual courtroom transcripts."),
        ("Fusion", "MemoCMT-style cross-modal transformer", "Integrate face, speech, and transcript evidence."),
    ]
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    add_body(
        slide,
        [
            "The committee should read this slide as the high-level project plan, not as a final legal system.",
            "The design starts from the Phase 1 backbone and adapts preprocessing, encoders, and evaluation targets.",
        ],
        0.75,
        5.95,
        11.6,
        0.75,
        font_size=13,
    )
    return slide


def add_diagram_slide(prs: Presentation):
    slide = prs.slides[11]  # existing Phase 2 pipeline slide
    clear_slide(slide)
    add_bg(slide, prs)
    add_title(
        slide,
        "12. Phase 2 Processing Pipeline",
        "Pipeline view for synchronized courtroom testimony analysis.",
    )
    add_body(
        slide,
        [
            "Segment testimony into aligned windows so face, audio, and text stay synchronized.",
            "Use a ViT-style visual encoder for witness face frames.",
            "Encode audio with XLS-R or IndicWav2Vec for prosody and stress-related features.",
            "Encode transcripts with IndicBERT v2 or MuRIL for multilingual courtroom language.",
            "Fuse the modality embeddings with a MemoCMT-style cross-modal transformer.",
            "Aggregate segment-level outputs into a continuous emotional timeline.",
        ],
        0.7,
        1.55,
        5.9,
        5.2,
        font_size=14,
    )
    add_diagram(slide, PHASE2_PIPELINE, 6.75, 1.55, 6.0)
    tx = slide.shapes.add_textbox(Inches(6.75), Inches(6.5), Inches(6.0), Inches(0.35))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Mermaid pipeline diagram exported as PNG"
    r.font.name = "Aptos"
    r.font.size = Pt(10.5)
    r.font.italic = True
    r.font.color.rgb = RGBColor(95, 95, 95)
    return slide


def add_checklist_slide(prs: Presentation):
    slide = prs.slides[12]  # existing Phase 2 adaptation slide
    clear_slide(slide)
    add_bg(slide, prs)
    add_title(
        slide,
        "13. Phase 2 Adaptation, Explainability, and Checklist",
        "Concise committee-facing implementation summary.",
    )
    left_bullets = [
        "Phase 2 starts from the Phase 1 backbone and adapts encoders, preprocessing, and evaluation targets.",
        "If courtroom data are limited, the practical strategy is pretraining on MELD and IEMOCAP, then fine-tuning on small legal-domain samples.",
        "Use partial freezing or transfer learning when data are sparse.",
    ]
    right_bullets = [
        "Expose attention maps or modality contribution summaries so the output is inspectable.",
        "Present results as emotional pattern analysis, not as a legal decision engine.",
        "Implementation checklist: reproduce Phase 1, align face/audio/text, integrate encoders, add fusion and temporal aggregation, then evaluate with metrics and timelines.",
    ]
    add_body(slide, left_bullets, 0.7, 1.55, 5.8, 4.9, font_size=14)
    add_body(slide, right_bullets, 6.9, 1.55, 5.8, 4.9, font_size=14)
    return slide


def renumber_titles(prs: Presentation):
    content_index = 1
    for i, slide in enumerate(prs.slides):
        # Skip title slide, keep it unnumbered.
        if i == 0:
            continue
        title_text = None
        title_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text.strip():
                txt = shape.text_frame.text.strip()
                if shape.name.startswith("TextBox"):
                    # use the first visible title-like text box
                    if title_text is None and txt[:1].isdigit():
                        title_text = txt
                        title_shape = shape
                        break
        if title_shape is None:
            continue
        # keep only the portion after the existing number, if any
        base = title_text
        if "." in title_text[:5]:
            base = title_text.split(".", 1)[1].strip()
        new_title = f"{content_index}. {base}"
        set_textbox_text(title_shape, new_title, font_size=24, bold=True)
        content_index += 1


def set_textbox_text(shape, text, font_size=15, bold=False, italic=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor(35, 35, 35)


def main():
    prs = Presentation(str(PPTX_PATH))
    add_table_slide(prs)
    add_diagram_slide(prs)
    add_checklist_slide(prs)
    renumber_titles(prs)
    prs.save(str(PPTX_PATH))
    print("Updated First Review PPTX")


if __name__ == "__main__":
    main()
