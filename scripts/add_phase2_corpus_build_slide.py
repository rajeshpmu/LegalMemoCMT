from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.before_corpus_build_slide.pptx"

NAVY = RGBColor(20, 48, 87)
TEXT = RGBColor(35, 35, 35)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = "Aptos"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def make_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.22))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    add_text(slide, "How the Corpus Was Built", 0.55, 0.38, 12.2, 0.55, 24, True, NAVY)
    add_text(slide, "The implementation converts public source material into traceable, MELD-style utterance rows.", 0.58, 0.95, 12.1, 0.35, 12, False, RGBColor(95, 95, 95))
    steps = [
        ("1. Source discovery", "data/clancy_urls.txt\ndata/tupac_trial_urls.txt\nUCR case and record pages", "Public URLs are recorded first; a URL is not yet corpus evidence."),
        ("2. Source manifest", "run_build_clancy_source_manifest.sh\nbuild_clancy_source_manifest.py", "yt-dlp probes title, duration, subtitle and media metadata."),
        ("3. Transcript turns", "run_build_clancy_turn_manifest.sh\nbuild_clancy_turn_manifest.py", "Subtitle cues are grouped into turn rows; persistent rejection IDs are excluded."),
        ("4. Media clips", "run_build_clancy_turn_clips.sh\nbuild_clancy_turn_clips.py", "Timestamp boundaries and source offsets are passed to ffmpeg to create MP4 + WAV."),
        ("5. Quality gates", "filter_legalmeld_rows_by_use.py\nEDA and outlier reports", "Rows are separated into usable, review, reject and confidence categories."),
        ("6. Dataset split", "run_build_clancy_dataset_split.sh\nrun_build_clancy_dataset_validation.sh", "Complete source groups are assigned to train/dev/test and leakage is checked."),
    ]
    x_positions = [0.55, 4.48, 8.41]
    y_positions = [1.55, 4.2]
    for i, (heading, code, meaning) in enumerate(steps):
        x = x_positions[i % 3]; y = y_positions[i // 3]
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.55), Inches(2.0))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(238, 244, 250)
        shape.line.color.rgb = RGBColor(150, 175, 200)
        add_text(slide, heading, x + 0.14, y + 0.12, 3.25, 0.28, 14, True, NAVY)
        add_text(slide, code, x + 0.14, y + 0.48, 3.25, 0.62, 10, False, TEXT)
        add_text(slide, meaning, x + 0.14, y + 1.22, 3.25, 0.58, 10, False, TEXT)
    add_text(slide, "Important: the pipeline creates candidate rows first. Training eligibility still requires transcript correctness, valid audio/video, acceptable duration, non-duplicate media, and reliable labels.", 0.7, 6.55, 11.8, 0.42, 11, True, NAVY)
    return slide


def move_after(prs: Presentation, slide, after_index: int) -> None:
    slide_list = prs.slides._sldIdLst
    # python-pptx stores slide references in sldIdLst, not as direct slide XML nodes.
    slide_id = slide_list[-1]
    slide_list.remove(slide_id)
    slide_list.insert(after_index, slide_id)


def main() -> None:
    shutil.copy2(DECK, BACKUP)
    prs = Presentation(DECK)
    slide = make_slide(prs)
    # Insert after the existing Source Sites and Download Flow slide (slide 7).
    move_after(prs, slide, 7)
    prs.save(DECK)
    print(f"Added corpus-build slide; wrote {len(prs.slides)} slides to {DECK}")
    print(f"Backup: {BACKUP}")


if __name__ == "__main__":
    main()
