from __future__ import annotations

import json
import subprocess
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from docx import Document
from docx.enum.text import WD_BREAK
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCX_PATH = ROOT / "implementation_docments" / "Second_Review_LegalMemoCMT.docx"
PPTX_PATH = ROOT / "implementation_docments" / "Second_Review_LegalMemoCMT.pptx"
FIG_DIR = ROOT / "implementation_docments" / "figures" / "second_review"
FIG_DIR.mkdir(parents=True, exist_ok=True)

MELD_SUMMARY = ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "summary.json"
MELD_FOLD2_CM = ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "fold_2" / "analysis_test" / "confusion_matrix.csv"
MELD_FOLD4_CM = ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "fold_4" / "analysis_test" / "confusion_matrix.csv"
CREMA_METRICS = ROOT / "results" / "paper_aligned_crema_d" / "cmt_min" / "metrics.json"
CREMA_CM = ROOT / "results" / "paper_aligned_crema_d" / "cmt_min" / "analysis_test" / "confusion_matrix.csv"

BENCHMARK_SPLIT_MMD = FIG_DIR / "second_review_benchmark_split.mmd"
EXECUTION_ORDER_MMD = FIG_DIR / "second_review_execution_order.mmd"
BENCHMARK_SPLIT_SVG = FIG_DIR / "second_review_benchmark_split.svg"
BENCHMARK_SPLIT_PNG = FIG_DIR / "second_review_benchmark_split.png"
EXECUTION_ORDER_SVG = FIG_DIR / "second_review_execution_order.svg"
EXECUTION_ORDER_PNG = FIG_DIR / "second_review_execution_order.png"

MELD_CM_PNG = FIG_DIR / "meld_fold2_confusion_matrix.png"
MELD_FOLD4_CM_PNG = FIG_DIR / "meld_fold4_confusion_matrix.png"
CREMA_CM_PNG = FIG_DIR / "crema_d_confusion_matrix.png"


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def render_mermaid(code: str, mmd_path: Path, svg_path: Path, png_path: Path) -> None:
    mmd_path.write_text(code, encoding="utf-8")
    subprocess.run(
        ["npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(mmd_path), "-o", str(svg_path), "-b", "white"],
        check=True,
    )
    subprocess.run(
        ["npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(mmd_path), "-o", str(png_path), "-b", "white"],
        check=True,
    )


def _clean_label(value: str) -> str:
    s = str(value)
    return s.split(":", 1)[1] if ":" in s else s


def build_confusion_image(matrix_csv: Path, output_png: Path, title: str) -> None:
    df = pd.read_csv(matrix_csv, index_col=0)
    df.index = [_clean_label(v) for v in df.index]
    df.columns = [_clean_label(v) for v in df.columns]
    fig, ax = plt.subplots(figsize=(8.4, 6.4))
    im = ax.imshow(df.values, cmap="Blues", aspect="auto")
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            ax.text(j, i, f"{int(df.values[i, j])}", ha="center", va="center", fontsize=8)
    ax.set_xticks(range(len(df.columns)))
    ax.set_xticklabels(df.columns, rotation=45, ha="right")
    ax.set_yticks(range(len(df.index)))
    ax.set_yticklabels(df.index)
    ax.set_title(title, fontsize=13)
    ax.set_xlabel("Predicted label")
    ax.set_ylabel("Actual label")
    ax.set_xlim(-0.5, len(df.columns) - 0.5)
    ax.set_ylim(len(df.index) - 0.5, -0.5)
    fig.tight_layout()
    fig.savefig(output_png, dpi=220, bbox_inches="tight")
    plt.close(fig)


def ppt_style(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_bg(slide, prs: Presentation):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None, y=0.35):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(y + 0.52), Inches(12.1), Inches(0.55))
        tf2 = tx2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x=0.75, y=1.35, w=11.95, h=5.65, font_size=18, indent=0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.08
        p.left_margin = Inches(indent)
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_picture(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        return slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    return None


def add_caption(slide, text: str, left: float, top: float, width: float):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(10.5)
    r.font.italic = True
    r.font.color.rgb = RGBColor(95, 95, 95)


def existing_titles(prs: Presentation) -> set[str]:
    titles = set()
    for slide in prs.slides:
        for sh in slide.shapes:
            if hasattr(sh, "text") and sh.text.strip():
                titles.add(sh.text.strip().splitlines()[0])
                break
    return titles


def remove_slides_from(prs: Presentation, start_index: int) -> None:
    """Remove all slides with 0-based index >= start_index."""
    sldIdLst = prs.slides._sldIdLst
    for idx in range(len(prs.slides) - 1, start_index - 1, -1):
        rel = sldIdLst[idx]
        prs.part.drop_rel(rel.rId)
        del sldIdLst[idx]


def add_para(doc: Document, text: str, *, italic: bool = False) -> None:
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.font.name = "Times New Roman"
    r.font.size = Pt(12)
    r.italic = italic


def add_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        doc.add_paragraph(item, style="List Bullet")


def add_content_slide(prs: Presentation, title: str, subtitle: str, bullets, font_size=16):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    add_body(slide, bullets, x=0.75, y=1.5, w=11.9, h=5.55, font_size=font_size)
    return slide


def add_text_diagram_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    bullets,
    image_path: Path,
    caption: str,
    font_size=16,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    add_body(slide, bullets, x=0.7, y=1.5, w=6.0, h=5.45, font_size=font_size)
    add_picture(slide, image_path, 7.05, 1.75, 5.55)
    add_caption(slide, caption, 7.05, 6.45, 5.55)
    return slide


def add_table_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    footer: str | None = None,
    widths: list[float] | None = None,
    font_size: float = 10.25,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.55), Inches(1.55), Inches(12.25), Inches(4.8)).table
    if widths is None:
        widths = [0.8, 1.5, 1.45, 1.45, 1.55, 5.5]
    for i, header in enumerate(headers):
        table.columns[i].width = Inches(widths[i])
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(20, 48, 87)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Aptos"
                r.font.size = Pt(10.5)
                r.font.bold = True
                r.font.color.rgb = RGBColor(255, 255, 255)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_idx < len(row) - 1 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Aptos"
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = RGBColor(35, 35, 35)
    if footer:
        tx = slide.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.4))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = footer
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.italic = True
        r.font.color.rgb = RGBColor(95, 95, 95)
    return slide


def set_slide_title(slide, new_title: str) -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            shape.text = new_title
            return


def move_appended_slides_before_anchor(prs: Presentation, anchor_slide, appended_count: int) -> None:
    """Move the last appended_count slides so they appear before anchor_slide."""
    if appended_count <= 0:
        return
    sldIdLst = prs.slides._sldIdLst
    slide_ids = list(sldIdLst)
    try:
        anchor_id = next(sldId for sldId in slide_ids if prs.slides[slide_ids.index(sldId)] is anchor_slide)
    except StopIteration:
        return

    new_ids = slide_ids[-appended_count:]
    for sldId in new_ids:
        sldIdLst.remove(sldId)

    anchor_index = list(sldIdLst).index(anchor_id)
    for offset, sldId in enumerate(new_ids):
        sldIdLst.insert(anchor_index + offset, sldId)


def append_pptx() -> None:
    meld_summary = read_json(MELD_SUMMARY)["metrics"]
    meld_fold2 = read_json(ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "fold_2" / "metrics.json")
    crema = read_json(CREMA_METRICS)
    prs = Presentation(PPTX_PATH)
    ppt_style(prs)
    remove_slides_from(prs, 24)

    add_content_slide(
        prs,
        "25. Phase 1 Completion Snapshot",
        "A concise status slide for the second review discussion.",
        [
            "Code implementation is roughly 90% complete.",
            "Paper-adaptation fidelity is roughly 75% complete.",
            "The core multimodal model, data plumbing, training, evaluation, and analysis tools are already present.",
            "The remaining work is mostly about benchmark finalization, metric clarity, and tighter paper-style alignment.",
        ],
        font_size=16,
    )

    add_content_slide(
        prs,
        "26. Phase 1 Result Summary",
        "Weighted cross-entropy is the strongest Phase 1 baseline and the right stopping point before adding video.",
        [
            "The MELD 5-fold CV summary gives mean accuracy 0.6247, mean weighted F1 0.6195, and mean macro F1 0.4395.",
            "Fold 2 is the strongest single-fold anchor with accuracy 0.6375, weighted F1 0.6254, and macro F1 0.4430.",
            "The focal-loss alternatives were worse, so the weighted-CE baseline remains the best Phase 1 result.",
            "The thesis novelty is therefore the legal-domain adaptation and the planned video-frame extension, not a small MELD gain.",
        ],
        font_size=15,
    )

    add_text_diagram_slide(
        prs,
        "27. Benchmark Split",
        "CREMA-D is now the speech-emotion track, while MELD remains the conversational track.",
        [
            "The new benchmark split does not change the implementation code.",
            "It changes how the same model is organized for discussion and reporting.",
            "CREMA-D is treated as the speech-emotion benchmark and uses speaker-independent CV.",
            "MELD remains the paper-aligned conversational benchmark and uses dialogue-based CV.",
        ],
        BENCHMARK_SPLIT_PNG,
        "Figure: benchmark split for the new project framing.",
        font_size=15,
    )

    add_text_diagram_slide(
        prs,
        "28. MELD Fold 2 Error Analysis",
        f"This slide is grounded in LegalMemoCMT_MELD_CV_5Fold_Analysis_Report.docx: mean Acc {meld_summary['accuracy']['mean']:.4f}, Fold 2 Acc {read_json(ROOT / 'results/paper_aligned_meld_cv/cmt_min/fold_2/metrics.json')['accuracy']:.4f}.",
        [
            "The MELD 5-fold report says Fold 2 is the strongest option if weighted aggregate performance is the priority.",
            "The same report says Fold 4 is the strongest option if balanced class behavior is the priority.",
            "The confusion matrix shows neutral-collapse behavior, especially between neutral, joy, anger, and surprise.",
            "This is good enough to say the implementation is in the right direction, but minority-class separation is still not fully solved.",
        ],
        MELD_CM_PNG,
        "Figure: MELD Fold 2 confusion matrix.",
        font_size=15,
    )

    add_text_diagram_slide(
        prs,
        "29. MELD Fold 4 Confusion Matrix",
        "Fold 4 is the best balanced-class anchor from the MELD 5-fold report.",
        [
            "Fold 4 has the strongest unweighted accuracy and macro F1 among the MELD folds.",
            "That makes it a useful complement to Fold 2, which is the strongest weighted-aggregate anchor.",
            "The matrix still shows confusion, but it is the clearest place to see the model's more balanced behavior.",
        ],
        MELD_FOLD4_CM_PNG,
        "Figure: MELD Fold 4 confusion matrix.",
        font_size=15,
    )

    add_text_diagram_slide(
        prs,
        "30. CREMA-D Status",
        f"This slide is grounded in LegalMemoCMT_CREMA_D_Analysis_Report.docx; the completed CREMA-D held-out run is still weak: Acc {crema['accuracy']:.4f}, Macro F1 {crema['macro_f1']:.4f}.",
        [
            "The CREMA-D analysis report describes the current held-out output as a diagnostic failure rather than a strong benchmark result.",
            "The new speaker-independent CREMA-D CV workflow exists to make the speech-emotion track more paper-like than the single held-out run.",
            "The current risk is class collapse, so the CV summary and fold consistency matter more than any single split result.",
        ],
        CREMA_CM_PNG,
        "Figure: current CREMA-D confusion matrix baseline.",
        font_size=15,
    )

    add_table_slide(
        prs,
        "31. Fold Status Overview",
        "This is the fold-level summary from LegalMemoCMT_MELD_CV_5Fold_Analysis_Report.docx.",
        ["Fold", "Acc / W-Acc", "UW-Acc", "Macro F1", "Weighted F1", "Reading"],
        [
            ["0", "0.6257", "0.4322", "0.4334", "0.6184", "Stable, neutral-heavy"],
            ["1", "0.6238", "0.4425", "0.4395", "0.6221", "Stable, sadness + neutral"],
            ["2", "0.6375", "0.4369", "0.4430", "0.6254", "Best weighted aggregate"],
            ["3", "0.6165", "0.4295", "0.4209", "0.6122", "Lowest overall fold"],
            ["4", "0.6199", "0.4672", "0.4606", "0.6194", "Best balanced-class fold"],
        ],
        footer="Fold 2 is the strongest option if weighted aggregate performance is the priority; Fold 4 is the strongest option if balanced class behavior is the priority.",
    )

    add_table_slide(
        prs,
        "32. Fold 4 Three-Way Training Comparison",
        "Fold 4 is the balanced-class anchor, so it is useful for checking whether the loss change helps minority behavior.",
        ["Run", "Script", "Strategy", "Acc / W-Acc", "UW-Acc", "Macro F1", "Weighted F1"],
        [
            [
                "Weighted-CE baseline",
                "run_paper_aligned_meld_cv.sh",
                "Fresh weighted-CE training from the paper-aligned path",
                f"{read_json(ROOT / 'results' / 'paper_aligned_meld_cv' / 'cmt_min' / 'fold_4' / 'metrics.json')['accuracy']:.4f}",
                f"{read_json(ROOT / 'results' / 'paper_aligned_meld_cv' / 'cmt_min' / 'fold_4' / 'metrics.json')['unweighted_accuracy']:.4f}",
                f"{read_json(ROOT / 'results' / 'paper_aligned_meld_cv' / 'cmt_min' / 'fold_4' / 'metrics.json')['macro_f1']:.4f}",
                f"{read_json(ROOT / 'results' / 'paper_aligned_meld_cv' / 'cmt_min' / 'fold_4' / 'metrics.json')['weighted_f1']:.4f}",
            ],
            [
                "Warm-start focal",
                "run_improvement_class_balanced_focal_meld_selected.sh",
                "Warm-start focal-loss run on the same Fold 4 split",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmstart_focal' / 'meld_selected' / 'cmt_min' / 'fold_4' / 'metrics.json')['accuracy']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmstart_focal' / 'meld_selected' / 'cmt_min' / 'fold_4' / 'metrics.json')['unweighted_accuracy']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmstart_focal' / 'meld_selected' / 'cmt_min' / 'fold_4' / 'metrics.json')['macro_f1']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmstart_focal' / 'meld_selected' / 'cmt_min' / 'fold_4' / 'metrics.json')['weighted_f1']:.4f}",
            ],
        ],
        widths=[1.35, 1.55, 2.65, 1.15, 1.0, 1.05, 1.05],
        font_size=9.4,
        footer="Fold 4 is the best balanced-class anchor in the paper-aligned MELD report; the warm-start focal run does not improve it, so the baseline still wins here.",
    )

    add_content_slide(
        prs,
        "33. W-Acc and UW-Acc Explained",
        "The speech-emotion CV track uses two related metrics so class imbalance is visible instead of hidden.",
        [
            "W-Acc means weighted accuracy in the base-paper naming style. In this codebase, it is the ordinary sample-level accuracy over all examples.",
            "UW-Acc means unweighted accuracy. It is the mean of the per-class accuracies, so each class contributes equally even when the dataset is imbalanced.",
            "A small example: if 8 of 10 samples are correct overall, W-Acc = 0.80. But if class accuracies are 1.0, 0.5, and 0.0, then UW-Acc = (1.0 + 0.5 + 0.0) / 3 = 0.50.",
            "This is why both numbers matter: W-Acc shows overall correctness, while UW-Acc shows whether minority classes are being learned fairly.",
        ],
        font_size=15,
    )

    add_content_slide(
        prs,
        "34. Guidance Call Results Summary",
        "The guidance discussion is based on the completed MELD 5-fold report and the CREMA-D analysis report.",
        [
            f"The MELD 5-fold report says Fold 2 is the strongest weighted-aggregate option and Fold 4 is the strongest balanced-class option; Fold 2 still gives the clearest paper-comparison anchor with accuracy {read_json(ROOT / 'results' / 'paper_aligned_meld_cv' / 'cmt_min' / 'fold_2' / 'metrics.json')['accuracy']:.4f} versus the paper's CMT+MIN test accuracy of 0.6418.",
            f"The MELD errors are structured rather than random, which is why the report says the implementation is in the right direction even though the match is not exact.",
            f"The CREMA-D analysis report says the current held-out run is weak: accuracy {crema['accuracy']:.4f}, macro F1 {crema['macro_f1']:.4f}, with a strong collapse toward a single class.",
            "The overall picture is that the MELD analysis report supports the paper-aligned conversational claim, while the CREMA-D report shows a weak baseline that still needs work.",
        ],
        font_size=15,
    )

    add_content_slide(
        prs,
        "35. Guidance Call Next Steps",
        "The next discussion should focus on the next benchmark decision implied by the two analysis reports.",
        [
            "Keep MELD as the main conversational result because the 5-fold report already gives a coherent paper-aligned story.",
            "Treat CREMA-D CV as the primary speech-emotion benchmark, but improve it before presenting it as a strong result.",
            "Use the CREMA-D report to decide whether the next improvement should target imbalance handling, training stability, or a stricter speech-emotion protocol.",
            "Use the mentor call to confirm whether the paper-aligned MELD report should stay the final conversational authority or whether a stricter paper-exact branch is still needed.",
        ],
        font_size=15,
    )

    add_content_slide(
        prs,
        "36. Second Review Questions",
        "Questions to steer the next step without over-committing to the wrong benchmark story.",
        [
            "Should CREMA-D CV remain the primary speech-emotion benchmark, or should a stricter paper-exact speech-emotion branch be pursued?",
            "Should the completed MELD paper-aligned CV result remain the main conversational result for now, while the paper-exact MELD template stays a future protocol option unless it is also run to completion?",
            "What should be improved first: imbalance handling, tuning schedule, or the weak benchmark itself?",
            "Is the current codebase complete enough for thesis refinement, or is one more protocol-level refinement needed?",
        ],
        font_size=16,
    )

    add_table_slide(
        prs,
        "37. Fold 2 Three-Way Training Comparison",
        "Same fold, same paper-aligned architecture, different training objective and checkpoint strategy.",
        ["Run", "Script", "Strategy", "Acc / W-Acc", "UW-Acc", "Macro F1", "Weighted F1"],
        [
            [
                "Weighted-CE baseline",
                "run_paper_aligned_meld_cv.sh",
                "Fresh weighted-CE training from the paper-aligned path",
                f"{meld_fold2['accuracy']:.4f}",
                f"{meld_fold2['unweighted_accuracy']:.4f}",
                f"{meld_fold2['macro_f1']:.4f}",
                f"{meld_fold2['weighted_f1']:.4f}",
            ],
            [
                "Focal from scratch",
                "run_improvement_class_balanced_focal_meld_selected.sh",
                "Fresh focal-loss run on the same Fold 2 split",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmstart_focal' / 'meld_selected' / 'cmt_min' / 'fold_2' / 'metrics.json')['accuracy']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmstart_focal' / 'meld_selected' / 'cmt_min' / 'fold_2' / 'metrics.json')['unweighted_accuracy']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmstart_focal' / 'meld_selected' / 'cmt_min' / 'fold_2' / 'metrics.json')['macro_f1']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmstart_focal' / 'meld_selected' / 'cmt_min' / 'fold_2' / 'metrics.json')['weighted_f1']:.4f}",
            ],
            [
                "Resume warm-start",
                "run_resume_warmstart_focal_meld_fold2.sh --epochs 5",
                "Load weighted-CE checkpoint and continue with focal loss",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmresume_focal' / 'meld_fold_2' / 'metrics.json')['accuracy']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmresume_focal' / 'meld_fold_2' / 'metrics.json')['unweighted_accuracy']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmresume_focal' / 'meld_fold_2' / 'metrics.json')['macro_f1']:.4f}",
                f"{read_json(ROOT / 'results' / 'improvement' / 'warmresume_focal' / 'meld_fold_2' / 'metrics.json')['weighted_f1']:.4f}",
            ],
        ],
        widths=[1.35, 1.55, 2.65, 1.15, 1.0, 1.05, 1.05],
        font_size=9.3,
        footer="The baseline weighted-CE run remains the strongest Fold 2 result; focal from scratch is the weakest; warm-start focal is a partial recovery but still not enough to beat the baseline.",
    )

    add_content_slide(
        prs,
        "38. Second Review Answers",
        "The evidence so far supports moving from baseline tuning to the video-enhanced LegalMemoCMT stage.",
        [
            "Keep the weighted-CE MELD Phase 1 result as the backbone baseline because it is still the best current conversational result.",
            "Treat the focal-loss experiments as diagnostics: they showed the loss is not enough by itself and did not beat the baseline.",
            "Proceed to HuBERT + BERT with Cross-Modal Transformer fusion plus ViT facial cues for the next stage.",
            "Also keep the HuBERT + BERT with Cross-Modal Transformer fusion text-audio path so the video stage can be compared cleanly against a non-visual baseline.",
            "The novelty should now be framed around legal-domain adaptation and multimodal video integration, not a small MELD gain.",
        ],
        font_size=15,
    )

    prs.save(PPTX_PATH)


def append_docx() -> None:
    meld_summary = read_json(MELD_SUMMARY)["metrics"]
    meld_fold2 = read_json(ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "fold_2" / "metrics.json")
    crema = read_json(CREMA_METRICS)

    doc = Document(DOCX_PATH)
    existing = {p.text.strip() for p in doc.paragraphs if p.text.strip()}
    if "20. Visual Summary for Second Review" in existing:
        return
    styles = doc.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"].font.size = Pt(12)
    for name in ["Heading 1", "Heading 2", "Heading 3", "Title"]:
        if name in styles:
            styles[name].font.name = "Times New Roman"

    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    doc.add_heading("20. Visual Summary for Second Review", level=1)
    add_para(
        doc,
        "This appendix-style section mirrors the slide deck. It gives the second review a visual summary of the new benchmark split, the MELD error structure, and the current CREMA-D status, so the discussion can stay aligned across the document and the presentation.",
    )

    doc.add_heading("20.1 Completion Snapshot", level=2)
    add_bullets(
        doc,
        [
            "Code implementation: about 90% complete.",
            "Paper-adaptation fidelity: about 75% complete.",
            "The remaining work is mostly benchmark finalization and protocol refinement.",
        ],
    )

    doc.add_heading("20.2 Benchmark Split Diagram", level=2)
    doc.add_picture(str(BENCHMARK_SPLIT_PNG), width=Inches(6.6))
    add_para(
        doc,
        "CREMA-D is the primary speech-emotion benchmark in the new split, while MELD is the primary conversational benchmark. The diagram is meant to make that separation visible at a glance.",
        italic=True,
    )

    doc.add_heading("20.3 Execution Order Diagram", level=2)
    doc.add_picture(str(EXECUTION_ORDER_PNG), width=Inches(6.6))
    add_para(
        doc,
        "The clean order is to complete or inspect the CREMA-D CV track first, then run or inspect the MELD CV track, and finally bring those results into the mentor discussion.",
        italic=True,
    )

    doc.add_heading("20.4 MELD Fold 2 Confusion Matrix", level=2)
    doc.add_picture(str(MELD_CM_PNG), width=Inches(6.6))
    add_para(
        doc,
        f"Fold 2 is the strongest overall held-out MELD fold in the CV summary, with accuracy {meld_fold2['accuracy']:.4f} and weighted F1 {meld_fold2['weighted_f1']:.4f}. The confusion pattern still shows a neutral-heavy bias, so the result is good but not fully solved.",
        italic=True,
    )

    doc.add_heading("20.5 CREMA-D Confusion Matrix", level=2)
    doc.add_picture(str(CREMA_CM_PNG), width=Inches(6.6))
    add_para(
        doc,
        f"The current CREMA-D single-split baseline remains weak, with accuracy {crema['accuracy']:.4f} and macro F1 {crema['macro_f1']:.4f}. The new CV workflow is intended to replace this as the main speech-emotion summary.",
        italic=True,
    )

    doc.add_heading("20.6 Second Review Talking Points", level=2)
    add_bullets(
        doc,
        [
            f"MELD is close to the base paper on the main accuracy range: mean CV accuracy {meld_summary['accuracy']['mean']:.4f}.",
            "CREMA-D is now the new speech-emotion track, but the current single-split baseline is very weak.",
            "The next question is whether to treat CREMA-D CV as the final speech-emotion story or keep it as a stepping stone.",
            "The codebase is already substantial; the remaining work is mainly about tightening the benchmark narrative and the strongest evaluation path.",
        ],
    )

    doc.add_heading("20.7 Second Review Answers After the Mentor Discussion", level=2)
    add_bullets(
        doc,
        [
            "Keep the weighted-CE MELD Phase 1 result as the backbone baseline, because it is still the strongest current conversational result.",
            "Treat the focal-loss experiments as diagnostics rather than the final answer, because they did not beat the baseline and exposed optimization instability.",
            "Proceed to HuBERT + BERT with Cross-Modal Transformer fusion plus ViT facial cues for the next stage of the thesis.",
            "Keep the HuBERT + BERT with Cross-Modal Transformer fusion text-audio path as the non-visual comparison backbone, so the video stage can be evaluated cleanly.",
            "Frame the thesis novelty around legal-domain adaptation and multimodal video integration, not around a small MELD gain.",
        ],
    )
    add_para(
        doc,
        "The practical conclusion is that Phase 1 is good enough to stop optimizing for the sake of MELD alone. The best use of the remaining effort is to extend the model with visual facial cues and compare the video-enhanced system against the current text-audio backbone.",
        italic=True,
    )

    DOCX_PATH.write_bytes(b"")
    doc.save(DOCX_PATH)


def main() -> None:
    render_mermaid(
        """flowchart LR
  A[Phase 1 Codebase] --> B[Primary Speech-Emotion Track]
  A --> C[Primary Conversational Track]
  B --> D[CREMA-D 5-fold CV]
  C --> E[MELD 5-fold CV]
  D --> F[W-Acc / UW-Acc Summary]
  E --> G[Accuracy / F1 / Confusion Analysis]
""",
        BENCHMARK_SPLIT_MMD,
        BENCHMARK_SPLIT_SVG,
        BENCHMARK_SPLIT_PNG,
    )
    render_mermaid(
        """flowchart TD
  A[Run CREMA-D CV] --> B[Inspect CREMA-D summary]
  B --> C[Run MELD CV]
  C --> D[Inspect MELD fold 2 analysis]
  D --> E[Discuss next steps with mentor]
""",
        EXECUTION_ORDER_MMD,
        EXECUTION_ORDER_SVG,
        EXECUTION_ORDER_PNG,
    )
    build_confusion_image(MELD_FOLD2_CM, MELD_CM_PNG, "MELD Fold 2 Confusion Matrix")
    build_confusion_image(MELD_FOLD4_CM, MELD_FOLD4_CM_PNG, "MELD Fold 4 Confusion Matrix")
    build_confusion_image(CREMA_CM, CREMA_CM_PNG, "CREMA-D Confusion Matrix")
    append_pptx()
    append_docx()
    print(PPTX_PATH)
    print(DOCX_PATH)


if __name__ == "__main__":
    main()
