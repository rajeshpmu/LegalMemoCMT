from __future__ import annotations

from pathlib import Path
import shutil

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "implementation_docments"
ESA_PPTX = OUT_DIR / "LegalMemoCMT_Phase1_ESA_Presentation.pptx"
ESA_DOCX = OUT_DIR / "LegalMemoCMT_Phase1_ESA_Prep_Guide.docx"
LOGO_SRC = Path("/Users/rajeshpmu/Downloads/PES_Mtech_Project_report_Sample_Latex/PES_Mtech_Project_report_latex_format/pes_logo.png")
LOGO_DST = OUT_DIR / "pes_logo.png"
FIG_DIR = ROOT / "implementation_docments" / "phase1_project_report_pes" / "figures"

COLOR_NAVY = RGBColor(18, 47, 85)
COLOR_TEAL = RGBColor(38, 115, 119)
COLOR_AMBER = RGBColor(204, 124, 0)
COLOR_LIGHT = RGBColor(245, 247, 250)
COLOR_BODY = RGBColor(40, 40, 40)
COLOR_MUTED = RGBColor(94, 94, 94)


def ensure_assets() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    if LOGO_SRC.exists():
        shutil.copy2(LOGO_SRC, LOGO_DST)


def fig(name: str) -> Path:
    p = FIG_DIR / name
    if not p.exists():
        raise FileNotFoundError(f"Missing expected figure asset: {p}")
    return p


def add_slide_background(slide, title: str, slide_no: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, PptInches(13.333), PptInches(0.58))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_NAVY
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, PptInches(0.58), PptInches(13.333), PptInches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_TEAL
    accent.line.fill.background()
    tx = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.12), PptInches(11.8), PptInches(0.3))
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos"
    run.font.size = PptPt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    tx2 = slide.shapes.add_textbox(PptInches(12.0), PptInches(0.12), PptInches(0.9), PptInches(0.3))
    p2 = tx2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = f"{slide_no:02d}"
    r2.font.name = "Aptos"
    r2.font.size = PptPt(14)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(255, 255, 255)


def add_body_box(slide, x, y, w, h, fill=COLOR_LIGHT):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(215, 220, 225)
    return shp


def add_textbox(slide, x, y, w, h, text, font_size=18, color=COLOR_BODY, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, font_size=18, bullet_color=COLOR_BODY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    first = True
    for item in items:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.text = f"• {item}"
        para.level = 0
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.name = "Aptos"
            run.font.size = PptPt(font_size)
            run.font.color.rgb = bullet_color
    return tb


def add_picture_fit(slide, path: Path, x, y, w, h):
    from PIL import Image as PILImage

    with PILImage.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box = w / h
    if aspect >= box:
        pw = w
        ph = w / aspect
    else:
        ph = h
        pw = h * aspect
    slide.shapes.add_picture(str(path), x + (w - pw) / 2, y + (h - ph) / 2, width=pw, height=ph)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, PptInches(4.1), PptInches(7.5))
    left.fill.solid()
    left.fill.fore_color.rgb = COLOR_NAVY
    left.line.fill.background()
    slim = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(4.1), 0, PptInches(0.18), PptInches(7.5))
    slim.fill.solid()
    slim.fill.fore_color.rgb = COLOR_TEAL
    slim.line.fill.background()
    if LOGO_DST.exists():
        slide.shapes.add_picture(str(LOGO_DST), PptInches(0.45), PptInches(0.45), width=PptInches(1.1))
    add_textbox(slide, PptInches(0.45), PptInches(1.55), PptInches(3.15), PptInches(2.0),
                "UE20CS971\nProject Phase 1\nEnd Semester Assessment", font_size=20, color=RGBColor(255, 255, 255), bold=True)
    add_textbox(slide, PptInches(4.55), PptInches(0.65), PptInches(8.1), PptInches(1.3),
                "LegalMemoCMT: An Explainable Multilingual Multimodal Emotional Cue Analysis Framework for Indian Courtroom Testimony Using Cross-Modal Transformers",
                font_size=22, color=COLOR_NAVY, bold=True)
    add_textbox(slide, PptInches(4.55), PptInches(2.0), PptInches(7.8), PptInches(1.0),
                "Phase 1 ESA presentation and technical defense pack", font_size=18, color=COLOR_TEAL, bold=True)
    add_textbox(slide, PptInches(4.55), PptInches(3.0), PptInches(7.8), PptInches(1.35),
                "Rajesh Upadhyaya  |  PES2PGE24DS200\nGuide: Ramesh Prakash Guledgudd", font_size=16, color=COLOR_BODY)
    add_textbox(slide, PptInches(4.55), PptInches(4.55), PptInches(7.9), PptInches(1.25),
                "This deck is built for the ESA review and keeps the Phase 1 story focused on the reproducible MELD baseline, the facial-cue extension, the demo path, and the error analysis that explains the current model behavior.", font_size=15, color=COLOR_MUTED)
    add_textbox(slide, PptInches(4.55), PptInches(6.35), PptInches(7.5), PptInches(0.4),
                "LegalMemoCMT Phase 1", font_size=13, color=COLOR_TEAL, bold=True)


def add_simple_slide(prs, title, slide_no, bullets, right_text=None, right_image=None, table=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, title, slide_no)
    add_body_box(slide, PptInches(0.3), PptInches(0.82), PptInches(12.7), PptInches(6.35), fill=RGBColor(255, 255, 255))
    add_bullets(slide, PptInches(0.55), PptInches(1.1), PptInches(5.8), PptInches(5.75), bullets, font_size=16)
    if right_text:
        right_text_h = 1.55 if table else 2.0
        add_body_box(slide, PptInches(6.55), PptInches(1.1), PptInches(6.2), PptInches(right_text_h), fill=RGBColor(246, 248, 250))
        add_textbox(slide, PptInches(6.75), PptInches(1.25), PptInches(5.7), PptInches(right_text_h - 0.2), right_text, font_size=14, color=COLOR_BODY)
    if right_image:
        add_body_box(slide, PptInches(6.6), PptInches(1.1), PptInches(6.1), PptInches(4.8), fill=RGBColor(246, 248, 250))
        add_picture_fit(slide, right_image, PptInches(6.7), PptInches(1.22), PptInches(5.9), PptInches(4.55))
    if table:
        rows, cols, data = table
        table_y = PptInches(2.85) if right_text else PptInches(1.15)
        table_h = PptInches(3.6) if right_text else PptInches(4.8)
        shape = slide.shapes.add_table(rows, cols, PptInches(6.55), table_y, PptInches(6.2), table_h)
        tbl = shape.table
        for r in range(rows):
            for c in range(cols):
                cell = tbl.cell(r, c)
                cell.text = data[r][c]
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.name = "Aptos"
                        run.font.size = PptPt(11)
                        run.font.color.rgb = COLOR_BODY
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_NAVY
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.bold = True
        for r in range(rows):
            for c in range(cols):
                tbl.cell(r, c).margin_left = PptInches(0.05)
                tbl.cell(r, c).margin_right = PptInches(0.05)
    return slide


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    add_title_slide(prs)
    add_simple_slide(
        prs, "Agenda", 2,
        [
            "Problem statement and scope",
            "Literature and why the baseline is paper-aligned",
            "Design approach, constraints, and methodology",
            "Architecture, implementation, and technologies",
            "Progress, results, references, and demo discussion",
        ],
        right_text="Use this slide to set the story: reproduce the baseline first, then show where facial cues help, and finally show the errors and current limits honestly."
    )
    add_simple_slide(
        prs, "Problem Statement", 3,
        [
            "The long-term project goal is courtroom-testimony understanding: when a witness or speaker talks in a legal setting, the system should help identify the emotion or affective state that is being expressed from the available evidence.",
            "Courtroom testimony is difficult because the same statement can sound neutral, tense, doubtful, or emotional depending on the speaking style, the audio tone, and the visible facial cue.",
            "Phase 1 uses MELD as the benchmark to build a reproducible multimodal pipeline that can learn from text, audio, and optionally video while keeping every sample traceable from manifest to checkpoint.",
            "MemoCMT is the closest paper-style reference because it gives a strong multimodal fusion baseline that I can reproduce first before extending the system toward the courtroom-use case.",
        ],
        right_text="One sentence defense: the real-world problem is legal-testimony emotion understanding, and MELD Phase 1 is the controlled benchmark step that proves the pipeline can learn and explain emotion patterns before moving to courtroom adaptation."
    )
    add_simple_slide(
        prs, "Abstract and Scope", 4,
        [
            "LegalMemoCMT is a cross-modal emotional cue analysis framework for courtroom-testimony use.",
            "Phase 1 validates the core implementation on MELD before later courtroom adaptation.",
            "Paper-aligned baseline: BERT + HuBERT + cross-modal fusion + MIN pooling.",
            "ViT facial-cue path: raw mp4 clips -> sampled frames -> cached embeddings.",
            "Scope: reproduce the benchmark, study the error pattern, and prepare the next adaptation step.",
        ],
        right_text="Say the abstract in two parts: first, reproduce the paper-aligned MELD baseline; second, add the visual support path and use the results to prepare for courtroom-testimony adaptation."
    )
    add_simple_slide(
        prs, "Literature Survey", 5,
        [
            "MemoCMT: cross-modal transformer feature fusion for multimodal emotion recognition.",
            "MELD: dialogue benchmark with seven emotion labels and strong class imbalance.",
            "BERT and HuBERT: pretrained text and speech encoders for the baseline.",
            "Vision Transformer: frame-based visual encoder for facial-cue analysis.",
            "Weighted cross-entropy and focal loss: imbalance-aware training objectives.",
        ],
        right_text="The literature slide should show that the design is not random: it follows established multimodal emotion-recognition ideas and keeps the visual branch as a support experiment."
    )
    add_simple_slide(
        prs, "Suggestions from Review - 3", 6,
        [
            "Keep the paper-aligned baseline as the anchor reference.",
            "Show fold-level metrics, confusion matrices, and top confusions.",
            "Explain the raw-video demo with confidence and top-3 probabilities.",
            "Be honest that the visual branch helps some cases but not every case.",
        ],
        right_text="Use this slide to show that you listened to the panel: you kept the baseline strong, added inspectable visual support, and improved the explanation around errors and confidence."
    )
    add_simple_slide(
        prs, "Design Approach", 7,
        [
            "Stage 1: reproduce the paper-aligned baseline on MELD.",
            "Stage 2: attach face-crop / full-frame ViT support to raw videos.",
            "Stage 3: test gated fusion and auxiliary loss as controlled changes.",
            "Stage 4: compare metrics, confusion matrices, and demo cases.",
        ],
        right_text="Explain why this approach is good: each step changes only one part of the system, so the result can be interpreted instead of guessed."
    )
    add_simple_slide(
        prs, "Design Constraints, Assumptions & Dependencies", 8,
        [
            "Constraint: MELD is class imbalanced, so accuracy alone is not enough.",
            "Assumption: sample_id, split, and label remain aligned in every manifest row.",
            "Dependency: pretrained backbones, ffmpeg, and cached video/audio files must be present.",
            "Dependency: raw MP4 processing needs consistent frame sampling and feature caching.",
        ],
        right_text="The key technical point is that the pipeline only works if the sample mapping stays stable. If the manifest or cache is off, then the prediction explanation is no longer trustworthy."
    )
    add_simple_slide(
        prs, "Design Details", 9,
        [
            "Manifest row -> transcript, audio path, video path, label, split.",
            "Preprocessing -> sample frames, crop face, encode frames with ViT, save .npy.",
            "Model -> text/audio encoders + fusion block + pooling + classifier.",
            "Analysis -> predictions CSV, metrics JSON, confusion matrix, top confusions.",
        ],
        right_image=fig("phase1_replication_pipeline.png")
    )
    add_simple_slide(
        prs, "Proposed Methodology / Approach", 10,
        [
            "Train the baseline and select the best checkpoint by validation metric.",
            "Run the facial-cue support path and compare with the same test examples.",
            "Check if the visual branch changes the error pattern or only the confidence.",
            "Use confusion analysis to decide whether the next step should be gating or auxiliary supervision.",
        ],
        right_image=fig("vit_support_pipeline.png")
    )
    add_simple_slide(
        prs, "Architecture (if applicable)", 11,
        [
            "Input streams: text, audio, and optionally video.",
            "Text branch: BERT-based utterance representation.",
            "Audio branch: HuBERT-based paralinguistic representation.",
            "Fusion: cross-modal transformer over modality embeddings.",
            "Output: emotion class probabilities and top-3 confidence values.",
        ],
        right_image=fig("phase1_metric_comparison.png")
    )
    add_simple_slide(
        prs, "Design Description (if applicable)", 12,
        [
            "src/data/preprocessing.py: frame sampling, face crop, audio loading, text normalization.",
            "src/models/model.py: encoders, fusion, pooling, classifier, and optional video support.",
            "src/train/train.py: training loop, checkpoint selection, and imbalance-aware losses.",
            "src/train/evaluate.py: metrics, predictions export, and held-out evaluation.",
            "scripts/*: reproducible entry points for training, analysis, and demo.",
        ],
        right_text="If asked where the result comes from, say: preprocessing builds the evidence, model.py performs fusion, train.py decides the checkpoint, and evaluate.py records the final numbers."
    )
    add_simple_slide(
        prs, "Technologies Used", 13,
        [
            "Python, PyTorch, Hugging Face Transformers",
            "NumPy, Pandas, scikit-learn, Matplotlib",
            "OpenCV / ffmpeg for video handling",
            "python-docx / python-pptx for report and slide generation",
            "RunPod or local macOS for execution and demo",
        ],
        right_text="Do not list tools only. For each tool, know its role: PyTorch trains the model, ffmpeg opens video, pandas and scikit-learn support analysis, and python-docx/pptx produce the review material."
    )
    progress_table = (
        6, 3, [
            ["Item", "Status", "What to say"],
            ["Baseline MELD replication", "Done", "This is the anchor result."],
            ["Fold analysis and confusion matrices", "Done", "This explains the class-level behavior."],
            ["Raw MP4 demo pipeline", "Done", "This proves end-to-end inference."],
            ["Face-crop ViT support", "Done", "This is the visual extension."],
            ["Gated / aux comparisons", "Done", "These are controlled improvement steps."],
        ]
    )
    add_simple_slide(
        prs, "Project Progress", 14,
        [
            "The Phase 1 core benchmark is complete and review-ready.",
            "The visual support path is implemented and has been analyzed on representative clips.",
            "The current story is not 'perfect accuracy'; it is 'traceable implementation plus meaningful error analysis'.",
        ],
        right_text="Use the progress slide to say that the baseline, analysis, demo path, and visual support branch are all in place. The only thing that remains is broader adaptation and refinement, not missing plumbing.",
        table=progress_table
    )
    add_simple_slide(
        prs, "References", 15,
        [
            "MemoCMT paper",
            "MELD dataset paper",
            "CREMA-D dataset paper",
            "BERT, HuBERT, and Vision Transformer references",
            "Focal loss / weighted CE reference",
        ],
        right_text="When the examiner asks for references, keep it short: cite the base paper, the datasets, the encoders, and the imbalance-loss literature."
    )
    add_simple_slide(
        prs, "Any other information", 16,
        [
            "Demo order: correct neutral -> correct non-neutral -> near-miss -> confident wrong -> summary metrics.",
            "Top-3 values are softmax probabilities, not accuracy.",
            "If the prediction is wrong, check sampling, modality bias, label ambiguity, and fusion behavior.",
            "Always relate the error to the confusion matrix rather than to a single output line.",
        ],
        right_text="This slide is your viva safety net. It tells you how to answer the 'why wrong?' question without overclaiming."
    )
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, "Thank You", 17)
    add_textbox(slide, PptInches(0.9), PptInches(2.0), PptInches(11.5), PptInches(0.8),
                "Thank you", font_size=28, color=COLOR_NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, PptInches(1.15), PptInches(3.0), PptInches(11.0), PptInches(1.2),
                "ESA review focus: reproducible Phase 1 baseline, explainable facial-cue support, and honest error analysis.", font_size=18, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    prs.save(str(ESA_PPTX))


def doc_paragraph(doc: Document, text: str, bold_prefix: str | None = None):
    p = doc.add_paragraph()
    p.style = doc.styles["Normal"]
    if bold_prefix and text.startswith(bold_prefix):
        r1 = p.add_run(bold_prefix)
        r1.bold = True
        r1.font.name = "Times New Roman"
        r1.font.size = Pt(11.5)
        r2 = p.add_run(text[len(bold_prefix):])
        r2.font.name = "Times New Roman"
        r2.font.size = Pt(11.5)
    else:
        r = p.add_run(text)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    return p


def add_heading(doc: Document, text: str, level: int = 1):
    p = doc.add_heading(text, level=level)
    for r in p.runs:
        r.font.name = "Times New Roman"
    return p


def build_docx() -> None:
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.7)
    sec.bottom_margin = Inches(0.7)
    sec.left_margin = Inches(0.8)
    sec.right_margin = Inches(0.8)
    for style_name in ["Normal", "Title", "Heading 1", "Heading 2", "Heading 3"]:
        if style_name in doc.styles:
            doc.styles[style_name].font.name = "Times New Roman"
    doc.styles["Normal"].font.size = Pt(11.5)

    if LOGO_DST.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(str(LOGO_DST), width=Inches(0.85))
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("ESA Preparation Guide")
    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(18)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("LegalMemoCMT Phase 1")
    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(16)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("For the UE20CS971 End Semester Assessment")
    r.italic = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(12)
    doc.add_paragraph()
    intro = (
        "This guide combines the most useful parts of the raw-video demo SOP, the three-video ESA demo guide, the code explanation, "
        "the benchmark execution guide, the RunPod SOP, and the paper-aligned technical reference. It is written to help you explain "
        "the project in Phase 1 clearly during the ESA review and to give you a single place to practise the demo flow, the code flow, "
        "and the error-analysis story."
    )
    doc_paragraph(doc, intro)
    add_heading(doc, "1. How to Present the ESA Story", 1)
    for t in [
        "Start with the problem statement: the project needs a reproducible multimodal emotion pipeline for MELD with clear traceability from raw data to predictions.",
        "Explain that Phase 1 is a controlled replication effort: the baseline is the main result, and the facial-cue branch is a support path for analysis and future courtroom-style adaptation.",
        "Use the demo to prove three things: the model accepts raw MP4 input, the outputs are interpretable, and the error analysis is honest rather than decorative.",
    ]:
        doc_paragraph(doc, t)
    add_heading(doc, "2. Slide-by-Slide Speaking Notes", 1)
    slide_notes = [
        ("Slide 1", "Say the title clearly and mention that Phase 1 focuses on reproducibility, not overclaiming. Give the full project title and the current scope in one sentence."),
        ("Slide 2", "Tell the examiner that the deck is organized from problem statement to progress, ending with references and a demo-friendly discussion."),
        ("Slide 3", "State the end goal first: courtroom-testimony emotion understanding. Then explain that Phase 1 uses MELD as the reproducible benchmark to build the multimodal pipeline, learn the class-imbalance behavior, and prepare the MemoCMT-style baseline for later legal-domain adaptation."),
        ("Slide 4", "Explain verbally that LegalMemoCMT is a cross-modal emotional cue analysis framework aimed at courtroom-testimony understanding. Phase 1 validates the MELD implementation first, using a paper-aligned baseline of BERT + HuBERT + cross-modal fusion + MIN pooling. The ViT support path then converts raw MP4 clips into sampled frames and cached embeddings. End with the scope statement: reproduce the benchmark, understand the error pattern, and prepare the path for courtroom-testimony adaptation later."),
        ("Slide 5", "Justify the design using MemoCMT, MELD, BERT, HuBERT, ViT, and imbalance-loss literature."),
        ("Slide 6", "Summarize review suggestions as action items: keep the baseline stable, show fold-level analysis, and explain confidence plus confusion matrix."),
        ("Slide 7", "Walk through the staged design: baseline -> face-crop/video -> gated fusion -> auxiliary loss -> analysis."),
        ("Slide 8", "Discuss dependencies and failure points: raw video quality, HF backbones, ffmpeg, and stable manifest rows."),
        ("Slide 9", "Use the pipeline figure to explain how raw data becomes cached .npy embeddings and then model inputs."),
        ("Slide 10", "Explain the experimental ladder and why each stage changes only one factor at a time."),
        ("Slide 11", "Describe the architecture as branches for text, audio, and optionally video that meet at fusion and classification."),
        ("Slide 12", "Be ready to name the files: preprocessing.py, model.py, train.py, evaluate.py, and the script wrappers."),
        ("Slide 13", "Talk about tools in terms of roles, not just names: what each library does in the pipeline."),
        ("Slide 14", "State the progress honestly: baseline done, analysis done, demo done, video support done, Phase 2 still pending."),
        ("Slide 15", "Cite the base paper, datasets, encoders, and imbalance loss papers only; do not overload this slide."),
        ("Slide 16", "Use this as your viva answer slide: top-3 are softmax probabilities, and errors come from sampling, modality bias, ambiguity, or fusion."),
        ("Slide 17", "Close confidently and be ready to move to questions or a short live demo."),
    ]
    tbl = doc.add_table(rows=1, cols=3)
    hdr = tbl.rows[0].cells
    hdr[0].text = "Slide"
    hdr[1].text = "What to say"
    hdr[2].text = "What to emphasize"
    for c in hdr:
        for p in c.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.name = "Times New Roman"
    for slide, note in slide_notes:
        row = tbl.add_row().cells
        row[0].text = slide
        row[1].text = note
        row[2].text = "Keep it technical, concise, and traceable."
    add_heading(doc, "2.1 Slide 4 Extended Talking Points", 2)
    for t in [
        "LegalMemoCMT is the broader framework name, and it is being developed as a cross-modal emotional cue analysis system for courtroom-testimony use.",
        "Phase 1 is the benchmark validation step: it proves that the MELD implementation works, is reproducible, and produces stable metrics before any legal-domain adaptation is attempted.",
        "The paper-aligned baseline uses BERT for text, HuBERT for audio, cross-modal transformer fusion, and MIN pooling, because that is the closest point of comparison to the reference-style implementation.",
        "The ViT support path starts from raw mp4 clips, samples frames, extracts facial cues, and stores the result as cached embeddings so the same clip can be analyzed again without repeating extraction.",
        "The scope statement should stay short and direct: reproduce the benchmark, understand the error pattern, and prepare for courtroom-testimony adaptation later.",
    ]:
        doc_paragraph(doc, t)
    add_heading(doc, "3. Demo SOP You Should Practise", 1)
    demo_steps = [
        "Run the reviewer bundle first so the metrics summary, confusion matrix, and top confusions are already ready.",
        "Open the baseline raw-mp4 demo command and show one correct neutral clip first.",
        "Show one correct non-neutral clip to prove the model is not simply predicting neutral everywhere.",
        "Show one near-miss or confident-wrong clip and explain why the top-3 probabilities matter.",
        "Finish with the summary files: metrics.json, confusion matrix.csv, top_confusions.csv, and the predicted-vs-actual table.",
    ]
    for t in demo_steps:
        doc_paragraph(doc, t)
    add_heading(doc, "4. Commands to Practise", 1)
    commands = [
        ("Prepare reviewer bundle", "bash scripts/run_phase1_review_demo_bundle.sh"),
        ("Baseline raw-MP4 demo", "DEVICE=cuda bash scripts/run_demo_paper_aligned_raw_mp4.sh test_dia279_utt9 data/MELD/raw/MELD.Raw/test/output_repeated_splits_test/dia279_utt9.mp4"),
        ("Gated+aux raw-MP4 demo", "DEVICE=cuda CHECKPOINT=results/facial_cues/meld_vit_facecrop_gated_video_aux/fold_4/best_model.pt bash scripts/run_demo_meld_vit_facecrop_gated_video_aux_raw_mp4.sh test_dia279_utt9 data/MELD/raw/MELD.Raw/test/output_repeated_splits_test/dia279_utt9.mp4"),
        ("Inspect cached embedding", "python3 scripts/read_meld_vit_facecue_npy.py --file results/phase1_review_demo/raw_mp4_cache/test_dia278_utt5_facecrop.npy"),
    ]
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Script / Command"
    table.rows[0].cells[1].text = "Why it matters"
    for c in table.rows[0].cells:
        for p in c.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.name = "Times New Roman"
    for cmd, why in commands:
        row = table.add_row().cells
        row[0].text = cmd
        row[1].text = why
    add_heading(doc, "5. Code Defense Cheat Sheet", 1)
    code_items = [
        ("src/data/preprocessing.py", "Samples frames, crops faces, loads audio, normalizes text, and creates the visual feature tensor that later gets cached as .npy."),
        ("src/models/model.py", "Defines the text/audio encoders, fusion block, pooling, and classifier. This is where modality interaction is decided."),
        ("src/train/train.py", "Loads the manifest, builds the train/validation split, trains the model, computes metrics, and saves the best checkpoint."),
        ("src/train/evaluate.py", "Loads a checkpoint and turns predictions into accuracy, weighted F1, macro F1, unweighted accuracy, and CSV output."),
        ("scripts/predict_phase1_raw_mp4_demo.py", "Drives the raw-video demo: mp4 input -> cached features -> model inference -> softmax probabilities -> top-3 output."),
    ]
    for path, desc in code_items:
        p = doc.add_paragraph()
        p.style = doc.styles["Normal"]
        r = p.add_run(f"{path}: ")
        r.bold = True
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        r2 = p.add_run(desc)
        r2.font.name = "Times New Roman"
        r2.font.size = Pt(11.5)
    add_heading(doc, "6. How to Explain a Wrong Prediction", 1)
    for t in [
        "Say whether the error is a near-miss or a confident wrong case.",
        "Check whether the transcript sounds biased toward another class.",
        "Check whether the sampled frames likely missed the peak facial expression.",
        "Check whether the visual branch is being overruled by stronger text/audio evidence.",
        "Point to the confusion matrix and top confusions, not just the single prediction line.",
    ]:
        doc_paragraph(doc, t)
    add_heading(doc, "7. Recommended Practice Order", 1)
    practice = [
        "Read the slide-by-slide notes once without running anything.",
        "Practise the baseline and gated demo commands on a single video clip.",
        "Open the confusion matrix and predicted-vs-actual table while speaking.",
        "Answer one likely viva question: why can a prediction be confident and still wrong?",
        "Answer another likely viva question: why is macro F1 lower than weighted F1?",
    ]
    for t in practice:
        doc_paragraph(doc, t)
    add_heading(doc, "Appendix: Best Source Documents to Keep Open During Prep", 1)
    refs = [
        "Phase1_Raw_MP4_Demo_SOP_Student_Guide.docx",
        "Phase1_ESA_Three_Video_Demo_Guide.docx",
        "LegalMemoCMT_Phase1_Code_Explanation.docx",
        "LegalMemoCMT_Benchmark_Execution_and_Analysis_Guide.docx",
        "LegalMemoCMT_RunPod_SOP.docx",
        "LegalMemoCMT_Paper_Aligned_Technical_Details_Expanded.docx",
    ]
    for ref in refs:
        doc_paragraph(doc, ref)
    doc.save(str(ESA_DOCX))


def main():
    ensure_assets()
    build_pptx()
    build_docx()
    print(f"Wrote {ESA_PPTX}")
    print(f"Wrote {ESA_DOCX}")


if __name__ == "__main__":
    main()
