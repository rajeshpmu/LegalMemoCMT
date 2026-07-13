from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.util import Inches as PptInches


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
READING_DOCX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Reading_Script.docx"
GENERATOR_PATH = ROOT / "scripts/build_phase1_esa_reading_script.py"


ROWS = [
    {
        "sample_id": "test_dia279_utt9",
        "cue": "Why don’t you phase it out?",
        "gt": "neutral",
        "base_pred": "anger",
        "base_conf": "0.4422",
        "base_top3": "anger 0.4422 | neutral 0.4204 | surprise 0.0500",
        "gate_pred": "neutral",
        "gate_conf": "0.9428",
        "gate_top3": "neutral 0.9428 | anger 0.0187 | joy 0.0142",
        "result": "Wrong -> Correct",
        "reason": "Near-boundary neutral case; gated+aux pulls the decision strongly to neutral.",
    },
    {
        "sample_id": "test_dia4_utt6",
        "cue": "Oh wait, Joey, you can’t go like that! You stink!",
        "gt": "disgust",
        "base_pred": "disgust",
        "base_conf": "0.5126",
        "base_top3": "disgust 0.5126 | anger 0.3356 | surprise 0.0572",
        "gate_pred": "disgust",
        "gate_conf": "0.9670",
        "gate_top3": "disgust 0.9670 | anger 0.0142 | sadness 0.0119",
        "result": "Correct -> Correct",
        "reason": "Both are correct; gated+aux sharply separates disgust from the next best class.",
    },
    {
        "sample_id": "test_dia278_utt5",
        "cue": "That’s not true, there are great pictures of us!",
        "gt": "surprise",
        "base_pred": "joy",
        "base_conf": "0.4809",
        "base_top3": "joy 0.4809 | surprise 0.3564 | neutral 0.0646",
        "gate_pred": "joy",
        "gate_conf": "0.9435",
        "gate_top3": "joy 0.9435 | surprise 0.0284 | anger 0.0092",
        "result": "Wrong -> Wrong",
        "reason": "The transcript itself looks positive/playful, so both checkpoints lean to joy; gated+aux becomes more confident in the wrong class.",
    },
    {
        "sample_id": "test_dia153_utt5",
        "cue": "If you're afraid of bugs.....get a bug.",
        "gt": "neutral",
        "base_pred": "fear",
        "base_conf": "0.6990",
        "base_top3": "fear 0.6990 | neutral 0.1666 | anger 0.0633",
        "gate_pred": "disgust",
        "gate_conf": "0.7620",
        "gate_top3": "disgust 0.7620 | anger 0.0775 | neutral 0.0763",
        "result": "Wrong -> Wrong",
        "reason": "The wording carries fear/disgust cues, so neutral is hard; gated+aux changes the wrong class but does not recover neutral.",
    },
    {
        "sample_id": "test_dia244_utt14",
        "cue": "Ross , foot on the floor or come over no more!",
        "gt": "neutral",
        "base_pred": "anger",
        "base_conf": "0.8735",
        "base_top3": "anger 0.8735 | neutral 0.0352 | fear 0.0339",
        "gate_pred": "anger",
        "gate_conf": "0.9486",
        "gate_top3": "anger 0.9486 | disgust 0.0139 | fear 0.0120",
        "result": "Wrong -> Wrong",
        "reason": "The utterance sounds commanding/angry, so both checkpoints prefer anger; gated+aux only increases confidence in the wrong class.",
    },
]


def ppt_style_run(run, *, size=11, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color)


def doc_style_run(run, *, size=11, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_box(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor.from_string(fill)
    if line:
        shape.line.color.rgb = PptRGBColor.from_string(line)
    else:
        shape.line.color.rgb = PptRGBColor.from_string(fill)
    return shape


def add_textbox(slide, left, top, width, height, text, *, size=11, bold=False, color="000000", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    ppt_style_run(r, size=size, bold=bold, color=color)
    return tb


def slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            t = shape.text_frame.text.strip()
            if t:
                return t.splitlines()[0].strip()
    return "Untitled"


def insert_doc_paragraph_before(anchor_paragraph, text, *, size=11, bold=False, color="000000", indent_pts=0):
    new_p = OxmlElement("w:p")
    anchor_paragraph._p.addprevious(new_p)
    paragraph = Paragraph(new_p, anchor_paragraph._parent)
    paragraph.text = text
    if indent_pts:
        paragraph.paragraph_format.left_indent = Pt(indent_pts)
    paragraph.paragraph_format.space_after = Pt(4)
    paragraph.paragraph_format.space_before = Pt(0)
    for run in paragraph.runs:
        doc_style_run(run, size=size, bold=bold, color=color)
    return paragraph


def update_pptx():
    prs = Presentation(str(PPTX_PATH))
    if any(
        "Live Demo Results Summary" in (shape.text_frame.text if getattr(shape, "has_text_frame", False) else "")
        for slide in prs.slides
        for shape in slide.shapes
    ):
        print("PPTX already has the live demo summary slide; skipping.")
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0.0, 0.0, 13.333, 0.58, "122F55")
    add_box(slide, 0.0, 0.58, 13.333, 0.08, "267377")
    add_textbox(slide, 0.45, 0.12, 11.8, 0.3, "Live Demo Results Summary", size=18, bold=True, color="FFFFFF")
    add_textbox(slide, 12.0, 0.12, 0.9, 0.3, str(len(prs.slides)), size=12, color="FFFFFF", align=PP_ALIGN.RIGHT)

    # Table
    headers = [
        "sample_id",
        "Transcript cue",
        "GT",
        "Baseline pred / conf",
        "Baseline top-3",
        "Gated+aux pred / conf",
        "Gated+aux top-3",
        "Result / reason",
    ]
    n_rows = len(ROWS) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(
        n_rows,
        n_cols,
        PptInches(0.18),
        PptInches(0.88),
        PptInches(12.75),
        PptInches(5.55),
    ).table

    col_widths = [1.38, 2.18, 0.72, 1.35, 1.96, 1.35, 1.96, 1.85]
    for i, w in enumerate(col_widths):
        table.columns[i].width = PptInches(w)

    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptRGBColor.from_string("DDEBF7")
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                ppt_style_run(r, size=8, bold=True, color="122F55")

    for i, row in enumerate(ROWS, start=1):
        values = [
            row["sample_id"],
            row["cue"],
            row["gt"],
            f'{row["base_pred"]} ({row["base_conf"]})',
            row["base_top3"],
            f'{row["gate_pred"]} ({row["gate_conf"]})',
            row["gate_top3"],
            f'{row["result"]}\n{row["reason"]}',
        ]
        for j, value in enumerate(values):
            cell = table.cell(i, j)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j not in [2] else PP_ALIGN.CENTER
                for r in p.runs:
                    ppt_style_run(r, size=7.3 if j != 7 else 7.0, color="202020")

    # Summary box
    add_box(slide, 0.35, 6.55, 12.55, 0.72, "EEF4FA", "C9D7E6")
    add_textbox(
        slide,
        0.52,
        6.66,
        12.2,
        0.45,
        "Summary: gated+aux improves 1/5 clips (test_dia279_utt9), keeps 1/5 correct, and remains wrong on 3/5 clips. The honest conclusion is selective improvement, not universal superiority.",
        size=10.5,
        bold=True,
        color="122F55",
    )
    add_textbox(
        slide,
        0.52,
        7.02,
        12.2,
        0.25,
        "Note: the baseline paper-aligned run shows missing gated_fusion/video_aux keys because those layers are not present in the baseline checkpoint; this is expected partial loading, not a demo failure.",
        size=8.5,
        color="5E5E5E",
    )

    prs.save(str(PPTX_PATH))
    print(f"Updated PPTX: {PPTX_PATH}")


def update_reading_doc():
    doc = Document(str(READING_DOCX_PATH))
    heading = "Slide 26: Live Demo Results Summary"
    if any(p.text.strip() == heading for p in doc.paragraphs):
        print("Reading script already has slide 26 section; skipping.")
        return

    # Append to the end.
    h = doc.add_paragraph()
    h.style = "Heading 1"
    r = h.add_run(heading)
    doc_style_run(r, size=14, bold=True, color="122F55")

    paras = [
        "Use this slide as the live evidence slide. Explain that each row uses the same raw mp4 clip for both checkpoints, so the only thing changing is the model checkpoint and its learned fusion behavior.",
        "Start by telling the examiner what the columns mean: sample_id identifies the clip, the transcript cue shows the wording that the text branch sees, GT is the MELD label, the baseline and gated columns show the predicted label plus confidence, and the top-3 columns show the three highest softmax probabilities for that sample.",
        "The first row, test_dia279_utt9, is the clearest improvement case. The transcript is a short neutral question, 'Why don’t you phase it out?'. The baseline is almost a tie between anger and neutral, with anger only slightly above neutral, so this is a true boundary case rather than a collapse. The gated+aux checkpoint moves neutral to 0.9428 and drops anger to 0.0187, which means the later model learned a much cleaner neutral boundary on this clip. That is the actual reason this row is better: the same input becomes confidently correct after the checkpoint change.",
        "The second row, test_dia4_utt6, is a correct non-neutral example in both checkpoints. The transcript, 'Oh wait, Joey, you can’t go like that! You stink!', is strongly negative and fits disgust much better than a neutral reading. The gated+aux checkpoint does not change the class, but it sharply increases the confidence from 0.5126 to 0.9670 and pushes the next-best class, anger, almost to zero. That is useful because it shows the later model is not only correct, but more separated from competing labels.",
        "The third row, test_dia278_utt5, is the important honest failure case. The transcript, 'That’s not true, there are great pictures of us!', sounds positive or playful, so joy is an understandable bias from the text branch. Both checkpoints predict joy, but the true label is surprise. The gated+aux checkpoint is worse in the calibration sense because it becomes much more confident in the wrong class, and the surprise probability collapses from 0.3564 to 0.0284. The actual reason to say here is not that the model is broken; it is that the transcript itself is misleading for this label, and the later checkpoint locked onto that misleading cue more strongly.",
        "The fourth row, test_dia153_utt5, is a hard neutral failure. The transcript, 'If you're afraid of bugs.....get a bug.', contains emotional lexical cues like afraid and bugs, so the model has a reason to drift away from neutral. The baseline chooses fear, while the gated+aux checkpoint shifts to disgust. Both are wrong, which tells you the later model changed the wrong class instead of recovering the neutral label. That is an important point because it shows selective improvement: the gated model did not solve every ambiguous neutral case.",
        "The fifth row, test_dia244_utt14, is another neutral failure where the wording sounds commanding or angry: 'Ross , foot on the floor or come over no more!'. Both checkpoints choose anger, and the gated+aux checkpoint becomes even more confident. The actual reading is that the utterance itself carries a strong anger-like surface form, so the model is not inventing a random class; it is responding to a cue that conflicts with the neutral ground truth.",
        "The correct overall conclusion from these five clips is that gated+aux is not universally better. It is clearly better on test_dia279_utt9 because it converts a boundary mistake into the correct neutral class. It is also stronger on the already-correct disgust example because it separates the correct class more cleanly. But it remains wrong on the other three clips, and in two of them it becomes more confident in the wrong answer. That is the honest result to present.",
        "When the examiner asks why the gated model helps in some cases, the technical answer is that the later checkpoint has learned different modality weighting and a sharper decision boundary. The same raw clip and the same cached face-crop feature are used in both runs, so the change is checkpoint-level rather than input-level. The gated fusion and auxiliary video loss can help when the visual or contextual evidence aligns with the true class, but they can also amplify a misleading transcript cue when the sample is ambiguous.",
        "Also explain the baseline missing-key message if asked. The baseline paper-aligned run prints missing gated_fusion and video_aux_classifier keys because the baseline checkpoint does not contain those extra modules. That is expected partial loading into a shared model class and does not mean the demo failed. The gated+aux run matches the richer checkpoint, so it loads without that mismatch.",
        "A concise viva sentence for this slide is: the gated+aux checkpoint improves the neutral boundary case and strengthens one correct non-neutral case, but it does not fix every hard sample. So the result is selective improvement, not a universal gain over the baseline.",
    ]
    for para in paras:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        r = p.add_run(para)
        doc_style_run(r, size=11)

    doc.save(str(READING_DOCX_PATH))
    print(f"Updated reading script: {READING_DOCX_PATH}")


def patch_generator_for_future_rebuilds():
    # Keep the generator in sync so a later rebuild preserves the new slide notes.
    text = GENERATOR_PATH.read_text()
    marker = '        25: [\\n'
    if '        26: [' in text:
        return
    insert_after = text.rfind('        25: [')
    if insert_after == -1:
        return
    # Minimal patch: add a 26 entry before the closing brace of the dict.
    insertion = '''
        26: [
            "Use this slide as the live evidence slide. Explain that each row uses the same raw mp4 clip for both checkpoints, so the only thing changing is the model checkpoint and its learned fusion behavior.",
            "The slide shows the transcript cue, ground truth, baseline prediction, baseline confidence, baseline top-3 probabilities, gated+aux prediction, gated+aux confidence, gated+aux top-3 probabilities, and a short reason for the outcome.",
            "The technical conclusion is selective improvement: the gated+aux checkpoint fixes the near-boundary neutral case, keeps one non-neutral success, but remains wrong on the other hard clips and can become more confident in the wrong answer.",
        ],\n'''
    text = text.replace('\n    }\n\n\nif __name__ == "__main__":', f'{insertion}    }}\\n\\n\\nif __name__ == "__main__":')
    GENERATOR_PATH.write_text(text)
    print(f"Updated generator: {GENERATOR_PATH}")


if __name__ == "__main__":
    update_pptx()
    update_reading_doc()
    patch_generator_for_future_rebuilds()
