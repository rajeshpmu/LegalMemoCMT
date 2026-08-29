"""Add acceptance-gate failure-check slides and speaking notes."""
from pathlib import Path
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document

ROOT=Path(__file__).resolve().parents[1]
PPTX=ROOT/'implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Presentation_v1.pptx'
DOCX=ROOT/'implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Speaking_Doc.docx'
IMAGE=ROOT/'implementation_docments/figures/phase2_acceptance_gate_failure_checks.png'
MARKER='GATE_FAILURE_CHECKS_V1'
NAVY=RGBColor(20,48,87); TEXT=RGBColor(35,35,35); GREY=RGBColor(95,95,95)

def text(slide,value,x,y,w,h,size=16,color=TEXT,bold=False):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); box.text_frame.word_wrap=True; box.text_frame.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p=box.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=value; r.font.name='Aptos'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color

def frame(prs,heading,subtitle):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor(255,255,255)
    b=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.22)); b.fill.solid(); b.fill.fore_color.rgb=NAVY; b.line.fill.background()
    text(s,heading,.55,.38,12.2,.55,24,NAVY,True); text(s,subtitle,.58,.96,12,.35,12,GREY); return s

def before(slides,slide,index):
    ids=slides._sldIdLst; item=ids[-1]; ids.remove(item); ids.insert(index,item)

def update_ppt():
    prs=Presentation(PPTX)
    if any(MARKER in sh.text for s in prs.slides for sh in s.shapes if hasattr(sh,'text')): return
    backup=PPTX.with_name(PPTX.name+'.before_gate_failure_checks.pptx')
    if not backup.exists(): shutil.copy2(PPTX,backup)
    insertion=len(prs.slides)-1
    s=frame(prs,'Acceptance Gate: Two Critical-Failure Checks','The gate distinguishes a materially ambiguous row from an ordinary low-confidence row')
    s.shapes.add_picture(str(IMAGE),Inches(.7),Inches(1.35),width=Inches(12.0),height=Inches(5.5)); before(prs.slides,s,insertion); insertion+=1
    s=frame(prs,'Critical Conflicts Versus Ordinary Weak Rows','Current 200-row pilot: 7 critical conflicts within 75 unresolved rows')
    text(s,'CRITICAL CONFLICT\n\n• Phase 1 is non-neutral but the basic candidate changes to neutral without safe quoted/non-self scope plus CALM_COMPOSED evidence.\n• DISTRESSED is proposed without distress corroboration.\n• These rows require focused review because the evidence supports materially different interpretations.',.75,1.4,5.9,5.2,17)
    text(s,'ORDINARY GATE FAILURE\n\n• Basic confidence < 0.70.\n• Affect confidence < 0.60.\n• Basic or affect candidate is UNKNOWN/UNRESOLVED.\n\nThese rows are also UNRESOLVED/WEAK, but the reason is insufficient support rather than a specific cross-signal conflict.',6.9,1.4,5.7,5.2,17)
    before(prs.slides,s,insertion); prs.save(PPTX)

def update_doc():
    doc=Document(DOCX)
    if any(MARKER in p.text for p in doc.paragraphs): return
    backup=DOCX.with_name(DOCX.name+'.before_gate_failure_checks.docx')
    if not backup.exists(): shutil.copy2(DOCX,backup)
    doc.add_heading('Acceptance Gate Failure Checks: Speaking Notes',1)
    p=doc.add_paragraph(); p.add_run(MARKER).font.size=Pt(8)
    doc.add_paragraph('The acceptance gate performs two different kinds of protection. A critical conflict means the available signals support materially different interpretations, so the row needs focused human inspection. An ordinary gate failure means the row does not meet a confidence or resolution requirement, even if no specific contradiction has been found. Both outcomes are UNRESOLVED/WEAK, but their review priorities and explanations differ.')
    doc.add_heading('Critical check 1: Phase 1 versus integrated basic emotion',2)
    doc.add_paragraph('The code first detects a non-neutral Phase 1 prediction that changes to a neutral basic-emotion candidate. That change is allowed only when the transcript scope is OTHER_PERSON_DESCRIBED, EVENT_DESCRIBED, or QUOTED_SPEECH, speaker_emotion_evidence_present is NO, and the behavioral affect is CALM_COMPOSED. Otherwise the change is treated as critical because the model and the integrated interpretation disagree about the witness’s own emotion without enough attribution evidence.')
    doc.add_heading('Critical check 2: unsupported distress',2)
    doc.add_paragraph('The second check catches proposed_courtroom_affect=DISTRESSED when distress_corroboration_present is not YES. Low valence and elevated excitement can indicate negative activation, but they do not prove distress. Corroboration may come from a categorical audio cue, distress-specific language, or separately recorded visual evidence. Without it, the row must not be automatically accepted as distressed.')
    doc.add_heading('Ordinary unresolved checks',2)
    doc.add_paragraph('A row can be UNRESOLVED/WEAK without being critical. This occurs when basic confidence is below 0.70, affect confidence is below 0.60, or either candidate is UNKNOWN/UNRESOLVED. These rows are not necessarily contradictory; the system simply lacks enough support for automatic acceptance.')
    doc.add_heading('How to explain the pilot counts',2)
    doc.add_paragraph('In the current 200-row pilot, 125 rows are AUTO_ADJUDICATED/SILVER and 75 are UNRESOLVED/WEAK. Seven of the 75 unresolved rows have critical_conflict=YES. The critical rows are a subset of the unresolved rows, so the counts must not be added together. I inspect the seven first, then review the remaining unresolved rows by their acceptance_gate_reason. SILVER is machine-assisted acceptance, not human gold truth.')
    doc.add_paragraph('The Mermaid diagram shows the control flow: candidates are generated first, then the gate checks critical ambiguity and candidate validity before assigning SILVER or WEAK. This makes the implementation reproducible and prevents a high confidence in one dimension from hiding missing evidence in another.')
    doc.save(DOCX)

def main():
    update_ppt(); update_doc(); print({'pptx':str(PPTX),'docx':str(DOCX),'marker':MARKER})

if __name__=='__main__': main()
