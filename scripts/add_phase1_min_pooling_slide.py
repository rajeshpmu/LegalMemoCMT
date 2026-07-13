from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from docx.oxml import OxmlElement
from docx.shared import Pt, RGBColor
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Reading_Script.docx"


def style_run(run, *, size=10.5, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def style_ppt_run(run, *, size=10.5, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color)


def insert_paragraph_before(paragraph: Paragraph, text: str = "", style: str | None = None) -> Paragraph:
    new_p = OxmlElement("w:p")
    paragraph._p.addprevious(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        run = new_para.add_run(text)
        style_run(run)
    return new_para


def insert_paragraph_after(paragraph: Paragraph, text: str = "", style: str | None = None) -> Paragraph:
    new_p = OxmlElement("w:p")
    paragraph._p.addnext(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        run = new_para.add_run(text)
        style_run(run)
    return new_para


def add_box(slide, left, top, width, height, fill, line=None, radius=False):
    shp_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shp_type,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = PptRGBColor.from_string(fill)
    shp.line.color.rgb = PptRGBColor.from_string(line or fill)
    return shp


def add_text(slide, left, top, width, height, text, *, size=11, bold=False, color="000000", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    style_ppt_run(r, size=size, bold=bold, color=color)
    return tb


def insert_slide_before(prs: Presentation, before_index_0_based: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(before_index_0_based, sldId)
    return slide


def title_shape(slide):
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if float(shape.top) > 1200000:
            continue
        candidates.append(shape)
    if not candidates:
        return None

    def score(shape):
        max_font = 0.0
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    max_font = max(max_font, float(run.font.size))
        return (max_font, -float(shape.top))

    return max(candidates, key=score)


def number_shape(slide):
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = shp.text_frame.text.strip()
        if not txt or not re.fullmatch(r"\d+", txt):
            continue
        if float(shp.top) <= 350000 and float(shp.left) >= 10000000:
            return shp
    return None


def renumber_titles(prs: Presentation) -> None:
    seq = 1
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        ts = title_shape(slide)
        if ts is None:
            continue
        tf = ts.text_frame
        old = tf.text.strip()
        base = re.sub(r"^\d+\.\s*", "", old)
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{seq}. {base}"
        r.font.name = "Aptos Display"
        r.font.size = PptPt(24)
        r.font.bold = True
        r.font.color.rgb = PptRGBColor.from_string("122F55")
        seq += 1


def renumber_slide_numbers(prs: Presentation) -> None:
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        nshape = number_shape(slide)
        if nshape is not None:
            nshape.text = str(idx)
            for p in nshape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
                for run in p.runs:
                    style_ppt_run(run, size=12, color="FFFFFF")


def build_pptx() -> None:
    prs = Presentation(str(PPTX_PATH))

    if any(
        "Why MIN Pooling for MELD" in (shp.text_frame.text if getattr(shp, "has_text_frame", False) else "")
        for slide in prs.slides
        for shp in slide.shapes
    ):
        renumber_titles(prs)
        renumber_slide_numbers(prs)
        prs.save(str(PPTX_PATH))
        print(f"Updated numbering only in {PPTX_PATH}")
        return

    slide = insert_slide_before(prs, 5)  # insert before current Slide 6

    add_box(slide, 0, 0, 13.333, 0.58, "122F55")
    add_box(slide, 0, 0.58, 13.333, 0.08, "267377")
    add_text(
        slide,
        0.42,
        0.12,
        11.4,
        0.3,
        "6. Why MIN Pooling for MELD?",
        size=17.5,
        bold=True,
        color="FFFFFF",
    )
    add_text(slide, 12.0, 0.12, 0.8, 0.3, "6", size=12, color="FFFFFF", align=PP_ALIGN.RIGHT)
    add_text(
        slide,
        0.5,
        0.72,
        12.2,
        0.42,
        "MIN is not arbitrary here: it is the pooling choice used by the MemoCMT MELD reference path, and the saved Phase 1 results show it as the strongest stable cross-check on MELD.",
        size=10.5,
        color="5E5E5E",
    )

    # Short evidence strip.
    left = add_box(slide, 0.55, 1.24, 6.0, 0.72, "EAF4F8", line="9EC3D1", radius=True)
    add_text(
        slide,
        0.75,
        1.38,
        5.65,
        0.38,
        "Base paper cross-check\nMemoCMT MELD CMT+MIN: 64.18% Acc, 62.52% F1",
        size=11.0,
        bold=True,
        color="122F55",
    )
    right = add_box(slide, 6.7, 1.24, 6.1, 0.72, "F7FBF5", line="B8D7BA", radius=True)
    add_text(
        slide,
        6.9,
        1.38,
        5.7,
        0.38,
        "Phase 1 paper-aligned MIN run\nAcc 0.6247 | W-F1 0.6195 | Macro F1 0.4395 | U-Acc 0.4417",
        size=11.0,
        bold=True,
        color="2F5A38",
    )

    rows = [
        ["MIN", "Weakest per-dimension activation after fusion", "0.6247 / 0.6195 / 0.4395 / 0.4417", "Paper-aligned anchor and strongest saved MELD aggregate"],
        ["MEAN", "Average fused token evidence", "0.6199 / 0.6103 / 0.4264 / 0.4209", "Close, but slightly below the MIN anchor"],
        ["MAX", "Strongest fused token activation", "0.6180 / 0.6106 / 0.4272 / 0.4238", "Useful comparator, but not the chosen MELD anchor"],
        ["CLS", "First valid token as summary vector", "0.6149 / 0.6045 / 0.4217 / 0.4144", "Saved case-study variant, but weakest of the four"],
    ]
    headers = ["Pooling", "What it summarizes", "Saved metric snapshot", "Interpretation"]

    table = slide.shapes.add_table(5, 4, PptInches(0.55), PptInches(2.15), PptInches(12.15), PptInches(3.1)).table
    widths = [1.1, 3.0, 3.2, 4.85]
    for idx, w in enumerate(widths):
        table.columns[idx].width = PptInches(w)

    for c, txt in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = txt
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                style_ppt_run(run, size=10.5, bold=True, color="FFFFFF")
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptRGBColor.from_string("122F55")

    for r, row in enumerate(rows, start=1):
        for c, txt in enumerate(row):
            cell = table.cell(r, c)
            cell.text = txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptRGBColor.from_string("F8FAFC" if r % 2 else "EEF3F8")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c else PP_ALIGN.CENTER
                for run in p.runs:
                    style_ppt_run(run, size=9.5, color="1F1F1F", bold=(c == 0))

    add_text(
        slide,
        0.6,
        5.45,
        12.0,
        0.55,
        "Why MIN? It matches the paper-aligned MemoCMT MELD setup, keeps the reference comparison honest, and the saved Phase 1 metrics make it the most defensible aggregation choice in this repository.",
        size=10.0,
        color="5E5E5E",
        align=PP_ALIGN.LEFT,
    )
    add_text(
        slide,
        0.6,
        6.05,
        12.0,
        0.3,
        "Important caveat: MIN is benchmark-specific here. The project uses it because MELD and the MemoCMT baseline are being compared on the same design axis, not because MIN is always better for every task.",
        size=9.4,
        color="6B7280",
    )

    renumber_titles(prs)
    renumber_slide_numbers(prs)
    prs.save(str(PPTX_PATH))
    print(f"Inserted MIN pooling slide and renumbered {PPTX_PATH}")


def find_heading(doc: Document, text: str):
    for p in doc.paragraphs:
        if p.text.strip() == text:
            return p
    return None


def remove_paragraph(paragraph: Paragraph) -> None:
    paragraph._element.getparent().remove(paragraph._element)


def block_between(doc: Document, start_text: str, end_text: str):
    paragraphs = doc.paragraphs
    start = end = None
    for idx, p in enumerate(paragraphs):
        txt = p.text.strip()
        if start is None and txt == start_text:
            start = idx
        elif start is not None and txt == end_text:
            end = idx
            break
    if start is None or end is None or end <= start:
        return []
    return paragraphs[start + 1 : end]


def renumber_docx_headings(doc: Document, preserve_heading: str, start_from: int = 6) -> None:
    seen_new = False
    pattern = re.compile(r"^Slide (\d+):\s*(.*)$")
    for p in doc.paragraphs:
        txt = p.text.strip()
        if txt == preserve_heading:
            seen_new = True
            continue
        m = pattern.match(txt)
        if not m:
            continue
        num = int(m.group(1))
        if num >= start_from and seen_new:
            base = m.group(2)
            p.text = f"Slide {num + 1}: {base}"


def update_docx() -> None:
    doc = Document(str(DOCX_PATH))
    new_heading = "Slide 6: Why MIN Pooling for MELD"
    anchor = find_heading(doc, "Slide 6: MemoCMT vs Phase 1 Implementation Mapping")
    if anchor is None:
        raise RuntimeError("Could not find the existing Slide 6 heading in the reading script DOCX.")

    desired_lines = [
        "Use this slide to explain why MIN is the pooling choice for the paper-aligned MELD path. The key point is that MIN is not an arbitrary preference. It is the choice that keeps the implementation aligned with the MemoCMT reference design and gives the strongest stable MELD summary in the repository.",
        "Start with the base-paper cross-check: MemoCMT reports MELD CMT+MIN test results of 64.18% accuracy and 62.52% F1. That gives you the benchmark anchor. Then move to the repository numbers: the paper-aligned 5-fold MIN run reports mean accuracy 0.6247, weighted F1 0.6195, macro F1 0.4395, and unweighted accuracy 0.4417. Those values show that Phase 1 is tracking the paper in the right range.",
        "Then compare the saved case-study pooling variants. CLS is 0.6149 accuracy, 0.6045 weighted F1, and 0.4217 macro F1. MEAN is 0.6199 accuracy, 0.6103 weighted F1, and 0.4264 macro F1. MAX is 0.6180 accuracy, 0.6106 weighted F1, and 0.4272 macro F1. The spread is not huge, but the saved MIN run is the cleanest paper-aligned anchor and the strongest stable conversational summary in the repository.",
        "Explain that the comparison is benchmark-specific. You are not claiming that MIN is universally superior for every task. You are saying that for the MemoCMT-style MELD setting, MIN is the safest choice because it cross-checks best with the base paper and preserves comparability across the Phase 1 story.",
        "A viva-friendly line is: MIN was chosen because it matches the MemoCMT MELD reference design and the saved Phase 1 results show it as the most defensible pooling choice for the paper-aligned conversational baseline.",
    ]

    existing_new_heading = find_heading(doc, new_heading)
    if existing_new_heading is not None:
        between = block_between(doc, new_heading, "Slide 7: MemoCMT Paper vs Phase 1 Implementation Mapping")
        for p in between:
            remove_paragraph(p)
        cur = existing_new_heading
        for line in desired_lines:
            cur = insert_paragraph_after(cur, line)
        renumber_docx_headings(doc, new_heading)
        doc.save(str(DOCX_PATH))
        print(f"Rebuilt MIN pooling section in {DOCX_PATH}")
        return

    # Insert the new section before the existing Slide 6 section.
    current = anchor
    p = insert_paragraph_before(current, new_heading, style="Heading 1")
    rr = p.runs[0]
    rr.font.name = "Aptos"
    rr.font.size = Pt(14)
    rr.font.bold = True
    rr.font.color.rgb = RGBColor.from_string("122F55")
    current = p
    for line in desired_lines:
        current = insert_paragraph_after(current, line)

    renumber_docx_headings(doc, new_heading)
    doc.save(str(DOCX_PATH))
    print(f"Inserted MIN pooling section into {DOCX_PATH}")


def main() -> None:
    build_pptx()
    update_docx()


if __name__ == "__main__":
    main()
