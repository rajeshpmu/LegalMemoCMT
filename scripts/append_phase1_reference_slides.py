from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(20, 48, 87)
SUB_COLOR = RGBColor(72, 72, 72)
BODY_COLOR = RGBColor(35, 35, 35)
BAR_COLOR = RGBColor(20, 48, 87)
WHITE = RGBColor(255, 255, 255)


PHASE1_SLIDES = {
    "Zeroth_Review_LegalMemoCMT.pptx": [
        (
            "13. Phase 1 Implementation References",
            "Core papers that directly support the current Phase 1 architecture and benchmark path.",
            [
                ["Base paper (primary reference)", "MemoCMT: multimodal emotion recognition using cross-modal transformer-based feature fusion", "https://www.nature.com/articles/s41598-025-89202-x", "Main Phase 1 architecture reference and direct source for the CMT design."],
                ["Trimodal baseline", "Multimodal Emotion Recognition Based on Facial Expressions, Speech, and Body Gestures", "https://www.mdpi.com/2079-9292/13/18/3756", "Strong trimodal baseline for face, speech, and body-gesture fusion."],
                ["Speech review", "Speech emotion recognition in conversations using artificial intelligence: a systematic review and meta-analysis", "https://link.springer.com/article/10.1007/s10462-025-11197-8", "Supports the speech branch, metrics, and conversational evaluation framing."],
                ["Conversation fusion", "Multimodal Emotion Recognition in Conversations Using Transformer and Graph Neural Networks", "https://www.mdpi.com/2076-3417/15/22/11971", "Dialogue structure, transformer fusion, and graph reasoning."],
                ["Survey backbone", "Recent Trends of Multimodal Affective Computing: A Survey from NLP Perspective", "https://arxiv.org/abs/2409.07388", "Task taxonomy, datasets, and metric framing for the thesis."],
                ["Broader system", "A Multimodal Emotion Recognition System: Integrating Facial Expressions, Body Movement, Speech, and Spoken Language", "https://arxiv.org/abs/2412.17907", "Applied multimodal example beyond the MemoCMT core."],
            ],
        ),
        (
            "14. Phase 1 Supporting References",
            "Supporting papers that help with alternative fusion ideas and lightweight comparison baselines.",
            [
                ["Capsule/graph fusion", "Meaningful Multimodal Emotion Recognition Based on Capsule Graph Transformer Architecture", "https://www.mdpi.com/2078-2489/16/1/40", "Useful for graph-transformer comparisons and richer fusion design."],
                ["Dual attention fusion", "HyFusER: Hybrid Multimodal Transformer for Emotion Recognition Using Dual Cross Modal Attention", "https://www.mdpi.com/2076-3417/15/3/1053", "Useful for dual attention and hybrid transformer comparisons."],
                ["Speech-text fusion", "Multimodal Affective Communication Analysis: Fusing Speech Emotion and Text Sentiment Using Machine Learning", "https://www.mdpi.com/2076-3417/14/15/6631", "Good for simpler speech-plus-text baseline framing."],
                ["Adaptive exchange fusion", "Adaptive multimodal transformer based on exchanging for multimodal sentiment analysis", "https://www.nature.com/articles/s41598-025-11848-4", "Shows stronger adaptive fusion than fixed pooling baselines."],
                ["Hierarchical attention", "Hierarchical cross-modal attention and dual audio pathways for enhanced multimodal sentiment analysis", "https://www.nature.com/articles/s41598-025-09000-3", "Useful for hierarchical attention and audio-path design ideas."],
                ["Conversation survey", "Multimodal Emotion Recognition in Conversations: A Survey of Methods, Trends, Challenges and Prospects", "https://aclanthology.org/2025.findings-emnlp.332/", "Newest high-level MERC survey for background and comparison."],
            ],
        ),
    ],
    "First_Review_LegalMemoCMT.pptx": [
        (
            "14. Phase 1 Implementation References",
            "Core papers that directly support the current Phase 1 architecture and benchmark path.",
            [
                ["Base paper (primary reference)", "MemoCMT: multimodal emotion recognition using cross-modal transformer-based feature fusion", "https://www.nature.com/articles/s41598-025-89202-x", "Main Phase 1 architecture reference and direct source for the CMT design."],
                ["Trimodal baseline", "Multimodal Emotion Recognition Based on Facial Expressions, Speech, and Body Gestures", "https://www.mdpi.com/2079-9292/13/18/3756", "Strong trimodal baseline for face, speech, and body-gesture fusion."],
                ["Speech review", "Speech emotion recognition in conversations using artificial intelligence: a systematic review and meta-analysis", "https://link.springer.com/article/10.1007/s10462-025-11197-8", "Supports the speech branch, metrics, and conversational evaluation framing."],
                ["Conversation fusion", "Multimodal Emotion Recognition in Conversations Using Transformer and Graph Neural Networks", "https://www.mdpi.com/2076-3417/15/22/11971", "Dialogue structure, transformer fusion, and graph reasoning."],
                ["Survey backbone", "Recent Trends of Multimodal Affective Computing: A Survey from NLP Perspective", "https://arxiv.org/abs/2409.07388", "Task taxonomy, datasets, and metric framing for the thesis."],
                ["Broader system", "A Multimodal Emotion Recognition System: Integrating Facial Expressions, Body Movement, Speech, and Spoken Language", "https://arxiv.org/abs/2412.17907", "Applied multimodal example beyond the MemoCMT core."],
            ],
        ),
        (
            "15. Phase 1 Supporting References",
            "Supporting papers that help with alternative fusion ideas and lightweight comparison baselines.",
            [
                ["Capsule/graph fusion", "Meaningful Multimodal Emotion Recognition Based on Capsule Graph Transformer Architecture", "https://www.mdpi.com/2078-2489/16/1/40", "Useful for graph-transformer comparisons and richer fusion design."],
                ["Dual attention fusion", "HyFusER: Hybrid Multimodal Transformer for Emotion Recognition Using Dual Cross Modal Attention", "https://www.mdpi.com/2076-3417/15/3/1053", "Useful for dual attention and hybrid transformer comparisons."],
                ["Speech-text fusion", "Multimodal Affective Communication Analysis: Fusing Speech Emotion and Text Sentiment Using Machine Learning", "https://www.mdpi.com/2076-3417/14/15/6631", "Good for simpler speech-plus-text baseline framing."],
                ["Adaptive exchange fusion", "Adaptive multimodal transformer based on exchanging for multimodal sentiment analysis", "https://www.nature.com/articles/s41598-025-11848-4", "Shows stronger adaptive fusion than fixed pooling baselines."],
                ["Hierarchical attention", "Hierarchical cross-modal attention and dual audio pathways for enhanced multimodal sentiment analysis", "https://www.nature.com/articles/s41598-025-09000-3", "Useful for hierarchical attention and audio-path design ideas."],
                ["Conversation survey", "Multimodal Emotion Recognition in Conversations: A Survey of Methods, Trends, Challenges and Prospects", "https://aclanthology.org/2025.findings-emnlp.332/", "Newest high-level MERC survey for background and comparison."],
            ],
        ),
    ],
}


def add_bg(slide, prs):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BAR_COLOR
    bar.line.fill.background()


def add_title(slide, title, subtitle):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(0.8), Inches(12.1), Inches(0.55))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Aptos"
    r2.font.size = Pt(11.5)
    r2.font.italic = True
    r2.font.color.rgb = SUB_COLOR


def add_table_slide(prs, title, subtitle, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)

    table = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.35), Inches(1.35), Inches(12.6), Inches(5.8)).table
    widths = [1.35, 4.0, 3.0, 4.25]
    for idx, w in enumerate(widths):
        table.columns[idx].width = Inches(w)

    headers = ["Phase 1 role", "Article", "Link", "Why it matters"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Aptos"
                r.font.size = Pt(10)
                r.font.bold = True
                r.font.color.rgb = TITLE_COLOR

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.space_after = Pt(0)
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(8.4)
                    run.font.color.rgb = BODY_COLOR
        if r_idx % 2 == 0:
            for c_idx in range(4):
                table.cell(r_idx, c_idx).fill.solid()
                table.cell(r_idx, c_idx).fill.fore_color.rgb = RGBColor(245, 247, 250)


def main():
    for name, slidespecs in PHASE1_SLIDES.items():
        path = Path("implementation_docments") / name
        prs = Presentation(str(path))
        for title, subtitle, rows in slidespecs:
            add_table_slide(prs, title, subtitle, rows)
        prs.save(str(path))
        print("updated", path)


if __name__ == "__main__":
    main()
