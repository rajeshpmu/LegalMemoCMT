from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT_DIR = Path("implementation_docments")
SLIDE_WIDE = (13.333, 7.5)


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(SLIDE_WIDE[0])
    prs.slide_height = Inches(SLIDE_WIDE[1])


def add_bg(slide, prs: Presentation):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None, y=0.35):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(y + 0.52), Inches(12.1), Inches(0.55))
        tf2 = tx2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x=0.75, y=1.35, w=11.95, h=5.65, font_size=18, indent=0.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.08
        p.left_margin = Inches(indent)
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_diagram(slide, image_path: Path, left: float, top: float, width: float):
    if not image_path.exists():
        return None
    return slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def add_numbered(slide, bullets, **kwargs):
    add_body(slide, [f"{i+1}. {b}" for i, b in enumerate(bullets)], **kwargs)


def add_title_slide(prs: Presentation, title: str, deck_tag: str, project_title: str, scope: str, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)

    # Main title
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(12.2), Inches(1.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)

    # Deck tag
    tx2 = slide.shapes.add_textbox(Inches(0.62), Inches(1.45), Inches(12.0), Inches(0.4))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = deck_tag
    r2.font.name = "Aptos"
    r2.font.size = Pt(14)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(66, 66, 66)

    # Project title and scope
    tx3 = slide.shapes.add_textbox(Inches(0.62), Inches(1.95), Inches(12.0), Inches(1.15))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    for i, line in enumerate([project_title, scope]):
        p3 = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p3.space_after = Pt(3)
        r3 = p3.add_run()
        r3.text = line
        r3.font.name = "Aptos"
        r3.font.size = Pt(15 if i == 0 else 12.5)
        r3.font.color.rgb = RGBColor(50, 50, 50)
        if i == 0:
            r3.font.bold = True

    # Bullets box
    add_body(slide, bullets, x=0.95, y=3.25, w=11.55, h=3.5, font_size=17)

    return slide


def add_content_slide(prs: Presentation, title: str, subtitle: str, bullets, font_size=17):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    add_body(slide, bullets, x=0.75, y=1.48, w=11.85, h=5.55, font_size=font_size)
    return slide


def add_two_column_slide(prs: Presentation, title: str, subtitle: str, left_title: str, left_bullets, right_title: str, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)

    for x, heading, bullets in [(0.6, left_title, left_bullets), (6.8, right_title, right_bullets)]:
        h = slide.shapes.add_textbox(Inches(x), Inches(1.5), Inches(5.6), Inches(0.35))
        tf = h.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = heading
        r.font.name = "Aptos"
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = RGBColor(20, 48, 87)
        add_body(slide, bullets, x=x, y=1.92, w=5.65, h=4.95, font_size=15)

    return slide


def add_text_diagram_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    bullets,
    image_path: Path,
    diagram_caption: str | None = None,
    font_size: int = 16,
    text_box=(0.7, 1.45, 6.0, 5.55),
    image_box=(7.0, 1.7, 5.7),
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)
    add_body(
        slide,
        bullets,
        x=text_box[0],
        y=text_box[1],
        w=text_box[2],
        h=text_box[3],
        font_size=font_size,
    )
    add_diagram(slide, image_path, image_box[0], image_box[1], image_box[2])
    if diagram_caption:
        tx = slide.shapes.add_textbox(Inches(image_box[0]), Inches(6.55), Inches(image_box[2]), Inches(0.35))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = diagram_caption
        r.font.name = "Aptos"
        r.font.size = Pt(10.5)
        r.font.italic = True
        r.font.color.rgb = RGBColor(95, 95, 95)
    return slide


def build_zeroth_presentation(output: Path):
    phase2_diagram = Path("artifacts/legalmemocmt_phase2.png")
    fusion_diagram = Path("artifacts/mermaid/paper_aligned_fusion_study.png")
    prs = Presentation()
    set_slide_size(prs)

    add_title_slide(
        prs,
        "Zeroth Review",
        "MTech Final Project",
        "Project Title: LegalMemoCMT: An Explainable Multilingual Multimodal Emotional Cue Analysis Framework for Indian Courtroom Testimony Using Cross-Modal Transformers",
        "Research Scope: The project focuses only on observable emotional cues from face, speech, and transcript. It does not predict guilt, innocence, intent, or deception.",
        [
            "Two-phase research plan: Phase 1 is the implemented benchmark scaffold, and Phase 2 is the longer-term courtroom extension.",
            "The current repository already includes a legacy baseline, a pretrained/paper-aligned path, and evaluation tooling.",
            "The deck emphasizes what is implemented now, what remains planned, and what must stay outside the legal decision boundary.",
        ],
    )

    add_content_slide(
        prs,
        "1. Introduction",
        "The introduction frames multimodal emotion analysis as the core ML problem and courtroom testimony as the high-stakes extension.",
        [
            "Emotion analysis is a core problem in multimodal machine learning because affective cues appear across speech, facial expression, and language.",
            "Courtroom testimony is a high-stakes application, so the project is intentionally limited to observable emotional state and does not attempt truth verification or legal inference.",
            "The current repository already supports that direction through the transformer baseline, pretrained/paper path, manifest loading, preprocessing utilities, training scripts, and prediction analysis tools.",
        ],
    )

    add_content_slide(
        prs,
        "2. Literature Survey",
        "The literature base stays focused on multimodal emotion recognition, speech emotion recognition, transformer fusion, and affective computing surveys.",
        [
            "The conceptual references remain valid for the thesis narrative.",
            "The implementation status in the repository has moved beyond a simple plan and now includes a working multimodal pipeline.",
            "The thesis should position the literature as justification for the current architecture and for the future courtroom adaptation.",
        ],
    )

    add_content_slide(
        prs,
        "3. Current Phase 1 Implementation",
        "This slide summarizes the present benchmark paths and the experimental controls already available in the codebase.",
        [
            "Phase 1 now supports MELD and CREMA-D as the main benchmark paths.",
            "IEMOCAP support exists in the codebase through a manifest builder, but it remains deferred in the project plan and is not the current primary benchmark.",
            "The model supports legacy transformer encoders and a pretrained/paper mode using HuggingFace backbones for text and speech.",
            "Training supports cross-entropy, weighted cross-entropy, and focal loss, plus modality subsets for ablation studies.",
            "Evaluation supports accuracy, weighted accuracy, unweighted accuracy, macro F1, weighted F1, JSON metrics export, and per-sample prediction CSV export.",
        ],
        font_size=16,
    )

    add_text_diagram_slide(
        prs,
        "4. Proposed System: Phase 1 and Phase 2",
        "The review separates the implemented benchmark path from the future courtroom adaptation.",
        [
            "Phase 1 is already a usable benchmark implementation with manifest-driven data loading, extracted-feature fallback, pretrained/paper-aligned inference, and held-out analysis.",
            "Phase 2 remains the courtroom extension concept: multilingual Indian testimony, interpretability support, and a strict affective-analysis-only scope.",
            "The diagram on the right makes the intended Phase 1 to Phase 2 progression explicit.",
        ],
        phase2_diagram,
        diagram_caption="Phase 1 to Phase 2 LegalMemoCMT direction",
        font_size=15,
        text_box=(0.7, 1.55, 6.1, 5.25),
        image_box=(7.0, 1.75, 5.8),
    )

    add_content_slide(
        prs,
        "5. Timeline and Feasibility",
        "The project is feasible because the data plumbing and model scaffolding already exist.",
        [
            "The remaining work is mainly benchmark execution, result curation, reporting, and any further courtroom-domain extension rather than core implementation from scratch.",
            "The Phase 1 path is already inspectable and reproducible, which lowers the risk for review and revision.",
            "Phase 2 can be treated as a scoped research extension instead of a full rebuild.",
        ],
    )

    add_content_slide(
        prs,
        "6. References",
        "This slide tells the committee how the reference list should be updated relative to the current implementation.",
        [
            "Keep the multimodal emotion recognition, speech emotion review, transformer fusion, and affective computing references.",
            "Update dataset references to reflect MELD, CREMA-D, and the deferred status of IEMOCAP in the implementation plan.",
            "The reference narrative should now support the implemented pipeline and the planned courtroom extension, not only the original proposal.",
        ],
    )

    add_text_diagram_slide(
        prs,
        "7. Operational Details",
        "This slide collects the appendix-style explanation that reviewers usually ask about during a presentation.",
        [
            "The fusion module first creates three modality vectors, one each for text, audio, and video, then applies a Transformer over that three-token sequence.",
            "cls keeps the first fused position as the summary vector and effectively makes the text-stream position the anchor representation after fusion.",
            "mean averages the three fused vectors and treats the modalities as equally important contributors.",
            "max and min are controlled comparison variants that expose different ways of compressing cross-modal evidence.",
            "The diagram on the right also reminds the committee that pooling is downstream of the fusion block, not a separate model.",
        ],
        fusion_diagram,
        diagram_caption="Fusion study variants used for pooling discussion",
        font_size=15,
        text_box=(0.7, 1.55, 6.15, 5.25),
        image_box=(6.95, 1.95, 5.95),
    )

    add_content_slide(
        prs,
        "8. Review Takeaways",
        "This final slide states the main committee-facing message for the Zeroth Review.",
        [
            "Phase 1 is now a usable benchmark implementation rather than only a paper plan.",
            "Phase 2 remains the courtroom extension concept with legal-domain safeguards and an affective-analysis-only boundary.",
            "The appendix-style implementation notes on pooling and manifest building should stay visible because they explain how the pipeline is actually wired.",
            "The thesis should keep the emphasis on explainable emotional interpretation, not legal judgment.",
        ],
    )

    prs.save(output)


def build_first_presentation(output: Path):
    architecture_diagram = Path("artifacts/mermaid/paper_aligned_architecture.png")
    fusion_diagram = Path("artifacts/mermaid/paper_aligned_fusion_study.png")
    prs = Presentation()
    set_slide_size(prs)

    add_title_slide(
        prs,
        "First Review",
        "MTech Final Project",
        "Project Title: LegalMemoCMT: An Explainable Multilingual Multimodal Emotional Cue Analysis Framework for Indian Courtroom Testimony Using Cross-Modal Transformers",
        "Research Scope: The project focuses only on observable emotional cues from face, speech, and transcript. It does not predict guilt, innocence, intent, or deception.",
        [
            "This review reflects the implemented system as well as the remaining thesis direction.",
            "Phase 1 has a working codebase for multimodal emotion recognition, including a legacy baseline and a pretrained/paper-aligned path.",
            "Phase 2 remains the courtroom-adaptation goal.",
        ],
    )

    add_text_diagram_slide(
        prs,
        "1. Architecture Design for the Current System",
        "The architecture is split into the implemented Phase 1 path and the planned Phase 2 courtroom extension.",
        [
            "Phase 1 is already a working multimodal pipeline with legacy and pretrained/paper-aligned branches.",
            "The current code supports transcript, audio, and video inputs, plus learned modality fusion and configurable pooling.",
            "Phase 2 extends the same design toward courtroom testimony with multilingual encoders, visual face cues, and explainability outputs.",
            "The architecture on the right matches the actual repository structure and the paper-aligned path.",
        ],
        architecture_diagram,
        diagram_caption="Paper-aligned architecture and fusion flow",
        font_size=15,
        text_box=(0.7, 1.55, 6.05, 5.25),
        image_box=(6.95, 1.72, 5.95),
    )

    add_content_slide(
        prs,
        "2. Algorithms and Techniques",
        "This slide compresses the current implementation stack into one operational view.",
        [
            "The implemented Phase 1 stack uses manifest-based preprocessing, sequence modeling, transformer fusion, and classification loss with optional imbalance handling.",
            "The current scripts support split-aware training, held-out evaluation, confusion analysis, and per-sample prediction export.",
            "Preprocessing covers transcript normalization, audio feature extraction or waveform loading, and video frame/feature handling.",
            "Fusion uses modality-token Transformer integration with configurable pooling.",
            "Training uses cross-entropy, weighted cross-entropy, or focal loss, and evaluation includes accuracy, macro F1, weighted F1, weighted accuracy, confusion analysis, and prediction CSVs.",
        ],
        font_size=16,
    )

    add_content_slide(
        prs,
        "3. Expected Outcomes",
        "The expected outcomes are now more concrete because the Phase 1 pipeline is already operating.",
        [
            "Phase 1 is now expected to support benchmark experiments on MELD and CREMA-D, plus a deferred IEMOCAP path in the codebase if needed later.",
            "Phase 2 still aims to demonstrate transfer to multilingual courtroom testimony with interpretable affective traces, but that remains a future extension rather than the current implementation state.",
            "The final system should provide interpretable emotional summaries rather than legal conclusions.",
            "Single-modality and multimodal fusion comparisons remain important for reporting.",
        ],
    )

    add_content_slide(
        prs,
        "4. Hardware Requirements and Local Feasibility",
        "This slide is about what can be run locally and what will still need more compute.",
        [
            "The current implementation is compatible with local development on a Mac for preprocessing, smoke tests, small runs, and script validation.",
            "Larger pretrained experiments are still the more expensive part of the pipeline and may benefit from a GPU workstation or cloud environment.",
            "The practical message for review is that the repository is already usable for controlled experimentation.",
        ],
    )

    add_content_slide(
        prs,
        "5. Datasets You Can Use",
        "The dataset section should clearly separate the active benchmark pair from the deferred path.",
        [
            "MELD for conversational multimodal emotion recognition.",
            "CREMA-D as the current public replacement benchmark for the deferred IEMOCAP slot.",
            "IEMOCAP remains supported in the repository through a manifest builder, but it is still deferred in the project plan.",
        ],
    )

    add_content_slide(
        prs,
        "6. Ethical and Legal Scope",
        "This review needs a clear ethical boundary because the project is explicitly aimed at courtroom-adjacent analysis.",
        [
            "The system must not infer guilt, innocence, deception, credibility, or intent.",
            "It should only analyze observable emotional cues and temporal affective patterns.",
            "That limitation should remain explicit in the thesis and in any future courtroom adaptation.",
        ],
    )

    add_content_slide(
        prs,
        "7. References",
        "This slide explains how the bibliography should support the now-implemented Phase 1 pipeline.",
        [
            "Keep the multimodal emotion recognition and survey references.",
            "Update the narrative so the citations support the now-implemented Phase 1 pipeline and the still-proposed courtroom adaptation.",
            "Use the references to justify the architecture, the fusion strategy, and the evaluation plan.",
        ],
    )

    add_text_diagram_slide(
        prs,
        "8. Appendices",
        "The last slide collects the appendix-style details that are often discussed after the main architecture slide.",
        [
            "The fusion module first creates three modality vectors, one each for text, audio, and video, then applies a Transformer over that three-token sequence.",
            "The pooling choice determines how those three fused modality vectors are collapsed into a single representation before classification.",
            "cls makes the text-stream position the anchor representation after fusion.",
            "mean, max, and min remain controlled comparison variants that change how cross-modal information is compressed.",
            "The manifest builder still matters because it normalizes dataset-specific folders into the shared CSV interface used by training and evaluation.",
        ],
        fusion_diagram,
        diagram_caption="Pooling comparison and held-out evaluation flow",
        font_size=15,
        text_box=(0.7, 1.55, 6.1, 5.25),
        image_box=(6.95, 1.95, 5.95),
    )

    prs.save(output)


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_zeroth_presentation(OUT_DIR / "Zeroth_Review_LegalMemoCMT.pptx")
    build_first_presentation(OUT_DIR / "First_Review_LegalMemoCMT.pptx")
    print("Wrote Zeroth_Review_LegalMemoCMT.pptx and First_Review_LegalMemoCMT.pptx")


if __name__ == "__main__":
    main()
