from __future__ import annotations

import copy
import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE_DECK = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Presentation.pptx"
OUTPUT = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
CLANCY = ROOT / "data/processed/phase2/clancy/clancy_turn_manifest_post_rejection.csv"
TRIBUNAL = ROOT / "data/processed/phase2/legalmeld_validated/legalmeld_metadata_validated.csv"
TUPAC_URLS = ROOT / "data/tupac_trial_urls.txt"

NAVY = RGBColor(20, 48, 87)
TEXT = RGBColor(35, 35, 35)
GREY = RGBColor(95, 95, 95)


def rows(path: Path) -> list[dict[str, str]]:
    if not path.exists():
        return []
    with path.open(newline="", encoding="utf-8-sig") as f:
        return list(csv.DictReader(f))


def duration(row: dict[str, str]) -> float:
    for key in ("clip_duration_seconds", "duration_seconds"):
        try:
            return float(row.get(key) or 0)
        except ValueError:
            pass
    return 0.0


def duration_stats(items: list[dict[str, str]]) -> dict[str, float | int]:
    values = sorted(duration(row) for row in items)
    if not values:
        return {"rows": 0, "hours": 0, "minutes": 0, "min": 0, "max": 0, "under": 0, "over": 0}
    return {
        "rows": len(values),
        "hours": round(sum(values) / 3600, 3),
        "minutes": round(sum(values) / 60, 2),
        "min": round(values[0], 2),
        "max": round(values[-1], 2),
        "under": sum(value < 0.8 for value in values),
        "over": sum(value > 30 for value in values),
    }


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def bg(slide, prs: Presentation) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()


def title(slide, text: str, subtitle: str = "") -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(12.2), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.name = "Aptos Display"; run.font.size = Pt(24); run.font.bold = True; run.font.color.rgb = NAVY
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.58), Inches(0.95), Inches(12.1), Inches(0.42))
        p2 = box2.text_frame.paragraphs[0]
        run2 = p2.add_run(); run2.text = subtitle
        run2.font.name = "Aptos"; run2.font.size = Pt(12.5); run2.font.italic = True; run2.font.color.rgb = GREY


def bullets(slide, items: list[str], x=0.75, y=1.42, w=11.8, h=5.65, size=16) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(5); p.line_spacing = 1.02
        run = p.add_run(); run.text = "• " + item
        run.font.name = "Aptos"; run.font.size = Pt(size); run.font.color.rgb = TEXT


def slide_with_bullets(prs, heading, subtitle, items, size=16):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide, prs); title(slide, heading, subtitle); bullets(slide, items, size=size); return slide


def table_slide(prs, heading, subtitle, headers, data, size=12):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide, prs); title(slide, heading, subtitle)
    table = slide.shapes.add_table(len(data) + 1, len(headers), Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.75)).table
    for c, value in enumerate(headers): table.cell(0, c).text = value
    for r, row in enumerate(data, 1):
        for c, value in enumerate(row): table.cell(r, c).text = str(value)
    for r_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
            tf = cell.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = "Aptos"; run.font.size = Pt(size); run.font.color.rgb = NAVY if r_idx == 0 else TEXT
                    if r_idx == 0: run.font.bold = True
    return slide


def two_col(prs, heading, subtitle, left_heading, left, right_heading, right, size=14):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide, prs); title(slide, heading, subtitle)
    for x, h, items in ((0.6, left_heading, left), (6.75, right_heading, right)):
        box = slide.shapes.add_textbox(Inches(x), Inches(1.38), Inches(5.8), Inches(0.35))
        p = box.text_frame.paragraphs[0]; r = p.add_run(); r.text = h; r.font.name = "Aptos"; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY
        bullets(slide, items, x=x, y=1.82, w=5.75, h=4.95, size=size)
    return slide


def clone_slide(src_prs: Presentation, dst_prs: Presentation, index: int):
    source = src_prs.slides[index - 1]
    slide = dst_prs.slides.add_slide(dst_prs.slide_layouts[6])
    for shape in source.shapes:
        slide.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")
    for rel in source.part.rels.values():
        if "notesSlide" not in rel.reltype and rel.is_external is False:
            try:
                slide.part.rels.get_or_add(rel.reltype, rel._target)
            except Exception:
                pass
    return slide


def replace_text(prs: Presentation, old: str, new: str) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def build() -> None:
    tribunal_rows = rows(TRIBUNAL)
    clancy_rows = rows(CLANCY)
    tribunal_stats = duration_stats(tribunal_rows)
    clancy_stats = duration_stats(clancy_rows)
    clancy_window = [r for r in clancy_rows if 0.8 <= duration(r) <= 30]
    clancy_window_stats = duration_stats(clancy_window)
    clancy_emotions = {}
    for row in clancy_window:
        label = (row.get("emotion_label") or "BLANK").strip() or "BLANK"
        clancy_emotions[label] = clancy_emotions.get(label, 0) + 1
    tupac_count = sum(1 for line in TUPAC_URLS.read_text().splitlines() if line.strip()) if TUPAC_URLS.exists() else 0

    src = Presentation(SOURCE_DECK)
    prs = Presentation(); set_slide_size(prs)
    for index in (1, 4, 5, 11):
        clone_slide(src, prs, index)
    replace_text(prs, "First Guidance Call", "Second Guidance Call")
    replace_text(prs, "first guidance call", "second guidance call")

    slide_with_bullets(prs, "Second Guidance Call: Corpus Expansion", "The discussion has moved from proving the pipeline to comparing corpus branches and controlling the next expansion step.", [
        "The first guidance call established the verification-first and utterance-level design.",
        "This call focuses on the actual corpus branches: IRMCT / ICTY, Lindsay Clancy, Tupac / Keffe D, and the planned Indian-SIM branch.",
        "The EDA separates what is currently processed from what is only a source plan.",
        "Duration and outlier checks are used to decide what can enter a MELD-style training set.",
    ], size=16)

    table_slide(prs, "LegalMemoCMT Courtroom Corpus", "Four branches with different roles in the research design.", ["Corpus branch", "Current role", "Witness / speech coverage", "Implementation status"], [
        ["IRMCT / ICTY", "Secondary bootstrap benchmark", "International tribunal witnesses", f"Processed: {tribunal_stats['rows']} rows; {tribunal_stats['hours']} h candidate"],
        ["Lindsay Clancy", "Primary courtroom benchmark", "Psychiatrists, family, medical witnesses", f"Processed: {clancy_stats['rows']} turn rows; {clancy_stats['hours']} h before duration window"],
        ["Tupac / Keffe D", "Planned expansion", "Hostile, gang-associated, police, forensic/medical witnesses", f"Source list only: {tupac_count} URLs; no processed EDA yet"],
        ["Indian-SIM", "Planned Indian adaptation", "Controlled Indian courtroom testimony", "Not acquired or processed yet"],
    ], size=11)

    table_slide(prs, "Source Sites and Download Flow", "The same evidence chain is used for existing and planned branches; a source URL is not itself a corpus row.", ["Branch", "Source / entry point", "Discovery", "Download / processing", "Evidence status"], [
        ["IRMCT / ICTY", "UCR: https://ucr.irmct.org/; ICTY cases: https://www.icty.org/en/cases", "Case ledger, UCR records, TAP/transcript metadata", "Strict UCR downloader → hearing/witness manifests → LegalMELD", "Processed tribunal pilot"],
        ["Lindsay Clancy", "Public YouTube URLs in data/clancy_urls.txt", "Source manifest and subtitle metadata", "Clancy source builder → turn manifest → ffmpeg turn clips → EDA", "Processed primary benchmark"],
        ["Tupac / Keffe D", "URLs in data/tupac_trial_urls.txt", "Manual source review and witness-type screening", "Planned downloader and same turn-level pipeline", "Source list only; not yet corpus evidence"],
        ["Indian-SIM", "Future public/mock-trial source pack", "Indian candidate ledger and source verification", "Reuse source, hearing, witness, alignment, and clip stages", "Planned adaptation branch"],
    ], size=10)

    table_slide(prs, "IRMCT / ICTY EDA", "Current tribunal branch: a small validated pilot, not a claim of a large finished corpus.", ["Measure", "Current result", "Interpretation"], [
        ["Metadata rows", tribunal_stats["rows"], "Utterance-level rows in legalmeld_metadata_validated.csv"],
        ["Summed candidate duration", f"{tribunal_stats['minutes']} min ({tribunal_stats['hours']} h)", "Current pilot duration"],
        ["Duration range", f"{tribunal_stats['min']}–{tribunal_stats['max']} sec", "No rows above 30 sec in this export"],
        ["Under 0.8 sec / over 30 sec", f"{tribunal_stats['under']} / {tribunal_stats['over']}", "Duration outlier flags for this pilot"],
        ["Source groups", len({r.get('hearing_id') for r in tribunal_rows}), "Hearing-level traceability"],
        ["Readiness", "Controlled pilot", "Alignment and witness visibility still require quality review"],
    ], size=12)

    table_slide(prs, "Lindsay Clancy EDA", "Primary benchmark branch after persistent rejection filtering.", ["Measure", "Current result", "Interpretation"], [
        ["Post-rejection rows", clancy_stats["rows"], "Rows remaining after persistent non-value exclusions"],
        ["Total candidate duration", f"{clancy_stats['minutes']} min ({clancy_stats['hours']} h)", "Includes short rows and long-turn candidates"],
        ["0.8–30 sec window", f"{clancy_window_stats['rows']} rows; {clancy_window_stats['minutes']} min ({clancy_window_stats['hours']} h)", "Initial MELD-style duration screen"],
        ["Below 0.8 sec", clancy_stats["under"], "Review or exclude as very short fragments"],
        ["Above 30 sec", clancy_stats["over"], "Split or manually reject; not direct MELD-style utterances"],
        ["Source videos", len({r.get('youtube_id') for r in clancy_rows}), "Source diversity currently represented"],
    ], size=11)

    emotion_rows = [[label, count, f"{count / len(clancy_window) * 100:.2f}%"] for label, count in sorted(clancy_emotions.items(), key=lambda item: (-item[1], item[0]))]
    table_slide(prs, "Clancy Duration Window and Emotion EDA", "The 0.8–30 second subset is a candidate training pool, while labels remain heuristic and imbalanced.", ["Emotion label", "Rows", "Share of 0.8–30 sec window"], emotion_rows, size=13)

    table_slide(prs, "Tupac / Keffe D: Current Status and EDA Boundary", "This branch should not be presented as processed corpus data yet.", ["Item", "Current state", "Next evidence required"], [
        ["Source URLs", tupac_count, "Open each URL and identify actual witness testimony segments"],
        ["Processed utterance rows", "0", "Download media and obtain aligned transcript/subtitle evidence"],
        ["Duration EDA", "Not available", "Run the same turn-manifest duration and outlier report after processing"],
        ["Witness coverage", "Planned categories", "Confirm hostile, gang-associated, police, and forensic/medical witnesses manually"],
        ["Training use", "Do not train yet", "Promote only after text/audio/video traceability and quality checks"],
    ], size=12)

    table_slide(prs, "Indian-SIM: Planned Adaptation Branch", "The Indian objective remains part of Phase 2, but no statistics are claimed before acquisition and verification.", ["Item", "Planned design", "Current status"], [
        ["Source type", "Controlled Indian courtroom testimony / mock-trial material", "No verified corpus rows yet"],
        ["Speaker coverage", "Witnesses, judge, prosecutor, defence", "To be defined by source pack"],
        ["Duration EDA", "0.8–30 sec first screen; outlier review above 30 sec", "Not available yet"],
        ["Labels", "Manual or adjudicated emotion/credibility labels", "Not available yet"],
        ["Pipeline reuse", "Candidate ledger → verification → manifests → alignment → clips", "Implementation pattern already exists"],
    ], size=12)

    table_slide(prs, "Cross-Corpus Duration and Outlier Comparison", "The EDA separates usable candidate windows from unresolved processing work.", ["Branch", "Rows", "Candidate hours", "0.8–30 sec rows", "Below 0.8 sec", "Above 30 sec", "Decision"], [
        ["IRMCT / ICTY", tribunal_stats["rows"], tribunal_stats["hours"], "Pilot rows within range", tribunal_stats["under"], tribunal_stats["over"], "Controlled pilot review"],
        ["Lindsay Clancy", clancy_stats["rows"], clancy_stats["hours"], f"{clancy_window_stats['rows']} ({clancy_window_stats['hours']} h)", clancy_stats["under"], clancy_stats["over"], "Primary duration-filtered pool"],
        ["Tupac / Keffe D", "0", "0", "Not available", "N/A", "N/A", "Acquire and process first"],
        ["Indian-SIM", "0", "0", "Not available", "N/A", "N/A", "Build adaptation source pack"],
    ], size=10)

    slide_with_bullets(prs, "What I Want to Confirm in Guidance", "The mentor discussion should focus on scope discipline and the evidence required before training.", [
        "Whether Clancy should remain the primary benchmark while tribunal remains the secondary bootstrap comparison.",
        "Whether Tupac / Keffe D should be processed next or remain a planned diversity branch until sources are manually verified.",
        "Whether the 0.8–30 second duration window is acceptable as the first training-selection rule.",
        "How much manual emotion-label review is required before using the current heuristic labels.",
        "What minimum Indian-SIM material is sufficient to demonstrate adaptation rather than only describe it.",
    ], size=16)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {len(prs.slides)} slides to {OUTPUT}")


if __name__ == "__main__":
    build()
