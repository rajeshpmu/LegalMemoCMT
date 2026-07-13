from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from docx.oxml import OxmlElement
from docx.shared import Pt, RGBColor
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Reading_Script.docx"


def insert_paragraph_before(paragraph: Paragraph, text: str = "", style: str | None = None) -> Paragraph:
    new_p = OxmlElement("w:p")
    paragraph._p.addprevious(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        new_para.add_run(text)
    return new_para


def insert_paragraph_after(paragraph: Paragraph, text: str = "", style: str | None = None) -> Paragraph:
    new_p = OxmlElement("w:p")
    paragraph._p.addnext(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        new_para.add_run(text)
    return new_para


def find_paragraph(doc: Document, text: str):
    for p in doc.paragraphs:
        if p.text.strip() == text:
            return p
    return None


def update_pptx() -> None:
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides[4]  # current Slide 5
    target = None
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False) and shp.has_text_frame:
            txt = shp.text_frame.text.strip()
            if txt.startswith("Closest published reference point for the Phase 1 baseline"):
                target = shp
                break
    if target is None:
        raise RuntimeError("Could not locate the MemoCMT body textbox on Slide 5.")

    lines = [
        "Closest published reference point for the Phase 1 baseline because it fuses pretrained text and audio encoders through a cross-modal transformer.",
        "Uses BERT for text and HuBERT for audio, so each modality is encoded first and then aligned through cross-attention instead of simple concatenation.",
        "CLS means the first valid fused token is treated as the summary vector; MEAN averages all valid fused tokens; MAX keeps the strongest hidden activation per dimension; MIN keeps the weakest hidden activation per dimension.",
        "The paper compares those aggregation choices after fusion, and the Phase 1 paper-aligned baseline follows the MIN choice for MELD.",
        "The paper also reports MELD as a case study, which makes it the right technical anchor for the conversational baseline story.",
        "This literature review maps directly to Phase 1 implementation: reproduce the text+audio baseline first, then extend with the visual branch.",
    ]
    target.text_frame.clear()
    for i, line in enumerate(lines):
        p = target.text_frame.paragraphs[0] if i == 0 else target.text_frame.add_paragraph()
        p.text = line if i == 0 else f"{line}"
        p.level = 0
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(13)
            run.font.color.rgb = PptRGBColor.from_string("1F1F1F")
    prs.save(str(PPTX_PATH))


def update_docx() -> None:
    doc = Document(str(DOCX_PATH))
    slide5 = find_paragraph(doc, "Slide 5: Literature Review 1: MemoCMT")
    slide6 = find_paragraph(doc, "Slide 6: Why MIN Pooling for MELD")
    if slide5 is None or slide6 is None:
        raise RuntimeError("Could not locate the Slide 5 / Slide 6 section in the reading script DOCX.")

    additions = [
        "Make the fusion logic explicit: the CMT block lets the two modalities attend to each other, and then the paper compares aggregation strategies such as CLS, MEAN, MAX, and MIN. The Phase 1 baseline follows the same idea and uses MIN as the strongest paper-aligned choice for MELD.",
        "CLS, MEAN, MAX, and MIN are not separate models. They are four ways of converting the fused token sequence into one fixed-size vector before the classifier sees it.",
        "CLS takes the first valid token as the summary. In transformer-style models, the first token is often used as the aggregate representation, so the model learns to route sequence-level information into that position.",
        "MEAN averages all valid fused tokens. That keeps the representation stable, but it can smooth away strong emotional cues if only a few tokens carry the useful signal.",
        "MAX keeps the largest activation in each hidden dimension. That can preserve strong evidence, but it can also overreact to a noisy spike in one dimension.",
        "MIN keeps the smallest activation in each hidden dimension. In this project, MIN is the paper-aligned choice for MELD because the MemoCMT reference path reports its best MELD result with MIN pooling.",
        "A strong viva explanation is that pooling is a summarization step, not the fusion mechanism itself. The fusion step creates the interaction between text and audio, while pooling decides how that interaction is compressed into one classifier input vector.",
        "In code, the pooling choices are implemented after cross-attention in `src/models/model.py:190-208`, where the fused sequence is reduced with `x[:, 0, :]` for CLS, `mean(dim=1)` for MEAN, `max(dim=1)` for MAX, and `min(dim=1)` for MIN.",
        "The higher-level CMT fusion itself is implemented in `src/models/model.py:160-246`, and the paper-aligned training command uses MIN explicitly in `scripts/run_paper_aligned_meld_cv.sh:38-62`.",
        "That is why the project first reasons about the paper’s aggregation choices and then keeps MIN for the paper-aligned MELD run: it preserves comparability with MemoCMT and keeps the Phase 1 baseline anchored to the published reference.",
    ]

    # If the block has already been inserted, rebuild the Slide 5 section cleanly so the order reads naturally.
    if any("CLS, MEAN, MAX, and MIN are not separate models." in p.text for p in doc.paragraphs):
        to_remove = []
        seen = False
        for p in doc.paragraphs:
            if p is slide5:
                seen = True
                continue
            if p is slide6:
                break
            if seen:
                to_remove.append(p)
        for p in to_remove:
            remove_paragraph(p)
        anchor = slide5
        for text in additions:
            anchor = insert_paragraph_after(anchor, text)
        doc.save(str(DOCX_PATH))
        return

    anchor = slide5
    for text in additions:
        anchor = insert_paragraph_after(anchor, text)
    doc.save(str(DOCX_PATH))


def main() -> None:
    update_pptx()
    update_docx()
    print("Updated Slide 5 pooling explanation in PPTX and DOCX.")


if __name__ == "__main__":
    main()
