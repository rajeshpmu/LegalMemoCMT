from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
import re


PPTX_PATH = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT/implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx")


def remove_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    prs.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def insert_slide_before(prs, slide, before_index_0_based):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(before_index_0_based, sldId)


def title_shape(slide):
    candidates = []
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        text = shp.text_frame.text.strip()
        if not text:
            continue
        if shp.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        # top title area
        if float(shp.top) > 1200000:  # ~1.3 in
            continue
        candidates.append(shp)
    if not candidates:
        return None
    return sorted(candidates, key=lambda s: (float(s.top), -float(s.left)))[0]


def number_shape(slide):
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = shp.text_frame.text.strip()
        if not txt or not re.fullmatch(r"\d+", txt):
            continue
        if float(shp.top) <= 350000 and float(shp.left) >= 10000000:
            return shp
    return None


def renumber(prs):
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        nshape = number_shape(slide)
        if nshape is not None:
            nshape.text = str(idx)


def clean_title(slide):
    tshape = title_shape(slide)
    if tshape is None:
        return
    txt = tshape.text_frame.text.strip()
    txt = re.sub(r"^\d+\.\s*", "", txt)
    tshape.text = txt


def main():
    prs = Presentation(str(PPTX_PATH))
    if len(prs.slides) < 27:
        raise SystemExit(f"Expected at least 27 slides, found {len(prs.slides)}")

    # Move the last slide (comparison slide) to after the live demo summary slide.
    comparison = prs.slides[-1]
    remove_slide(prs, len(prs.slides) - 1)
    insert_slide_before(prs, comparison, 24)  # before current Thank You slide

    # Clean title text and renumber visible slide-number boxes.
    clean_title(prs.slides[24])
    renumber(prs)

    prs.save(str(PPTX_PATH))
    print(f"Moved comparison slide to position 25 and renumbered visible slide numbers in {PPTX_PATH}")


if __name__ == "__main__":
    main()
