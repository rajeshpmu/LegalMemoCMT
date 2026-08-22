from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUT_DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Active_SemiSupervised_Annotation_Student_Guide.docx"
OUT_PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Active_SemiSupervised_Annotation_Guidance.pptx"


def add_code(doc: Document, text: str) -> None:
    p = doc.add_paragraph(style="No Spacing")
    r = p.add_run(text)
    r.font.name = "Courier New"
    r.font.size = Pt(8)


def add_box(slide, x, y, w, h, text, fill, size=16, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(70, 80, 95)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 1
    for run in p.runs:
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(25, 35, 50)


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = PptPt(19)
    return slide


# PPTX
prs = Presentation()
prs.slide_width = PptInches(13.333)
prs.slide_height = PptInches(7.5)
title = prs.slides.add_slide(prs.slide_layouts[0])
title.shapes.title.text = "LegalMemoCMT Phase 2"
title.placeholders[1].text = "Active and Semi-Supervised Corpus Construction\nStudent Technical Guide"
add_bullet_slide(prs, "1. Why This Stage Is Needed", [
    "Phase 1 already provides a MELD-trained tri-modal checkpoint for basic emotion.",
    "Courtroom data is different: speakers, witness roles, examination phases, and camera visibility vary.",
    "The new stage creates a human-in-the-loop path instead of treating model predictions as ground truth.",
    "Immediate priority remains reliable diarization, transcript alignment, witness extraction, and utterance clips.",
])
slide = prs.slides.add_slide(prs.slide_layouts[6])
tb = slide.shapes.add_textbox(PptInches(0.45), PptInches(0.2), PptInches(12.4), PptInches(0.5))
tb.text_frame.text = "2. End-to-End Active Annotation Loop"
tb.text_frame.paragraphs[0].runs[0].font.size = PptPt(26)
tb.text_frame.paragraphs[0].runs[0].font.bold = True
flow = ["Diarize", "Align", "Resolve roles", "Extract witness turns", "Build tri-modal clips", "Human seed", "Warm-start", "Pseudo-label", "Review uncertainty", "Retrain"]
for i, text in enumerate(flow):
    x = 0.35 + (i % 5) * 2.55
    y = 1.1 + (i // 5) * 1.55
    add_box(slide, x, y, 2.1, 0.75, text, ((220, 232, 245) if i < 5 else (231, 243, 235)), 14, True)
    if i % 5 < 4:
        line = slide.shapes.add_connector(1, PptInches(x + 2.1), PptInches(y + 0.37), PptInches(x + 2.5), PptInches(y + 0.37))
        line.line.end_arrowhead = True
    elif i == 4:
        line = slide.shapes.add_connector(1, PptInches(x + 1.05), PptInches(y + 0.75), PptInches(0.35 + 2.55 * 4 + 1.05), PptInches(2.65))
        line.line.end_arrowhead = True
add_bullet_slide(prs, "3. Implementation Boundary", [
    "Implemented now: schema extension, seed selection, reviewer queue, priority scoring, optional suggestions, human merge, and validation.",
    "Already present: pyannote setup, diarization worker, subtitle evidence, role review, role-map application, Phase 1 pseudo-label script.",
    "Deferred: automatic pseudo-label acceptance, iterative retraining, video affect inference, contradiction pairs, Indian multilingual adaptation.",
    "No diarization output is assumed until the one-source pilot completes.",
])
add_bullet_slide(prs, "4. Two Label Families", [
    "Basic emotion: neutral, anger, disgust, fear, joy, sadness, surprise.",
    "Courtroom affect: CALM_COMPOSED, HESITANT_UNCERTAIN, GUARDED, DEFENSIVE, ASSERTIVE, TENSE, DISTRESSED, AGITATED, UNCLEAR.",
    "These are separate tasks: anger must not be mechanically mapped to DEFENSIVE, and neutral must not be mapped to CALM_COMPOSED.",
    "Existing emotion_label remains compatible; canonical fields record source and annotation status.",
])
add_bullet_slide(prs, "5. Annotation Status and Provenance", [
    "UNLABELED: no accepted annotation exists.",
    "AUTO_SUGGESTED / PSEUDO_LABELED: machine output that is not human ground truth.",
    "HUMAN_SINGLE / HUMAN_MULTI / ADJUDICATED: reviewed labels with increasing agreement evidence.",
    "REJECTED: the row should not enter the target task; retain the reason.",
    "Suggestions never overwrite canonical human labels.",
])
add_bullet_slide(prs, "6. Model-Assisted Evidence", [
    "Pyannote: who spoke when; anonymous clusters only, not legal roles or emotion.",
    "Whisper/ASR and subtitles: what was said and when; supports alignment and context.",
    "SER-Odyssey or audeering audio models: provisional valence/arousal/dominance evidence.",
    "DeBERTa zero-shot: transcript-based candidate affect descriptions.",
    "Phase 1 checkpoint: warm-start and basic-emotion pseudo-label baseline, not courtroom-affect ground truth.",
])
add_bullet_slide(prs, "7. Context and Observable Cues", [
    "Retain examination_phase, question_type, challenge_level, response_stance, and previous_question_text.",
    "Record cues separately: hesitation, pause, correction, repetition, interruption, overlap, speech-rate change, visible distress.",
    "Numeric measures include response latency, mean pause, speech rate, and overlap duration.",
    "Observable cues describe behavior; they must never be converted into deception or credibility labels.",
])
add_bullet_slide(prs, "8. Human-Reviewed Seed Set", [
    "Start with approximately 500–1000 utterances, not only easy high-confidence rows.",
    "Cover Clancy and Tupac, witnesses, examination phases, professional/family/police witnesses, neutral and high-affect cases.",
    "Use a witness-disjoint and duplicate-source-disjoint benchmark before active learning begins.",
    "The seed set becomes the first trustworthy legal-domain annotation set.",
])
add_bullet_slide(prs, "9. Active-Learning Priority", [
    "Priority can combine entropy, top-1/top-2 margin, modality disagreement, class rarity, quality flags, and diversity.",
    "Keep the scoring weights in YAML so they can be inspected and changed without editing code.",
    "Compare uncertainty, disagreement, diversity, class-balance, witness-diversity, and random-control strategies.",
    "Disagreement is useful evidence for review; it is not a reason to force label agreement.",
])
add_bullet_slide(prs, "10. First Commands and Next Milestone", [
    "First: complete the one-source pyannote pilot and save raw segments plus diarized manifest.",
    "Next: create the cluster review sheet, map roles, align clusters to turns, and validate witness-only clips.",
    "Then: extend the schema, select a seed, export the queue, and collect human annotations.",
    "Only after human seed labels exist: run model suggestions and plan warm-start adaptation.",
])
prs.save(OUT_PPTX)


# DOCX
doc = Document()
doc.add_heading("LegalMemoCMT Phase 2: Active and Semi-Supervised Corpus Construction", 0)
doc.add_paragraph("Student technical guide for the incremental next stage of the Clancy, Tupac, and future Indian corpus pipeline.")
doc.add_heading("1. Purpose and Scope", level=1)
doc.add_paragraph("The goal is to construct a reliable courtroom utterance corpus before applying machine-assisted annotation. Phase 2 must first solve the evidence problem: who spoke, what was said, when it was said, and whether text, audio, and video refer to the same utterance. Active learning is added only after these foundations exist.")
doc.add_heading("2. Existing Pipeline and Immediate Boundary", level=1)
doc.add_paragraph("The repository already contains the Phase 1 MELD tri-modal checkpoint, Clancy turn/clip preparation, subtitle evidence extraction, pyannote environment setup, diarization scripts, speaker-role review tools, and a Phase 1 pseudo-label script. The current diarization job must finish before its outputs are used. The new annotation package does not assume those output files already exist.")
doc.add_heading("3. Correct Execution Order", level=1)
for text in ["Run pyannote/speaker-diarization-3.1.", "Save raw diarization segments.", "Resolve anonymous clusters to courtroom roles where possible.", "Align diarization intervals to subtitles, transcripts, ASR, and existing turns.", "Build speaker-aware turns and identify witness utterances.", "Generate aligned text, WAV, and MP4 clips.", "Run text/audio/video quality checks.", "Extend the schema and prepare annotation queues.", "Create a human-reviewed seed set.", "Only then begin model suggestions, pseudo-labeling, and active-learning iterations."]:
    doc.add_paragraph(text, style="List Number")
doc.add_heading("4. New Annotation Package", level=1)
doc.add_paragraph("The additive package is under phase2/annotation/. It does not replace existing preprocessing or training behavior.")
for name, explanation in [
    ("extend_annotation_schema.py", "Adds canonical labels, statuses, suggestion fields, observable cues, priority fields, and iteration provenance. It preserves emotion_label and does not copy unsupported legacy placeholders such as unknown into basic_emotion."),
    ("build_seed_annotation_set.py", "Selects a configurable and diverse seed using source, witness, and examination-phase groups. It marks rows UNLABELED and requires review."),
    ("build_annotation_queue.py", "Exports a reviewer-friendly CSV containing clip paths, transcript context, model evidence, role fields, and priority reasons."),
    ("compute_annotation_priority.py", "Computes a configurable score from entropy, margin uncertainty, modality disagreement, class rarity, and quality flags. It is a selection aid, not a truth score."),
    ("generate_model_suggestions.py", "Optionally runs a text zero-shot model and/or audio classifier. It writes suggestion fields and AUTO_SUGGESTED status only; canonical labels are not overwritten."),
    ("merge_human_annotations.py", "Merges reviewed labels by utterance_id, validates vocabularies, writes manual provenance, and preserves machine suggestions."),
    ("validate_annotation_state.py", "Checks duplicate IDs, controlled statuses, allowed label vocabularies, and whether human statuses have proper provenance."),
]:
    doc.add_heading(name, level=2)
    doc.add_paragraph(explanation)
doc.add_heading("5. Label Families", level=1)
doc.add_paragraph("Basic emotion remains MELD-compatible: neutral, anger, disgust, fear, joy, sadness, and surprise. Canonical fields are basic_emotion, basic_emotion_source, basic_emotion_confidence, and basic_emotion_annotation_status. Existing emotion_label, emotion_label_source, and emotion_label_confidence remain available for older code.")
doc.add_paragraph("Courtroom affect is separate: CALM_COMPOSED, HESITANT_UNCERTAIN, GUARDED, DEFENSIVE, ASSERTIVE, TENSE, DISTRESSED, AGITATED, and UNCLEAR. The implementation never maps anger to DEFENSIVE or neutral to CALM_COMPOSED automatically.")
doc.add_heading("6. Annotation Status", level=1)
doc.add_paragraph("UNLABELED means no accepted label exists. AUTO_SUGGESTED records an output shown to a reviewer. PSEUDO_LABELED records an accepted machine label under a documented rule. HUMAN_SINGLE, HUMAN_MULTI, and ADJUDICATED record increasing human-review evidence. REJECTED means the row is excluded from the task. The final benchmark must not rely only on AUTO_SUGGESTED or unreviewed PSEUDO_LABELED rows.")
doc.add_heading("7. Model Stack and Limitations", level=1)
doc.add_paragraph("Pyannote/speaker-diarization-3.1 provides speech intervals and anonymous speaker clusters. It does not identify legal roles or emotion. Whisper/ASR and subtitles provide text and timing evidence. The SER-Odyssey multi-attribute audio model or audEERING dimensional model can provide provisional valence/arousal/dominance. A DeBERTa zero-shot NLI model can score transcript descriptions such as guarded or defensive. The Phase 1 MELD checkpoint provides a basic-emotion baseline and warm-start initialization, not courtroom-affect ground truth.")
doc.add_heading("8. Courtroom Context and Observable Cues", level=1)
doc.add_paragraph("The annotation context should retain the previous attorney question, target witness response, and optionally the following interaction. Store examination_phase, question_type, challenge_level, and response_stance. Store observable cues separately: hesitation, long pause, self-correction, repetition, interruption, overlap, voice rise/drop, speech-rate change, visible distress, gaze shift, head movement, and facial tension. Numeric fields can record response_latency_ms, mean_pause_ms, speech_rate_wpm, and overlap_duration_ms. These fields describe behavior and do not establish deception, credibility, truthfulness, or reliability.")
doc.add_heading("9. Seed Dataset Design", level=1)
doc.add_paragraph("The initial seed should contain approximately 500–1000 utterances distributed across Clancy and Tupac, witnesses, direct/cross/redirect examination, professional/family/police witnesses, neutral and high-affect examples, ambiguous examples, and different source videos. The selection must not only choose easy high-confidence rows. Reserve a witness-disjoint and duplicate-source-disjoint benchmark before active selection begins.")
doc.add_heading("10. Active-Learning Loop", level=1)
for text in ["Warm-start a legal-domain model from the preserved Phase 1 checkpoint.", "Predict unannotated courtroom utterances.", "Store prediction distributions, entropy, margin, modality disagreement, and model provenance.", "Select uncertain, disagreeing, rare, diverse, and random-control samples.", "Export an annotation queue with the original clip and context.", "Merge reviewed labels; human labels take precedence.", "Retrain and record annotation_iteration, training_iteration, model_checkpoint, selection_strategy, and review_timestamp."]:
    doc.add_paragraph(text, style="List Number")
doc.add_heading("11. Exact First Pilot Command", level=1)
add_code(doc, '''PYTHON_BIN=$PWD/.venv-diarization/bin/python \\
HF_TOKEN="$HF_TOKEN" \\
bash phase2/run_clancy_diarization.sh \\
  --input-csv data/processed/phase2/clancy/duration_0_8_to_20/clancy_dataset_manifest_vit_200_subtitle_evidence.csv \\
  --output-csv data/processed/phase2/clancy/duration_0_8_to_20/clancy_dataset_manifest_vit_200_diarized.csv \\
  --segments-csv data/processed/phase2/clancy/duration_0_8_to_20/clancy_diarization_segments_200.csv \\
  --device cpu \\
  --max-sources 1''')
doc.add_paragraph("Do not run annotation scripts against the diarized manifest until this pilot completes and the two output CSVs exist. Then create the cluster review sheet with phase2/create_clancy_speaker_role_review.py, manually map clusters, apply the role map, and validate witness-only tri-modal rows.")
doc.add_heading("12. Annotation Package Commands After Diarization", level=1)
add_code(doc, '''PYTHON_BIN=$PWD/.venv/bin/python \\
bash phase2/annotation/run_extend_annotation_schema.sh \\
  --input-csv <role-aware-manifest.csv> \\
  --output-csv <schema-extended-manifest.csv>

PYTHON_BIN=$PWD/.venv/bin/python \\
bash phase2/annotation/run_build_seed_annotation_set.sh \\
  --input-csv <schema-extended-manifest.csv> \\
  --output-csv <seed.csv> \\
  --summary-json <seed-summary.json> \\
  --config configs/annotation/active_learning.yaml

PYTHON_BIN=$PWD/.venv/bin/python \\
bash phase2/annotation/run_compute_annotation_priority.sh \\
  --input-csv <seed-or-suggestion-manifest.csv> \\
  --output-csv <priority-manifest.csv> \\
  --summary-json <priority-summary.json> \\
  --config configs/annotation/active_learning.yaml

PYTHON_BIN=$PWD/.venv/bin/python \\
bash phase2/annotation/run_build_annotation_queue.sh \\
  --input-csv <priority-manifest.csv> \\
  --output-csv <annotation_queue.csv> \\
  --max-rows 100''')
doc.add_heading("13. What Is Implemented and What Is Deferred", level=1)
doc.add_paragraph("Implemented: additive schema migration, seed selection, queue export, configurable priority scoring, optional text/audio suggestions, human merge, validation, YAML configuration, and an implementation report. Deferred: automatic pseudo-label acceptance, iterative active-learning retraining, video-affect model integration, contradiction-pair construction, and Indian multilingual adaptation. Deferring these is deliberate because diarization and the human seed set are prerequisites.")
doc.add_heading("14. Readiness Status", level=1)
doc.add_paragraph("Clancy: preprocessing and diarization prerequisites exist; the current diarization pilot must complete before its outputs are trusted. Tupac: source collection and processing remain pending. Indian corpus: planned adaptation branch, not prepared. The immediate milestone is completed pyannote output, role mapping, transcript/turn alignment, and validation of witness-only tri-modal clips.")
doc.add_heading("15. Scientific Description", level=1)
doc.add_paragraph("LegalMemoCMT uses a human-in-the-loop active and semi-supervised annotation strategy. Speaker-aware tri-modal courtroom utterances are first constructed using diarization, transcription alignment, and witness-role filtering. Pretrained modality-specific models provide candidate evidence but do not define ground truth. A manually reviewed seed corpus can warm-start the Phase-1 MELD model. The adapted model may later pseudo-label unannotated utterances, while uncertainty, modality disagreement, class rarity, and diversity drive later human review. Human labels supersede machine suggestions, and witness-disjoint evaluation data remain isolated from the loop.")
doc.save(OUT_DOCX)
print(f"Wrote {OUT_DOCX}")
print(f"Wrote {OUT_PPTX}")
