from pathlib import Path
import shutil

from docx import Document
from docx.enum.section import WD_SECTION
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
PPTX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.before_hf_diarization.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Speaking_Doc.docx"
DOCX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Speaking_Doc.before_hf_diarization.docx"


def add_code(document: Document, text: str) -> None:
    paragraph = document.add_paragraph(style="No Spacing")
    run = paragraph.add_run(text)
    run.font.name = "Courier New"
    run.font.size = Pt(8)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptPt(20)
    # Keep the slide readable on the existing 16:9 deck.
    body.margin_left = Inches(0.25)
    body.margin_right = Inches(0.25)


if not PPTX_BACKUP.exists():
    shutil.copy2(PPTX, PPTX_BACKUP)
prs = Presentation(PPTX)
add_bullet_slide(
    prs,
    "Hugging Face Models: Why I Use Them",
    [
        "pyannote/speaker-diarization-3.1 is used to find who spoke when in the source courtroom audio.",
        "Its segmentation dependency detects speech regions; the speaker-embedding component groups similar voices into anonymous cluster IDs.",
        "The models solve a timing and voice-grouping problem, not a legal-role problem. A cluster is not automatically a Witness, Judge, or Prosecutor.",
        "I use a separate .venv-diarization with torch 2.2.2, torchaudio 2.2.2, pyannote.audio 3.4.0, and the compatible Hugging Face Hub version.",
        "Model access is provenance-controlled: the token is used only at runtime and is never stored in the corpus manifest.",
    ],
)
add_bullet_slide(
    prs,
    "From Diarization Output to Phase 2 Rows",
    [
        "Input: source audio plus the subtitle/utterance manifest with timestamps.",
        "Output: source-level segments containing start time, end time, and anonymous speaker_cluster_id values.",
        "The pipeline overlaps each utterance interval with diarization segments and assigns the strongest matching cluster; this preserves source and timestamp traceability.",
        "Manual cluster review then maps selected clusters to Witness, Prosecutor, Defence, Judge, or Other with a confidence value.",
        "The enriched manifest records witness_in_segment, witness_speaking_status, speaker_role, speaker_role_source, speaker_role_confidence, visual_target_role, and visual_speaker_match.",
        "Only after this review can role-aware filtering or role-aware weak labels be used. Diarization alone does not create emotion, credibility, or deception labels.",
    ],
)
prs.save(PPTX)


if not DOCX_BACKUP.exists():
    shutil.copy2(DOCX, DOCX_BACKUP)
document = Document(DOCX)
document.add_section(WD_SECTION.NEW_PAGE)
document.add_heading("Additional Speaking Module: Hugging Face Diarization Models", level=1)
document.add_heading("1. Why I Need Diarization in a Courtroom Corpus", level=2)
document.add_paragraph(
    "A courtroom video may show the witness while a prosecutor is asking a question, or it may show a "
    "judge while another person is speaking. Therefore, the visible face cannot be treated as the speaking "
    "person by default. The transcript supplies words and approximate utterance boundaries, but it does not "
    "reliably identify the speaker in every Clancy subtitle. Speaker diarization adds an audio-based layer: "
    "it detects speech intervals and groups acoustically similar voice segments. This makes the corpus more "
    "aware of speaker changes before I assign courtroom roles."
)
document.add_heading("2. What the Models Do and Do Not Do", level=2)
document.add_paragraph(
    "The main pipeline is `pyannote/speaker-diarization-3.1`. It combines speech segmentation and speaker "
    "representation components. Segmentation estimates when speech is present. Speaker embeddings convert "
    "short voice regions into numerical representations, and clustering groups regions that sound like the "
    "same speaker. The resulting labels are anonymous identifiers such as `SPEAKER_00` and `SPEAKER_01`."
)
document.add_paragraph(
    "These models do not know that `SPEAKER_00` is a witness or that `SPEAKER_01` is a prosecutor. They also "
    "do not determine emotion, truthfulness, deception, credibility, or legal responsibility. Those meanings "
    "must not be inferred from a cluster label or a face. The model output is evidence for a later role-review "
    "step, not a final legal interpretation."
)
document.add_heading("3. Why the Gated Model Dependencies Matter", level=2)
document.add_paragraph(
    "The diarization pipeline downloads a configuration and its dependent models. In this environment the "
    "important dependency chain includes `pyannote/speaker-diarization-3.1` and the gated segmentation model "
    "`pyannote/segmentation-3.0`, with any linked speaker-embedding model requested by the pipeline. A valid "
    "Hugging Face identity is not sufficient: the account must also accept the model conditions and use a "
    "read-enabled token."
)
document.add_paragraph(
    "The repository keeps this workload separate in `.venv-diarization` because pyannote.audio 3.4.0 needs "
    "the matched Torch audio stack. The validated versions are torch 2.2.2, torchaudio 2.2.2 with "
    "`AudioMetaData`, pyannote.audio 3.4.0, and a compatible Hugging Face Hub release. The token is supplied "
    "through `HF_TOKEN` at runtime and is not written to source files or manifests."
)
document.add_heading("4. How I Run the Model", level=2)
add_code(document, '''export HF_TOKEN="<your Hugging Face read token>"
./.venv-diarization/bin/python \\
  phase2/check_clancy_diarization_prerequisites.py \\
  --model pyannote/speaker-diarization-3.1 \\
  --load-model''')
document.add_paragraph(
    "The preflight performs three different checks. First, it checks local Torch, torchaudio, and pyannote "
    "compatibility. Second, it checks repository access using Hugging Face model information. Third, it loads "
    "the pipeline and therefore tests access to dependent gated models. I proceed only when both repository "
    "access and model loading pass."
)
add_code(document, '''PYTHON_BIN=$PWD/.venv-diarization/bin/python \\
HF_TOKEN="$HF_TOKEN" \\
bash phase2/run_clancy_diarization.sh \\
  --input-csv data/processed/phase2/clancy/duration_0_8_to_20/clancy_dataset_manifest_vit_200_subtitle_evidence.csv \\
  --output-csv data/processed/phase2/clancy/duration_0_8_to_20/clancy_dataset_manifest_vit_200_diarized.csv \\
  --segments-csv data/processed/phase2/clancy/duration_0_8_to_20/clancy_diarization_segments_200.csv \\
  --device cpu \\
  --max-sources 1''')
document.add_heading("5. How the Output Enters the Phase 2 Pipeline", level=2)
for item in [
    "`clancy_diarization_segments_200.csv` is the source-level audit file. It records the source audio path, anonymous cluster ID, segment start, segment end, and diarization model.",
    "`clancy_dataset_manifest_vit_200_diarized.csv` preserves the original utterance row and adds the best-overlapping cluster assignment where an overlap exists.",
    "The overlap is a temporal join: an utterance receives a cluster when the diarization interval and utterance interval share time. The chosen cluster is the candidate with the greatest overlap score.",
    "A manual review sheet maps stable clusters to legal roles. The mapping must record the sample listened to, the role decision, and confidence; it must not deanonymize protected witnesses.",
    "The final role-aware manifest can support witness-only filtering, role-aware analysis, and later weak-label conditioning while retaining the original transcript, audio, video, and source paths.",
]:
    document.add_paragraph(item, style="List Bullet")
document.add_heading("6. What the Outcome Means for Training", level=2)
document.add_paragraph(
    "The immediate outcome is not a trained model. It is a better structured dataset in which each utterance "
    "has stronger evidence about who was speaking and whether the witness was speaking or being addressed. "
    "This helps prevent a visible-witness clip containing a lawyer's question from being incorrectly treated as "
    "witness speech. It also lets later experiments compare all courtroom speech with witness-only speech."
)
document.add_paragraph(
    "For multimodal emotion work, diarization can be used as a conditioning or filtering signal, while the "
    "Phase 1 MELD model remains only a source of provisional basic-emotion pseudo-labels. Courtroom-specific "
    "affect labels require a separate annotation policy and confidence tracking. No diarization output should "
    "be converted automatically into `deceptive`, `truthful`, `credible`, or `unreliable`."
)
document.add_heading("7. Current Result and Honest Limitation", level=2)
document.add_paragraph(
    "The local prerequisites and Hugging Face model access now pass, including the gated segmentation download. "
    "The next controlled experiment is one source with a small manifest. A successful run will demonstrate "
    "cluster extraction, not perfect speaker identification. Accented speech, overlapping courtroom speech, "
    "microphones, camera cuts, and short utterances can produce incorrect or fragmented clusters, so manual "
    "role verification remains a required quality gate."
)
document.add_heading("8. One-Sentence Guidance-Call Explanation", level=2)
document.add_paragraph(
    "I use Hugging Face pyannote models to add an auditable audio-based speaker-turn layer, then combine it "
    "with subtitle evidence and manual cluster-role verification so the Phase 2 manifest can distinguish who "
    "was speaking without pretending that an automatic model knows courtroom roles or emotional truth."
)
document.save(DOCX)
print(f"Updated {PPTX}")
print(f"Updated {DOCX}")
