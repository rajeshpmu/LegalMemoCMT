"""Add gated-manifest lineage and example explanations to the manually edited v1 deck."""
from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Presentation_v1.pptx"
IMAGE = ROOT / "implementation_docments/figures/phase2_gated_annotation_flow.png"
MARKER = "GATED_MANIFEST_LINEAGE_V1"
NAVY = RGBColor(20,48,87); TEXT = RGBColor(35,35,35); GREY = RGBColor(95,95,95)

def text(slide, value, x, y, w, h, size=16, color=TEXT, bold=False):
    box=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); box.text_frame.word_wrap=True; box.text_frame.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p=box.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=value; r.font.name='Aptos'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color

def frame(prs, heading, subtitle):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=RGBColor(255,255,255)
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.22)); bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
    text(slide,heading,.55,.38,12.2,.55,24,NAVY,True); text(slide,subtitle,.58,.96,12,.35,12,GREY)
    return slide

def insert_before(slide_list, slide, before_index):
    ids=slide_list._sldIdLst; new_id=ids[-1]; ids.remove(new_id); ids.insert(before_index, new_id)

def main():
    if not PPTX.exists() or not IMAGE.exists(): raise SystemExit('Missing deck or rendered diagram')
    prs=Presentation(PPTX)
    if any(MARKER in sh.text for s in prs.slides for sh in s.shapes if hasattr(sh,'text')): print('Already updated'); return
    backup=PPTX.with_name(PPTX.name+'.before_gated_flow_v1.pptx')
    if not backup.exists(): shutil.copy2(PPTX,backup)
    before=len(prs.slides)-1
    s=frame(prs,'How the Gated CSV Was Produced','Reproducible data lineage from scope evidence to final machine-assisted fields')
    s.shapes.add_picture(str(IMAGE), Inches(.65), Inches(1.35), width=Inches(12.0), height=Inches(4.95))
    text(s,'The diagram separates candidate generation from acceptance. The first script combines the existing scope-aware CSV with audio, transcript, and visual evidence. The second script applies thresholds and conflict rules; it does not overwrite Phase 1 predictions.',.85,6.35,11.7,.55,13)
    insert_before(prs.slides,s,before); before+=1
    s=frame(prs,'Slide 17 Explained: Why the SILVER Row Passed','DCBWoWhsTpA_turn06801')
    text(s,'Input lineage\n\n1. emotion_scope_review_200_scope_aware.csv preserves Phase 1 output and adds scope/audio evidence.\n2. propose_clancy_courtroom_affect.py proposes neutral basic emotion and CALM_COMPOSED delivery.\n3. apply_clancy_annotation_acceptance_gate.py checks thresholds and conflicts.',.75,1.4,5.9,4.9,17)
    text(s,'Gate result\n\nBasic candidate: neutral @ 0.75\nAffect candidate: CALM_COMPOSED @ 0.60\nCritical conflict: NO\n\nTherefore:\nfinal_basic_emotion=neutral\nfinal_courtroom_affect=CALM_COMPOSED\nannotation_status=AUTO_ADJUDICATED\nannotation_tier=SILVER\n\nThis is machine-assisted acceptance, not human gold truth.',6.9,1.4,5.7,4.9,17)
    insert_before(prs.slides,s,before); before+=1
    s=frame(prs,'Slide 18 Explained: Why the Row Stayed WEAK','DCBWoWhsTpA_turn06570')
    text(s,'Evidence\n\nPhase 1: disgust @ 0.472638\nBasic candidate confidence: 0.47\nCourtroom candidate: HESITANT_UNCERTAIN @ 0.75\nTranscript: “Um, [snorts] I don’t believe...”\n\nThe affect candidate is useful, but the basic-emotion evidence is below the 0.70 gate.',.75,1.4,5.9,4.9,17)
    text(s,'Gate result\n\nfinal_basic_emotion=UNRESOLVED\nfinal_courtroom_affect=HESITANT_UNCERTAIN\nannotation_status=UNRESOLVED\nannotation_tier=WEAK\n\nThe gate prevents the strong affect score from hiding weak basic-emotion evidence. The bracketed non-speech marker is context only and is not converted into disgust.',6.9,1.4,5.7,4.9,17)
    insert_before(prs.slides,s,before)
    prs.save(PPTX); print('Updated',PPTX,'slides',len(prs.slides))

if __name__=='__main__': main()
