from pathlib import Path
import shutil

from docx import Document
from docx.enum.section import WD_SECTION
from docx.shared import Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
PPTX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.before_visible_witness_next_stage.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Speaking_Doc.docx"
DOCX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Speaking_Doc.before_visible_witness_next_stage.docx"
TITLES = {
    "After Visible-Witness Selection: The Next Stage",
    "Pool A Readiness Gate Before Adaptation",
}


def box(slide, x, y, w, h, text, color, size=15, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.color.rgb = RGBColor(70, 80, 95)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 1
    for run in p.runs:
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(25, 35, 50)


def arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(70, 80, 95)
    line.line.end_arrowhead = True


prs = Presentation(PPTX)
if not PPTX_BACKUP.exists():
    shutil.copy2(PPTX, PPTX_BACKUP)

if not any(any(title in sh.text for sh in slide.shapes if hasattr(sh, "text")) for slide in prs.slides for title in TITLES):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.5))
    title.text_frame.text = "After Visible-Witness Selection: The Next Stage"
    title.text_frame.paragraphs[0].runs[0].font.size = PptPt(26)
    title.text_frame.paragraphs[0].runs[0].font.bold = True
    steps = [
        ("1. Freeze\nrole-aware manifest", (220, 232, 245)),
        ("2. Filter\nPool A", (231, 243, 235)),
        ("3. Validate\nmodalities + provenance", (248, 237, 215)),
        ("4. Split by\nsource/hearing", (232, 224, 245)),
        ("5. Audit\nPhase 1 labels", (208, 224, 220)),
    ]
    x = 0.35
    for i, (text, color) in enumerate(steps):
        box(slide, x, 1.25, 2.25, 0.95, text, color, 16, True)
        if i < len(steps) - 1:
            arrow(slide, x + 2.25, 1.72, x + 2.5, 1.72)
        x += 2.55
    box(slide, 2.1, 3.05, 3.7, 1.0, "Controlled pilot\n100–300 Pool A rows", (220, 232, 245), 18, True)
    box(slide, 6.2, 3.05, 3.7, 1.0, "Warm-start adaptation\nonly after readiness passes", (231, 243, 235), 18, True)
    arrow(slide, 4.0, 2.2, 3.95, 3.05)
    arrow(slide, 5.8, 3.55, 6.2, 3.55)
    note = slide.shapes.add_textbox(Inches(0.7), Inches(4.65), Inches(11.8), Inches(1.9))
    tf = note.text_frame
    tf.clear()
    tf.word_wrap = True
    lines = [
        "The first goal is a trustworthy subset, not the largest possible row count.",
        "Phase 1 pseudo-labels remain provisional and retain model provenance.",
        "Courtroom-specific affect is a separate annotation task; it is not inferred from diarization.",
        "Rows from other pools remain available for later interaction and sensitivity experiments.",
    ]
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = PptPt(17)
    prs.save(PPTX)
    ids = prs.slides._sldIdLst
    new_id = ids[-1]
    ids.remove(new_id)
    anchor = "Selecting the Visible-Witness Speaking Corpus"
    insert_at = next(
        (
            i + 1
            for i, sid in enumerate(ids)
            if any(anchor in sh.text for sh in prs.part.related_part(sid.rId).slide.shapes if hasattr(sh, "text"))
        ),
        len(ids),
    )
    ids.insert(insert_at, new_id)
    prs.save(PPTX)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.5))
    title.text_frame.text = "Pool A Readiness Gate Before Adaptation"
    title.text_frame.paragraphs[0].runs[0].font.size = PptPt(26)
    title.text_frame.paragraphs[0].runs[0].font.bold = True
    box(slide, 0.65, 1.1, 3.7, 1.15, "Required row evidence\nWitness • SPEAKING • WITNESS_VISIBLE", (208, 224, 220), 17, True)
    box(slide, 4.8, 1.1, 3.7, 1.15, "Required modality checks\ntext • audio • video • provenance", (220, 232, 245), 17, True)
    box(slide, 8.95, 1.1, 3.7, 1.15, "Required evaluation controls\ngrouped splits • no leakage", (248, 237, 215), 17, True)
    box(slide, 2.3, 3.0, 3.7, 1.15, "Then: pseudo-label\na 100–300 row audit", (231, 243, 235), 17, True)
    box(slide, 6.4, 3.0, 3.7, 1.15, "Finally: compare\nbaseline vs adaptation", (232, 224, 245), 17, True)
    arrow(slide, 2.5, 2.25, 4.0, 3.0)
    arrow(slide, 6.65, 2.25, 4.8, 3.0)
    arrow(slide, 10.8, 2.25, 8.2, 3.0)
    arrow(slide, 6.0, 3.58, 6.4, 3.58)
    note = slide.shapes.add_textbox(Inches(0.85), Inches(4.75), Inches(11.6), Inches(1.45))
    note.text_frame.text = "Do not fine-tune on unverified pseudo-labels or leakage-prone splits. The readiness report must show what passed, what was excluded, and why."
    note.text_frame.paragraphs[0].runs[0].font.size = PptPt(19)
    note.text_frame.paragraphs[0].runs[0].font.bold = True
    prs.save(PPTX)
    ids = prs.slides._sldIdLst
    new_id = ids[-1]
    ids.remove(new_id)
    anchor = "After Visible-Witness Selection: The Next Stage"
    insert_at = next(
        (
            i + 1
            for i, sid in enumerate(ids)
            if any(anchor in sh.text for sh in prs.part.related_part(sid.rId).slide.shapes if hasattr(sh, "text"))
        ),
        len(ids),
    )
    ids.insert(insert_at, new_id)
    prs.save(PPTX)


if not DOCX_BACKUP.exists():
    shutil.copy2(DOCX, DOCX_BACKUP)
document = Document(DOCX)
document.add_section(WD_SECTION.NEW_PAGE)
document.add_heading("Additional Speaking Module: What Happens After Pool A Exists", level=1)
document.add_heading("1. Why Visible-Witness Selection Is Not the Final Dataset", level=2)
document.add_paragraph(
    "Visible-witness selection removes an important source of label noise, but it does not automatically prove "
    "that every row is ready for model training. A row can have a witness face and still have incorrect transcript "
    "boundaries, a silent audio clip, a corrupted video, a weak diarization assignment, or a duplicate source "
    "segment. The next stage converts the selected rows into a controlled and auditable Pool A rather than treating "
    "the visual selection as a final quality certificate."
)
document.add_heading("2. Step One: Freeze the Role-Aware Manifest", level=2)
document.add_paragraph(
    "First I preserve the role-aware manifest as a versioned intermediate artifact. I do not overwrite the source "
    "manifest or remove rejected rows. Every selected row must retain its utterance ID, source video/audio paths, "
    "subtitle path, timestamps, speaker cluster, role source, confidence, and visual review status. This allows a "
    "future reviewer to trace a training row back to the original recording and the decision that admitted it."
)
document.add_heading("3. Step Two: Construct Pool A", level=2)
document.add_paragraph(
    "Pool A is the conservative witness-emotion subset. The selection rule requires `speaker_role=Witness`, "
    "`witness_speaking_status=SPEAKING`, and `speaker_visible=WITNESS_VISIBLE`. I also require valid audio, a "
    "resolved video path, acceptable transcript alignment, and no persistent rejection-list match. The rejected "
    "rows are written to an exclusion manifest with a reason such as unresolved speaker, witness listening, poor "
    "visibility, break footage, invalid media, or weak alignment."
)
document.add_heading("4. Step Three: Validate the Three Modalities", level=2)
document.add_paragraph(
    "The text check confirms that the utterance text and transcript source are present and that the timestamp is "
    "valid. The audio check confirms that the WAV exists, can be decoded, has a non-trivial duration, and is not "
    "silent or heavily clipped. The video check confirms that the MP4 exists, can be decoded, and has the expected "
    "time interval. These checks are independent: a valid audio file does not prove that the video shows the witness, "
    "and a valid MP4 does not prove that its speech matches the transcript."
)
document.add_heading("5. Step Four: Create Leakage-Aware Splits", level=2)
document.add_paragraph(
    "The split must use a source-level or hearing-level group. Adjacent utterances from one video share the same "
    "room acoustics, camera, microphone, legal vocabulary, and often the same continuous speaker turn. Randomly "
    "splitting rows would let the model see nearly identical context in training and testing. Therefore, complete "
    "source videos or hearings are assigned to only one of train, development, or test. The validation report must "
    "show zero group leakage before training."
)
document.add_heading("6. Step Five: Use Phase 1 for Provisional Basic-Emotion Labels", level=2)
document.add_paragraph(
    "After Pool A passes the structural checks, I can run the Phase 1 three-modality checkpoint to produce weak "
    "basic-emotion predictions. These predictions are not human ground truth. The manifest must preserve the model "
    "checkpoint, modality configuration, inference date, confidence or probability information, and source row ID. "
    "The original `emotion_label` field remains for compatibility, while `basic_emotion_source` and "
    "`basic_emotion_confidence` explain how the value was obtained."
)
document.add_paragraph(
    "I should review a stratified sample from every split, not only high-confidence examples. I should inspect "
    "whether predictions are dominated by neutral, whether certain speakers or videos receive abnormal labels, and "
    "whether visible courtroom expressions are being confused with MELD-style emotion categories. Low-confidence or "
    "obviously incorrect predictions should be flagged rather than silently promoted."
)
document.add_heading("7. Step Six: Keep Courtroom Affect Separate", level=2)
document.add_paragraph(
    "Basic emotion and courtroom affect answer different research questions. The Phase 1 model can provide a "
    "provisional MELD-compatible label such as sadness or anger. A separate annotation layer can describe courtroom "
    "interactional behavior such as guarded, hesitant, defensive, tense, or distressed. Neither layer should infer "
    "deception, truthfulness, credibility, or unreliability from facial behavior or voice."
)
document.add_heading("8. Step Seven: Run a Controlled Pilot", level=2)
document.add_paragraph(
    "Before processing the entire Pool A, I should select approximately 100–300 rows across multiple source videos, "
    "speakers, and courtroom contexts. I compare the Phase 1 checkpoint before adaptation with the checkpoint after a "
    "small warm-start experiment. The held-out source group is used for evaluation. This pilot tests whether the "
    "data preparation and labels are coherent before I spend compute on the full corpus."
)
document.add_heading("9. What I Report to the Guide", level=2)
for item in [
    "Pool A row count and total duration",
    "unique source videos, hearings, and manually verified witness clusters",
    "text, audio, and video validation counts",
    "alignment-confidence distribution and manual-review count",
    "emotion pseudo-label distribution and confidence distribution",
    "train/dev/test group counts and leakage result",
    "number of rows excluded from Pool A and the reasons",
    "baseline versus adaptation-pilot results",
]:
    document.add_paragraph(item, style="List Bullet")
document.add_heading("10. Student-Level Summary", level=2)
document.add_paragraph(
    "After visible-witness selection, I freeze the evidence, build Pool A, validate all three modalities, create "
    "grouped splits, audit Phase 1 pseudo-labels, and only then run a small adaptation pilot. This order protects the "
    "project from training on lawyer speech, duplicated context, invalid clips, or unexamined pseudo-labels."
)
document.save(DOCX)
print(f"Updated {PPTX}")
print(f"Updated {DOCX}")
