from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Presentation.pptx"
ASSET_DIR = ROOT / "implementation_docments" / "phase1_esa_assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)
MMD_PATH = ASSET_DIR / "system_architecture_flow.mmd"
SVG_PATH = ASSET_DIR / "system_architecture_flow.svg"
PNG_PATH = ASSET_DIR / "system_architecture_flow.png"

MERMAID = """%%{init: {'themeVariables': {'fontSize': '34px', 'fontFamily': 'Aptos'}}}%%
flowchart LR
  subgraph A[Input layer]
    direction TB
    A1[MP4 video]
    A2[Transcript text]
    A3[Audio stream]
  end
  subgraph B[Preprocessing layer]
    direction TB
    B1[Frame sampling]
    B2[Face crop]
    B3[FFmpeg audio extraction]
    B4[Text normalization]
  end
  subgraph C[Manifest and cache layer]
    direction TB
    C1[CSV manifest]
    C2[Cached .npy features]
  end
  subgraph D[Model layer]
    direction TB
    D1[BERT text encoder]
    D2[HuBERT audio encoder]
    D3[ViT visual encoder]
    D4[Gated fusion]
    D5[Classifier]
  end
  subgraph E[Analysis layer]
    direction TB
    E1[Metrics JSON]
    E2[Confusion matrix]
    E3[Top confusions]
    E4[Predictions CSV]
  end
  subgraph F[Future Phase 2]
    direction TB
    F1[Courtroom metadata]
    F2[Multilingual transcripts]
    F3[Explainability summaries]
  end
  A1 --> B1 --> B2 --> C1
  A2 --> B4 --> C1
  A3 --> B3 --> C2
  B2 --> C2
  C1 --> D1
  C1 --> D2
  C1 --> D3
  C2 --> D3
  D1 --> D4
  D2 --> D4
  D3 --> D4
  D4 --> D5 --> E1 --> E2 --> E3 --> E4 --> F1 --> F2 --> F3
"""


def render_diagram() -> None:
    MMD_PATH.write_text(MERMAID, encoding="utf-8")
    for out in (SVG_PATH, PNG_PATH):
        subprocess.run(
            [
                "npx",
                "-y",
                "@mermaid-js/mermaid-cli",
                "-i",
                str(MMD_PATH),
                "-o",
                str(out),
                "-b",
                "white",
                "-s",
                "3",
            ],
            cwd=str(ROOT),
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )


def replace_slide16_image() -> None:
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides[15]
    picture = None
    for shp in slide.shapes:
        if shp.shape_type == 13:
            picture = shp
            break
    if picture is None:
        raise RuntimeError("Slide 16 picture not found")
    left, top, width, height = picture.left, picture.top, picture.width, picture.height
    picture._element.getparent().remove(picture._element)
    # Slightly reduce height to avoid clipping while preserving readability.
    slide.shapes.add_picture(str(PNG_PATH), left, top, width=width, height=height)
    prs.save(str(PPTX_PATH))


def main() -> None:
    render_diagram()
    replace_slide16_image()
    print(f"Updated {PPTX_PATH}")


if __name__ == "__main__":
    main()
