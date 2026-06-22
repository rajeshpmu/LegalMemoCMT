from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.docx"


def add_bg(slide, prs: Presentation):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), prs.slide_width, PptInches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None):
    tx = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.35), PptInches(12.2), PptInches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = PptPt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)

    if subtitle:
        tx2 = slide.shapes.add_textbox(PptInches(0.58), PptInches(0.87), PptInches(12.1), PptInches(0.45))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = PptPt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x, y, w, h, font_size=14):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        p.line_spacing = 1.06
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = PptPt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_note_box(slide, left, top, width, height, text):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 244, 214)
    shape.line.color.rgb = RGBColor(184, 145, 42)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = PptPt(11)
    r.font.color.rgb = RGBColor(95, 72, 0)
    return shape


def insert_slide_before(prs: Presentation, slide, before_index_0_based: int):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(before_index_0_based, sldId)


def renumber_titles(prs: Presentation):
    seq = 1
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        candidates = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            text = shape.text_frame.text.strip()
            if not text or shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            if float(shape.top) > float(PptInches(1.25)):
                continue
            candidates.append(shape)
        if not candidates:
            continue
        def score(shape):
            max_font = 0.0
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.size:
                        max_font = max(max_font, float(run.font.size))
            return (max_font, -float(shape.top))

        title_shape = max(candidates, key=score)
        tf = title_shape.text_frame
        base = tf.text.strip()
        if "." in base[:6]:
            base = base.split(".", 1)[1].strip()
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{seq}. {base}"
        r.font.name = "Aptos Display"
        r.font.size = PptPt(24)
        r.font.bold = True
        r.font.color.rgb = RGBColor(20, 48, 87)
        seq += 1


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def add_docx_after(paragraph, text: str = "", style: str | None = None):
    new_p = OxmlElement("w:p")
    paragraph._p.addnext(new_p)
    para = paragraph._parent.paragraphs[0]._p.__class__(new_p)  # type: ignore[attr-defined]
    # wrap in a python-docx Paragraph object by using the document API
    return para


def insert_paragraph_after(paragraph, text: str = "", style: str | None = None):
    new_p = OxmlElement("w:p")
    paragraph._p.addnext(new_p)
    from docx.text.paragraph import Paragraph

    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        new_para.add_run(text)
    return new_para


def add_docx_section():
    doc = Document(str(DOCX_PATH))
    if any("1.1 What a Modality Means" in p.text for p in doc.paragraphs):
        return

    anchor = None
    for p in doc.paragraphs:
        if p.text.strip() == "1. Title and Abstract":
            anchor = p
            break
    if anchor is None:
        raise RuntimeError("Could not find the Title and Abstract heading in the DOCX.")

    heading = insert_paragraph_after(anchor, "1.1 What a Modality Means", style="Heading 2")
    p1 = insert_paragraph_after(
        heading,
        "A modality is one type of input signal. In this project, the three main modalities are text, audio, and video. Text carries the transcript, audio carries tone and stress, and video carries facial behavior and visible cues.",
    )
    p2 = insert_paragraph_after(
        p1,
        "The system is multimodal because it combines more than one modality and checks how they agree, disagree, or complement each other when analyzing emotional cues.",
    )
    p3 = insert_paragraph_after(p2, "Simple examples:")
    for text in [
        "Text: what was said in the testimony.",
        "Audio: how it was said, including pitch, pause, and stress.",
        "Video: what was visible on the face and in the upper-body behavior.",
    ]:
        p3 = insert_paragraph_after(p3, text, style="List Bullet")

    table = doc.add_table(rows=1, cols=1)
    table.style = "Table Grid"
    cell = table.cell(0, 0)
    set_cell_shading(cell, "FFF2CC")
    cell.text = ""
    cell.paragraphs[0].add_run("Review note: ").bold = True
    cell.paragraphs[0].add_run(
        "when the thesis says multimodal, it means the model uses more than one kind of evidence rather than relying on only text or only speech."
    )
    cell.paragraphs[0].runs[0].font.name = "Times New Roman"
    cell.paragraphs[0].runs[0].font.size = Pt(10.5)
    cell.paragraphs[0].runs[0].font.color.rgb = DocxRGBColor(90, 68, 0)
    cell.paragraphs[0].runs[1].font.name = "Times New Roman"
    cell.paragraphs[0].runs[1].font.size = Pt(10.5)
    cell.paragraphs[0].runs[1].font.color.rgb = DocxRGBColor(90, 68, 0)
    p3._p.addnext(table._tbl)

    doc.save(str(DOCX_PATH))


def add_pptx_slide():
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(
        slide,
        "4. What a Modality Means",
        "A short glossary note to keep the terminology clear before the architecture slides.",
    )
    add_body(
        slide,
        [
            "A modality is one type of input signal.",
            "In this project, the main modalities are text, audio, and video.",
            "Text carries the transcript, audio carries tone and stress, and video carries facial behavior and visible cues.",
            "The model is multimodal because it combines more than one modality and studies how they agree or disagree.",
        ],
        0.7,
        1.5,
        6.0,
        3.0,
        font_size=14,
    )
    add_note_box(
        slide,
        6.95,
        1.65,
        5.75,
        1.35,
        "Review note: multimodal means the model uses more than one kind of evidence instead of relying on only text or only speech.",
    )
    add_note_box(
        slide,
        6.95,
        3.25,
        5.75,
        1.55,
        "Examples:\n• Text = what was said.\n• Audio = how it was said.\n• Video = what was visible on the face and in the body.",
    )
    add_note_box(
        slide,
        6.95,
        5.05,
        5.75,
        1.05,
        "This definition helps explain why the model can compare different kinds of evidence instead of treating the transcript alone as enough.",
    )
    insert_slide_before(prs, slide, 3)
    renumber_titles(prs)
    prs.save(str(PPTX_PATH))


def main():
    add_pptx_slide()
    add_docx_section()
    print("Added modality glossary to First Review PPTX and DOCX.")


if __name__ == "__main__":
    main()
