from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"


def title_shape(slide):
    candidates = []
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = shp.text_frame.text.strip()
        if not txt:
            continue
        if shp.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if float(shp.top) > 1200000:
            continue
        candidates.append(shp)
    if not candidates:
        return None

    def score(shape):
        max_font = 0.0
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    max_font = max(max_font, float(run.font.size))
        return (max_font, -float(shape.top), -float(shape.left))

    return max(candidates, key=score)


def renumber_titles(prs: Presentation) -> None:
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        tshape = title_shape(slide)
        if tshape is None:
            continue
        raw = tshape.text_frame.text.strip()
        if not raw:
            continue
        base = re.sub(r"^\d+\.\s*", "", raw)
        # Avoid changing the title slide or obvious section markers only.
        tshape.text_frame.clear()
        p = tshape.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"{idx}. {base}"
        r.font.name = "Aptos Display"
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string("122F55")


def main() -> None:
    prs = Presentation(str(PPTX_PATH))
    renumber_titles(prs)
    prs.save(str(PPTX_PATH))
    print(f"Renumbered visible titles in {PPTX_PATH}")


if __name__ == "__main__":
    main()
