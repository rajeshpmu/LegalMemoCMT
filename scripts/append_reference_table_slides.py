from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(20, 48, 87)
SUB_COLOR = RGBColor(72, 72, 72)
BODY_COLOR = RGBColor(35, 35, 35)
BAR_COLOR = RGBColor(20, 48, 87)
WHITE = RGBColor(255, 255, 255)


REFERENCE_SETS = {
    "Zeroth_Review_LegalMemoCMT.pptx": [
        (
            "10. Reading Links for Gap Coverage",
            "These articles close the missing-modality and conversational MER gaps identified in the teaching draft.",
            [
                ["Missing modalities", "Dynamic Modality and View Selection for Multimodal Emotion Recognition with Missing Modalities", "https://arxiv.org/abs/2404.12251", "Dynamic modality selection when one stream is missing or degraded."],
                ["Missing modalities", "Multimodal Prompt Learning with Missing Modalities for Sentiment Analysis and Emotion Recognition", "https://arxiv.org/abs/2407.05374", "Prompt-based missing-modality generation and compact adaptation."],
                ["Missing modalities", "Leveraging Retrieval Augment Approach for Multimodal Emotion Recognition Under Missing Modalities", "https://arxiv.org/abs/2410.02804", "Retrieval-assisted recovery when one modality is absent."],
                ["Missing modalities", "Enhancing Emotion Recognition in Incomplete Data: CM-ARR", "https://arxiv.org/abs/2407.09029", "Cross-modal alignment, reconstruction, and refinement under incomplete inputs."],
                ["Conversational MER", "Multimodal Emotion Recognition in Conversations: A Survey of Methods, Trends, Challenges and Prospects", "https://aclanthology.org/2025.findings-emnlp.332/", "Newest high-level survey for MERC."],
                ["Conversational MER", "Multi-modal emotion recognition in conversation based on prompt learning with text-audio fusion features", "https://www.nature.com/articles/s41598-025-89758-8", "Prompt-based text-audio fusion in conversational ERC."],
            ],
        ),
        (
            "11. Reading Links for Explainability, Bias, and Fusion",
            "These articles support interpretability, fairness, and stronger fusion design for Phase 2.",
            [
                ["Conversational MER", "Cross-modal gated feature enhancement for multimodal emotion recognition in conversations", "https://www.nature.com/articles/s41598-025-11989-6", "Gated fusion and cross-modal interaction in ERC."],
                ["Fusion / adaptation", "Adaptive multimodal transformer based on exchanging for multimodal sentiment analysis", "https://www.nature.com/articles/s41598-025-11848-4", "Adaptive exchange fusion beyond fixed pooling baselines."],
                ["Fusion / adaptation", "Hierarchical cross-modal attention and dual audio pathways for enhanced multimodal sentiment analysis", "https://www.nature.com/articles/s41598-025-09000-3", "Hierarchical attention and stronger audio processing."],
                ["Explainability", "EMERSK -- Explainable Multimodal Emotion Recognition with Situational Knowledge", "https://arxiv.org/abs/2306.08657", "Modular multimodal explanation and situational context."],
                ["Bias / fairness", "Emo-bias: A Large Scale Evaluation of Social Bias on Speech Emotion Recognition", "https://arxiv.org/abs/2406.05065", "Gender bias, dataset bias, and model bias in SER."],
                ["Explainability", "A Review on Explainability in Multimodal Deep Neural Nets", "https://arxiv.org/abs/2105.07878", "General explainability toolkit for multimodal deep models."],
                ["Explainability", "Interpretable Multimodal Emotion Recognition using Facial Features and Physiological Signals", "https://arxiv.org/abs/2306.02845", "Modality contribution analysis and interpretability methods."],
            ],
        ),
    ],
    "First_Review_LegalMemoCMT.pptx": [
        (
            "11. Reading Links for Gap Coverage",
            "These articles close the missing-modality and conversational MER gaps identified in the teaching draft.",
            [
                ["Missing modalities", "Dynamic Modality and View Selection for Multimodal Emotion Recognition with Missing Modalities", "https://arxiv.org/abs/2404.12251", "Dynamic modality selection when one stream is missing or degraded."],
                ["Missing modalities", "Multimodal Prompt Learning with Missing Modalities for Sentiment Analysis and Emotion Recognition", "https://arxiv.org/abs/2407.05374", "Prompt-based missing-modality generation and compact adaptation."],
                ["Missing modalities", "Leveraging Retrieval Augment Approach for Multimodal Emotion Recognition Under Missing Modalities", "https://arxiv.org/abs/2410.02804", "Retrieval-assisted recovery when one modality is absent."],
                ["Missing modalities", "Enhancing Emotion Recognition in Incomplete Data: CM-ARR", "https://arxiv.org/abs/2407.09029", "Cross-modal alignment, reconstruction, and refinement under incomplete inputs."],
                ["Conversational MER", "Multimodal Emotion Recognition in Conversations: A Survey of Methods, Trends, Challenges and Prospects", "https://aclanthology.org/2025.findings-emnlp.332/", "Newest high-level survey for MERC."],
                ["Conversational MER", "Multi-modal emotion recognition in conversation based on prompt learning with text-audio fusion features", "https://www.nature.com/articles/s41598-025-89758-8", "Prompt-based text-audio fusion in conversational ERC."],
            ],
        ),
        (
            "12. Reading Links for Explainability, Bias, and Fusion",
            "These articles support interpretability, fairness, and stronger fusion design for Phase 2.",
            [
                ["Conversational MER", "Cross-modal gated feature enhancement for multimodal emotion recognition in conversations", "https://www.nature.com/articles/s41598-025-11989-6", "Gated fusion and cross-modal interaction in ERC."],
                ["Fusion / adaptation", "Adaptive multimodal transformer based on exchanging for multimodal sentiment analysis", "https://www.nature.com/articles/s41598-025-11848-4", "Adaptive exchange fusion beyond fixed pooling baselines."],
                ["Fusion / adaptation", "Hierarchical cross-modal attention and dual audio pathways for enhanced multimodal sentiment analysis", "https://www.nature.com/articles/s41598-025-09000-3", "Hierarchical attention and stronger audio processing."],
                ["Explainability", "EMERSK -- Explainable Multimodal Emotion Recognition with Situational Knowledge", "https://arxiv.org/abs/2306.08657", "Modular multimodal explanation and situational context."],
                ["Bias / fairness", "Emo-bias: A Large Scale Evaluation of Social Bias on Speech Emotion Recognition", "https://arxiv.org/abs/2406.05065", "Gender bias, dataset bias, and model bias in SER."],
                ["Explainability", "A Review on Explainability in Multimodal Deep Neural Nets", "https://arxiv.org/abs/2105.07878", "General explainability toolkit for multimodal deep models."],
                ["Explainability", "Interpretable Multimodal Emotion Recognition using Facial Features and Physiological Signals", "https://arxiv.org/abs/2306.02845", "Modality contribution analysis and interpretability methods."],
            ],
        ),
    ],
}


def add_bg(slide, prs):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BAR_COLOR
    bar.line.fill.background()


def add_title(slide, title, subtitle):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(0.8), Inches(12.1), Inches(0.55))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Aptos"
    r2.font.size = Pt(11.5)
    r2.font.italic = True
    r2.font.color.rgb = SUB_COLOR


def add_table_slide(prs, title, subtitle, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, title, subtitle)

    table = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.35), Inches(1.35), Inches(12.6), Inches(5.8)).table
    widths = [1.25, 4.0, 3.0, 4.35]
    for idx, w in enumerate(widths):
        table.columns[idx].width = Inches(w)

    headers = ["Gap area", "Article", "Link", "Why to read"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Aptos"
                r.font.size = Pt(10)
                r.font.bold = True
                r.font.color.rgb = TITLE_COLOR

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.space_after = Pt(0)
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(8.5)
                    run.font.color.rgb = BODY_COLOR
        if r_idx % 2 == 0:
            for c_idx in range(4):
                table.cell(r_idx, c_idx).fill.solid()
                table.cell(r_idx, c_idx).fill.fore_color.rgb = RGBColor(245, 247, 250)


def main():
    for name, slidespecs in REFERENCE_SETS.items():
        path = Path("implementation_docments") / name
        prs = Presentation(str(path))
        for title, subtitle, rows in slidespecs:
            add_table_slide(prs, title, subtitle, rows)
        prs.save(str(path))
        print("updated", path)


if __name__ == "__main__":
    main()
