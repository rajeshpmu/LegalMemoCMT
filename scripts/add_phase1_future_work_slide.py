from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocRGBColor
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Prep_Guide.docx"


def _set_text_style(run, *, name="Aptos", size=18, bold=False, color="000000"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color)


def _add_box(slide, left, top, width, height, fill="FFFFFF", line="FFFFFF", transparency=0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor.from_string(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = PptRGBColor.from_string(line)
    return shape


def _add_textbox(slide, left, top, width, height, text, *, font_size=18, bold=False, color="000000",
                 align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _set_text_style(r, size=font_size, bold=bold, color=color)
    return tx


def _add_bullets(slide, left, top, width, height, bullets, *, title=None):
    box = _add_box(slide, left, top, width, height, fill="F7F9FC", line="D7E1EE")
    tf = slide.shapes.add_textbox(
        PptInches(left + 0.15),
        PptInches(top + 0.12),
        PptInches(width - 0.3),
        PptInches(height - 0.24),
    ).text_frame
    tf.word_wrap = True
    tf.clear()
    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        _set_text_style(r, size=15, bold=True, color="122F55")
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(4)
            if p.runs:
                _set_text_style(p.runs[0], size=13, color="404040")
    else:
        first = True
        for bullet in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = bullet
            p.level = 0
            p.space_after = Pt(4)
            if p.runs:
                _set_text_style(p.runs[0], size=13, color="404040")
    return box


def update_pptx():
    prs = Presentation(str(PPTX_PATH))
    future_title = "Potential Improvements / Next Steps"
    if any(
        future_title.lower() in (shape.text_frame.text if getattr(shape, "has_text_frame", False) else "").lower()
        for slide in prs.slides
        for shape in slide.shapes
    ):
        print("PPTX already contains the future-work slide; skipping add.")
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Header band, matching the rest of the deck.
    _add_box(slide, 0.0, 0.0, 13.333, 0.58, fill="122F55", line="122F55")
    _add_box(slide, 0.0, 0.58, 13.333, 0.08, fill="267377", line="267377")

    title = _add_textbox(
        slide, 0.45, 0.12, 11.8, 0.3, future_title, font_size=18, bold=True, color="FFFFFF"
    )
    number = _add_textbox(slide, 12.0, 0.12, 0.9, 0.3, str(len(prs.slides)), font_size=12, color="FFFFFF")
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    number.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    _add_bullets(
        slide,
        0.45,
        1.05,
        5.95,
        5.6,
        [
            "Helpful because the current model still shows neutral bias, long-range dialogue dependence, and speaker-to-speaker ambiguity.",
            "Implement after the Phase 1 benchmark is frozen; keep Phase 1 claims stable before making structural changes.",
            "If time remains in Phase 1, only prototype a small dialogue-context ablation. Keep the full redesign for Phase 2 groundwork.",
        ],
        title="Why this belongs next",
    )

    _add_bullets(
        slide,
        6.65,
        1.05,
        6.15,
        5.6,
        [
            "Add speaker_id, dialogue_id, and turn_index to the manifest or dataset object.",
            "Create a local context window over previous turns and add a speaker embedding or dialogue-state encoder in src/models/model.py.",
            "For Phase 2, add courtroom-role metadata (witness, judge, attorney, defendant) before the context encoder.",
            "Evaluate with the same folds, weighted F1, macro F1, confusion matrix, and top-confusion analysis so the comparison stays fair.",
        ],
        title="How to start technically",
    )

    callout = _add_box(slide, 0.9, 6.9, 11.6, 0.75, fill="EEF4FA", line="C9D7E6")
    tx = slide.shapes.add_textbox(PptInches(1.1), PptInches(7.02), PptInches(11.2), PptInches(0.45))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run()
    r.text = "Best timing: lock Phase 1 now, then use this as the first context-aware extension before Phase 2 courtroom adaptation."
    _set_text_style(r, size=14, bold=True, color="122F55")

    prs.save(str(PPTX_PATH))
    print(f"Updated PPTX: {PPTX_PATH}")


def _add_paragraph(doc, text, *, bold=False, size=11.5, color="000000", style=None):
    p = doc.add_paragraph(style=style)
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.space_before = Pt(0)
    r = p.add_run(text)
    _set_doc_run_style(r, bold=bold, size=size, color=color)
    return p


def _set_doc_run_style(run, *, bold=False, size=11.5, color="000000"):
    run.font.name = "Aptos"
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = DocRGBColor.from_string(color)


def _insert_paragraph_before(target_paragraph, text, style=None):
    new_p = OxmlElement("w:p")
    target_paragraph._p.addprevious(new_p)
    paragraph = Paragraph(new_p, target_paragraph._parent)
    if style is not None:
        paragraph.style = style
    paragraph.text = text
    return paragraph


def update_docx():
    doc = Document(str(DOCX_PATH))
    heading_text = "2.11 Slide 25 Potential Improvements / Next Steps"
    existing_indices = [i for i, p in enumerate(doc.paragraphs) if heading_text in p.text]

    # If the section already exists from an earlier run, remove the full block so
    # we can reinsert it in the correct order.
    if existing_indices:
        heading_idx = existing_indices[0]
        start_idx = max(0, heading_idx - 7)
        end_idx = heading_idx
        for idx in range(end_idx, start_idx - 1, -1):
            p = doc.paragraphs[idx]
            p._element.getparent().remove(p._element)

    anchor = None
    for p in doc.paragraphs:
        if p.text.strip().startswith("3. Demo SOP You Should Practise"):
            anchor = p
            break
    if anchor is None:
        anchor = doc.add_paragraph("")

    paragraphs = [
        (heading_text, True, 13.0, "122F55"),
        (
            "This slide should be presented as future work, not as part of the already-claimed Phase 1 result. "
            "The current Phase 1 baseline is strongest when it is frozen and defended with its own metrics. "
            "These improvements are useful because they target the same failure pattern still visible in MELD: "
            "neutral bias, speaker-context dependence, and emotion transitions across turns.",
            False, 11.5, "000000",
        ),
        ("When to implement:", True, 11.5, "122F55"),
        (
            "If Phase 1 still has buffer, prototype only the first two items as small ablations on MELD. "
            "The full version should be implemented after Phase 1 is locked and before Phase 2 courtroom adaptation begins. "
            "Courtroom-role metadata belongs mainly to Phase 2 because MELD does not contain legal roles.",
            False, 11.5, "000000",
        ),
        ("How to start technically:", True, 11.5, "122F55"),
        (
            "1) Extend the manifest or dataset record with speaker_id, dialogue_id, turn_index, and a role placeholder where available.\n"
            "2) Build a short dialogue-history window so each utterance can see previous turns.\n"
            "3) Add a speaker embedding or dialogue-context encoder in src/models/model.py.\n"
            "4) Keep the existing text/audio/video branches unchanged for the first ablation, then compare baseline vs +speaker context vs +shift modeling.\n"
            "5) For Phase 2, add courtroom-role parsing and role-aware segmentation before the context encoder.",
            False, 11.5, "000000",
        ),
        ("Where the code changes would start:", True, 11.5, "122F55"),
        (
            "scripts/build_meld_cv_folds.py already preserves dialogue-safe grouping, but it does not yet build explicit speaker memory.\n"
            "src/data/preprocessing.py and the manifest builders are where speaker / turn metadata would be injected.\n"
            "src/models/model.py is where a context encoder, speaker embedding, or role-aware branch would live.\n"
            "src/train/train.py would control the ablations and loss configuration, while src/train/evaluate.py would compare confusion matrices and per-class changes.",
            False, 11.5, "000000",
        ),
        (
            "Best timing: lock Phase 1 now, then use this as the first context-aware extension before Phase 2 courtroom adaptation.",
            True, 11.5, "122F55",
        ),
    ]

    for text, bold, size, color in paragraphs:
        p = _insert_paragraph_before(anchor, text)
        p.paragraph_format.space_after = Pt(4)
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.line_spacing = 1.0
        if p.runs:
            _set_doc_run_style(p.runs[0], bold=bold, size=size, color=color)
        if "\n" in text:
            # Apply font to subsequent runs if Word splits the paragraph text.
            for run in p.runs[1:]:
                _set_doc_run_style(run, bold=bold, size=size, color=color)

    doc.save(str(DOCX_PATH))
    print(f"Updated DOCX: {DOCX_PATH}")


if __name__ == "__main__":
    update_pptx()
    update_docx()
