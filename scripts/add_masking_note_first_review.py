from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.docx"

NOTE_TEXT = (
    "Masking note: padding tokens are ignored so mean, max, and min only summarize valid text or audio positions."
)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_margins(cell, top=90, start=110, bottom=90, end=110) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcMar = tcPr.first_child_found_in("w:tcMar")
    if tcMar is None:
        tcMar = OxmlElement("w:tcMar")
        tcPr.append(tcMar)
    for m, val in [("top", top), ("start", start), ("bottom", bottom), ("end", end)]:
        node = tcMar.find(qn(f"w:{m}"))
        if node is None:
            node = OxmlElement(f"w:{m}")
            tcMar.append(node)
        node.set(qn("w:w"), str(val))
        node.set(qn("w:type"), "dxa")


def add_docx_note() -> None:
    doc = Document(str(DOCX_PATH))
    if any(NOTE_TEXT in p.text for p in doc.paragraphs):
        return

    target_para = None
    for p in doc.paragraphs:
        if "For the text-plus-audio paper-style branch, masking matters" in p.text:
            target_para = p
            break
    if target_para is None:
        raise RuntimeError("Could not find the masking explanation paragraph in DOCX.")

    table = doc.add_table(rows=1, cols=1)
    table.style = "Table Grid"
    cell = table.cell(0, 0)
    set_cell_shading(cell, "FFF2CC")
    set_cell_margins(cell)
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run("Masking note: ")
    r1.bold = True
    r1.font.name = "Times New Roman"
    r1.font.size = Pt(10.5)
    r1.font.color.rgb = DocxRGBColor(70, 54, 0)
    r2 = p.add_run("padding tokens are ignored so mean, max, and min only summarize valid text or audio positions.")
    r2.font.name = "Times New Roman"
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = DocxRGBColor(70, 54, 0)

    target_para._p.addnext(table._tbl)
    doc.save(str(DOCX_PATH))


def add_pptx_note() -> None:
    prs = Presentation(str(PPTX_PATH))
    target_slide = None
    for slide in prs.slides:
        title = ""
        for sh in slide.shapes:
            if hasattr(sh, "text_frame") and sh.text_frame and sh.text_frame.text.strip():
                txt = sh.text_frame.text.strip().split("\n", 1)[0]
                if txt[:1].isdigit() and "Pretrained Paper-Style Pooling Explained" in txt:
                    title = txt
                    break
        if title:
            target_slide = slide
            break
    if target_slide is None:
        raise RuntimeError("Could not find the pooling slide.")

    if any(NOTE_TEXT in sh.text_frame.text for sh in target_slide.shapes if hasattr(sh, "text_frame") and sh.text_frame):
        prs.save(str(PPTX_PATH))
        return

    bullet_box = None
    for sh in target_slide.shapes:
        if hasattr(sh, "text_frame") and sh.text_frame and "In the paper-style branch, pooling is applied" in sh.text_frame.text:
            bullet_box = sh
            break
    if bullet_box is None:
        raise RuntimeError("Could not locate the bullet text box on the pooling slide.")

    note = target_slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PptInches(0.78),
        PptInches(1.34),
        PptInches(5.85),
        PptInches(0.54),
    )
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(255, 242, 204)
    note.line.color.rgb = RGBColor(185, 150, 50)

    tf = note.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "Masking note: "
    r1.font.name = "Aptos"
    r1.font.size = PptPt(11)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(90, 68, 0)
    r2 = p.add_run()
    r2.text = "padding tokens are ignored so mean, max, and min only summarize valid text or audio positions."
    r2.font.name = "Aptos"
    r2.font.size = PptPt(11)
    r2.font.color.rgb = RGBColor(90, 68, 0)

    # Make room for the note without changing the slide structure.
    bullet_box.top = PptInches(1.96)
    bullet_box.height = PptInches(4.55)

    prs.save(str(PPTX_PATH))


def main() -> None:
    add_pptx_note()
    add_docx_note()
    print("Added masking note to First Review PPTX and DOCX.")


if __name__ == "__main__":
    main()
