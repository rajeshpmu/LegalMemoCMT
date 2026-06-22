from __future__ import annotations

import json
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "implementation_docments"
FIG_DIR = OUT_DIR / "figures" / "guidance_advanced_ai_ml"
FIG_DIR.mkdir(parents=True, exist_ok=True)

PPTX_PATH = OUT_DIR / "Guidance_Call_Advanced_AI_ML_Recommendations.pptx"
DOCX_PATH = OUT_DIR / "Guidance_Call_Advanced_AI_ML_Recommendations.docx"

FOLD2_METRICS = ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "fold_2" / "metrics.json"
FOLD4_METRICS = ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "fold_4" / "metrics.json"
CREMA_METRICS = ROOT / "results" / "paper_aligned_crema_d" / "cmt_min" / "metrics.json"
FOLD2_CM_CSV = ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "fold_2" / "analysis_test" / "confusion_matrix.csv"
FOLD4_CM_CSV = ROOT / "results" / "paper_aligned_meld_cv" / "cmt_min" / "fold_4" / "analysis_test" / "confusion_matrix.csv"
CREMA_CM_CSV = ROOT / "results" / "paper_aligned_crema_d" / "cmt_min" / "analysis_test" / "confusion_matrix.csv"

FOLD2_CM_PNG = FIG_DIR / "meld_fold2_confusion_matrix.png"
FOLD4_CM_PNG = FIG_DIR / "meld_fold4_confusion_matrix.png"
CREMA_CM_PNG = FIG_DIR / "crema_d_confusion_matrix.png"


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def clean_label(value: str) -> str:
    s = str(value)
    return s.split(":", 1)[1] if ":" in s else s


def build_confusion_image(csv_path: Path, png_path: Path, title: str) -> None:
    df = pd.read_csv(csv_path, index_col=0)
    df.index = [clean_label(v) for v in df.index]
    df.columns = [clean_label(v) for v in df.columns]

    fig, ax = plt.subplots(figsize=(8.5, 6.5))
    im = ax.imshow(df.values, cmap="Blues", aspect="auto")
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            ax.text(j, i, f"{int(df.values[i, j])}", ha="center", va="center", fontsize=8)
    ax.set_xticks(range(len(df.columns)))
    ax.set_xticklabels(df.columns, rotation=45, ha="right")
    ax.set_yticks(range(len(df.index)))
    ax.set_yticklabels(df.index)
    ax.set_title(title, fontsize=13)
    ax.set_xlabel("Predicted label")
    ax.set_ylabel("Actual label")
    ax.set_xlim(-0.5, len(df.columns) - 0.5)
    ax.set_ylim(len(df.index) - 0.5, -0.5)
    fig.tight_layout()
    fig.savefig(png_path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def set_ppt_style(prs: Presentation) -> None:
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)


def ppt_bg(slide, prs: Presentation):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), prs.slide_width, PptInches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def ppt_title(slide, title: str, subtitle: str | None = None, y: float = 0.35):
    tx = slide.shapes.add_textbox(PptInches(0.55), PptInches(y), PptInches(12.2), PptInches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = PptPt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)
    if subtitle:
        tx2 = slide.shapes.add_textbox(PptInches(0.58), PptInches(y + 0.52), PptInches(12.1), PptInches(0.55))
        tf2 = tx2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = PptPt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def ppt_bullets(slide, bullets, x=0.75, y=1.35, w=11.95, h=5.65, font_size=18):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = PptPt(6)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = PptPt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)


def ppt_add_picture(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), PptInches(left), PptInches(top), width=PptInches(width))


def add_slide(prs: Presentation, title: str, subtitle: str, bullets, image: Path | None = None, caption: str | None = None, font_size=16):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ppt_bg(slide, prs)
    ppt_title(slide, title, subtitle)
    if image is None:
        ppt_bullets(slide, bullets, x=0.75, y=1.5, w=11.9, h=5.55, font_size=font_size)
    else:
        ppt_bullets(slide, bullets, x=0.7, y=1.5, w=6.1, h=5.25, font_size=font_size)
        ppt_add_picture(slide, image, 7.05, 1.7, 5.55)
        if caption:
            tx = slide.shapes.add_textbox(PptInches(7.05), PptInches(6.4), PptInches(5.55), PptInches(0.35))
            tf = tx.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = caption
            r.font.name = "Aptos"
            r.font.size = PptPt(10.5)
            r.font.italic = True
            r.font.color.rgb = RGBColor(95, 95, 95)
    return slide


def add_table_slide(prs: Presentation, title: str, subtitle: str, headers: list[str], rows: list[list[str]], footer: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ppt_bg(slide, prs)
    ppt_title(slide, title, subtitle)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), PptInches(0.55), PptInches(1.55), PptInches(12.25), PptInches(4.8)).table
    widths = [0.8, 1.5, 1.45, 1.45, 1.55, 5.5]
    for i, header in enumerate(headers):
        table.columns[i].width = PptInches(widths[i])
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(20, 48, 87)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Aptos"
                r.font.size = PptPt(10.5)
                r.font.bold = True
                r.font.color.rgb = RGBColor(255, 255, 255)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_idx < len(row) - 1 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Aptos"
                    r.font.size = PptPt(10.25)
                    r.font.color.rgb = RGBColor(35, 35, 35)
    if footer:
        tx = slide.shapes.add_textbox(PptInches(0.7), PptInches(6.4), PptInches(11.9), PptInches(0.4))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = footer
        r.font.name = "Aptos"
        r.font.size = PptPt(10.5)
        r.font.italic = True
        r.font.color.rgb = RGBColor(95, 95, 95)


def configure_doc(doc: Document) -> None:
    styles = doc.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"].font.size = Pt(12)
    for name in ["Title", "Heading 1", "Heading 2", "Heading 3"]:
        if name in styles:
            styles[name].font.name = "Times New Roman"
    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.8)
    section.left_margin = Inches(0.95)
    section.right_margin = Inches(0.95)


def dpara(doc: Document, text: str, *, italic: bool = False) -> None:
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.font.name = "Times New Roman"
    r.font.size = Pt(12)
    r.italic = italic


def dbullets(doc: Document, items: list[str]) -> None:
    for item in items:
        doc.add_paragraph(item, style="List Bullet")


def dtable(doc: Document, headers: list[str], rows: list[list[str]]) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = value
    doc.add_paragraph()


def build_ppt() -> None:
    fold2 = read_json(FOLD2_METRICS)
    fold4 = read_json(FOLD4_METRICS)
    crema = read_json(CREMA_METRICS)
    prs = Presentation()
    set_ppt_style(prs)

    add_slide(
        prs,
        "Guidance Call: Advanced AI/ML Next Steps",
        "Student-level recommendation deck based on the current MELD and CREMA-D error patterns.",
        [
            "The current results are stable enough to trust, but the model still confuses the same classes repeatedly.",
            "MELD shows a strong but biased pattern, while CREMA-D is much closer to collapse.",
            "The next step should change the loss, the representation, and the context modeling instead of only tuning hyperparameters.",
        ],
        font_size=16,
    )

    add_table_slide(
        prs,
        "Current Evidence",
        "These measured results are the basis for the recommendations.",
        ["Source", "Acc / W-Acc", "UW-Acc", "Macro F1", "Weighted F1", "Interpretation"],
        [
            ["MELD Fold 2", f"{fold2['accuracy']:.4f}", f"{fold2['unweighted_accuracy']:.4f}", f"{fold2['macro_f1']:.4f}", f"{fold2['weighted_f1']:.4f}", "Best weighted anchor; still neutral-heavy"],
            ["MELD Fold 4", f"{fold4['accuracy']:.4f}", f"{fold4['unweighted_accuracy']:.4f}", f"{fold4['macro_f1']:.4f}", f"{fold4['weighted_f1']:.4f}", "Best balanced-class anchor"],
            ["CREMA-D", f"{crema['accuracy']:.4f}", f"{crema['unweighted_accuracy']:.4f}", f"{crema['macro_f1']:.4f}", f"{crema['weighted_f1']:.4f}", "Near-collapse into one class"],
        ],
        footer="Fold 2 gives the best weighted anchor; Fold 4 gives the best balanced-class anchor. CREMA-D still needs stabilization before it can be trusted as a strong secondary benchmark.",
    )

    add_slide(
        prs,
        "MELD Fold 2: What the Matrix Means",
        "Fold 2 is the strongest weighted-comparison fold, but it still over-predicts neutral.",
        [
            "The model learns broad affect, but it still merges neutral, joy, anger, and surprise too often.",
            "That means the embedding space is not separating fine-grained emotions strongly enough.",
            "This is the best place to test class-balanced focal loss and supervised contrastive learning.",
        ],
        image=FOLD2_CM_PNG,
        caption="MELD Fold 2 confusion matrix",
        font_size=15,
    )

    add_slide(
        prs,
        "MELD Fold 4: Why It Matters",
        "Fold 4 is the strongest balanced-class fold, so it is the best evidence that the model can improve UW-Acc-like behavior.",
        [
            "Fold 4 gives the best UW-Acc and macro F1 in the 5-fold report.",
            "That means the model can behave more fairly when the metric gives each class equal importance.",
            "The goal should be to keep this balanced behavior while recovering stronger weighted performance.",
        ],
        image=FOLD4_CM_PNG,
        caption="MELD Fold 4 confusion matrix",
        font_size=15,
    )

    add_slide(
        prs,
        "CREMA-D: Why It Needs Stabilization",
        "CREMA-D is currently near collapse, so the first fix should be conservative fine-tuning and class-balanced optimization.",
        [
            "The matrix suggests the model is overconfident in a small set of classes.",
            "Before changing architecture, use gradual unfreezing, discriminative learning rates, and class-balanced focal loss.",
            "If collapse remains, add supervised contrastive learning and class-aware sampling.",
        ],
        image=CREMA_CM_PNG,
        caption="CREMA-D confusion matrix",
        font_size=15,
    )

    add_slide(
        prs,
        "Recommended Next-Step Stack",
        "This is the order I would implement the improvements in.",
        [
            "1. Class-balanced focal loss or LDAM-DRW to fix the loss function.",
            "2. Supervised contrastive learning on the fused embedding to separate emotion clusters.",
            "3. Dialogue-context modeling for MELD so turns are interpreted with local context.",
            "4. Gradual unfreezing and smaller learning rates for CREMA-D.",
            "5. Calibration and hard-example mining after the main objective changes.",
        ],
        font_size=15,
    )

    add_slide(
        prs,
        "1. Class-Balanced Focal Loss",
        "This is the best first change because both MELD and CREMA-D are showing imbalance-driven collapse.",
        [
            "Focal loss downweights easy examples and focuses learning on hard ones.",
            "Add class weights so minority emotions matter more in the gradient.",
            "Implementation: cross-entropy multiplied by (1 - p_t)^gamma, plus alpha weights from class frequency.",
            "Expected effect: fewer false neutral predictions and better macro F1 / UW-Acc.",
        ],
        font_size=15,
    )

    add_slide(
        prs,
        "2. Supervised Contrastive Learning",
        "This is the best representation-learning upgrade because the confusion matrix suggests the embedding space is not separating classes well.",
        [
            "Add a projection head on top of the fused embedding.",
            "Use supervised contrastive loss so same-label examples are pulled together and different-label examples are pushed apart.",
            "Keep classification loss as the main objective and add the contrastive term as an auxiliary loss.",
            "Expected effect: clearer class clusters, especially between neutral and the nearby emotions joy, anger, and surprise.",
        ],
        font_size=15,
    )

    add_slide(
        prs,
        "3. Dialogue Context Modeling for MELD",
        "MELD errors are often context errors, not just feature errors.",
        [
            "Use a local context window around each utterance.",
            "Add speaker embeddings so repeated speakers are represented consistently.",
            "A hierarchical Transformer or dialogue-level attention block can encode turn history before classification.",
            "Expected effect: fewer neutral-to-emotion confusions caused by missing context.",
        ],
        font_size=15,
    )

    add_slide(
        prs,
        "4. CREMA-D Stabilization Plan",
        "CREMA-D needs a conservative fine-tuning regime before more advanced tricks will help.",
        [
            "Freeze the pretrained backbones for a warm-up period, then unfreeze gradually.",
            "Use a smaller learning rate for BERT/HuBERT than for the classifier head.",
            "Use balanced sampling so the minority classes are not starved during training.",
            "Expected effect: stop the collapse and recover a wider class distribution.",
        ],
        font_size=15,
    )

    add_slide(
        prs,
        "Implementation Order",
        "Do not change everything at once; the confusion matrices should guide each ablation.",
        [
            "First run class-balanced focal loss on MELD and CREMA-D.",
            "Then add supervised contrastive learning to the fused representation.",
            "Next add dialogue context for MELD only.",
            "Finally stabilize CREMA-D with gradual unfreezing and discriminative learning rates.",
        ],
        font_size=15,
    )

    add_slide(
        prs,
        "What to Ask in the Guidance Call",
        "The mentor conversation should focus on the best next algorithmic move, not on changing everything at once.",
        [
            "Should the project prioritize class-balanced focal loss or LDAM-DRW first?",
            "Is supervised contrastive learning worth the extra complexity for MELD and CREMA-D in this timeline?",
            "Should MELD context modeling be added before any further tuning of the pretrained encoders?",
            "Is CREMA-D worth stabilizing further, or should it stay a secondary diagnostic benchmark?",
        ],
        font_size=15,
    )

    prs.save(PPTX_PATH)


def build_doc() -> None:
    fold2 = read_json(FOLD2_METRICS)
    fold4 = read_json(FOLD4_METRICS)
    crema = read_json(CREMA_METRICS)
    doc = Document()
    configure_doc(doc)

    doc.add_heading("Guidance Call: Advanced AI/ML Next Steps", level=0)
    dpara(
        doc,
        "Purpose: this document explains what the current confusion matrices are telling us, why the model still makes the same kinds of mistakes, and which advanced AI/ML methods are most likely to reduce those errors. It is written at student level but includes enough implementation detail to discuss the next steps with a mentor or guide.",
    )

    doc.add_heading("1. Why the Current Results Point to a Next-Step Change", level=1)
    dpara(
        doc,
        "The current results are not random noise. MELD Fold 2 and Fold 4 show a stable trade-off: Fold 2 gives the strongest weighted aggregate result, while Fold 4 gives the strongest class-balanced result. CREMA-D is weaker still, with a near-collapse pattern. That tells us the model is learning something useful, but the objective and representation are still not strong enough to separate emotion classes cleanly.",
    )
    dtable(
        doc,
        ["Source", "Acc / W-Acc", "UW-Acc", "Macro F1", "Weighted F1", "Reading"],
        [
            ["MELD Fold 2", f"{fold2['accuracy']:.4f}", f"{fold2['unweighted_accuracy']:.4f}", f"{fold2['macro_f1']:.4f}", f"{fold2['weighted_f1']:.4f}", "Best weighted anchor; still neutral-heavy"],
            ["MELD Fold 4", f"{fold4['accuracy']:.4f}", f"{fold4['unweighted_accuracy']:.4f}", f"{fold4['macro_f1']:.4f}", f"{fold4['weighted_f1']:.4f}", "Best balanced-class anchor"],
            ["CREMA-D", f"{crema['accuracy']:.4f}", f"{crema['unweighted_accuracy']:.4f}", f"{crema['macro_f1']:.4f}", f"{crema['weighted_f1']:.4f}", "Near-collapse into one class"],
        ],
    )
    dpara(
        doc,
        "The practical implication is that better tuning alone is unlikely to solve the confusion matrices. The next step should change the learning signal and the way the model represents emotion. That is why the recommendations below combine loss design, representation learning, and context modeling.",
    )

    doc.add_heading("2. What the Confusion Matrices Are Telling Us", level=1)
    dpara(
        doc,
        "MELD Fold 2 is the paper-comparison fold. It still confuses neutral with joy, anger, and surprise, which means the model is not separating subtle emotion boundaries strongly enough. MELD Fold 4 shows that the model can behave more fairly across classes, so the problem is not complete failure. CREMA-D is stronger evidence of under-separation: the model collapses into one class far too often, which is a sign that the optimization regime and class balance are not yet right.",
    )
    dbullets(
        doc,
        [
            "Fold 2: best overall / weighted performance, but still neutral-heavy.",
            "Fold 4: best balanced-class behavior, so class-wise learning is possible.",
            "CREMA-D: collapse means the model is overfitting one dominant rule or adapting too aggressively.",
        ],
    )

    doc.add_heading("3. Best Advanced AI/ML Recommendations", level=1)
    doc.add_heading("3.1 Class-Balanced Focal Loss", level=2)
    dpara(
        doc,
        "This is the first and most practical change. Focal loss downweights easy examples and focuses the gradient on hard examples. On an imbalanced emotion dataset, that means the model spends less effort on easy neutral samples and more effort on the confusing minority emotions. In code, you start with cross-entropy, multiply by (1 - p_t)^gamma, and add class weights alpha based on class frequency. A gamma value around 2 is a common starting point.",
    )
    dbullets(
        doc,
        [
            "Implementation idea: compute class frequencies from the training manifest.",
            "Use alpha weights so rare classes contribute more to the loss.",
            "Use gamma to reduce the gradient from easy examples.",
            "Expected outcome: fewer false neutral predictions and better macro F1 / UW-Acc.",
        ],
    )
    dpara(
        doc,
        "A stronger option is LDAM-DRW if you want a margin-based imbalance fix. LDAM increases the decision margin required for minority classes, and DRW switches on re-weighting later in training. If the goal is a cleaner confusion matrix, class-balanced focal loss is easier to add first, while LDAM-DRW is the more advanced follow-up.",
    )

    doc.add_heading("3.2 Supervised Contrastive Learning", level=2)
    dpara(
        doc,
        "The confusion matrix suggests the fused representation does not separate classes enough. Supervised contrastive learning fixes that by pulling same-label samples together and pushing different-label samples apart in the embedding space. The simplest implementation is to add a projection head after the fused CMT vector and compute a contrastive loss on the normalized embeddings before the final classifier. This can be combined with the standard classification loss.",
    )
    dbullets(
        doc,
        [
            "Add a small projection head on top of the fused representation.",
            "Normalize the projected embeddings and compute contrastive loss by label.",
            "Keep classification loss as the main objective and contrastive loss as an auxiliary term.",
            "Expected outcome: cleaner clusters for neutral, joy, anger, sadness, surprise, and fewer overlaps.",
        ],
    )
    dpara(
        doc,
        "This method is especially useful when the model is learning broad affect but not fine-grained class identity. That is exactly what the MELD fold analyses show. If the representation becomes more separable, the confusion matrix usually improves before the headline accuracy does.",
    )

    doc.add_heading("3.3 Dialogue Context Modeling for MELD", level=2)
    dpara(
        doc,
        "MELD is a dialogue dataset, so the current utterance is often not enough. A sentence that looks neutral in isolation may be anger or joy when read with the previous turn. A context-aware model should therefore encode a short window of surrounding turns. The cleanest implementation is a hierarchical model: utterance encoder first, then dialogue-level context encoder over a sliding window, then fusion/classification.",
    )
    dbullets(
        doc,
        [
            "Use a local context window of previous and next turns where available.",
            "Add speaker embeddings so repeated speakers are represented consistently.",
            "Use a dialogue Transformer or graph attention layer over turn embeddings.",
            "Expected outcome: fewer neutral-to-emotion confusions caused by missing context.",
        ],
    )
    dpara(
        doc,
        "This is the most natural way to reduce MELD-specific confusion. The Fold 2 and Fold 4 errors are not just random misclassifications; they are often context-sensitive mistakes. That means dialogue-level modeling should improve the model in a way that pure loss tuning cannot.",
    )

    doc.add_heading("3.4 CREMA-D Stabilization", level=2)
    dpara(
        doc,
        "CREMA-D is currently too close to a one-class collapse, so the first priority is to make fine-tuning conservative and stable. Use a lower learning rate for the pretrained encoders, a higher learning rate for the classifier head, and gradual unfreezing. If the backbones are changing too quickly, the classifier never gets a stable representation to learn from.",
    )
    dbullets(
        doc,
        [
            "Start with the encoder frozen and train only the classifier head.",
            "Unfreeze one backbone stage at a time and use a smaller encoder learning rate.",
            "Use class-balanced focal loss or balanced softmax to prevent class collapse.",
            "Expected outcome: predictions stop collapsing into one class and class separation improves.",
        ],
    )
    dpara(
        doc,
        "Balanced sampling can help here as well, especially if some CREMA-D classes are underrepresented in the current fold splits. The point is not to make the model memorize the minority class. The point is to make the optimization stable enough that the minority classes stay visible during training.",
    )

    doc.add_heading("3.5 Calibration and Hard-Example Mining", level=2)
    dpara(
        doc,
        "After the main objective and representation are improved, calibration and hard-example mining are the next helpful steps. Temperature scaling can reduce overconfident wrong predictions, and hard-example mining can focus training on the exact examples that cause confusion. These are usually refinement steps, not the first fix, but they are valuable once the model is no longer collapsing.",
    )
    dbullets(
        doc,
        [
            "Temperature scaling helps when predictions are too confident in the wrong class.",
            "Hard-example mining can upweight the utterances repeatedly confused in the matrix.",
            "Curriculum learning can be used if you want to move from easier to harder emotion distinctions.",
        ],
    )

    doc.add_heading("4. Recommended Implementation Order", level=1)
    for item in [
        "Add class-balanced focal loss first and re-run MELD Fold 2 and Fold 4 plus the CREMA-D run.",
        "If the confusion matrix is still too mixed, add supervised contrastive learning on the fused embedding.",
        "For MELD, add dialogue context modeling so the model sees neighboring turns.",
        "For CREMA-D, stabilize fine-tuning with gradual unfreezing and discriminative learning rates.",
        "If needed, add calibration and hard-example mining as a final refinement pass.",
    ]:
        doc.add_paragraph(item, style="List Number")
    dpara(
        doc,
        "This order is important. It starts with the simplest high-impact change and moves toward more structural changes only if the confusion matrix still does not improve. That is the safest way to iterate in a research project because it keeps the ablation path readable.",
    )

    doc.add_heading("5. What Improvement Should Look Like", level=1)
    dpara(
        doc,
        "The main sign of success is not only higher accuracy. The confusion matrix should show less absorption into neutral on MELD and less collapse into a single class on CREMA-D. A good result would preserve Fold 2-level accuracy, raise Fold 4-level balanced behavior, and improve CREMA-D from near collapse toward meaningful multi-class separation.",
    )
    dbullets(
        doc,
        [
            "MELD: fewer neutral-heavy mistakes, especially on joy, anger, sadness, and surprise.",
            "CREMA-D: predictions should spread across the classes instead of collapsing.",
            "Metrics: macro F1 and UW-Acc should improve, not just sample accuracy.",
        ],
    )

    doc.add_heading("6. Questions for the Guidance Call", level=1)
    dbullets(
        doc,
        [
            "Should class-balanced focal loss be the first change, or should I jump directly to LDAM-DRW?",
            "Is supervised contrastive learning worth adding now, or only after the loss function is improved?",
            "For MELD, should I implement dialogue context modeling before changing CREMA-D again?",
            "For CREMA-D, should I prioritize gradual unfreezing and discriminative learning rates before trying a larger architecture change?",
        ],
    )

    doc.save(DOCX_PATH)


def main() -> None:
    build_confusion_image(FOLD2_CM_CSV, FOLD2_CM_PNG, "MELD Fold 2 Confusion Matrix")
    build_confusion_image(FOLD4_CM_CSV, FOLD4_CM_PNG, "MELD Fold 4 Confusion Matrix")
    build_confusion_image(CREMA_CM_CSV, CREMA_CM_PNG, "CREMA-D Confusion Matrix")
    build_ppt()
    build_doc()
    print(PPTX_PATH)
    print(DOCX_PATH)


if __name__ == "__main__":
    main()
