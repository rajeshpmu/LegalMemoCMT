"""Add SILVER training and evaluation slides plus detailed speaking notes."""
from pathlib import Path
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document

ROOT=Path(__file__).resolve().parents[1]
PPTX=ROOT/'implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Presentation_v1.pptx'
DOCX=ROOT/'implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Speaking_Doc.docx'
FIG=ROOT/'implementation_docments/figures/phase2_silver_training_evaluation_flow.png'
MARK='SILVER_TRAINING_EVALUATION_V1'

def add_text(s,v,x,y,w,h,size=15,color=(35,35,35),bold=False):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); b.text_frame.word_wrap=True
    p=b.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT; r=p.add_run(); r.text=v
    r.font.name='Aptos'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=RGBColor(*color)

def slide_base(prs,title,sub):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=RGBColor(255,255,255)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.22)); bar.fill.solid(); bar.fill.fore_color.rgb=RGBColor(20,48,87); bar.line.fill.background()
    add_text(s,title,.55,.38,12.2,.45,23,(20,48,87),True); add_text(s,sub,.58,.9,12,.3,12,(95,95,95)); return s

def add_table(s,headers,rows,x=.55,y=1.35,w=12.2,h=5.55,size=11):
    t=s.shapes.add_table(len(rows)+1,len(headers),Inches(x),Inches(y),Inches(w),Inches(h)).table
    for j,v in enumerate(headers): t.cell(0,j).text=str(v)
    for i,row in enumerate(rows,1):
        for j,v in enumerate(row): t.cell(i,j).text=str(v)
    for i,row in enumerate(t.rows):
        for c in row.cells:
            c.text_frame.word_wrap=True
            for p in c.text_frame.paragraphs:
                for r in p.runs: r.font.name='Aptos'; r.font.size=Pt(size); r.font.color.rgb=RGBColor(20,48,87) if i==0 else RGBColor(35,35,35); r.font.bold=i==0

def insert_after(prs,index,slide):
    ids=prs.slides._sldIdLst; item=ids[-1]; ids.remove(item); ids.insert(index,item)

def update_ppt():
    prs=Presentation(PPTX)
    if any(MARK in getattr(sh,'text','') for s in prs.slides for sh in s.shapes): return
    backup=PPTX.with_name(PPTX.name+'.before_silver_training_evaluation.pptx')
    if not backup.exists(): shutil.copy2(PPTX,backup)
    # Insert immediately before the current final decision slide.
    index=len(prs.slides)-1
    s=slide_base(prs,'SILVER Corpus: Fine-Tuning and Target Labels','Use machine-assisted labels for a controlled adaptation experiment, not as courtroom ground truth')
    add_table(s,['Decision','Phase 2 plan','Reason'],[
        ['Training target','`final_basic_emotion` for AUTO_ADJUDICATED/SILVER rows; seven MELD labels','This is the existing compatible classification head and preserves `emotion_label` compatibility.'],
        ['Courtroom affect','Do not train as gold from current machine candidates','CALM_COMPOSED, TENSE, and related labels require a human-validated seed set; current candidates are evidence/review suggestions.'],
        ['Warm start','`results/facial_cues/meld_vit_facecrop_gated_video_aux/fold_4/best_model.pt`','Available trimodal face-crop checkpoint with compatible 768-D video features; fold 4 metrics: accuracy 0.599, macro-F1 0.433.'],
        ['Data use','Start with SILVER rows only after group split and provenance checks','The current 125 SILVER rows come from one controlled pilot source, so they are useful for a smoke test, not a final generalization claim.'],
        ['Controls','Keep Phase 1 baseline, compare adapted model, and preserve all original fields','This exposes confirmation bias and domain-shift effects rather than hiding them.'],
    ],size=10); add_text(s,'Initial experiment: adapt basic emotion first; build courtroom-affect supervision only after human review.',.8,6.5,11.7,.35,13,(20,48,87),True); insert_after(prs,index,s); index+=1
    s=slide_base(prs,'Evaluation Plan: Does Adaptation Actually Help?','Evaluate on groups and labels that were not used to create the SILVER training candidates')
    add_table(s,['Evaluation component','How it will be done','What success means'],[
        ['Gold test target','Human-validated basic emotion on held-out witness/source groups','No Silver-derived target is treated as gold.'],
        ['Primary metrics','Macro-F1, balanced accuracy, per-class precision/recall/F1, confusion matrix','Improvement is judged across classes, not only by majority-class accuracy.'],
        ['Calibration','Confidence reliability, expected calibration error, and prediction entropy','A model should be less overconfident on courtroom domain-shift cases.'],
        ['Ablations','Text-only, audio-only, video-only, and trimodal comparisons','Shows which modality contributes and whether video actually helps.'],
        ['Leakage checks','Split by source video/hearing/witness; no group appears in multiple partitions','A gain is credible only if it survives group-disjoint evaluation.'],
        ['Courtroom-affect evaluation','Separate human-validated affect test with label agreement and abstention rate','Measures the new task independently; basic emotion and affect are not conflated.'],
    ],size=10); add_text(s,'Expected outcome: better calibrated courtroom adaptation is possible, but the experiment may also reveal that 125 SILVER rows are insufficient.',.75,6.5,11.9,.35,12,(20,48,87),True); insert_after(prs,index,s)
    prs.save(PPTX); print('Updated',PPTX,'slides',len(prs.slides))

def update_doc():
    doc=Document(DOCX)
    if any(MARK in p.text for p in doc.paragraphs): return
    backup=DOCX.with_name(DOCX.name+'.before_silver_training_evaluation.docx')
    if not backup.exists(): shutil.copy2(DOCX,backup)
    doc.add_heading('Using AUTO_ADJUDICATED / SILVER Rows for Training and Evaluation',1)
    p=doc.add_paragraph(); p.add_run(MARK).font.size=Pt(8)
    doc.add_heading('Recommended target labels',2)
    doc.add_paragraph('For the first controlled adaptation experiment, the supervised target should be final_basic_emotion on AUTO_ADJUDICATED/SILVER rows. Its vocabulary is the seven MELD-compatible labels: neutral, anger, disgust, fear, joy, sadness, and surprise. The original Phase 1 prediction and confidence remain in separate fields, while the selected training target is recorded explicitly. Existing emotion_label compatibility should not be broken.')
    doc.add_paragraph('I should not use the current proposed_courtroom_affect or final_courtroom_affect machine candidates as if they were human gold labels. Courtroom affect is a second task describing presentation, such as CALM_COMPOSED or HESITANT_UNCERTAIN, and needs a human-validated seed set. It can be retained as machine evidence or used for review prioritization, but not presented as a reliable supervised target yet.')
    doc.add_heading('Warm-start checkpoint',2)
    doc.add_paragraph('The executable warm start is results/facial_cues/meld_vit_facecrop_gated_video_aux/fold_4/best_model.pt. It is the available trimodal Phase 1 checkpoint using the face-crop video representation expected by the Clancy ViT feature paths. Its recorded MELD metrics are accuracy 0.599, macro-F1 0.433, and weighted-F1 0.606 on 2,610 samples. Fold 2 has a slightly higher recorded macro-F1 of 0.435, but the expected fold_2 best_model.pt is not present at the audited path, so I must not claim that it is executable until the checkpoint is located and verified.')
    doc.add_paragraph('This is warm-start transfer learning: initialize the LegalMemoCMT model with the MELD-trained parameters, then continue training with a small learning rate on carefully selected courtroom rows. It is not training from scratch. I should begin with the classification head or conservative partial fine-tuning, use early stopping, and compare against the unchanged Phase 1 checkpoint.')
    doc.add_heading('Why the current SILVER set is only a pilot',2)
    doc.add_paragraph('The current gated pilot contains 125 AUTO_ADJUDICATED/SILVER rows and 75 UNRESOLVED/WEAK rows. The SILVER rows are machine-assisted, not human gold. They are useful for a pipeline smoke test and for identifying domain-shift patterns, but 125 rows from one controlled source are too small and too correlated for a strong generalization claim. A group split by source video or witness is required, and a single-source pilot may not support a meaningful independent test.')
    doc.add_heading('Training procedure',2)
    for x in [
        'Freeze the raw evidence fields and create an explicit training-target column from final_basic_emotion only for approved SILVER rows.',
        'Split by source video/hearing/witness before optimization; never randomly split adjacent utterances from the same source.',
        'Load the fold_4 trimodal checkpoint, verify model_cfg, especially use_video=true and video_dim=768, and confirm all feature paths exist.',
        'Run a small warm-start experiment with a reduced learning rate and early stopping; retain the baseline checkpoint for direct comparison.',
        'Keep unresolved rows out of supervised loss, but retain them for later human review and active learning.',
    ]: doc.add_paragraph(x,style='List Bullet')
    doc.add_heading('Evaluation design',2)
    doc.add_paragraph('The primary evaluation target is a held-out human-validated basic-emotion test set. I will report macro-F1, balanced accuracy, per-class precision/recall/F1, confusion matrices, confidence calibration, and abstention or unresolved rates. Accuracy alone is unsafe because neutral may dominate the distribution. The evaluation must compare three systems: the unchanged Phase 1 checkpoint, the warm-start adapted model, and simple modality ablations such as text-only, audio-only, video-only, and trimodal.')
    doc.add_paragraph('For courtroom affect, I will create a separate human-validated test set with the courtroom-affect vocabulary and intensity labels. I will not evaluate courtroom affect by comparing machine candidates to themselves. I will also report agreement between reviewers, disagreement by witness/source, and errors caused by quoted or other-person-described emotion content.')
    doc.add_heading('Expected outcomes and interpretation',2)
    doc.add_paragraph('A successful result would show improved macro-F1 or calibration on held-out courtroom data without increasing leakage or collapsing into the neutral majority class. It may also show that the MELD warm start transfers poorly, that video features are unreliable when the witness is not visible, or that 125 SILVER rows are insufficient. Those are scientifically useful outcomes. The adaptation should not be declared successful merely because training loss falls or Silver-set accuracy rises.')
    doc.add_heading('Student explanation for the guidance call',2)
    doc.add_paragraph('I will explain that SILVER rows are a controlled machine-assisted training experiment. I use final_basic_emotion as the initial target because it is compatible with the Phase 1 MELD head. I warm-start from the available trimodal face-crop fold_4 checkpoint, preserve the original predictions, and test whether adaptation improves on a separate human-validated group-disjoint set. Courtroom affect remains a separate annotation task until enough human labels exist. The result I want is not simply a larger score; it is evidence that the model generalizes to courtroom presentation without confusing described emotion with the witness’s own emotion.')
    doc.save(DOCX); print('Updated',DOCX)

if __name__=='__main__': update_ppt(); update_doc()
