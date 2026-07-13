from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PPTX_PATH = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT/implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx")

ROWS = [
    ("MemoCMT paper CMT+MIN", "64.18%", "62.52%", "—", "—"),
    ("Paper-aligned MELD CV mean", "62.47%", "61.95%", "43.95%", "44.17%"),
    ("Best baseline fold 2", "63.75%", "62.54%", "44.30%", "43.69%"),
    ("Gated + aux fold 2", "60.54%", "60.22%", "43.51%", "44.92%"),
    ("Gated + aux fold 4", "59.92%", "60.56%", "43.30%", "46.38%"),
]


def remove_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    rId = slide_id.rId
    prs.part.drop_rel(rId)
    slide_id_list.remove(slide_id)


def add_box(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line:
        shp.line.color.rgb = RGBColor.from_string(line)
    else:
        shp.line.color.rgb = RGBColor.from_string(fill)
    return shp


def style_run(run, size=10, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_text(slide, left, top, width, height, text, size=11, bold=False, color="000000", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    style_run(r, size=size, bold=bold, color=color)
    return tb


def rebuild():
    prs = Presentation(str(PPTX_PATH))
    if len(prs.slides) >= 27:
        remove_slide(prs, len(prs.slides) - 1)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_no = len(prs.slides)

    add_box(slide, 0.0, 0.0, 13.333, 0.58, "122F55")
    add_box(slide, 0.0, 0.58, 13.333, 0.08, "267377")
    add_text(slide, 0.45, 0.12, 11.4, 0.3, f"{slide_no}. MELD Result Comparison: MemoCMT vs Baseline vs Gated+Aux", size=18, bold=True, color="FFFFFF")
    add_text(slide, 11.95, 0.12, 0.95, 0.3, str(slide_no), size=12, color="FFFFFF", align=PP_ALIGN.RIGHT)
    add_text(slide, 0.52, 0.72, 12.15, 0.32, "Directional comparison only: the paper uses a different MELD evaluation framing, so the trend is the important part.", size=10.5, color="5E5E5E")

    add_box(slide, 0.35, 1.15, 12.6, 1.0, "EEF4FA", "C9D7E6")
    add_text(
        slide,
        0.55,
        1.28,
        12.2,
        0.22,
        "Key readout: paper 64.18/62.52, baseline mean 62.47/61.95, best baseline fold 2 63.75/62.54, gated+aux fold 2 60.54/60.22, gated+aux fold 4 59.92/60.56.",
        size=12,
        bold=True,
        color="122F55",
    )
    add_text(
        slide,
        0.55,
        1.60,
        12.1,
        0.28,
        "Interpretation: the baseline remains the closest stable reproduction of MemoCMT on MELD; gated+aux helps some clips but is not yet a universal improvement.",
        size=11,
        color="3A3A3A",
    )

    headers = ["Run", "Acc / W-Acc", "W-F1", "Macro F1", "UW-Acc"]
    table = slide.shapes.add_table(
        len(ROWS) + 1,
        len(headers),
        Inches(0.35),
        Inches(2.35),
        Inches(12.6),
        Inches(3.65),
    ).table
    widths = [3.2, 1.6, 1.3, 1.3, 1.3]
    for idx, w in enumerate(widths):
        table.columns[idx].width = Inches(w)
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string("DDEBF7")
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                style_run(r, size=10.5, bold=True, color="122F55")
    for i, row in enumerate(ROWS, start=1):
        values = row
        for j, value in enumerate(values):
            cell = table.cell(i, j)
            cell.text = value
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor.from_string("F8FBFF")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    style_run(r, size=9.4 if j == 0 else 10, color="202020")

    add_box(slide, 0.35, 6.15, 12.6, 0.86, "EEF4FA", "C9D7E6")
    add_text(
        slide,
        0.55,
        6.28,
        12.2,
        0.25,
        "Paper gap on the mean baseline is about 1.71 points in accuracy and 0.57 points in weighted F1.",
        size=11,
        bold=True,
        color="122F55",
    )
    add_text(
        slide,
        0.55,
        6.56,
        12.0,
        0.22,
        "Use this slide to say that gated+aux is a selective branch, not a universal replacement for the baseline.",
        size=10,
        color="4B4B4B",
    )

    prs.save(str(PPTX_PATH))
    print(f"Rebuilt slide {slide_no} in {PPTX_PATH}")


if __name__ == "__main__":
    rebuild()
