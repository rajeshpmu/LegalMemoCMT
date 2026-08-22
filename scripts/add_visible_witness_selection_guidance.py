from pathlib import Path
import shutil

from docx import Document
from docx.enum.section import WD_SECTION
from docx.shared import Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
PPTX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.before_visible_witness_selection.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Speaking_Doc.docx"
DOCX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Speaking_Doc.before_visible_witness_selection.docx"
TITLE = "Selecting the Visible-Witness Speaking Corpus"


def add_box(slide, x, y, w, h, text, color, size=15, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.color.rgb = RGBColor(70, 80, 95)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 1
    for run in p.runs:
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(25, 35, 50)


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(70, 80, 95)
    line.line.end_arrowhead = True


prs = Presentation(PPTX)
if not PPTX_BACKUP.exists():
    shutil.copy2(PPTX, PPTX_BACKUP)
if not any(any(TITLE in sh.text for sh in slide.shapes if hasattr(sh, "text")) for slide in prs.slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.5))
    title.text_frame.text = TITLE
    title.text_frame.paragraphs[0].runs[0].font.size = PptPt(26)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

    add_box(slide, 0.55, 1.05, 2.2, 0.7, "Pyannote\naudio speaking intervals", (248, 237, 215), 16, True)
    add_box(slide, 3.2, 1.05, 2.2, 0.7, "Manual cluster\nrole mapping", (231, 243, 235), 16, True)
    add_box(slide, 5.85, 1.05, 2.2, 0.7, "Visual review\nwitness visible?", (220, 232, 245), 16, True)
    add_box(slide, 8.5, 1.05, 3.9, 0.7, "Primary witness-emotion\ntraining pool", (208, 224, 220), 16, True)
    add_arrow(slide, 2.75, 1.4, 3.2, 1.4)
    add_arrow(slide, 5.4, 1.4, 5.85, 1.4)
    add_arrow(slide, 8.05, 1.4, 8.5, 1.4)

    add_box(slide, 0.7, 2.35, 3.6, 1.15, "Pool A\nWitness + SPEAKING\nWITNESS_VISIBLE", (208, 224, 220), 17, True)
    add_box(slide, 4.85, 2.35, 3.6, 1.15, "Pool B\nWitness speaking\nvisibility uncertain", (242, 235, 205), 17, True)
    add_box(slide, 9.0, 2.35, 3.1, 1.15, "Pool C/D\nlistening, non-witness,\nor unresolved", (240, 220, 220), 16, True)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(4.05), Inches(11.8), Inches(2.4))
    tf = note.text_frame
    tf.clear()
    tf.word_wrap = True
    bullets = [
        "Pyannote can support audio-speaking evidence, but cannot identify a visible witness by itself.",
        "A visible face is not proof that the person is speaking; the camera may show the witness while counsel speaks.",
        "Use Pool A for the first witness-emotion model; retain Pools B–D for later context and review.",
        "Do not convert speaking status into emotion, credibility, deception, or truthfulness automatically.",
    ]
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = PptPt(17)
    prs.save(PPTX)

    # Place after the architecture slide.
    ids = prs.slides._sldIdLst
    new_id = ids[-1]
    ids.remove(new_id)
    anchor = "Whisper + Pyannote: Building a Witness-Aware Utterance"
    insert_at = next(
        (
            i + 1
            for i, sid in enumerate(ids)
            if any(anchor in sh.text for sh in prs.part.related_part(sid.rId).slide.shapes if hasattr(sh, "text"))
        ),
        len(ids),
    )
    ids.insert(insert_at, new_id)
    prs.save(PPTX)


if not DOCX_BACKUP.exists():
    shutil.copy2(DOCX, DOCX_BACKUP)
document = Document(DOCX)
document.add_section(WD_SECTION.NEW_PAGE)
document.add_heading("Additional Speaking Module: Selecting Visible Witness Speech", level=1)
document.add_heading("1. The Recommended Primary Pool", level=2)
document.add_paragraph(
    "For the first witness-emotion experiments, I should use a conservative Pool A rather than every row in a "
    "courtroom recording. Pool A contains an utterance whose role is verified as Witness, whose diarization "
    "evidence indicates that the mapped witness cluster is speaking, and whose visual review shows the witness "
    "is visible. The practical rule is: `speaker_role=Witness`, `witness_speaking_status=SPEAKING`, and "
    "`speaker_visible=WITNESS_VISIBLE`, together with valid audio and high or medium transcript alignment."
)
document.add_heading("2. What Pyannote Contributes", level=2)
document.add_paragraph(
    "Pyannote contributes the audio side of the decision. It detects speech intervals and assigns anonymous "
    "speaker clusters. When an utterance timestamp overlaps a cluster interval, the pipeline records the cluster "
    "and `speaker_cluster_overlap_seconds`. After manual review maps that cluster to the witness, the row can be "
    "marked `witness_speaking_status=SPEAKING`. This is stronger than assuming that every subtitle row is spoken "
    "by the person visible in the frame."
)
document.add_paragraph(
    "Pyannote does not inspect the video face and does not know legal roles. `SPEAKER_02` is only an anonymous "
    "voice grouping. It must be mapped to Witness, Prosecutor, Defence, Judge, or Other by listening to samples "
    "and recording the confidence and review notes."
)
document.add_heading("3. Why Visual Review Is Still Required", level=2)
document.add_paragraph(
    "Courtroom cameras often remain on the witness while counsel asks a question. In that situation the witness "
    "may be visible but listening, while the lawyer is the speaking person. Conversely, the witness may speak off "
    "camera. Therefore, audio diarization and face visibility answer different questions: Pyannote estimates which "
    "voice is speaking, while visual review estimates whether the target person is visible and plausibly the active "
    "speaker. A dedicated active-speaker model could assist later, but it is not currently treated as a validated "
    "source of truth."
)
document.add_heading("4. Required Processing Sequence", level=2)
for item in [
    "Run subtitle evidence extraction so the manifest retains cue boundaries, question patterns, testimony-block evidence, and transcript text.",
    "Run Pyannote once per source audio and retain the source-level segments CSV as the audit record.",
    "Generate the cluster review sheet with one or more representative audio/video clips and total cluster row counts.",
    "Listen to cluster samples and map anonymous clusters to courtroom roles; record role confidence and notes.",
    "Inspect the representative video clips and fill `visual_target_role` and `visual_speaker_match` conservatively.",
    "Apply the role map to create the role-aware manifest.",
    "Filter Pool A only after the role, speaking-status, visual, alignment, and audio-quality conditions pass.",
]:
    document.add_paragraph(item, style="List Number")
document.add_heading("5. Meaning of the Pools", level=2)
for item in [
    "Pool A: Witness + SPEAKING + WITNESS_VISIBLE + valid audio + high/medium alignment. Use for the initial witness-emotion training experiment.",
    "Pool B: Witness is speaking but visual visibility is uncertain. Retain for later sensitivity analysis or manual expansion.",
    "Pool C: Witness is visible but listening while counsel or judge speaks. Useful for courtroom interaction analysis, not as a witness facial-emotion sample.",
    "Pool D: Non-witness speech, unresolved clusters, invalid media, or poor alignment. Exclude from the initial witness corpus but retain an exclusion reason.",
]:
    document.add_paragraph(item, style="List Bullet")
document.add_heading("6. Labels That Must Not Be Inferred", level=2)
document.add_paragraph(
    "The speaking-status decision is not an emotion label. A tense witness voice is not automatically anger, "
    "fear, deception, or unreliability. Similarly, a lawyer's emotional argument must not enter the witness-emotion "
    "pool merely because the camera shows the witness. Basic emotion pseudo-labels from a Phase 1 MELD model and "
    "courtroom-specific affect labels must remain separate fields with their own source and confidence."
)
document.add_heading("7. Guidance-Call Explanation", level=2)
document.add_paragraph(
    "I use Pyannote to establish audio-based speaking evidence, manual review to map anonymous clusters to courtroom "
    "roles, and visual inspection to determine whether the witness is visible. Only the intersection of verified "
    "witness role, witness speaking status, and visible-witness evidence enters the first witness-emotion training "
    "pool; the other categories remain available for controlled secondary experiments."
)
document.save(DOCX)
print(f"Updated {PPTX}")
print(f"Updated {DOCX}")
