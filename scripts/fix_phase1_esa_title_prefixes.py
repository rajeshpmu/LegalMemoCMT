from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"


def title_shape(slide):
    candidates = []
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = shp.text_frame.text.strip()
        if not txt:
            continue
        if shp.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if float(shp.top) > 1200000:
            continue
        candidates.append(shp)
    if not candidates:
        return None
    # Prefer the larger, upper-left title text box.
    return sorted(candidates, key=lambda s: (float(s.top), -float(s.width)))[0]


def set_title(slide, prefix: int) -> bool:
    tshape = title_shape(slide)
    if tshape is None:
        return False
    raw = tshape.text_frame.text.strip()
    if not raw:
        return False
    base = raw
    if "." in raw[:4]:
        base = raw.split(".", 1)[1].strip()
    if raw.startswith(f"{prefix}."):
        return True
    tshape.text_frame.clear()
    p = tshape.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"{prefix}. {base}"
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("122F55")
    return True


def main() -> None:
    prs = Presentation(str(PPTX_PATH))
    fixes = {
        5: 5,
        15: 15,
        18: 18,
        20: 20,
        21: 21,
    }
    changed = []
    for idx, prefix in fixes.items():
        if idx - 1 < len(prs.slides) and set_title(prs.slides[idx - 1], prefix):
            changed.append(idx)
    prs.save(str(PPTX_PATH))
    print(f"Updated title prefixes for slides: {changed}")


if __name__ == "__main__":
    main()
