from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.shared import Pt
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Prep_Guide.docx"

BODY = RGBColor(42, 42, 42)
MUTED = RGBColor(92, 92, 92)


def style_text(run, size=12, bold=False, color=BODY):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def insert_paragraph_before(paragraph, text="", style=None):
    new_p = OxmlElement("w:p")
    paragraph._p.addprevious(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        new_para.add_run(text)
    return new_para


def update_pptx():
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides[11]

    # Keep the title and slide number as-is, only refresh the text boxes.
    bullet_box = slide.shapes[5]
    bullet_box.text = "\n".join(
        [
            "• Constraint: MELD is class imbalanced, so macro F1 matters more than accuracy.",
            "• Assumption: sample_id, split, and label stay aligned across manifest, cache, and checkpoint.",
            "• Dependency: pretrained backbones, ffmpeg, and cached audio/video features must exist.",
            "• Impact: any missing or misaligned asset can break reproducibility and demo traceability.",
        ]
    )
    for p in bullet_box.text_frame.paragraphs:
        for run in p.runs:
            style_text(run, 13.2, False, BODY)

    note_box = slide.shapes[7]
    note_box.text = (
        "This slide is about why the chosen design can be trusted. "
        "The project only makes sense if the data labels, cached features, and model checkpoint all refer to the same utterance. "
        "If one dependency is missing or mismatched, then the result may still produce a number, but the number is no longer a valid comparison for Phase 1."
    )
    for p in note_box.text_frame.paragraphs:
        for run in p.runs:
            style_text(run, 11.2, False, MUTED)

    prs.save(str(PPTX_PATH))


def update_docx():
    doc = Document(str(DOCX_PATH))
    table = doc.tables[0]
    for row in table.rows:
        if row.cells[0].text.strip() == "Slide 12":
            row.cells[1].text = (
                "Discuss the design constraints and assumptions behind the chosen design approach. "
                "Explain the key dependencies such as pretrained backbones, ffmpeg, and stable manifest-to-cache alignment, "
                "and say how they affect reproducibility and interpretability."
            )
            row.cells[2].text = "Keep it technical, concise, and traceable."
            break

    paras = list(doc.paragraphs)
    start = next((i for i, p in enumerate(paras) if p.text.strip() == "2.4 Slide 12 Extended Talking Points"), None)
    if start is not None:
        end = next(
            (i for i in range(start + 1, len(paras)) if paras[i].text.strip().startswith("3. Demo SOP You Should Practise")),
            len(paras),
        )
        for i in range(start, end):
            el = paras[i]._element
            el.getparent().remove(el)

    anchor = next((p for p in doc.paragraphs if p.text.strip().startswith("3. Demo SOP You Should Practise")), None)
    if anchor is None:
        anchor = doc.paragraphs[-1]
    heading = insert_paragraph_before(anchor, style="Heading 2")
    heading.add_run("2.4 Slide 12 Extended Talking Points")
    sections = [
        "Design constraints: the MELD label distribution is heavily skewed toward neutral, so accuracy alone can look better than the model really is on minority emotions. For that reason the project must always discuss macro F1, weighted F1, and the confusion matrix together.",
        "Design assumptions: every manifest row must keep the same sample_id, split, transcript, and cached feature file. This is what allows a raw clip, a .npy cache, and a checkpointed prediction to refer to the same utterance during analysis.",
        "Dependencies: the pipeline depends on pretrained text and audio backbones, ffmpeg for video extraction, and locally cached media/features so that the same input can be replayed during review without reprocessing the full raw dataset every time.",
        "Impact of dependency failures: if frame sampling changes, a file path is wrong, or a cache is stale, the reported prediction may still be printed but it no longer supports a reliable comparison. That is why the project treats these dependencies as part of the scientific method, not just implementation detail.",
        "Why this matters for the design choice: the controlled experimental ladder only works when the data plumbing is stable. If the inputs are not aligned, then any improvement or regression could be caused by the pipeline error rather than the model design itself.",
    ]
    for t in sections:
        p = insert_paragraph_before(anchor)
        r = p.add_run(t)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        p.alignment = 3  # JUSTIFY
    doc.save(str(DOCX_PATH))


def main():
    update_pptx()
    update_docx()
    print(f"Updated {PPTX_PATH}")
    print(f"Updated {DOCX_PATH}")


if __name__ == "__main__":
    main()
