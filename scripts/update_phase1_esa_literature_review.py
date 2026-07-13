from __future__ import annotations

from pathlib import Path
import re

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Prep_Guide.docx"

NAVY = RGBColor(18, 47, 85)
TEAL = RGBColor(38, 115, 119)
LIGHT = RGBColor(246, 248, 250)
BODY = RGBColor(42, 42, 42)
MUTED = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)


LIT = [
    {
        "title": "Literature Review 1: MemoCMT",
        "paper": "MemoCMT: multimodal emotion recognition using cross-modal transformer-based feature fusion",
        "link": "https://doi.org/10.1038/s41598-025-89202-x",
        "bullets": [
            "Paper-level reference for cross-modal transformer fusion in multimodal emotion recognition.",
            "Shows why text and audio should be fused with attention instead of simple concatenation.",
            "Supports the Phase 1 choice of keeping the conversational baseline as the main comparison point.",
            "Most important because the whole project is anchored to this implementation style.",
        ],
        "note": (
            "MemoCMT is the closest technical reference for the Phase 1 baseline. "
            "In speaking terms, explain that the project is not inventing a new fusion idea from scratch; it is reproducing a MemoCMT-style text+audio pipeline first so the later facial-cue extension has a strong reference point. "
            "Mention that the cross-modal transformer is what lets one modality influence the other rather than just appending features."
        ),
    },
    {
        "title": "Literature Review 2: MELD",
        "paper": "MELD: A Multimodal Multi-Party Dataset for Emotion Recognition in Conversations",
        "link": "https://aclanthology.org/P19-1050/",
        "bullets": [
            "Primary benchmark dataset for dialogue-based emotion recognition in this project.",
            "Multi-party conversations make speaker context and turn history important.",
            "Seven emotion labels and heavy class imbalance make macro F1 essential.",
            "Explains why fold-safe dialogue grouping is needed in Phase 1.",
        ],
        "note": (
            "MELD is the benchmark that makes the Phase 1 problem concrete. "
            "When you speak about it, emphasize that the project is not just classifying emotions in isolation; it is handling utterances inside conversation, where surrounding turns and speaker history matter. "
            "The dataset also explains the neutral-heavy bias and the need for weighted metrics."
        ),
    },
    {
        "title": "Literature Review 3: BERT + HuBERT",
        "paper": "BERT: Pre-training of Deep Bidirectional Transformers for Language Understanding; HuBERT: Self-Supervised Speech Representation Learning by Masked Prediction of Hidden Units",
        "link": "https://arxiv.org/abs/1810.04805 ; https://arxiv.org/abs/2106.07447",
        "bullets": [
            "BERT turns transcripts into contextual token embeddings for the text branch.",
            "HuBERT turns waveform audio into self-supervised speech representations for the audio branch.",
            "Both are pretrained, so Phase 1 can focus on adaptation rather than training encoders from scratch.",
            "Together they give the baseline its strongest conversational signal.",
        ],
        "note": (
            "This slide should answer the question: why text and audio first? "
            "Explain that BERT captures semantic context from the transcript, while HuBERT captures speech delivery and paralinguistic patterns. "
            "The important defense point is that both encoders are pretrained, so Phase 1 leverages existing language and speech knowledge instead of spending compute on learning those representations from zero."
        ),
    },
    {
        "title": "Literature Review 4: Vision Transformer",
        "paper": "An Image is Worth 16x16 Words: Transformers for Image Recognition at Scale",
        "link": "https://arxiv.org/abs/2010.11929",
        "bullets": [
            "Used as the facial-cue encoder for sampled RGB frames from raw MP4 clips.",
            "Converts each frame into patch tokens and learns global visual attention.",
            "Works well for face-crop or full-frame image embeddings before caching.",
            "Supports the visual path without changing the text/audio baseline design.",
        ],
        "note": (
            "ViT is the visual support branch, not the main benchmark engine. "
            "Explain that the project samples RGB frames from mp4, optionally crops the face region, and then feeds the frames to ViT so that the facial signal becomes a compact embedding. "
            "The key technical message is that the video is not fed as a raw motion stream; it is turned into frame embeddings that the multimodal model can fuse."
        ),
    },
    {
        "title": "Literature Review 5: Weighted CE + Focal Loss",
        "paper": "Focal Loss for Dense Object Detection (reference for imbalance-aware training)",
        "link": "https://arxiv.org/abs/1708.02002",
        "bullets": [
            "Class imbalance is one of the central issues in MELD, especially because neutral dominates.",
            "Weighted cross entropy gives more learning weight to minority classes.",
            "Focal loss reduces the impact of easy majority examples and can improve boundary behavior.",
            "These losses are important because accuracy alone can hide minority-class failure.",
        ],
        "note": (
            "Weighted CE is the baseline imbalance fix in code, while focal loss is the more aggressive imbalance-aware reference. "
            "In the ESA explanation, stress that these losses do not magically solve class imbalance; they just change how much the model learns from minority examples versus easy majority examples. "
            "Also remind the examiner that weighted F1 and macro F1 are the metrics that reveal whether the loss actually helps the hard classes."
        ),
    },
]


def insert_slide(prs: Presentation, index: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    sldIdLst.insert(index, sldIdLst[-1])
    return slide


def add_title_bar(slide, title: str, num: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.58))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.58), Inches(13.333), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(11.8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    tb2 = slide.shapes.add_textbox(Inches(12.0), Inches(0.12), Inches(0.9), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{num:02d}"
    r2.font.name = "Aptos"
    r2.font.size = Pt(14)
    r2.font.bold = True
    r2.font.color.rgb = WHITE


def add_box(slide, x, y, w, h, fill=LIGHT):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(214, 220, 225)
    return shp


def add_bullets(slide, x, y, w, h, bullets):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(16)
            run.font.color.rgb = BODY


def add_citation(slide, x, y, w, h, paper, link):
    add_box(slide, x, y, w, h, fill=RGBColor(246, 248, 250))
    tb = slide.shapes.add_textbox(x + Inches(0.14), y + Inches(0.12), w - Inches(0.24), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    parts = [f"Paper: {paper}", f"Link: {link}", "Why it matters: this is the reference that justifies the design choice on this slide."]
    for i, line in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(12 if i < 2 else 11)
            run.font.color.rgb = BODY if i < 2 else MUTED


def set_slide_number(slide, num: int):
    for shape in slide.shapes:
        if hasattr(shape, "text") and re.fullmatch(r"\d{1,2}", shape.text.strip()):
            shape.text = f"{num:02d}"
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(14)
                    run.font.bold = True
                    run.font.color.rgb = WHITE


def update_slide_5(slide, item):
    # shape 2 title, shape 5 bullets, shape 7 right text in the current deck
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            txt = shape.text.strip()
            if txt == "Literature Survey":
                shape.text = item["title"]
            elif txt.startswith("• "):
                shape.text = "\n".join([f"• {b}" for b in item["bullets"]])
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.name = "Aptos"
                        run.font.size = Pt(15)
                        run.font.color.rgb = BODY
            elif txt.startswith("The literature slide should show"):
                shape.text = (
                    "This slide is the key paper-level reference for the Phase 1 literature path.\n"
                    f"{item['paper']}\n"
                    f"{item['link']}"
                )
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.name = "Aptos"
                        run.font.size = Pt(12)
                        run.font.color.rgb = MUTED if p.text.startswith("This slide") else BODY


def build_pptx():
    prs = Presentation(str(PPTX_PATH))
    update_slide_5(prs.slides[4], LIT[0])

    # Insert four new slides after current slide 5.
    for idx, item in enumerate(LIT[1:], start=5):
        slide = insert_slide(prs, idx)
        add_title_bar(slide, item["title"], idx + 1)
        add_box(slide, Inches(0.3), Inches(0.82), Inches(12.7), Inches(6.35), fill=WHITE)
        add_box(slide, Inches(0.55), Inches(1.08), Inches(5.85), Inches(5.62), fill=WHITE)
        add_bullets(slide, Inches(0.72), Inches(1.3), Inches(5.45), Inches(5.1), item["bullets"])
        add_citation(slide, Inches(6.6), Inches(1.15), Inches(5.95), Inches(2.0), item["paper"], item["link"])
        add_box(slide, Inches(6.6), Inches(3.35), Inches(5.95), Inches(3.15), fill=LIGHT)
        tb = slide.shapes.add_textbox(Inches(6.8), Inches(3.5), Inches(5.45), Inches(2.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item["note"]
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(13)
            run.font.color.rgb = BODY
        set_slide_number(slide, idx + 1)

    # Renumber all visible slide number boxes after insertion.
    for i, slide in enumerate(prs.slides, start=1):
        set_slide_number(slide, i)

    prs.save(str(PPTX_PATH))


def append_docx_section(doc: Document):
    doc.add_heading("2.2 Literature Review Slides (Slides 5-9)", level=2)
    doc.add_paragraph(
        "The ESA deck now treats the literature survey as five separate slides so each reference can be explained with enough detail. "
        "The slide itself stays compact, while these notes carry the deeper student-level explanation and the exact links."
    )
    for item in LIT:
        doc.add_heading(item["title"], level=3)
        p = doc.add_paragraph()
        p.add_run(item["paper"] + "\n").bold = True
        p.add_run(item["link"])
        for bullet in item["bullets"]:
            doc.add_paragraph(bullet, style="List Bullet")
        doc.add_paragraph(
            item["note"],
        )


def build_docx():
    doc = Document(str(DOCX_PATH))
    append_docx_section(doc)
    doc.save(str(DOCX_PATH))


def main():
    build_pptx()
    build_docx()
    print(f"Updated {PPTX_PATH}")
    print(f"Updated {DOCX_PATH}")


if __name__ == "__main__":
    main()
