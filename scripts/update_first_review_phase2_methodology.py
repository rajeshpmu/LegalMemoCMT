from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.docx"
PHASE2_DIAGRAM = ROOT / "artifacts" / "legalmemocmt_phase2.png"
PIPELINE_DIAGRAM = ROOT / "artifacts" / "mermaid" / "phase2_methodology_pipeline.png"


def add_bg(slide, prs: Presentation):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), prs.slide_width, PptInches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None, y: float = 0.35):
    tx = slide.shapes.add_textbox(PptInches(0.55), PptInches(y), PptInches(12.2), PptInches(0.7))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = PptPt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)

    if subtitle:
        tx2 = slide.shapes.add_textbox(PptInches(0.58), PptInches(y + 0.52), PptInches(12.1), PptInches(0.55))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = PptPt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x=0.75, y=1.45, w=6.0, h=5.45, font_size=15):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = PptPt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_diagram(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), PptInches(left), PptInches(top), width=PptInches(width))


def create_phase2_slides(prs: Presentation):
    slide_specs = [
        (
            "Proposed Phase 2 Methodology",
            "This section expands the courtroom-adaptation plan described in the Phase 2 methodology draft.",
            [
                "Phase 2 adapts the Phase 1 MemoCMT-style pipeline to Indian courtroom testimony.",
                "The model analyzes observable emotional cues only and does not predict guilt, innocence, intent, or deception.",
                "Target outputs are fear, anxiety, anger, stress timeline, and emotional transitions over time.",
                "Inputs are courtroom video testimony, audio testimony, and multilingual transcripts.",
            ],
            PHASE2_DIAGRAM,
            "Phase 2 adaptation overview",
            (0.7, 1.55, 6.05, 5.2),
            (6.95, 1.72, 5.95),
        ),
        (
            "Phase 2 Processing Pipeline",
            "The pipeline keeps the modalities aligned by time-window segmentation before fusion.",
            [
                "Segment testimony into aligned windows so face, audio, and text stay synchronized.",
                "Use a ViT-style visual encoder for witness face frames.",
                "Encode audio with XLS-R or IndicWav2Vec for prosody and stress-related features.",
                "Encode transcripts with IndicBERT v2 or MuRIL for multilingual courtroom language.",
                "Fuse the modality embeddings with a MemoCMT-style cross-modal transformer.",
                "Aggregate segment-level outputs into a continuous emotional timeline.",
            ],
            PIPELINE_DIAGRAM,
            "Phase 2 pipeline and temporal aggregation",
            (0.7, 1.55, 6.05, 5.2),
            (6.95, 1.72, 5.95),
        ),
        (
            "Phase 2 Adaptation, Explainability, and Checklist",
            "The Phase 2 plan extends the baseline without rebuilding the stack from zero.",
            [
                "Start from the Phase 1 backbone and adapt encoders, preprocessing, and evaluation targets.",
                "Pretrain on MELD and IEMOCAP, then fine-tune on small legal-domain samples if available.",
                "Use partial freezing or transfer learning when courtroom data are limited.",
                "Expose attention maps or modality contribution summaries for explanation.",
                "Present outputs as emotional pattern analysis, not as a legal decision engine.",
                "Implementation checklist: reproduce Phase 1, build face/audio/text alignment, integrate encoders, add fusion and temporal aggregation, then evaluate with emotion metrics and timelines.",
            ],
            None,
            None,
            (0.7, 1.45, 12.0, 5.45),
            None,
        ),
    ]

    new_slides = []
    for title, subtitle, bullets, image_path, caption, text_box, image_box in slide_specs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, prs)
        add_title(slide, title, subtitle)
        if image_path is None:
            add_body(slide, bullets, x=text_box[0], y=text_box[1], w=text_box[2], h=text_box[3], font_size=15)
        else:
            add_body(slide, bullets, x=text_box[0], y=text_box[1], w=text_box[2], h=text_box[3], font_size=15)
            add_diagram(slide, image_path, image_box[0], image_box[1], image_box[2])
            if caption:
                tx = slide.shapes.add_textbox(PptInches(image_box[0]), PptInches(6.55), PptInches(image_box[2]), PptInches(0.35))
                tf = tx.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = caption
                r.font.name = "Aptos"
                r.font.size = PptPt(10.5)
                r.font.italic = True
                r.font.color.rgb = RGBColor(95, 95, 95)
        new_slides.append(slide)
    return new_slides


def move_slide_ids(prs: Presentation, slide_ids, insert_at: int):
    sldIdLst = prs.slides._sldIdLst
    for sldId in reversed(slide_ids):
        sldIdLst.remove(sldId)
        sldIdLst.insert(insert_at, sldId)


def update_pptx():
    prs = Presentation(str(PPTX_PATH))
    new_slides = create_phase2_slides(prs)
    new_ids = list(prs.slides._sldIdLst[-len(new_slides):])
    # Move the new slides to appear after the existing "Ethical and Legal Scope" slide.
    move_slide_ids(prs, new_ids, 10)
    prs.save(str(PPTX_PATH))


def update_docx():
    doc = Document(str(DOCX_PATH))
    doc.add_heading("11. Proposed Phase 2 Methodology", level=1)
    doc.add_paragraph(
        "This section mirrors the proposed Phase 2 plan used in the updated First Review presentation. It keeps the scope narrow: the model analyzes emotional patterns in courtroom testimony, but it does not make legal conclusions."
    )
    doc.add_paragraph(
        "Phase 2 is not a rebuild from scratch. It starts from the Phase 1 backbone, then adapts the encoders, preprocessing, alignment, fusion, and evaluation so the system can operate on multilingual courtroom testimony."
    )

    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "Phase 2 component"
    hdr[1].text = "Planned method"
    hdr[2].text = "Purpose"
    rows = [
        ("Problem framing", "Observable emotion only", "Keep the legal scope bounded and safe."),
        ("Inputs", "Video, audio, transcript", "Use all available testimony evidence streams."),
        ("Visual encoder", "ViT-style face feature extractor", "Capture witness facial cues and motion."),
        ("Speech encoder", "XLS-R or IndicWav2Vec", "Capture prosody, pauses, and stress-related vocal patterns."),
        ("Text encoder", "IndicBERT v2 or MuRIL", "Handle multilingual courtroom transcripts."),
        ("Fusion", "MemoCMT-style cross-modal transformer", "Learn interactions across face, speech, and transcript."),
        ("Temporal modeling", "Segment-level aggregation", "Convert local outputs into a continuous timeline."),
        ("Explainability", "Attention and modality contribution summaries", "Make the output inspectable."),
    ]
    for row in rows:
        cells = table.add_row().cells
        for i, val in enumerate(row):
            cells[i].text = val
    doc.add_paragraph()

    doc.add_heading("11.1 Problem Framing", level=2)
    doc.add_paragraph(
        "Courtroom testimony differs from benchmark emotion datasets because it is longer, more formal, multilingual, and often emotionally restrained. The research problem is therefore not simple emotion classification. It is emotionally cautious pattern analysis over structured testimony."
    )
    doc.add_paragraph(
        "The target outputs are fear score, anxiety score, anger score, stress timeline, and emotional transitions over time."
    )

    doc.add_heading("11.2 Input Modalities", level=2)
    for text in [
        "Video testimony provides facial expressions and visible behavioral cues.",
        "Audio testimony provides prosody, pause structure, pitch variation, and stress-related vocal patterns.",
        "Transcript provides lexical and semantic context in English, Hindi, Tamil, Telugu, and Kannada.",
    ]:
        doc.add_paragraph(text, style="List Bullet")

    doc.add_heading("11.3 Processing Pipeline", level=2)
    for text in [
        "First, the testimony is segmented into aligned time windows so that face, audio, and text remain synchronized.",
        "Face frames are passed through a Vision Transformer to extract visual emotion embeddings.",
        "Audio segments are encoded with XLS-R or IndicWav2Vec to capture speech-level affective features.",
        "Text segments are encoded with IndicBERT v2 or MuRIL to handle multilingual courtroom language.",
        "A MemoCMT-style cross-modal transformer fuses the modality embeddings and learns cross-attention between face, speech, and transcript.",
        "A temporal aggregation layer converts per-segment predictions into a continuous emotional timeline.",
    ]:
        doc.add_paragraph(text, style="List Bullet")

    doc.add_paragraph()
    doc.add_picture(str(PHASE2_DIAGRAM), width=Inches(5.8))
    cap = doc.add_paragraph("Figure: LegalMemoCMT Phase 2 pipeline overview")
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("11.4 Adaptation Strategy", level=2)
    for text in [
        "Phase 2 starts from the Phase 1 backbone. The architecture is not rebuilt from zero; it is adapted by changing the encoders, input preprocessing, and evaluation target.",
        "If courtroom data is limited, the practical strategy is to pretrain on MELD and IEMOCAP, then fine-tune on small legal-domain samples using transfer learning, partial freezing, or lightweight adaptation layers.",
        "Language handling is done at the transcript level, while the audio branch remains language-agnostic and the fusion layer learns modality interactions across languages.",
    ]:
        doc.add_paragraph(text)

    doc.add_heading("11.5 Explainability and Legal Safety", level=2)
    for text in [
        "The system should expose attention maps or modality contribution summaries so a reviewer can see why a segment received a certain emotional score.",
        "The thesis must clearly state that the model is a research tool for emotional pattern analysis, not a legal decision engine.",
    ]:
        doc.add_paragraph(text, style="List Bullet")

    doc.add_heading("11.6 Expected Practical Output", level=2)
    for text in [
        "A working prototype that accepts testimony audio, video, and text and returns emotion scores per segment.",
        "A timeline plot showing how emotional intensity changes during testimony.",
        "A comparison between benchmark-domain performance and courtroom-domain behavior.",
    ]:
        doc.add_paragraph(text, style="List Bullet")

    doc.add_heading("11.7 Implementation Checklist", level=2)
    for text in [
        "Reproduce MemoCMT-like results on MELD and IEMOCAP first.",
        "Build preprocessing for face extraction, audio segmentation, and transcript alignment.",
        "Integrate ViT, XLS-R / IndicWav2Vec, and IndicBERT v2 / MuRIL encoders.",
        "Add cross-modal fusion and temporal aggregation.",
        "Evaluate with emotion metrics and visual timelines.",
    ]:
        doc.add_paragraph(text, style="List Bullet")

    doc.save(str(DOCX_PATH))


def main():
    if not PIPELINE_DIAGRAM.exists():
        raise FileNotFoundError(f"Missing diagram: {PIPELINE_DIAGRAM}")
    update_pptx()
    update_docx()
    print("Updated First Review PPTX and DOCX with Phase 2 methodology content.")


if __name__ == "__main__":
    main()
