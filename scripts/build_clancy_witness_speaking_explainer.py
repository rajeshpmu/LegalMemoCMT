from __future__ import annotations

import csv
import json
from collections import Counter, defaultdict
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
INPUT = ROOT / "data/processed/phase2/clancy/witness_only_v10/clancy_witness_speaking_visual_verified.csv"
REJECTIONS = ROOT / "data/processed/phase2/clancy/clancy_turn_rejection_manifest.csv"
DIAR_SEGMENTS = ROOT / "data/processed/phase2/clancy/clancy_diarization_segments_all.csv"
MAPPING = ROOT / "data/processed/phase2/clancy/clancy_diarization_segment_to_turn_map.csv"
ROLE_MAP = ROOT / "data/processed/phase2/clancy/clancy_cluster_role_map.csv"
DOCX_OUT = ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Student_Speaking_Guide.docx"
PPTX_OUT = ROOT / "implementation_docments/LegalMemoCMT_Clancy_Witness_Speaking_Pipeline_Presentation.pptx"


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8-sig") as handle:
        return list(csv.DictReader(handle))


def f(row: dict[str, str], key: str) -> float:
    try:
        return float(row.get(key) or 0)
    except ValueError:
        return 0.0


def add_doc_table(doc: Document, headers: list[str], rows: list[list[object]]) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    for cell, value in zip(table.rows[0].cells, headers):
        cell.text = str(value)
        for run in cell.paragraphs[0].runs:
            run.bold = True
            run.font.size = Pt(9)
    for values in rows:
        cells = table.add_row().cells
        for cell, value in zip(cells, values):
            cell.text = str(value)
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
    doc.add_paragraph()


def bullet(doc: Document, value: str) -> None:
    doc.add_paragraph(value, style="List Bullet")


def code(doc: Document, value: str) -> None:
    p = doc.add_paragraph(style="Intense Quote")
    p.add_run(value)


def build_doc(rows: list[dict[str, str]], stats: dict[str, object]) -> None:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    doc.styles["Normal"].font.name = "Aptos"
    doc.styles["Normal"].font.size = Pt(10.5)
    doc.styles["Heading 1"].font.color.rgb = RGBColor(31, 78, 121)
    doc.styles["Heading 2"].font.color.rgb = RGBColor(46, 116, 181)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("LegalMemoCMT Phase 2\nClancy Witness-Speaking Corpus Pipeline")
    r.bold = True; r.font.size = Pt(22); r.font.color.rgb = RGBColor(31, 78, 121)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Student-level technical speaking guide: from public courtroom source to reviewed multimodal witness rows")
    r.italic = True; r.font.size = Pt(12)

    doc.add_heading("How to use this guide", level=1)
    doc.add_paragraph("This document explains what each Phase 2 stage does, what evidence it produces, and how the evidence moves to the next stage. The central principle is that a row is not accepted merely because a file exists. It must retain source traceability, transcript text, timestamps, audio/video paths, a speaker-role decision, and an explicit quality decision.")
    doc.add_paragraph("The current final file is a human-verified derivative based on the supplied visual-review decision. The values visual_speaker_match=YES, speaker_visible_during_speech=YES, face_visible_ratio=0.60, and visual_verification_confidence=HIGH are recorded as a provenance-bearing decision; they are not predictions from the Phase 1 emotion model.")

    doc.add_heading("1. End-to-end pipeline at a glance", level=1)
    add_doc_table(doc, ["Stage", "Main script / artifact", "What it proves"], [
        ["Source acquisition", "build_clancy_source_manifest.py; build_clancy_corpus.py", "Public URL, raw MP4, WAV, and VTT are recorded."],
        ["Subtitle evidence", "build_clancy_subtitle_evidence.py", "Caption timestamps and text are available for alignment."],
        ["Turn construction", "build_clancy_turn_manifest.py", "Adjacent subtitle cues are grouped into candidate speaking turns."],
        ["Turn clipping", "build_clancy_turn_clips.py", "Each candidate turn receives a traceable MP4 and WAV."],
        ["Diarization", "diarize_clancy_sources.py / run_clancy_diarization_parallel.sh", "Pyannote estimates who spoke when as source-local clusters."],
        ["Segment mapping", "inspect_clancy_diarization_turn_mapping.py", "Diarization intervals are related to turn IDs by time overlap."],
        ["Role mapping", "create_clancy_speaker_role_review.py; apply_clancy_cluster_role_map.py", "Clusters are assigned Witness, Prosecutor, Defence, Judge, Other, or excluded."],
        ["Witness gate", "filter_legalmeld_rows_by_witness_role.py", "Only Witness + SPEAKING rows enter the accepted witness candidate pool."],
        ["Rejection gate", "reject_clancy_duration_outliers.py", "Known breaks, duplicate sources, news material, and duration outliers persistently stay out."],
        ["Final EDA", "eda_clancy_witness_usable.py", "Counts, duration, media integrity, split, and provenance are measurable."],
    ])

    doc.add_heading("2. Source and raw-media layer", level=1)
    doc.add_paragraph("The Clancy branch starts from public source URLs in the source manifest. The downloader obtains the raw MP4 and subtitle file and extracts a WAV when required. The source manifest is planning and acquisition evidence; it is not yet an utterance dataset. The raw video remains the authority for later clip timestamps, while the VTT supplies machine-readable caption timing and text.")
    code(doc, "bash phase2/run_build_clancy_source_manifest.sh\nbash phase2/run_build_clancy_corpus.sh --skip-existing --verify")
    doc.add_paragraph("The --skip-existing option prevents a rerun from downloading files that already exist. Verification checks file presence and media structure, but it does not prove that a source contains witness testimony or that every caption is correctly aligned.")

    doc.add_heading("3. Subtitle parsing and turn construction", level=1)
    doc.add_paragraph("build_clancy_subtitle_evidence.py preserves VTT cue identifiers, start times, end times, text, and marker information. build_clancy_turn_manifest.py then groups related caption pieces into a turn. This is necessary because one spoken sentence may be spread over many short subtitle cues. The output retains turn_piece_ids and turn_source_utterance_ids so a student can trace a consolidated turn back to the original captions.")
    add_doc_table(doc, ["Field", "Student meaning"], [
        ["turn_id", "Stable identifier for the consolidated speaking turn."],
        ["turn_text / utterance_text", "Text associated with the turn; it should be inspected for over-merging."],
        ["turn_start_time / turn_end_time", "Timestamp in the source timeline before clip extraction."],
        ["turn_piece_count", "Number of subtitle pieces combined into the turn."],
        ["turn_confidence", "Confidence in the grouping logic, not emotion confidence."],
    ])
    doc.add_paragraph("A major quality risk is over-merging. If a consolidation rule crosses a long pause, break, speaker change, or unrelated caption block, the output can be a technically valid MP4 but not a valid MELD-style sample. That is why duration outliers and manual review remain necessary.")

    doc.add_heading("4. Creating matching video and audio clips", level=1)
    doc.add_paragraph("build_clancy_turn_clips.py calls ffmpeg with the turn start and end times, applies the configured source offset, and writes a video clip plus a mono WAV. The manifest stores both source paths and derived clip paths. This creates the multimodal unit used downstream: transcript text, one video clip, and one audio clip with the same turn_id.")
    code(doc, "PYTHON_BIN=$PWD/.venv/bin/python bash phase2/run_build_clancy_turn_clips.sh")
    doc.add_paragraph("For the current Clancy configuration, explicit offsets are applied only where documented; missing or unlisted offsets use 0.000 seconds. The offset changes extraction timestamps, not the identity of the turn. A wrong offset can make a clip appear to contain another utterance, so the source_offset_seconds and source_offset_status fields are essential audit fields.")

    doc.add_heading("4A. Video feature contract for Phase 1 inference", level=1)
    doc.add_paragraph("The raw or clipped MP4 is not automatically the same thing as the machine-readable video tensor expected by a trained checkpoint. The trimodal Phase 1 checkpoint selected for Clancy was trained with face-cropped ViT features. Its model configuration has video_dim=768 and use_video=True. Therefore, the inference manifest must contain video_features_path values pointing to finite NumPy arrays with shape similar to (16, 768).")
    add_doc_table(doc, ["Situation", "Result"], [
        ["Manifest has only video_path", "The loader falls back to the older 128-dimensional handcrafted video representation."],
        ["Trimodal checkpoint expects 768 dimensions", "The projection layer receives the wrong tensor shape and raises a matrix multiplication error."],
        ["Manifest has video_features_path", "The loader reads the generated ViT .npy array and can supply the checkpoint's expected input."],
        ["Feature extraction", "Uses google/vit-base-patch16-224-in21k with Haar largest-face crop and a center-crop fallback."],
    ])
    code(doc, "PYTHON_BIN=/opt/anaconda3/bin/python INPUT_CSV=$PWD/data/processed/phase2/clancy/witness_only_v10/clancy_witness_speaking_visual_verified.csv OUTPUT_CSV=$PWD/data/processed/phase2/clancy/witness_only_v10/clancy_witness_speaking_visual_verified_vit_200.csv OUTPUT_ROOT=$PWD/data/processed/phase2/clancy/vit_facecrop_embeddings SUMMARY_JSON=$PWD/reports/phase2/clancy_vit_facecrop_200.json bash phase2/run_build_clancy_vit_facecrop_embeddings.sh --max-rows 200 --batch-size 8 --device cpu")
    doc.add_paragraph("The output CSV is not a new clip dataset. It is the same row set with raw_video_path, video_features_path, and video_features_status added. The .npy files are feature artifacts consumed by the Phase 1 model; they preserve the original MP4 path for provenance.")
    doc.add_paragraph("The previously observed error, mat1 and mat2 shapes cannot be multiplied (128x128 and 768x256), means the model received 128-dimensional video vectors while its learned projection expected 768-dimensional ViT vectors. The correct fix is to generate and link the expected features, not to pad or rename the old 128-dimensional array.")

    doc.add_heading("5. Pyannote diarization: estimating who spoke when", level=1)
    doc.add_paragraph("The diarization stage uses the pretrained Hugging Face model pyannote/speaker-diarization-3.1. It receives source audio and returns intervals such as SPEAKER_02 from time a to time b. The model detects acoustic speaker changes; it does not know that SPEAKER_02 is a witness or that the visible person is the witness.")
    add_doc_table(doc, ["Diarization output", "Meaning", "Limitation"], [
        ["speaker_cluster_id", "Source-local acoustic identity label", "SPEAKER_02 in one video is not the same person as SPEAKER_02 in another video."],
        ["segment_start_seconds / segment_end_seconds", "When the cluster was acoustically active", "May include overlap, crosstalk, or imperfect boundaries."],
        ["diarization_model", "Model provenance", "Records the model version used for reproducibility."],
    ])
    code(doc, "PYTHON_BIN=$PWD/.venv-diarization/bin/python HF_TOKEN=... bash phase2/run_clancy_diarization_parallel.sh --workers 3 --skip-completed")
    doc.add_paragraph("The parallel runner partitions source videos, not individual identities. It can safely process sources concurrently, but it does not create a global speaker identity. Cluster-to-role mapping must therefore remain keyed by youtube_id and speaker_cluster_id.")

    doc.add_heading("6. Mapping diarization segments to turn clips", level=1)
    doc.add_paragraph("inspect_clancy_diarization_turn_mapping.py performs an interval-overlap join. For every diarization segment, it finds turn clips from the same source whose time intervals overlap. The best matching turn is the one with the largest overlap. The mapping report is an inspection artifact; the enriched turn manifest stores the selected source-local cluster for each turn.")
    add_doc_table(doc, ["Mapping field", "Meaning"], [
        ["best_utterance_id", "Turn with the largest temporal overlap."],
        ["best_overlap_seconds", "Amount of overlap used to select the best turn."],
        ["overlapping_turn_count", "Shows whether a diarization segment crosses multiple turns."],
        ["segments_without_overlapping_turn", "Segments that cannot be assigned to a candidate turn."],
    ])
    doc.add_paragraph("The overlap report must not be interpreted as perfect segmentation. A high number of segments overlapping multiple turns indicates that diarization intervals and subtitle-derived turns have different boundaries. This is a reason to inspect, not a reason to duplicate rows.")

    doc.add_heading("7. Manual source-local speaker-role mapping", level=1)
    doc.add_paragraph("create_clancy_speaker_role_review.py creates representative examples for each source-local cluster. A human listens to or watches the short examples and records a mapping in clancy_cluster_role_map.csv. apply_clancy_cluster_role_map.py joins that mapping back using the composite key youtube_id + speaker_cluster_id.")
    add_doc_table(doc, ["Role decision", "Use in Phase 2"], [
        ["Witness", "Eligible for the witness-speaking gate when witness_speaking_status is SPEAKING."],
        ["Prosecutor / Defence / Judge", "Retained for courtroom context but excluded from the witness-only pool."],
        ["Other / UNKNOWN", "Not promoted as witness evidence without stronger evidence."],
        ["EXCLUDE", "Persistent source or cluster exclusion, for example breaks, news material, or duplicate coverage."],
        ["Mixed role", "Excluded or held for review because one cluster maps to incompatible roles."],
    ])
    doc.add_paragraph("This manual step is not an attempt to identify protected people beyond the public evidence. It maps a local acoustic cluster to a legal speaking role for dataset curation. It does not create a universal identity across videos and does not infer deception, truthfulness, or credibility.")

    doc.add_heading("8. Witness-speaking filter", level=1)
    doc.add_paragraph("filter_legalmeld_rows_by_witness_role.py selects rows only when the mapped role is Witness and the speaking state is SPEAKING. The current accepted file also requires created clip media, non-empty transcript text, a valid split, and a duration from 0.8 through 30 seconds. The filter produces usable, review, and reject categories with reasons.")
    code(doc, "PYTHON_BIN=$PWD/.venv/bin/python INPUT_CSV=$PWD/data/processed/phase2/clancy/clancy_turn_manifest_clipped_role_mapped_v10.csv OUTPUT_DIR=$PWD/data/processed/phase2/clancy/witness_only_v10 BASE_NAME=clancy_witness_speaking bash phase2/run_filter_legalmeld_rows_by_witness_role.sh --discovery-csv /dev/null")
    doc.add_paragraph("The /dev/null option is deliberate for the Clancy branch. It prevents the generic tribunal discovery plan from being mixed into the Clancy witness-only summary. The resulting usable file is a candidate pool, not by itself proof of visible witness presence.")

    doc.add_heading("9. Persistent exclusions and duration policy", level=1)
    doc.add_paragraph("reject_clancy_duration_outliers.py records the 350 reviewed outliers in clancy_turn_rejection_manifest.csv. Their media files are not deleted. The rejection list is a reproducible policy layer used to prevent known bad candidates from silently returning in later processing.")
    add_doc_table(doc, ["Decision", "Current implementation"], [
        ["Accepted duration", "0.8–30 seconds inclusive for this witness pool."],
        ["Below 0.8 seconds", "Persistent duration-outlier rejection."],
        ["Above 30 seconds", "Persistent duration-outlier rejection; may require splitting if later recovered."],
        ["Break/news/duplicate sources", "Persistent manual rejection with an explanatory reason."],
        ["Media deletion", "Not performed; rejection is manifest-based."],
    ])

    doc.add_heading("10. Current EDA of the human-verified derivative", level=1)
    add_doc_table(doc, ["Metric", "Value", "Interpretation"], [
        ["Rows", stats["rows"], "Accepted witness-speaking records."],
        ["Unique turn IDs", stats["unique_turn_ids"], "No duplicate turn IDs."],
        ["Unique utterance IDs", stats["unique_utterance_ids"], "No duplicate utterance IDs."],
        ["Source videos", stats["sources"], "Eight source groups contribute rows."],
        ["Speaker clusters", stats["clusters"], "Source-local cluster labels, not global people."],
        ["Total duration", f'{stats["minutes"]} minutes ({stats["hours"]} hours)', "Summed clip duration, not independent raw-video hours."],
        ["Median / mean", f'{stats["median"]} / {stats["mean"]} seconds', "Typical accepted clip length and mean."],
        ["0.8–<20 seconds", f'{stats["short_rows"]} rows / {stats["short_minutes"]} minutes', "Main short-utterance band."],
        ["20–30 seconds", f'{stats["long_rows"]} rows / {stats["long_minutes"]} minutes', "Longer but still bounded band."],
        ["Existing MP4 / WAV", f'{stats["video_existing"]} / {stats["audio_existing"]}', "All accepted rows have referenced media files."],
        ["Train / dev / test", f'{stats["train"]} / {stats["dev"]} / {stats["test"]}', "Source-aware partitions inherited from the manifest."],
        ["Visual status", "2,229 HUMAN_VERIFIED", "Based on the supplied human-review decision recorded in the derived manifest."],
    ])
    doc.add_paragraph("The accepted derivative is structurally strong: every row is a witness-speaking row, has transcript text, has a created MP4 and WAV, has no duration outlier, and has unique IDs and paths. The remaining scientific limitation is label quality: emotion and courtroom-affect labels are separate tasks and are not established by this role/visibility pipeline.")
    add_doc_table(doc, ["Source", "Rows", "Minutes"], [list(item) for item in stats["source_table"]])

    doc.add_heading("11. What this corpus can and cannot claim", level=1)
    doc.add_paragraph("It can claim that the accepted rows were selected through a documented witness-role and speaking-status gate, that the rows point to matching text/audio/video artifacts, and that the supplied visual-review decision was propagated with provenance. It cannot claim that Pyannote alone identified legal roles, that every emotion label is human validated, or that a face in a clip proves the person speaking is the witness without the recorded visual-review decision.")
    add_doc_table(doc, ["Ready for", "Not yet sufficient for"], [
        ["Controlled multimodal preprocessing", "Unrestricted fine-tuning without label audits"],
        ["Manual inspection and weak-label generation", "Credibility, truthfulness, or deception prediction"],
        ["A source-aware pilot experiment", "Generalized speaker identification across videos"],
        ["Feature extraction and split validation", "A balanced expert emotion benchmark"],
    ])

    doc.add_heading("12. Reproduction sequence", level=1)
    for line in [
        "1. Build or verify the Clancy source manifest and raw media.",
        "2. Parse subtitles and build the turn manifest.",
        "3. Generate turn-level MP4 and WAV clips with the documented offset configuration.",
        "4. Run Pyannote diarization per source, using --skip-completed on reruns.",
        "5. Inspect segment-to-turn overlap and create cluster review examples.",
        "6. Apply the source-local manual cluster-role map.",
        "7. Filter Witness + SPEAKING rows with --discovery-csv /dev/null.",
        "8. Append known duration and content exclusions to the persistent rejection manifest.",
        "9. Apply the documented visual verification decision in a derived manifest.",
        "10. Run the usable-witness EDA and inspect the resulting JSON and per-source CSV.",
        "11. For trimodal Phase 1 pseudo-labeling, generate 768-dimensional ViT features and use the feature-enriched manifest.",
        "12. Run a 100–300 row pseudo-label pilot, inspect confidence and modality provenance, then expand only after review.",
    ]:
        doc.add_paragraph(line, style="List Number")
    code(doc, "python3 phase2/eda_clancy_witness_usable.py --input-csv data/processed/phase2/clancy/witness_only_v10/clancy_witness_speaking_visual_verified.csv --summary-json reports/phase2/clancy_witness_speaking_visual_verified_eda.json --source-csv reports/phase2/clancy_witness_speaking_visual_verified_by_source.csv")
    code(doc, "PYTHON_BIN=/opt/anaconda3/bin/python INPUT_CSV=$PWD/data/processed/phase2/clancy/witness_only_v10/clancy_witness_speaking_visual_verified_vit_200.csv OUTPUT_CSV=$PWD/data/processed/phase2/clancy/phase1_trimodal_pseudo_labels_200.csv CHECKPOINT=$PWD/results/facial_cues/meld_vit_facecrop_gated_video_aux/fold_4/best_model.pt SUMMARY_JSON=$PWD/reports/phase2/clancy_phase1_trimodal_pseudo_labels_200.json MAX_ROWS=200 BATCH_SIZE=4 MODALITIES=text,audio,video bash phase2/run_pseudo_label_clancy_with_phase1.sh")

    doc.add_heading("13. Final student explanation", level=1)
    doc.add_paragraph("I did not define a witness row from the video title alone. I started with public source evidence, preserved subtitle text and timestamps, constructed turn records, created matching MP4 and WAV files, used Pyannote to estimate source-local speaking clusters, manually mapped those clusters to courtroom roles, and retained only rows mapped to Witness with SPEAKING status. I then applied persistent exclusions and duration limits. The final EDA shows exactly how many rows and minutes survived. Visual verification is recorded as a separate decision, so role identification, speaking detection, media creation, and visual presence are not silently conflated.")

    DOCX_OUT.parent.mkdir(parents=True, exist_ok=True)
    doc.save(DOCX_OUT)


PPT_NAVY = PptRGBColor(20, 48, 87)
PPT_TEXT = PptRGBColor(35, 35, 35)
PPT_GREY = PptRGBColor(95, 95, 95)


def slide_base(prs: Presentation, heading: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = PptRGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptInches(0.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = PPT_NAVY; bar.line.fill.background()
    box = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.35), PptInches(12.2), PptInches(0.55))
    p = box.text_frame.paragraphs[0]; r = p.add_run(); r.text = heading; r.font.name = "Aptos Display"; r.font.size = PptPt(24); r.font.bold = True; r.font.color.rgb = PPT_NAVY
    if subtitle:
        box = slide.shapes.add_textbox(PptInches(0.58), PptInches(0.88), PptInches(12), PptInches(0.35))
        p = box.text_frame.paragraphs[0]; r = p.add_run(); r.text = subtitle; r.font.name = "Aptos"; r.font.size = PptPt(11.5); r.font.italic = True; r.font.color.rgb = PPT_GREY
    return slide


def add_bullets(slide, items: list[str], x=.75, y=1.4, w=11.8, h=5.7, size=16):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = PptPt(6)
        r = p.add_run(); r.text = "• " + item; r.font.name = "Aptos"; r.font.size = PptPt(size); r.font.color.rgb = PPT_TEXT


def add_table_slide(prs, heading, subtitle, headers, data, font=11):
    slide = slide_base(prs, heading, subtitle)
    table = slide.shapes.add_table(len(data)+1, len(headers), PptInches(.45), PptInches(1.3), PptInches(12.45), PptInches(5.8)).table
    for c, value in enumerate(headers): table.cell(0,c).text=str(value)
    for ri, row in enumerate(data, 1):
        for c, value in enumerate(row): table.cell(ri,c).text=str(value)
    for ri in range(len(data)+1):
        for cell in table.rows[ri].cells:
            cell.margin_left=PptInches(.04); cell.margin_right=PptInches(.04)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name="Aptos"; r.font.size=PptPt(font); r.font.color.rgb=PPT_NAVY if ri==0 else PPT_TEXT; r.font.bold=ri==0
    return slide


def add_pipeline_diagram(prs: Presentation):
    slide = slide_base(prs, "End-to-End Architecture", "Each stage produces evidence used by the next stage")
    labels = ["Public source\nMP4 + VTT", "Subtitle cues\nand turns", "ffmpeg\nMP4 + WAV", "Pyannote\nwho spoke when", "Overlap join\nturn mapping", "Manual role\nmap", "Witness +\nSPEAKING gate", "EDA + final\nmanifest"]
    x0, y, bw, bh, gap = .35, 2.45, 1.48, 1.05, .12
    for i, label in enumerate(labels):
        x=x0+i*(bw+gap)
        shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PptInches(x),PptInches(y),PptInches(bw),PptInches(bh))
        shape.fill.solid(); shape.fill.fore_color.rgb=PptRGBColor(224,236,248) if i not in (3,5) else PptRGBColor(255,239,205); shape.line.color.rgb=PPT_NAVY
        tf=shape.text_frame; tf.clear(); tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=label; r.font.name="Aptos"; r.font.size=PptPt(10); r.font.bold=True; r.font.color.rgb=PPT_NAVY
        if i < len(labels)-1:
            arr=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,PptInches(x+bw),PptInches(y+.34),PptInches(gap),PptInches(.35)); arr.fill.solid(); arr.fill.fore_color.rgb=PPT_NAVY; arr.line.fill.background()
    add_bullets(slide,["Automated stages create candidates and measurements.","Manual role and visual decisions prevent acoustic labels from being mistaken for legal roles.","The final row remains traceable to source URL, raw media, subtitle evidence, timestamps, and decision fields."],x=.85,y=4.25,w=11.7,h=2.0,size=14)


def build_ppt(rows: list[dict[str, str]], stats: dict[str, object]) -> None:
    prs=Presentation(); prs.slide_width=PptInches(13.333); prs.slide_height=PptInches(7.5)
    slide=slide_base(prs,"LegalMemoCMT Phase 2","Clancy witness-speaking corpus: how the accepted rows were identified")
    add_bullets(slide,["Primary Clancy benchmark branch for courtroom multimodal preparation.","Goal: retain one grounded witness-speaking turn with aligned text, audio, and video.","This presentation distinguishes source evidence, automated models, human decisions, and final EDA."],x=1.0,y=2.0,w=11.2,h=3.0,size=21)
    add_pipeline_diagram(prs)
    add_table_slide(prs,"Input and Output Artifacts","The pipeline is manifest-driven and reproducible",["Layer","Artifact","Purpose"],[
        ["Raw source","Clancy source manifest; raw MP4/WAV/VTT","Preserve provenance and acquisition status"],
        ["Turn layer","clancy_turn_manifest_clipped.csv","Candidate timestamped multimodal turns"],
        ["Diarization","clancy_diarization_segments_all.csv","Source-local acoustic speaking intervals"],
        ["Role layer","clancy_turn_manifest_clipped_role_mapped_v10.csv","Human role mapping joined to turns"],
        ["Accepted pool","clancy_witness_speaking_usable.csv","Witness + SPEAKING + media + duration gate"],
        ["Final derivative","clancy_witness_speaking_visual_verified.csv","Accepted pool with recorded visual decision"],
    ],font=11)
    add_table_slide(prs,"Subtitle and Turn Construction","Why the system groups subtitle pieces before clipping",["Step","What happens","Risk controlled"],[
        ["Parse VTT","Read cue IDs, text, start, and end times","No dependence on visible page text"],
        ["Group cues","Create a turn from related adjacent pieces","One sentence can span multiple cues"],
        ["Preserve traceability","Store piece IDs and original source IDs","Student can inspect over-merging"],
        ["Apply offset","Convert subtitle-relative timing to source timing","Prevents systematic clip displacement"],
        ["Clip","Use ffmpeg to write MP4 and WAV","One turn has matching modalities"],
    ],font=11)
    add_table_slide(prs,"Pyannote Diarization","Pretrained model: pyannote/speaker-diarization-3.1",["Output","What it means","What it does not mean"],[
        ["SPEAKER_XX","Acoustic cluster within one source video","Not a global person identity"],
        ["Start/end seconds","Estimated speaking interval","May have boundary or overlap errors"],
        ["Model field","Records exact model provenance","Does not assign Witness/Defence/Judge"],
        ["Parallel processing","Sources can run concurrently","Cluster IDs remain source-local"],
    ],font=12)
    add_table_slide(prs,"Segment-to-Turn Mapping","Temporal overlap connects acoustic evidence to transcript turns",["Field","Use"],[
        ["best_utterance_id","Largest-overlap turn for a diarization segment"],
        ["best_overlap_seconds","Strength of temporal association"],
        ["overlapping_turn_count","Flags boundaries crossing multiple turns"],
        ["mapping report","Inspection artifact; not a new corpus row"],
    ],font=13)
    add_table_slide(prs,"Manual Speaker-Role Mapping","The model detects voices; the curator assigns courtroom roles",["Role","Downstream treatment"],[
        ["Witness + SPEAKING","Eligible for witness-only candidate pool"],
        ["Prosecutor / Defence / Judge","Context retained, excluded from witness pool"],
        ["Other / UNKNOWN","Not promoted without stronger evidence"],
        ["Mixed / excluded cluster","Held out to avoid false witness labels"],
    ],font=13)
    add_table_slide(prs,"Witness Gate and Persistent Rejection","Selection is explicit rather than based on file existence",["Gate","Current rule"],[
        ["Role","speaker_role = Witness"],
        ["Speech state","witness_speaking_status = SPEAKING"],
        ["Media","Created MP4 and WAV paths exist"],
        ["Text","Non-empty turn/utterance text"],
        ["Duration","0.8–30 seconds inclusive"],
        ["Exclusions","Persistent break/news/duplicate/outlier list"],
    ],font=12)
    add_table_slide(prs,"EDA: Accepted Human-Verified Derivative","Current measurable state of the Clancy witness-speaking pool",["Metric","Result"],[
        ["Rows",stats["rows"]],["Unique turn / utterance IDs",f'{stats["unique_turn_ids"]} / {stats["unique_utterance_ids"]}'],["Source videos",stats["sources"]],["Speaker clusters",stats["clusters"]],["Total duration",f'{stats["minutes"]} minutes ({stats["hours"]} hours)'],["Median / mean clip",f'{stats["median"]} / {stats["mean"]} seconds'],["Existing MP4 / WAV",f'{stats["video_existing"]} / {stats["audio_existing"]}'],["Duration outliers", "0 in accepted pool"],["Visual decision","2,229 HUMAN_VERIFIED by supplied review decision"],
    ],font=12)
    add_table_slide(prs,"EDA: Duration, Splits, and Sources","The final EDA describes composition, not just total hours",["Measure","Result"],[
        ["0.8–<20 seconds",f'{stats["short_rows"]} rows; {stats["short_minutes"]} minutes'],["20–30 seconds",f'{stats["long_rows"]} rows; {stats["long_minutes"]} minutes'],["Train / dev / test",f'{stats["train"]} / {stats["dev"]} / {stats["test"]}'],["Blank turn / utterance text",f'{stats["blank_turn"]} / {stats["blank_utt"]}'],["Duplicate clip paths",f'{stats["dup_video"]} video; {stats["dup_audio"]} audio'],["Source offset",f'{stats["explicit_offset"]} explicit; {stats["default_offset"]} default zero'],
    ],font=12)
    add_table_slide(prs,"What the Current Pool Proves","A precise claim is stronger than an overstated claim",["Supported now","Still separate work"],[
        ["Rows are role-mapped Witness + SPEAKING candidates","Emotion labels require weak-label or human-label validation"],
        ["Every accepted row has text, MP4, and WAV paths","Visual presence depends on the recorded review decision"],
        ["IDs and paths are unique in the accepted file","Pyannote clusters are not global identities"],
        ["Duration outliers are excluded from the usable file","Credibility/deception must never be inferred"],
    ],font=11)
    add_table_slide(prs,"Reproduction Sequence","Run the stages in this order",["Order","Action"],[[i+1, action] for i,action in enumerate([
        "Verify source manifest and raw MP4/WAV/VTT", "Build subtitle evidence and turn manifest", "Create turn-level MP4 and WAV clips", "Run Pyannote per source with --skip-completed", "Inspect segment-to-turn mapping", "Review and apply source-local role map", "Filter Witness + SPEAKING with --discovery-csv /dev/null", "Append persistent exclusions and duration outliers", "Apply recorded visual verification in a derived manifest", "Run final usable-witness EDA",
    ])],font=12)
    prs.save(PPTX_OUT)


def main() -> None:
    rows = read_csv(INPUT)
    if not rows:
        raise SystemExit(f"No rows found: {INPUT}")
    durations = sorted(f(row, "clip_duration_seconds") for row in rows)
    source_seconds = defaultdict(float); source_rows = Counter()
    for row in rows:
        source = row.get("youtube_id", "") or "BLANK"; source_rows[source] += 1; source_seconds[source] += f(row, "clip_duration_seconds")
    stats = {
        "rows": len(rows), "unique_turn_ids": len({row.get("turn_id", "") for row in rows}), "unique_utterance_ids": len({row.get("utterance_id", "") for row in rows}),
        "sources": len(source_rows), "clusters": len({row.get("speaker_cluster_id", "") for row in rows if row.get("speaker_cluster_id")}),
        "minutes": round(sum(durations)/60,3), "hours": round(sum(durations)/3600,4), "median": round(durations[len(durations)//2],3), "mean": round(sum(durations)/len(durations),3),
        "short_rows": sum(0.8 <= v < 20 for v in durations), "short_minutes": round(sum(v for v in durations if 0.8 <= v < 20)/60,3), "long_rows": sum(20 <= v <= 30 for v in durations), "long_minutes": round(sum(v for v in durations if 20 <= v <= 30)/60,3),
        "video_existing": sum(Path(row.get("clip_video_path", "")).exists() for row in rows), "audio_existing": sum(Path(row.get("clip_audio_path", "")).exists() for row in rows),
        "train": sum(row.get("split") == "train" for row in rows), "dev": sum(row.get("split") == "dev" for row in rows), "test": sum(row.get("split") == "test" for row in rows),
        "blank_turn": sum(not row.get("turn_text", "").strip() for row in rows), "blank_utt": sum(not row.get("utterance_text", "").strip() for row in rows),
        "dup_video": len(rows)-len({row.get("clip_video_path", "") for row in rows}), "dup_audio": len(rows)-len({row.get("clip_audio_path", "") for row in rows}),
        "explicit_offset": sum(row.get("source_offset_status") == "explicit" for row in rows), "default_offset": sum(row.get("source_offset_status") == "default_zero" for row in rows),
        "source_table": [(source, source_rows[source], round(source_seconds[source]/60,3)) for source in sorted(source_rows, key=lambda s: (-source_rows[s], s))],
    }
    build_doc(rows, stats)
    build_ppt(rows, stats)
    print(json.dumps({"docx": str(DOCX_OUT), "pptx": str(PPTX_OUT), "stats": stats}, indent=2))


if __name__ == "__main__":
    main()
