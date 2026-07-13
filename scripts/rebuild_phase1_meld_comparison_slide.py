from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PPTX_PATH = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT/implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx")


def style_run(run, *, size=10, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_box(slide, left, top, width, height, fill, line=None, radius=False):
    shp_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shp_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line:
        shp.line.color.rgb = RGBColor.from_string(line)
    else:
        shp.line.color.rgb = RGBColor.from_string(fill)
    return shp


def add_text(slide, left, top, width, height, text, *, size=11, bold=False, color="000000", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    style_run(r, size=size, bold=bold, color=color)
    return tb


def remove_last_slide(prs):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)


def rebuild():
    prs = Presentation(str(PPTX_PATH))
    # Replace the last slide (comparison slide) with a cleaner version.
    remove_last_slide(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_no = len(prs.slides)

    add_box(slide, 0, 0, 13.333, 0.58, "122F55")
    add_box(slide, 0, 0.58, 13.333, 0.08, "267377")
    add_text(slide, 0.42, 0.12, 11.9, 0.3, f"{slide_no}. MELD Result Comparison: MemoCMT vs Baseline vs Gated+Aux", size=18, bold=True, color="FFFFFF")
    add_text(slide, 12.0, 0.12, 0.8, 0.3, str(slide_no), size=12, color="FFFFFF", align=PP_ALIGN.RIGHT)
    add_text(slide, 0.48, 0.7, 12.2, 0.3, "Use this slide to show the numerical comparison clearly and read the story from left to right.", size=11, color="5E5E5E")

    # Top summary strip
    add_box(slide, 0.35, 1.0, 12.63, 0.95, "EEF4FA", "C9D7E6", radius=True)
    add_text(slide, 0.55, 1.16, 12.1, 0.2, "MemoCMT paper: 64.18% Acc / 62.52% W-F1", size=13.5, bold=True, color="122F55")
    add_text(slide, 0.55, 1.46, 12.1, 0.2, "Your paper-aligned baseline mean: 62.47% Acc / 61.95% W-F1  |  Best baseline fold 2: 63.75% / 62.54%", size=12, color="3A3A3A")

    # Three large cards
    cards = [
        (0.4, "Baseline", "Closest stable reproduction of MemoCMT on MELD.", "62.47%", "61.95%", "43.95%", "44.17%", "Best fold 2: 63.75% / 62.54%"),
        (4.45, "Gated + Aux fold 2", "Selective improvement on boundary clips, but not overall better.", "60.54%", "60.22%", "43.51%", "44.92%", "Fixes some neutral boundary cases."),
        (8.5, "Gated + Aux fold 4", "Still useful for some clips, but overall below the baseline.", "59.92%", "60.56%", "43.30%", "46.38%", "Higher UW-Acc, lower overall score."),
    ]
    for x, title, desc, acc, wf1, macro, uw, footer in cards:
        add_box(slide, x, 2.15, 3.95, 3.0, "FFFFFF", "C9D7E6", radius=True)
        add_box(slide, x, 2.15, 3.95, 0.45, "DDEBF7", "DDEBF7", radius=True)
        add_text(slide, x+0.15, 2.24, 3.6, 0.18, title, size=13.2, bold=True, color="122F55")
        add_text(slide, x+0.15, 2.72, 3.6, 0.38, desc, size=11, color="3A3A3A")
        add_box(slide, x+0.15, 3.25, 3.65, 0.55, "F8FBFF", "D9E4EE", radius=True)
        add_text(slide, x+0.2, 3.35, 1.0, 0.16, "Acc", size=10, bold=True, color="5E5E5E")
        add_text(slide, x+0.76, 3.3, 1.2, 0.22, acc, size=16, bold=True, color="122F55")
        add_text(slide, x+1.9, 3.35, 1.1, 0.16, "W-F1", size=10, bold=True, color="5E5E5E")
        add_text(slide, x+2.47, 3.3, 1.1, 0.22, wf1, size=16, bold=True, color="122F55")
        add_box(slide, x+0.15, 3.95, 3.65, 0.55, "F8FBFF", "D9E4EE", radius=True)
        add_text(slide, x+0.2, 4.05, 1.1, 0.16, "Macro", size=10, bold=True, color="5E5E5E")
        add_text(slide, x+0.76, 4.0, 1.2, 0.22, macro, size=16, bold=True, color="122F55")
        add_text(slide, x+1.9, 4.05, 1.1, 0.16, "UW", size=10, bold=True, color="5E5E5E")
        add_text(slide, x+2.47, 4.0, 1.1, 0.22, uw, size=16, bold=True, color="122F55")
        add_text(slide, x+0.15, 4.65, 3.6, 0.32, footer, size=10.5, color="4B4B4B")

    # Bottom conclusion bar
    add_box(slide, 0.35, 5.55, 12.63, 1.05, "EEF4FA", "C9D7E6", radius=True)
    add_text(slide, 0.55, 5.72, 12.15, 0.2, "Conclusion", size=13, bold=True, color="122F55")
    add_text(slide, 0.55, 6.0, 12.1, 0.24, "The paper-aligned baseline is still the closest overall reproduction of MemoCMT on MELD. Gated+aux is selective improvement, not a universal replacement.", size=12.2, bold=True, color="202020")
    add_text(slide, 0.55, 6.33, 12.1, 0.18, "Paper gap on the mean baseline: about 1.71 points in accuracy and 0.57 points in weighted F1.", size=11, color="5E5E5E")

    prs.save(str(PPTX_PATH))
    print(f"Rebuilt comparison slide in {PPTX_PATH}")


if __name__ == "__main__":
    rebuild()
