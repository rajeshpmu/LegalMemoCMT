from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from docx.oxml import OxmlElement
from docx.shared import Pt, RGBColor
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Reading_Script.docx"


def style_run(run, *, size=11, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def style_ppt_run(run, *, size=10, bold=False, color="000000"):
    run.font.name = "Aptos"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color)


def insert_paragraph_before(paragraph: Paragraph, text: str = "", style: str | None = None) -> Paragraph:
    new_p = OxmlElement("w:p")
    paragraph._p.addprevious(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        run = new_para.add_run(text)
        style_run(run)
    return new_para


def insert_paragraph_after(paragraph: Paragraph, text: str = "", style: str | None = None) -> Paragraph:
    new_p = OxmlElement("w:p")
    paragraph._p.addnext(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        run = new_para.add_run(text)
        style_run(run)
    return new_para


def remove_paragraph(paragraph: Paragraph) -> None:
    paragraph._element.getparent().remove(paragraph._element)


def add_box(slide, left, top, width, height, fill, line=None, radius=False):
    shp_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shp_type,
        PptInches(left),
        PptInches(top),
        PptInches(width),
        PptInches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = PptRGBColor.from_string(fill)
    shp.line.color.rgb = PptRGBColor.from_string(line or fill)
    return shp


def add_text(slide, left, top, width, height, text, *, size=11, bold=False, color="000000", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    style_ppt_run(r, size=size, bold=bold, color=color)
    return tb


def insert_slide_before(prs: Presentation, before_index_0_based: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(before_index_0_based, sldId)
    return slide


def title_shape(slide):
    candidates = []
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = shp.text_frame.text.strip()
        if not txt:
            continue
        if shp.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if float(shp.top) > 1200000:
            continue
        candidates.append(shp)
    if not candidates:
        return None

    def score(shape):
        max_font = 0.0
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    max_font = max(max_font, float(run.font.size))
        return (max_font, -float(shape.top), -float(shape.left))

    return max(candidates, key=score)


def number_shape(slide):
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = shp.text_frame.text.strip()
        if not txt or not re.fullmatch(r"\d+", txt):
            continue
        if float(shp.top) <= 350000 and float(shp.left) >= 10000000:
            return shp
    return None


def renumber_titles(prs: Presentation) -> None:
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        tshape = title_shape(slide)
        if tshape is None:
            continue
        raw = tshape.text_frame.text.strip()
        if not raw:
            continue
        base = re.sub(r"^\d+\.\s*", "", raw)
        tshape.text_frame.clear()
        p = tshape.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"{idx}. {base}"
        r.font.name = "Aptos Display"
        r.font.size = PptPt(24)
        r.font.bold = True
        r.font.color.rgb = PptRGBColor.from_string("122F55")


def renumber_slide_numbers(prs: Presentation) -> None:
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1:
            continue
        nshape = number_shape(slide)
        if nshape is not None:
            nshape.text = str(idx)
            for p in nshape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
                for run in p.runs:
                    style_ppt_run(run, size=12, color="FFFFFF")


def build_pptx() -> None:
    prs = Presentation(str(PPTX_PATH))
    if any("CLS / MEAN / MAX / MIN Pooling Example" in (sh.text_frame.text if getattr(sh, "has_text_frame", False) else "") for slide in prs.slides for sh in slide.shapes):
        renumber_titles(prs)
        renumber_slide_numbers(prs)
        prs.save(str(PPTX_PATH))
        return

    slide = insert_slide_before(prs, 5)  # before current Slide 6

    add_box(slide, 0, 0, 13.333, 0.58, "122F55")
    add_box(slide, 0, 0.58, 13.333, 0.08, "267377")
    add_text(slide, 0.42, 0.12, 11.4, 0.3, "6. CLS / MEAN / MAX / MIN Pooling Example", size=16.5, bold=True, color="FFFFFF")
    add_text(slide, 12.0, 0.12, 0.8, 0.3, "6", size=12, color="FFFFFF", align=PP_ALIGN.RIGHT)
    add_text(
        slide,
        0.5,
        0.7,
        12.2,
        0.38,
        "This example shows how a fused token sequence becomes one classifier vector. Dim means one hidden feature dimension in that vector.",
        size=10.2,
        color="5E5E5E",
    )

    headers = ["Token", "Dim 1", "Dim 2", "Dim 3", "Dim 4", "Pooling", "Result"]
    rows = [
        ["Token 1", "0.2", "-0.1", "0.5", "0.0", "Input", "[-]"],
        ["Token 2", "0.8", "0.3", "-0.2", "0.4", "Input", "[-]"],
        ["Token 3", "0.1", "0.7", "0.6", "-0.3", "Input", "[-]"],
        ["Summary", "-", "-", "-", "-", "CLS", "[0.2, -0.1, 0.5, 0.0]"],
        ["Summary", "-", "-", "-", "-", "MEAN", "[0.37, 0.30, 0.30, 0.03]"],
        ["Summary", "-", "-", "-", "-", "MAX", "[0.8, 0.7, 0.6, 0.4]"],
        ["Summary", "-", "-", "-", "-", "MIN", "[0.1, -0.1, -0.2, -0.3]"],
    ]

    table = slide.shapes.add_table(len(rows) + 1, len(headers), PptInches(0.45), PptInches(1.25), PptInches(12.45), PptInches(4.75)).table
    widths = [1.05, 0.88, 0.88, 0.88, 0.88, 1.25, 6.63]
    for idx, w in enumerate(widths):
        table.columns[idx].width = PptInches(w)

    for c, txt in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptRGBColor.from_string("122F55")
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                style_ppt_run(run, size=10, bold=True, color="FFFFFF")

    for r, row in enumerate(rows, start=1):
        for c, txt in enumerate(row):
            cell = table.cell(r, c)
            cell.text = txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptRGBColor.from_string("F8FAFC" if r % 2 else "EEF3F8")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c != 6 else PP_ALIGN.LEFT
                for run in p.runs:
                    style_ppt_run(run, size=9.2 if c != 6 else 8.8, color="1F1F1F", bold=(c == 0 or c == 5))

    add_box(slide, 0.45, 6.15, 12.45, 0.5, "F4F7FB", line="D7E0EA", radius=True)
    add_text(
        slide,
        0.62,
        6.25,
        12.0,
        0.3,
        "Interpretation: CLS uses one designated token, MEAN smooths the sequence, MAX keeps the strongest dimension-wise evidence, and MIN keeps the weakest dimension-wise evidence.",
        size=9.7,
        color="334155",
    )

    renumber_titles(prs)
    renumber_slide_numbers(prs)
    prs.save(str(PPTX_PATH))


def heading_paragraphs(doc: Document):
    for p in doc.paragraphs:
        if re.fullmatch(r"Slide \d+:.*", p.text.strip()):
            yield p


def insert_docx_section_before(doc: Document, anchor: Paragraph, heading: str, lines: list[str]) -> None:
    p = insert_paragraph_before(anchor, heading, style="Heading 1")
    rr = p.runs[0]
    rr.font.name = "Aptos"
    rr.font.size = Pt(14)
    rr.font.bold = True
    rr.font.color.rgb = RGBColor.from_string("122F55")
    cur = p
    for line in lines:
        cur = insert_paragraph_after(cur, line)


def shift_slide_headings(doc: Document, start_num: int, delta: int) -> None:
    pattern = re.compile(r"^Slide (\d+):\s*(.*)$")
    for p in doc.paragraphs:
        m = pattern.match(p.text.strip())
        if not m:
            continue
        num = int(m.group(1))
        if num >= start_num:
            p.text = f"Slide {num + delta}: {m.group(2)}"


def build_docx() -> None:
    doc = Document(str(DOCX_PATH))
    new_heading = "Slide 6: CLS / MEAN / MAX / MIN Pooling Example"
    if any(p.text.strip() == new_heading for p in doc.paragraphs):
        doc.save(str(DOCX_PATH))
        return

    anchor = None
    for p in doc.paragraphs:
        if p.text.strip() == "Slide 6: Why MIN Pooling for MELD":
            anchor = p
            break
    if anchor is None:
        raise RuntimeError("Could not find the Slide 6 Why MIN heading in the DOCX.")

    lines = [
        "Use this slide to show what CLS, MEAN, MAX, and MIN actually do to a fused feature sequence. The important concept is that the model does not classify the raw token list directly. It first creates a token sequence after cross-attention fusion, and then it compresses that sequence into one vector that the classifier can read.",
        "Start by explaining the meaning of Dim. Dim stands for feature dimension. Each token vector is not a single number; it is a hidden representation with multiple learned coordinates. In this example, each token is represented with four dimensions so that the pooling idea can be seen clearly. In the real model, the vector is much larger, but the logic is the same.",
        "The example table is intentionally small so the idea is visible. In the real model, a token vector could have 256, 768, or another hidden size depending on the encoder branch, but the pooling rule still acts dimension by dimension in exactly the same way.",
        "Then read the table row by row. Token 1, Token 2, and Token 3 are sample fused vectors. CLS uses the first valid vector as the summary, so it returns the Token 1 vector. MEAN averages all valid vectors, which gives a stable but smoother summary. MAX keeps the strongest activation in each hidden dimension. MIN keeps the weakest activation in each hidden dimension.",
        "Explain that none of these are different models. They are different summarization rules applied after the same fusion step. That is why the paper-style code keeps pooling separate from cross-attention. The fusion step learns interaction between text and audio, while pooling decides how to compress the fused sequence before classification.",
        "A good viva answer is to say that CLS is positional, MEAN is average-based, MAX is evidence-maximizing, and MIN is evidence-minimizing. In the MemoCMT MELD setting, the project keeps MIN because that is the paper-aligned choice and the Phase 1 MELD results are strongest and most defensible under that setting.",
        "If asked why CLS is not the default here, explain that CLS only reads the first fused token. That is useful when the architecture is designed around a special summary token, but the MemoCMT paper-style branch is being compared as a pooling choice after fusion, so MIN/MEAN/MAX/CLS are treated as alternative summarizers rather than the main architecture itself.",
        "The code path for this logic is in `src/models/model.py:190-208`, where `_masked_pool` implements CLS with `x[:, 0, :]`, MEAN with a masked average, MAX with masked maximum, and MIN with masked minimum. The broader CMT fusion that produces the sequence is in `src/models/model.py:160-246`, and the CLI flag that selects pooling is parsed in `src/train/train.py:32-35` and exposed in `src/train/train.py:215-218`.",
        "The training script `scripts/run_paper_aligned_meld_cv.sh` uses `--fusion-pooling min` because the goal is to stay aligned with the MemoCMT paper’s MELD result. So when you explain the slide, say that CLS, MEAN, MAX, and MIN are all valid options, but MIN is the one chosen for the paper-aligned baseline because the project is trying to match the published comparison point rather than invent a new pooling rule.",
        "One final clarification is useful in viva: pooling is not the same as fusion. Fusion decides how the modalities interact, while pooling decides how the resulting sequence is summarized. That distinction is important because the project later adds video and gated fusion, but the pooling step still remains the same kind of summarization layer.",
    ]
    insert_docx_section_before(doc, anchor, new_heading, lines)
    shift_slide_headings(doc, start_num=6, delta=1)
    doc.save(str(DOCX_PATH))


def main() -> None:
    build_pptx()
    build_docx()
    print("Added CLS/MEAN/MAX/MIN example slide and reading-script expansion.")


if __name__ == "__main__":
    main()
