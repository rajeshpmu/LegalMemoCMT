from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import subprocess

from docx import Document
from docx.shared import Pt
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Prep_Guide.docx"
ASSET_DIR = ROOT / "implementation_docments" / "phase1_esa_assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)
DIAGRAM_MMD = ASSET_DIR / "design_details_flow.mmd"
DIAGRAM_SVG = ASSET_DIR / "design_details_flow.svg"
DIAGRAM_PNG = ASSET_DIR / "design_details_flow.png"

BODY = RGBColor(42, 42, 42)
MUTED = RGBColor(92, 92, 92)

MERMAID = """%%{init: {'themeVariables': {'fontSize': '32px', 'fontFamily': 'Aptos'}}}%%
flowchart LR
  subgraph I[Inputs]
    direction TB
    A[MP4 video]
    B[WAV audio]
    C[Transcript text]
  end
  subgraph P[Preprocessing]
    direction TB
    D[Frame sampling]
    E[Face crop]
    F[FFmpeg audio extraction]
    G[Text normalization]
  end
  subgraph M[Manifests and cache]
    direction TB
    H[CSV manifest rows]
    J[Cached .npy features]
  end
  subgraph B1[Phase 1 backbone]
    direction TB
    K[BERT text]
    L[HuBERT audio]
    N[ViT visual]
    O[Gated fusion]
    Q[Weighted CE / aux loss]
  end
  subgraph O1[Outputs]
    direction TB
    R[Metrics]
    S[Confusion matrix]
    T[Top confusions]
  end
  subgraph F2[Future Phase 2]
    direction TB
    U[Courtroom metadata]
    V[Multilingual transcripts]
    W[Explainability summaries]
  end
  A --> D
  B --> F
  C --> G
  D --> E --> H
  F --> J
  G --> H
  H --> K
  H --> L
  H --> N
  K --> O
  L --> O
  N --> O
  O --> Q --> R --> S --> T --> U --> V --> W
"""


def style_text(run, size=12, bold=False, color=BODY):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def insert_paragraph_before(paragraph, text="", style=None):
    new_p = OxmlElement("w:p")
    paragraph._p.addprevious(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        new_para.add_run(text)
    return new_para


def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def clone_slide(prs, slide_idx):
    source = prs.slides[slide_idx]
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in source.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return new_slide


def insert_slide_after(prs, after_idx, new_slide):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(after_idx + 1, sldId)


def render_diagram():
    DIAGRAM_MMD.write_text(MERMAID, encoding="utf-8")
    for out in [DIAGRAM_SVG, DIAGRAM_PNG]:
        subprocess.run(
            ["npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(DIAGRAM_MMD), "-o", str(out), "-b", "white", "-s", "2.5"],
            cwd=str(ROOT),
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )


def update_pptx():
    render_diagram()
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides[12]
    summary = clone_slide(prs, 12)
    insert_slide_after(prs, 12, summary)
    summary = prs.slides[13]

    # Slide 13 becomes diagram-only. Remove the text-heavy boxes and enlarge the diagram.
    remove_shape(slide.shapes[5])
    remove_shape(slide.shapes[5])  # after removal, the old shape 6 becomes index 5
    pic = slide.shapes[5]
    pic.left = Inches(0.55)
    pic.top = Inches(1.55)
    pic.width = Inches(12.15)
    pic.height = Inches(5.45)

    summary.shapes[2].text = "Design Properties & Implications"
    summary.shapes[3].text = "14"
    # Remove the copied picture so this slide focuses on the key labels.
    for shp in list(summary.shapes):
        if shp.shape_type == 13:
            remove_shape(shp)
            break
    summary.shapes[5].text = "\n".join(
        [
            "• Novelty and innovativeness: courtroom-testimony adaptation is the main project direction.",
            "• Interoperability: MP4, WAV, CSV, .npy, and checkpoints must stay aligned across scripts.",
            "• Performance and reliability: pretrained encoders, fold checks, and confusion analysis will verify gains.",
            "• Maintainability and portability: the pipeline will stay modular for local GPU and RunPod CUDA execution.",
            "• Application compatibility: the same backbone will support legal transcripts, dashboards, and later courtroom studies.",
            "• Phase 2 direction and safety: courtroom metadata, multilingual transcripts, and explainability stay future work, and the model remains limited to observable emotional cues.",
        ]
    )
    for p in summary.shapes[5].text_frame.paragraphs:
        for run in p.runs:
            style_text(run, 11.2, False, BODY)
    summary.shapes[6].text = (
        "Use this slide to defend why the system is built as a modular pipeline rather than a single monolithic model. "
        "The point is that every property on the slide affects how believable the Phase 1 evidence is and how safely the later courtroom adaptation can be added."
    )
    for p in summary.shapes[6].text_frame.paragraphs:
        for run in p.runs:
            style_text(run, 10.8, False, MUTED)

    # Renumber slides from the inserted slide onward.
    for idx in range(13, len(prs.slides)):
        s = prs.slides[idx]
        for shp in s.shapes:
            if hasattr(shp, "text") and shp.text.strip().isdigit():
                shp.text = str(idx + 1)
                break
    prs.save(str(PPTX_PATH))


def update_docx():
    doc = Document(str(DOCX_PATH))
    table = doc.tables[0]
    for row in table.rows:
        if row.cells[0].text.strip() == "Slide 13":
            row.cells[1].text = (
                "Use the diagram to explain the pipeline: raw inputs, preprocessing, manifests/cache, backbone, fusion/loss, outputs, and future courtroom extension."
            )
            row.cells[2].text = "Keep it technical, concise, and traceable."
            break

    # Insert a new row for the added Slide 14, then shift later slide labels by +1.
    rows = table.rows
    insert_after = rows[13]
    new_tr = deepcopy(insert_after._tr)
    insert_after._tr.addnext(new_tr)
    from docx.table import _Row
    new_row = _Row(new_tr, table)
    new_row.cells[0].text = "Slide 14"
    new_row.cells[1].text = (
        "Summarize the design properties: novelty, interoperability, performance, reliability, maintainability, portability, "
        "application compatibility, and the future courtroom-adaptation path."
    )
    new_row.cells[2].text = "Keep it technical, concise, and traceable."
    # Update later row labels so the guide matches the new slide order.
    for idx in range(15, len(table.rows)):
        current = table.rows[idx]
        try:
            n = int(current.cells[0].text.strip().split()[-1])
            current.cells[0].text = f"Slide {n + 1}"
        except Exception:
            pass

    paras = list(doc.paragraphs)
    start = next((i for i, p in enumerate(paras) if p.text.strip() == "2.5 Slide 13 Extended Talking Points"), None)
    if start is not None:
        end = next((i for i in range(start + 1, len(paras)) if paras[i].text.strip().startswith("3. Demo SOP You Should Practise")), len(paras))
        for i in range(start, end):
            el = paras[i]._element
            el.getparent().remove(el)

    anchor = next((p for p in doc.paragraphs if p.text.strip().startswith("3. Demo SOP You Should Practise")), None)
    if anchor is None:
        anchor = doc.paragraphs[-1]
    heading = insert_paragraph_before(anchor, style="Heading 2")
    heading.add_run("2.5 Slide 13 Extended Talking Points")
    sections = [
        "Diagram walkthrough: the slide is now meant to be explained left to right. Raw inputs enter preprocessing, preprocessing creates manifests and cached features, the Phase 1 backbone consumes those cached artifacts, fusion and loss compute the prediction, and the output stage produces metrics and confusion analysis.",
    ]
    for t in sections:
        p = insert_paragraph_before(anchor)
        r = p.add_run(t)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        p.alignment = 3

    insert_after_26 = next((p for p in doc.paragraphs if p.text.strip().startswith("3. Demo SOP You Should Practise")), None)
    if insert_after_26 is None:
        insert_after_26 = doc.paragraphs[-1]
    h2 = insert_paragraph_before(insert_after_26, style="Heading 2")
    h2.add_run("2.6 Slide 14 Design Properties & Implications")
    sections_26 = [
        "Novelty and innovativeness: the project will be framed as a courtroom-testimony-oriented extension of multimodal emotion analysis rather than a generic benchmark model. That makes the contribution specific to the legal setting and the emotional cue timeline.",
        "Interoperability: the same utterance must move cleanly between MP4, WAV, CSV manifests, cached .npy features, and checkpoint files. This is what allows the training script, evaluation script, and demo script to all read the same sample consistently.",
        "Performance and reliability: Phase 1 will remain anchored to pretrained BERT, HuBERT, and ViT components, while fold-based evaluation, confusion matrices, and ablations will check whether the changes really help instead of only moving one aggregate metric.",
        "Maintainability and portability: the pipeline will stay modular so it can run on local hardware or RunPod CUDA and later be extended without rewriting the whole project.",
        "Application compatibility: the same backbone will stay compatible with public legal transcripts, research dashboards, and future courtroom-testimony studies.",
        "Phase 2 direction and safety: future courtroom segmentation, speaker metadata, multilingual transcripts, and explainability summaries will remain future work, and the model will stay limited to observable emotional cues only.",
    ]
    for t in sections_26:
        p = insert_paragraph_before(insert_after_26)
        r = p.add_run(t)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        p.alignment = 3
    doc.save(str(DOCX_PATH))


def main():
    update_pptx()
    update_docx()
    print(f"Updated {PPTX_PATH}")
    print(f"Updated {DOCX_PATH}")


if __name__ == "__main__":
    main()
