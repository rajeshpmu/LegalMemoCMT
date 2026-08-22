from pathlib import Path
import shutil

from docx import Document
from docx.enum.section import WD_SECTION
from docx.shared import Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.pptx"
PPTX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Guidance_Call_Presentation.before_whisper_pyannote_architecture.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Speaking_Doc.docx"
DOCX_BACKUP = ROOT / "implementation_docments/LegalMemoCMT_Phase2_First_Guidance_Call_Speaking_Doc.before_whisper_pyannote_architecture.docx"
TITLE = "Whisper + Pyannote: Building a Witness-Aware Utterance"


def add_box(slide, x, y, w, h, text, fill, font_size=16, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(70, 80, 95)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 1
    for run in p.runs:
        run.font.size = PptPt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(25, 35, 50)
    return shape


def arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(70, 80, 95)
    line.line.width = PptPt(1.5)
    line.line.end_arrowhead = True
    return line


if not PPTX_BACKUP.exists():
    shutil.copy2(PPTX, PPTX_BACKUP)
prs = Presentation(PPTX)
if not any(slide.shapes.title and slide.shapes.title.text == TITLE for slide in prs.slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.5))
    title.text_frame.text = TITLE
    title.text_frame.paragraphs[0].runs[0].font.size = PptPt(26)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

    add_box(slide, 5.0, 0.85, 3.3, 0.65, "Courtroom video", (220, 232, 245), 20, True)
    add_box(slide, 1.25, 1.95, 2.7, 0.7, "Whisper\nWHAT was said", (231, 243, 235), 17, True)
    add_box(slide, 9.05, 1.95, 2.7, 0.7, "Pyannote\nWHO spoke WHEN", (248, 237, 215), 17, True)
    add_box(slide, 4.65, 3.15, 4.0, 0.7, "Aligned transcript + speaker intervals", (220, 232, 245), 17, True)
    add_box(slide, 4.35, 4.35, 4.6, 0.7, "Witness utterance extraction", (232, 224, 245), 18, True)
    add_box(slide, 0.75, 5.55, 2.2, 0.75, "TEXT\ntranscript", (231, 243, 235), 17, True)
    add_box(slide, 5.05, 5.55, 2.2, 0.75, "AUDIO\nutterance WAV", (248, 237, 215), 17, True)
    add_box(slide, 9.35, 5.55, 2.2, 0.75, "VIDEO\nutterance MP4", (220, 232, 245), 17, True)
    add_box(slide, 4.45, 6.7, 4.8, 0.55, "LegalMemoCMT MELD-style row", (208, 224, 220), 18, True)

    arrow(slide, 6.65, 1.5, 2.6, 1.95)
    arrow(slide, 6.65, 1.5, 10.4, 1.95)
    arrow(slide, 2.6, 2.65, 5.8, 3.15)
    arrow(slide, 10.4, 2.65, 7.5, 3.15)
    arrow(slide, 6.65, 3.85, 6.65, 4.35)
    arrow(slide, 5.0, 5.05, 1.85, 5.55)
    arrow(slide, 6.65, 5.05, 6.15, 5.55)
    arrow(slide, 8.0, 5.05, 10.45, 5.55)
    arrow(slide, 1.85, 6.3, 6.0, 6.7)
    arrow(slide, 6.15, 6.3, 6.65, 6.7)
    arrow(slide, 10.45, 6.3, 7.3, 6.7)

    note = slide.shapes.add_textbox(Inches(0.65), Inches(7.35), Inches(12.0), Inches(0.35))
    note.text_frame.text = 'The alignment layer prevents a visible witness or emotional lawyer argument from being mislabeled as witness speech.'
    note.text_frame.paragraphs[0].runs[0].font.size = PptPt(13)
    note.text_frame.paragraphs[0].runs[0].font.italic = True

    # Move the new slide after the existing diarization-output slide.
    ids = prs.slides._sldIdLst
    new_id = ids[-1]
    ids.remove(new_id)
    insert_at = next(
        (i + 1 for i, sid in enumerate(ids) if prs.part.related_part(sid.rId).slide.shapes.title and prs.part.related_part(sid.rId).slide.shapes.title.text == "From Diarization Output to Phase 2 Rows"),
        len(ids),
    )
    ids.insert(insert_at, new_id)
    prs.save(PPTX)


if not DOCX_BACKUP.exists():
    shutil.copy2(DOCX, DOCX_BACKUP)
document = Document(DOCX)
document.add_section(WD_SECTION.NEW_PAGE)
document.add_heading("Additional Speaking Module: Whisper and Pyannote Together", level=1)
document.add_heading("1. The Technical Problem", level=2)
document.add_paragraph(
    "A courtroom video is a continuous recording, but the LegalMemoCMT training unit is one utterance. "
    "The system therefore needs two different kinds of evidence. It needs text evidence to know what was "
    "said, and it needs speaker evidence to estimate who spoke during each time interval. A transcript alone "
    "can contain questions, answers, captions, interruptions, and speaker changes without reliable role labels. "
    "A video frame alone can show the witness while a lawyer is speaking. Combining speech recognition and "
    "diarization addresses these two separate weaknesses."
)
document.add_heading("2. Whisper: What Was Said", level=2)
document.add_paragraph(
    "Whisper, used in the existing LegalMELD alignment path through the configured transcription backend, "
    "processes the hearing audio and produces recognized words with timestamps. These timestamps are used to "
    "compare the recorded audio with the transcript text and to create clip boundaries. Whisper is therefore "
    "primarily a content-and-time source: it helps answer what words occurred and approximately when they occurred."
)
document.add_paragraph(
    "Whisper does not reliably decide whether the speaker is a witness, prosecutor, defence lawyer, or judge. "
    "Its output can also differ from the official transcript because of accents, courtroom noise, crosstalk, "
    "or recognition errors. For that reason, `utterance_text`, `asr_text`, and the normalized transcript remain "
    "separate provenance fields rather than being silently replaced by one another."
)
document.add_heading("3. Pyannote: Who Spoke When", level=2)
document.add_paragraph(
    "The pyannote speaker-diarization pipeline processes the source audio and returns speech intervals with "
    "anonymous speaker cluster identifiers. The segmentation component estimates speech activity. The speaker "
    "representation and clustering stages group intervals that acoustically resemble one another. For example, "
    "several separated intervals may receive `SPEAKER_00`, while another voice receives `SPEAKER_01`."
)
document.add_paragraph(
    "The cluster is not a legal identity. It means only that the audio model considers the voice regions similar. "
    "I use short samples from each cluster for manual verification and then map the cluster to a role such as "
    "Witness, Prosecutor, Defence, Judge, or Other. This is why the manifest records both `speaker_cluster_id` "
    "and `speaker_role_source`: one records the machine grouping and the other records how the legal-role decision "
    "was obtained."
)
document.add_heading("4. The Alignment Join", level=2)
document.add_paragraph(
    "The two model outputs are joined by time. Whisper/transcript processing provides an utterance interval, "
    "for example 00:04:08.310 to 00:04:19.430. Pyannote provides overlapping speaker intervals. The pipeline "
    "finds the cluster with the strongest temporal overlap and writes that cluster to the utterance row. This is "
    "a temporal evidence join, not a claim that the model understood the courtroom dialogue."
)
document.add_paragraph(
    "The source-level diarization CSV remains important for auditing. It contains the original audio path, "
    "cluster ID, segment start, segment end, and model name. The enriched utterance manifest retains the original "
    "transcript, audio, video, and source paths, so a reviewer can move from a final row back to the source interval."
)
document.add_heading("5. Why This Matters for Clancy and Tupac", level=2)
document.add_paragraph(
    "A full trial recording can contain a defence question, a witness answer, a judge instruction, and an "
    "emotional closing argument in the same file. If I select rows only because the camera is pointing at the "
    "witness, I can incorrectly label a lawyer's speech as witness speech. The combined pipeline reduces this "
    "risk by requiring text/timestamp evidence and an audio speaker-cluster signal before role review."
)
document.add_heading("6. From Speaker-Aware Rows to MELD-Style Data", level=2)
for item in [
    "TEXT: the grounded transcript utterance and its normalized form.",
    "AUDIO: the WAV clip extracted from the same start and end boundaries.",
    "VIDEO: the MP4 clip extracted from the same boundaries, retaining visual context.",
    "ROLE METADATA: witness status, speaker role, cluster ID, and confidence fields.",
    "QUALITY METADATA: alignment method, alignment confidence, audio validation, video status, and manual-review status.",
]:
    document.add_paragraph(item, style="List Bullet")
document.add_paragraph(
    "Only after these fields are available can I create a defensible witness-only subset. The resulting row is "
    "still MELD-compatible because it has one text item, one audio path, and one video path, but it is more useful "
    "for courtroom research because it also records how the speaker-role evidence was obtained."
)
document.add_heading("7. What the Models Cannot Establish", level=2)
document.add_paragraph(
    "Neither Whisper nor pyannote can establish emotion, deception, credibility, truthfulness, or legal meaning. "
    "The Phase 1 MELD model may later provide provisional basic-emotion pseudo-labels, but those labels must retain "
    "their model source and confidence. Courtroom-specific affect requires a separate annotation layer. A witness "
    "speaking with a tense voice is not automatically deceptive, and a lawyer's emotional argument is not a witness "
    "emotion label."
)
document.add_heading("8. Student-Level Summary", level=2)
document.add_paragraph(
    "Whisper tells me what words occur and when; pyannote tells me which anonymous voice cluster overlaps that time; "
    "the transcript and diarization intervals are joined; and manual role review converts selected clusters into "
    "courtroom-role metadata. The final LegalMemoCMT row then connects text, audio, and video without confusing the "
    "person visible on camera with the person actually speaking."
)
document.save(DOCX)
print(f"Updated {PPTX}")
print(f"Updated {DOCX}")
