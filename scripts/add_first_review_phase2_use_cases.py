from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.docx"


def add_bg(slide, prs: Presentation):
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), prs.slide_width, PptInches(0.22)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str):
    tx = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.35), PptInches(12.2), PptInches(0.7))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = PptPt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)

    tx2 = slide.shapes.add_textbox(PptInches(0.58), PptInches(0.87), PptInches(12.1), PptInches(0.45))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Aptos"
    r2.font.size = PptPt(12.5)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x, y, w, h, font_size=14):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        p.line_spacing = 1.06
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = PptPt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_table(slide, left, top, width, height, rows):
    table = slide.shapes.add_table(
        len(rows), len(rows[0]), PptInches(left), PptInches(top), PptInches(width), PptInches(height)
    ).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = PptPt(11 if r == 0 else 10.25)
                    run.font.bold = r == 0
                    run.font.color.rgb = RGBColor(35, 35, 35)
    return table


def insert_slide_before(prs: Presentation, slide, before_index_0_based: int):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(before_index_0_based, sldId)


def _title_score(shape) -> tuple[float, float]:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return (0.0, 0.0)
    text = shape.text_frame.text.strip()
    if not text:
        return (0.0, 0.0)
    if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
        return (0.0, 0.0)
    top = float(shape.top)
    max_font = 0.0
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            if run.font.size:
                max_font = max(max_font, float(run.font.size))
    # Prefer the larger top-most textbox. Titles are at the top and use larger fonts.
    return (max_font, -top)


def renumber_titles(prs: Presentation):
    seq = 1
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue

        candidates = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            if float(shape.top) > float(PptInches(1.25)):
                continue
            candidates.append(shape)

        if not candidates:
            continue

        title_shape = sorted(candidates, key=_title_score, reverse=True)[0]
        tf = title_shape.text_frame
        old_text = tf.text.strip()
        base = re.sub(r"^\d+\.\s*", "", old_text)
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{seq}. {base}"
        r.font.name = "Aptos Display"
        r.font.size = PptPt(24)
        r.font.bold = True
        r.font.color.rgb = RGBColor(20, 48, 87)
        seq += 1


def update_pptx():
    prs = Presentation(str(PPTX_PATH))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(
        slide,
        "13. Phase 2 Use Cases and Interpretation",
        "Courtroom testimony + Indian languages + multimodal inputs + emotion timeline + explainability.",
    )
    add_body(
        slide,
        [
            "Phase 2 segments testimony into aligned windows, fuses face, audio, and transcript cues, and exposes the result as a timeline rather than a single verdict-like label.",
            "The goal is to help a reviewer inspect where emotional intensity changes and which topic or event may have contributed to the change.",
            "The system should remain a support tool for analysis, not a substitute for legal judgment.",
        ],
        0.7,
        1.42,
        5.95,
        2.15,
        font_size=14,
    )

    rows = [
        ["Use case", "Example pattern", "What the system can show", "What it must not claim"],
        [
            "Witness testimony review",
            "0-5 calm, 5-12 mild stress, 12-18 high anxiety, 18-22 stable, 22-30 anger indicators",
            "A topic-linked emotional timeline that lets lawyers ask what was being discussed during the spike",
            "It does not say the witness is lying or unreliable",
        ],
        [
            "Victim support",
            "Topic A low stress, Topic B fear response, Topic C high anxiety",
            "A structured emotional burden profile that may help support teams and psychologists",
            "It does not diagnose trauma or prove victimhood",
        ],
        [
            "Courtroom research",
            "English vs Tamil vs Hindi testimony",
            "Comparisons of emotion expression, language effects, and stress patterns across languages",
            "It does not rank languages or infer legal truth",
        ],
    ]
    add_table(slide, 6.75, 1.5, 6.1, 3.55, rows)

    add_body(
        slide,
        [
            "How it is achieved without overstating: segment the testimony, align the modalities, encode each stream, fuse the evidence, aggregate over time, and present the output with timestamps, confidence, and modality contribution summaries.",
        ],
        0.7,
        4.85,
        12.0,
        0.95,
        font_size=13,
    )
    add_body(
        slide,
        [
            "This framing is publishable as a multimodal analysis study if the data, annotation, and evaluation are done carefully; it should not be presented as a legal decision system.",
        ],
        0.7,
        5.75,
        12.0,
        0.55,
        font_size=12.5,
    )

    insert_slide_before(prs, slide, 13)
    renumber_titles(prs)
    prs.save(str(PPTX_PATH))


def update_docx():
    doc = Document(str(DOCX_PATH))

    doc.add_heading("11.10 Phase 2 Use Cases and Research Framing", level=2)
    doc.add_paragraph(
        "Phase 2 is intended to turn the system into a practical review aid for long courtroom testimony, not into a legal decision engine. The design goal is to segment testimony into time windows, align multimodal cues across face, audio, and text, and show how emotional intensity evolves over time. That makes the output easier to inspect, easier to explain, and safer to interpret."
    )
    doc.add_paragraph(
        "The key boundary should remain explicit: the system can highlight patterns such as calm, stress, fear, anxiety, or anger indicators, but it cannot determine whether someone is lying, guilty, innocent, or legally liable."
    )

    doc.add_heading("11.10.1 Witness Testimony Review", level=3)
    doc.add_paragraph(
        "A witness may speak for a long period, and the useful signal is often a sequence of changes rather than one static emotion label. In a 30-minute testimony, the system could summarize the timeline as calm at the start, mild stress during one topic, high anxiety during a sensitive section, a stable phase after that, and anger indicators near the end. The important point is that this is an emotional timeline, not a credibility verdict."
    )
    doc.add_paragraph("Example timeline:")
    for text in [
        "0-5 min: Calm",
        "5-12 min: Mild stress",
        "12-18 min: High anxiety",
        "18-22 min: Stable",
        "22-30 min: Anger indicators",
    ]:
        doc.add_paragraph(text, style="List Bullet")
    doc.add_paragraph(
        "The lawyer-facing question is: what topic was being discussed during the anxiety spike? That is a useful review question because it points to a segment that may deserve closer attention. The system should still avoid statements such as 'the witness is lying'; the technically defensible statement is simply that emotional intensity increased."
    )

    doc.add_heading("11.10.2 Victim Support", level=3)
    doc.add_paragraph(
        "A second use case is victim support. If a domestic violence survivor or another vulnerable witness records a statement, the model can surface sections where fear or anxiety appears stronger. That can help psychologists, victim support officers, and legal aid teams understand where the emotional burden is concentrated and where a slower, more careful review may be appropriate."
    )
    for text in [
        "Topic A: low stress",
        "Topic B: strong fear response",
        "Topic C: high anxiety",
    ]:
        doc.add_paragraph(text, style="List Bullet")
    doc.add_paragraph(
        "This is useful because the output is topic-linked and time-linked. It helps human reviewers see patterns in the statement, but it should not be treated as a diagnosis of trauma, proof of victimhood, or a substitute for professional judgment."
    )

    doc.add_heading("11.10.3 Courtroom Research", level=3)
    doc.add_paragraph(
        "A third use case is research across languages. Researchers may compare English testimony, Tamil testimony, and Hindi testimony to examine whether emotion expression, stress patterns, and multilingual mixing affect the multimodal outputs. If the data collection and annotation are careful, this can be framed as publishable multimodal research on cross-lingual emotional cue analysis."
    )
    for text in [
        "English testimony versus Tamil testimony versus Hindi testimony",
        "Emotion expression differences",
        "Stress pattern differences",
        "Language-specific interaction with audio and transcript cues",
    ]:
        doc.add_paragraph(text, style="List Bullet")
    doc.add_paragraph(
        "The correct framing is comparative analysis, not language ranking and not legal truth inference."
    )

    doc.add_heading("11.10.4 How This Can Be Achieved Without Overstating", level=3)
    for text in [
        "Segment the testimony into aligned time windows so that the transcript, audio, and visual cues refer to the same moment.",
        "Encode each modality separately using an appropriate model, then fuse the streams so the system can learn cross-modal interactions.",
        "Aggregate the segment-level outputs into a timeline so reviewers can inspect how emotion changes over time.",
        "Surface timestamps, confidence, and modality contribution summaries so the explanation remains visible and auditable.",
        "Use cautious wording in the thesis: the system supports emotional pattern analysis and does not infer lying, guilt, innocence, or credibility.",
    ]:
        doc.add_paragraph(text, style="List Bullet")

    doc.add_paragraph(
        "In other words, the implementation can be technically strong without making claims that the data cannot support. That is the correct way to present the work for a legal-AI context."
    )

    doc.add_heading("11.10.5 Interpretation Boundary", level=3)
    doc.add_paragraph(
        "The boundary should be repeated in every Phase 2 description. The pipeline can identify emotional patterns and their timing; it cannot determine truthfulness or legal liability. That limitation is not a weakness. It is what makes the project defensible, ethical, and technically credible."
    )

    doc.save(str(DOCX_PATH))


def main():
    update_pptx()
    update_docx()
    print("Updated First Review PPTX and DOCX with Phase 2 use cases.")


if __name__ == "__main__":
    main()
