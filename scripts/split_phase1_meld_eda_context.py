from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import subprocess

import matplotlib.pyplot as plt
from docx import Document
from docx.shared import Pt
from docx.oxml import OxmlElement
from docx.table import _Row
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Prep_Guide.docx"
ASSET_DIR = ROOT / "implementation_docments" / "phase1_esa_assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)
CONTEXT_MMD = ASSET_DIR / "meld_context_analysis.mmd"
CONTEXT_SVG = ASSET_DIR / "meld_context_analysis.svg"
CONTEXT_PNG = ASSET_DIR / "meld_context_analysis.png"

BODY = RGBColor(42, 42, 42)
MUTED = RGBColor(92, 92, 92)
PALE = RGBColor(249, 246, 222)
PALE2 = RGBColor(235, 231, 248)
ACCENT = RGBColor(112, 98, 184)

MERMAID = """%%{init: {'themeVariables': {'fontSize': '30px', 'fontFamily': 'Aptos'}}}%%
flowchart TB
  A[Inter-speaker influence] --> B[Emotion shifts]
  B --> C[Contextual distance]
  C --> D[Dialogue-aware Phase 1]
  D --> E[Weighted metrics + confusion analysis]
"""


def style_text(run, size=12, bold=False, color=BODY):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def insert_paragraph_before(paragraph, text="", style=None):
    new_p = OxmlElement("w:p")
    paragraph._p.addprevious(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        new_para.add_run(text)
    return new_para


def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def move_last_slide_to(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[-1]
    sldIdLst.remove(sldId)
    sldIdLst.insert(index, sldId)


def render_context_diagram():
    CONTEXT_MMD.write_text(MERMAID, encoding="utf-8")
    for out in [CONTEXT_SVG, CONTEXT_PNG]:
        subprocess.run(
            ["npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(CONTEXT_MMD), "-o", str(out), "-b", "white", "-s", "3"],
            cwd=str(ROOT),
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )


def update_pptx():
    render_context_diagram()
    prs = Presentation(str(PPTX_PATH))

    # Update Slide 7 to focus only on EDA + Table 10.
    slide7 = prs.slides[6]
    slide7.shapes[2].text = "MELD Dataset EDA & Table 10"
    slide7.shapes[3].text = "07"
    slide7.shapes[5].text = "\n".join(
        [
            "• MELD is a multi-party, multimodal dialogue benchmark with strong conversational context.",
            "• Table 10 shows shorter turns, a much larger neutral class, and a more dialogue-heavy profile than IEMOCAP.",
            "• The dataset statistics justify weighted metrics, fold-safe grouping, and context-aware analysis in Phase 1.",
            "• The paper EDA supports the decision to treat MELD as a conversational problem rather than isolated emotion classification.",
        ]
    )
    for p in slide7.shapes[5].text_frame.paragraphs:
        for run in p.runs:
            style_text(run, 12.0, False, BODY)
    slide7.shapes[6].text = (
        "Use Table 10 to explain why MELD needs dialogue-aware evaluation: the turns are shorter, the neutral class dominates, and context across speakers matters."
    )
    for p in slide7.shapes[6].text_frame.paragraphs:
        for run in p.runs:
            style_text(run, 10.8, False, MUTED)

    # Add new Slide 8 after current Slide 7, then shift later slides.
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])
    move_last_slide_to(prs, 7)
    slide8 = prs.slides[7]

    # Background panels.
    bg = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), Inches(0.82), Inches(5.95), Inches(5.80))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PALE
    bg.line.color.rgb = RGBColor(214, 205, 159)

    bg2 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.15), Inches(0.90), Inches(5.95), Inches(1.75))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = PALE2
    bg2.line.color.rgb = ACCENT

    bg3 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.15), Inches(2.78), Inches(5.95), Inches(3.65))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = PALE2
    bg3.line.color.rgb = ACCENT

    title = slide8.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(9.8), Inches(0.35))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "MELD Contextual Analysis"
    style_text(r, 18, True, BODY)

    num = slide8.shapes.add_textbox(Inches(11.35), Inches(0.12), Inches(0.8), Inches(0.35))
    p = num.text_frame.paragraphs[0]
    p.alignment = 2
    r = p.add_run()
    r.text = "08"
    style_text(r, 18, True, BODY)

    left = slide8.shapes.add_textbox(Inches(0.55), Inches(1.08), Inches(5.45), Inches(4.95))
    tf = left.text_frame
    tf.word_wrap = True
    bullets = [
        "Inter-speaker influence: another speaker’s previous turn often changes the target utterance’s emotion.",
        "Emotion shifts: the same speaker can move across emotions during a conversation, so one label is not enough to explain the whole dialogue.",
        "Contextual distance: nearby turns matter most, but farther turns still influence the model’s decision path.",
        "Project impact: Phase 1 must use context-aware metrics and confusion analysis, not isolated utterance reading.",
    ]
    for i, t in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.level = 0
        p.bullet = True
        for run in p.runs:
            style_text(run, 12.0, False, BODY)

    note = slide8.shapes.add_textbox(Inches(6.35), Inches(1.08), Inches(5.35), Inches(1.35))
    note_tf = note.text_frame
    note_tf.word_wrap = True
    p = note_tf.paragraphs[0]
    p.text = (
        "The contextual story explains why MELD is harder than a flat emotion task: speaker interactions, turn-to-turn shifts, and distance from surrounding context all affect the label."
    )
    for run in p.runs:
        style_text(run, 10.7, False, MUTED)

    slide8.shapes.add_picture(str(CONTEXT_PNG), Inches(6.25), Inches(2.95), width=Inches(5.55), height=Inches(3.15))

    # Renumber slides from the inserted slide onward.
    for idx in range(7, len(prs.slides)):
        s = prs.slides[idx]
        for shp in s.shapes:
            if hasattr(shp, "text") and shp.text.strip().isdigit():
                shp.text = f"{idx + 1:02d}"
                break

    prs.save(str(PPTX_PATH))


def update_docx():
    doc = Document(str(DOCX_PATH))
    table = doc.tables[0]

    # Update Slide 7 row and insert Slide 8 row.
    table.rows[7].cells[0].text = "Slide 7"
    table.rows[7].cells[1].text = (
        "Explain MELD as a conversational benchmark and show why Table 10 matters: short turns, strong neutral skew, multi-party context, and dialogue-heavy statistics."
    )
    table.rows[7].cells[2].text = "Keep it technical, concise, and traceable."

    new_tr = deepcopy(table.rows[7]._tr)
    table.rows[7]._tr.addnext(new_tr)
    new_row = _Row(new_tr, table)
    new_row.cells[0].text = "Slide 8"
    new_row.cells[1].text = (
        "Explain the contextual analysis: inter-speaker influence, emotion shifts, contextual distance, and why these factors matter for dialogue-aware evaluation in Phase 1."
    )
    new_row.cells[2].text = "Keep it technical, concise, and traceable."

    # Shift later slide labels by +1.
    for idx in range(9, len(table.rows)):
        current = table.rows[idx]
        txt = current.cells[0].text.strip()
        if txt.startswith("Slide "):
            try:
                n = int(txt.split()[-1])
                current.cells[0].text = f"Slide {n + 1}"
            except Exception:
                pass

    # Remove old MELD section paragraphs and insert the split versions.
    paras = list(doc.paragraphs)
    old_start = next((i for i, p in enumerate(paras) if p.text.strip() == "2.2 Slide 7 MELD Dataset EDA & Contextual Factors"), None)
    if old_start is not None:
        old_end = next((i for i in range(old_start + 1, len(paras)) if paras[i].text.strip().startswith("2.3 Literature Review Slides")), len(paras))
        for i in range(old_start, old_end):
            el = paras[i]._element
            el.getparent().remove(el)

    anchor = next((p for p in doc.paragraphs if p.text.strip().startswith("2.3 Literature Review Slides")), None)
    if anchor is None:
        anchor = doc.paragraphs[-1]

    h = insert_paragraph_before(anchor, style="Heading 2")
    h.add_run("2.2 Slide 7 MELD Dataset EDA & Table 10")
    for t in [
        "EDA summary: MELD is a multi-party, multimodal conversational dataset, so the project cannot treat each utterance as isolated text. The dataset is shorter-turn, more neutral-heavy, and dialogue-oriented, which is why Phase 1 must discuss weighted metrics and context-aware evaluation.",
        "Table 10 value: the MELD paper compares MELD with IEMOCAP and shows that MELD has more emotion-labeled utterances, more unique words, and much shorter average utterance and conversation lengths. This is useful for the ESA because it explains why dialogue history matters in the benchmark.",
        "Project relevance: the statistics justify fold-safe grouping, weighted F1, and confusion-matrix analysis in Phase 1, because accuracy alone can hide the neutral bias.",
    ]:
        p = insert_paragraph_before(anchor)
        r = p.add_run(t)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        p.alignment = 3

    h = insert_paragraph_before(anchor, style="Heading 2")
    h.add_run("2.3 Slide 8 MELD Contextual Analysis")
    for t in [
        "Inter-speaker influence: the MELD paper shows that many correct predictions depend on turns from different speakers, not only the target speaker. That means conversational context is a first-class signal in the project, not a side detail.",
        "Emotion shifts: the same speaker can shift emotion during the dialogue, and the paper notes that recognizing the shift and the exact emotion together is harder than detecting either one alone. This is why the project must explain errors as context-dependent events.",
        "Contextual distance: the model may attend to nearby turns most strongly, but farther historical and future turns also influence the decision. This supports the project’s idea of using the dialogue timeline when interpreting predictions.",
        "Phase 1 impact: these factors justify the use of weighted metrics, confusion matrices, and dialogue-aware analysis, because they reveal whether the model is learning conversation structure or just exploiting class imbalance.",
    ]:
        p = insert_paragraph_before(anchor)
        r = p.add_run(t)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        p.alignment = 3

    # Update later heading labels.
    replacements = {
        "2.3 Literature Review Slides": "2.4 Literature Review Slides",
        "2.4 Slide 11 Extended Talking Points": "2.5 Slide 11 Extended Talking Points",
        "2.5 Slide 12 Extended Talking Points": "2.6 Slide 12 Extended Talking Points",
        "2.6 Slide 13 Extended Talking Points": "2.7 Slide 13 Extended Talking Points",
        "2.7 Slide 14 Design Properties & Implications": "2.8 Slide 14 Design Properties & Implications",
        "2.8 Slide 16 System Architecture": "2.9 Slide 16 System Architecture",
        "2.9 Slide 17 Model Architecture": "2.10 Slide 17 Model Architecture",
    }
    for p in doc.paragraphs:
        t = p.text.strip()
        if t in replacements:
            p.text = replacements[t]

    doc.save(str(DOCX_PATH))


def main():
    update_pptx()
    update_docx()
    print(f"Updated {PPTX_PATH}")
    print(f"Updated {DOCX_PATH}")


if __name__ == "__main__":
    main()
