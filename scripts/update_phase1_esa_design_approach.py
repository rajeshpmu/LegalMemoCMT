from __future__ import annotations

from pathlib import Path
import shutil
import subprocess

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Prep_Guide.docx"
ASSET_DIR = ROOT / "implementation_docments" / "phase1_esa_assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)
DIAGRAM_MMD = ASSET_DIR / "design_approach_flow.mmd"
DIAGRAM_SVG = ASSET_DIR / "design_approach_flow.svg"
DIAGRAM_PNG = ASSET_DIR / "design_approach_flow.png"

NAVY = RGBColor(18, 47, 85)
TEAL = RGBColor(38, 115, 119)
LIGHT = RGBColor(246, 248, 250)
BODY = RGBColor(42, 42, 42)
MUTED = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)

MERMAID = """%%{init: {'themeVariables': {'fontSize': '30px', 'fontFamily': 'Aptos'}}}%%
flowchart TB
  A[Phase 1 goal] --> B[Baseline]
  B --> C[Add ViT cues]
  C --> D[Test gated / aux]
  D --> E[Use metrics]
  E --> F[Decide next step]

  B --- B1[Text + audio]
  C --- C1[Raw MP4 -> frames -> embeddings]
  D --- D1[One factor at a time]
"""


def render_mermaid():
    DIAGRAM_MMD.write_text(MERMAID, encoding="utf-8")
    for out in [DIAGRAM_SVG, DIAGRAM_PNG]:
        subprocess.run(
            ["npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(DIAGRAM_MMD), "-o", str(out), "-b", "white"],
            cwd=str(ROOT),
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )


def style_text(run, size=12, bold=False, color=BODY):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    if isinstance(color, RGBColor):
        run.font.color.rgb = color
    else:
        try:
            rgb = tuple(color)
            run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        except Exception:
            run.font.color.rgb = RGBColor(42, 42, 42)


def build_pptx():
    prs = Presentation(str(PPTX_PATH))
    slide = prs.slides[10]

    # Update title and slide number
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            txt = shape.text.strip()
            if txt == "Design Approach":
                shape.text = "Design Approach"
            elif txt == "11":
                shape.text = "11"

    # The slide layout has a left bullet box and a right explanation box.
    # Update those boxes directly so manual edits elsewhere remain intact.
    if len(slide.shapes) > 7:
        bullet_box = slide.shapes[5]
        bullet_box.text = "\n".join([
            "• What is the design approach followed? A controlled experimental ladder: baseline -> ViT -> gated/aux variants.",
            "• Benefits and drawbacks? It isolates each change for clean attribution, but it is slower than a single redesign.",
            "• Alternate design approaches? Early fusion, late fusion, or video-only are valid, but they blur which modality caused the effect.",
        ])
        for p in bullet_box.text_frame.paragraphs:
            for run in p.runs:
                style_text(run, 13.2, False, BODY)

        note_box = slide.shapes[7]
        note_box.text = (
            "The design is a controlled experimental ladder: the main benefit is interpretability because each change can be traced to one cause. "
            "The main drawback is speed, because every variant must be trained and compared separately before a conclusion is defensible."
        )
        for p in note_box.text_frame.paragraphs:
            for run in p.runs:
                style_text(run, 11.2, False, MUTED)

    # Add diagram on the right side, below the explanation box.
    # Remove a prior diagram if the script is rerun.
    for shp in list(slide.shapes):
        if getattr(shp, "name", "") == "DesignApproachDiagram":
            sp = shp._element
            sp.getparent().remove(sp)

    slide.shapes.add_picture(str(DIAGRAM_PNG), Inches(6.95), Inches(1.95), width=Inches(5.2), height=Inches(4.8)).name = "DesignApproachDiagram"

    # Re-save
    prs.save(str(PPTX_PATH))


def build_docx():
    doc = Document(str(DOCX_PATH))
    # update table row for Slide 11
    table = doc.tables[0]
    for row in table.rows:
        if row.cells[0].text.strip() == "Slide 11":
            row.cells[1].text = (
                "Explain the design approach as a controlled experimental ladder: baseline -> ViT -> gated/aux variants. "
                "Say why it is chosen, what the benefits and drawbacks are, and why alternate approaches such as early fusion, late fusion, or video-only would be less clean for attribution."
            )
            row.cells[2].text = "Keep it technical, concise, and traceable."
            break

    # Remove any prior Slide 11 extended sections so the guide stays clean.
    body = doc._element.body
    paras = list(doc.paragraphs)
    start_idxs = [i for i, p in enumerate(paras) if p.text.strip() == "2.3 Slide 11 Extended Talking Points"]
    if start_idxs:
        start = start_idxs[0]
        end = next((i for i in range(start + 1, len(paras)) if paras[i].text.strip().startswith("3. Demo SOP You Should Practise")), len(paras))
        for i in range(start, end):
            el = paras[i]._element
            el.getparent().remove(el)

    doc.add_heading("2.3 Slide 11 Extended Talking Points", level=2)
    sections = [
        "What is the design approach followed? It is a controlled experimental ladder: first validate the paper-aligned MELD baseline, then attach the ViT visual branch, and then test gated fusion and auxiliary loss as separate variants. This keeps every change isolated, so the observed shift in metrics or confusion behavior can be traced to one factor instead of several factors at once.",
        "Benefits and drawbacks: the main benefit is interpretability because each change can be traced to one cause. The main drawback is speed, because every variant must be trained and compared separately before a conclusion is defensible.",
        "What are the alternate approaches? Early fusion would mix modalities earlier in the network, late fusion would keep the branches separate for longer, and a video-only model would isolate the visual stream. Those alternatives are valid, but they weaken attribution for this project because they make it harder to tell whether the improvement came from text, audio, video, or the fusion rule itself.",
        "The evaluation loop is part of the design. The fold metrics, confusion matrix, and top-confusion table are what decide whether the design choice actually helped. If the confusion pattern does not improve, then the model change is not worth claiming as a better design even if one metric moves slightly upward.",
    ]
    for t in sections:
        p = doc.add_paragraph()
        r = p.add_run(t)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        p.alignment = PP_ALIGN.JUSTIFY
    doc.save(str(DOCX_PATH))


def main():
    render_mermaid()
    build_pptx()
    build_docx()
    print(f"Updated {PPTX_PATH}")
    print(f"Updated {DOCX_PATH}")


if __name__ == "__main__":
    main()
