"""Insert the exported 200-row pilot-processing diagram into review artifacts."""
from pathlib import Path
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Presentation_v1.pptx"
DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Speaking_Doc.docx"
FIGURE = ROOT / "implementation_docments/figures/phase2_clancy_200_row_pilot_processing.png"
MARKER = "CLANCY_200_ROW_PILOT_PROCESSING_V1"

def add_text(slide, value, x, y, w, h, size=16, color=(35,35,35), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = value
    run.font.name = "Aptos"; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

def update_ppt():
    prs = Presentation(PPTX)
    if any(MARKER in getattr(sh, "text", "") for s in prs.slides for sh in s.shapes): return
    backup = PPTX.with_name(PPTX.name + ".before_200_pilot_processing.pptx")
    if not backup.exists(): shutil.copy2(PPTX, backup)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255,255,255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.22))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(20,48,87); bar.line.fill.background()
    add_text(slide, "What Is Done in the First 200-Row Pilot?", .55, .38, 12.2, .45, 23, (20,48,87), True)
    add_text(slide, "Evidence generation, scope-aware review, and acceptance gating are tested before full-pool scaling", .58, .9, 12, .3, 12, (95,95,95))
    slide.shapes.add_picture(str(FIGURE), Inches(.45), Inches(1.22), width=Inches(12.45), height=Inches(5.85))
    add_text(slide, "DeBERTa NLI scope inference is the planned comparison stage; it must not overwrite existing fields automatically.", .7, 7.08, 12, .25, 11, (20,48,87), True)
    # Place after the existing corpus-to-pilot diagram (currently slide 14).
    ids = prs.slides._sldIdLst; item = ids[-1]; ids.remove(item); ids.insert(14, item)
    prs.save(PPTX); print("Inserted pilot-processing slide; slides=", len(prs.slides))

def update_doc():
    doc = Document(DOCX)
    if any(MARKER in p.text for p in doc.paragraphs): return
    backup = DOCX.with_name(DOCX.name + ".before_200_pilot_processing.docx")
    if not backup.exists(): shutil.copy2(DOCX, backup)
    doc.add_heading("How the First 200-Row Pilot Is Processed", 1)
    p = doc.add_paragraph(); p.add_run(MARKER).font.size = Pt(8)
    doc.add_paragraph(
        "This diagram is narrower than the full corpus diagram. It starts after the witness-visible speaking pool has been created and "
        "shows what happens to the controlled 200-row pilot. First, provenance is checked so every row retains its utterance ID, source "
        "video, subtitle/transcript path, and audio/video clip paths. This prevents a prediction from becoming detached from the evidence "
        "that produced it."
    )
    doc.add_paragraph(
        "The same rows then provide different evidence streams. The ViT face-crop stage creates 768-dimensional video features. The Phase 1 "
        "MELD checkpoint supplies an original basic-emotion prediction and confidence. Audio-SER supplies independent acoustic evidence, such "
        "as categorical emotion candidates and dimensional valence, excitement, and dominance. Transcript scope review detects whether the "
        "emotional content is self-expressed, quoted, about another person, or unclear. These signals are kept as separate fields rather than "
        "pretending that one model knows the complete courtroom interpretation."
    )
    doc.add_paragraph(
        "The evidence table feeds the scope-aware courtroom-affect candidate generator. Low-level evidence such as negative activation is not "
        "automatically promoted to labels such as DISTRESSED or TENSE. The acceptance gate then checks basic-emotion confidence, affect "
        "confidence, and critical conflicts. Rows that pass are marked AUTO_ADJUDICATED/SILVER as machine-assisted candidates; rows that do "
        "not pass remain UNRESOLVED/WEAK for human review. Neither category is a human gold label."
    )
    doc.add_paragraph(
        "The dashed DeBERTa branch is the planned enhancement. The NLI model will compare the transcript context with natural-language "
        "hypotheses for emotion-target and temporal scope. Its output should be compared with the existing scope fields and used as additional "
        "evidence or a review trigger. It must not silently overwrite Phase 1 predictions or create credibility, truthfulness, deception, or "
        "reliability labels."
    )
    doc.add_heading("Why this pilot is useful", 2)
    doc.add_paragraph(
        "A 200-row pilot is large enough to expose disagreement patterns but small enough for manual inspection. It lets me test whether the "
        "feature paths are valid, whether audio and text evidence are aligned, whether courtroom-affect rules are too deterministic, and whether "
        "the gate sends questionable cases to review. Only after these checks are stable should the process be rerun over the complete "
        "witness-visible pool and the final corpus statistics be recomputed."
    )
    doc.save(DOCX); print("Updated speaking document")

if __name__ == "__main__":
    update_ppt(); update_doc()
