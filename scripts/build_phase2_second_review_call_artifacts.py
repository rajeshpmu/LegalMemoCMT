"""Build a new Phase 2 Second Review Call deck and student speaking guide."""
from __future__ import annotations

import importlib.util
import json
import shutil
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt

ROOT = Path(__file__).resolve().parents[1]
SOURCE_BUILDER = ROOT / "scripts/build_phase2_second_guidance_call_presentation.py"
OUTPUT_PPTX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Presentation.pptx"
OUTPUT_DOCX = ROOT / "implementation_docments/LegalMemoCMT_Phase2_Second_Review_Call_Speaking_Doc.docx"
GATED = ROOT / "data/processed/phase2/clancy/courtroom_affect_candidates_200_v4_gated.csv"

NAVY = PptRGB(20, 48, 87); TEXT = PptRGB(35, 35, 35); GREY = PptRGB(95, 95, 95)


def load_base_builder():
    spec = importlib.util.spec_from_file_location("phase2_second_guidance_builder", SOURCE_BUILDER)
    module = importlib.util.module_from_spec(spec); assert spec.loader
    spec.loader.exec_module(module)
    module.OUTPUT = OUTPUT_PPTX
    module.CLANCY = GATED
    return module


def ptext(slide, text, x, y, w, h, size=16, color=TEXT, bold=False):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    box.text_frame.word_wrap = True; box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.name = "Aptos"; r.font.size = PptPt(size); r.font.bold = bold; r.font.color.rgb = color


def slide_frame(prs, heading, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = PptRGB(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptInches(.22)); bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    ptext(slide, heading, .55, .38, 12.2, .55, 24, NAVY, True); ptext(slide, subtitle, .58, .96, 12, .35, 12, GREY)
    return slide


def add_bullets(slide, lines, x=.8, y=1.45, w=11.8, h=5.4, size=17):
    ptext(slide, "\n".join("• " + line for line in lines), x, y, w, h, size)


def add_table(slide, headers, rows, x=.55, y=1.35, w=12.2, h=5.7, size=11):
    table = slide.shapes.add_table(len(rows)+1, len(headers), PptInches(x), PptInches(y), PptInches(w), PptInches(h)).table
    for c, value in enumerate(headers): table.cell(0,c).text = str(value)
    for ri, row in enumerate(rows, 1):
        for ci, value in enumerate(row): table.cell(ri,ci).text = str(value)
    for ri, row in enumerate(table.rows):
        for cell in row.cells:
            cell.text_frame.word_wrap = True
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name="Aptos"; run.font.size=PptPt(size); run.font.color.rgb=NAVY if ri == 0 else TEXT; run.font.bold=ri == 0
    return table


def build_pptx(df):
    builder = load_base_builder(); builder.build()
    prs = Presentation(OUTPUT_PPTX)
    # The base deck supplies title, abstract, corpus, source flow, and inherited
    # guidance material. These appended slides are generated from the current
    # gated manifest and are specific to this review call.
    s = slide_frame(prs, "Phase 2 Review: Gated Annotation EDA", "Current evidence from courtroom_affect_candidates_200_v4_gated.csv")
    add_table(s, ["Measure", "Observed result", "Technical interpretation"], [
        ["Rows", len(df), "200-row pilot, not the full Clancy corpus"],
        ["AUTO_ADJUDICATED / SILVER", int((df.annotation_status == 'AUTO_ADJUDICATED').sum()), "Machine-assisted accepted tier; not gold labels"],
        ["UNRESOLVED / WEAK", int((df.annotation_status == 'UNRESOLVED').sum()), "Requires human review before supervised training"],
        ["Critical conflicts", int((df.critical_conflict == 'YES').sum()), "Rows requiring special inspection"],
        ["Final basic emotion", ", ".join(f"{k}:{v}" for k,v in df.final_basic_emotion.value_counts().to_dict().items()), "Final field is unresolved for weak rows"],
        ["Final courtroom affect", ", ".join(f"{k}:{v}" for k,v in df.final_courtroom_affect.value_counts().to_dict().items()), "CALM_COMPOSED is behavior, not basic emotion"],
    ], size=10)

    s = slide_frame(prs, "Phase 2 Hugging Face Models", "Each model supplies a different evidence type; none independently creates courtroom ground truth")
    add_table(s, ["Model", "Phase 2 use", "Example output / limitation"], [
        ["faster-whisper tiny.en", "ASR word timestamps and transcript alignment", "Text timing; may misrecognize accents or courtroom overlap"],
        ["pyannote/speaker-diarization-3.1", "Anonymous speaker segments and cluster IDs", "SPEAKER_07 from 02:13:32–02:13:42; role still needs mapping"],
        ["google/vit-base-patch16-224-in21k", "768-D face-crop video embeddings", ".npy features for the Phase 1 video loader; not an emotion label"],
        ["3loi/SER-Odyssey-Baseline-WavLM-Multi-Attributes", "Continuous voice valence/excitement/dominance evidence", "Excitement is acoustic activation evidence, not distress"],
        ["speechbrain/emotion-recognition-wav2vec2-IEMOCAP", "Categorical audio cross-check", "neu/ang/hap/sad only; not comparable to all seven MELD classes"],
        ["MoritzLaurer/deberta-v3-large-zeroshot-v2.0", "NLI hypotheses for target and temporal scope", "SELF_EXPRESSED + PAST_SELF candidate; machine suggestion requiring review"],
    ], size=9)

    def row(uid):
        return df[df.utterance_id == uid].iloc[0]
    a = row("DCBWoWhsTpA_turn06801")
    s = slide_frame(prs, "Example: SILVER Auto-Adjudication", "DCBWoWhsTpA_turn06801: quoted content versus the witness's presentation")
    add_table(s, ["Field", "Value"], [[k, a.get(k, "")] for k in [
        "utterance_id","speaker_role","utterance_text","phase1_basic_emotion","phase1_basic_emotion_confidence",
        "basic_emotion_review_candidate","basic_emotion_review_candidate_confidence","proposed_courtroom_affect",
        "proposed_courtroom_affect_confidence","negative_activation_candidate","distress_corroboration_present",
        "speaker_emotion_evidence_present","emotion_target_scope","annotation_status","annotation_tier",
        "final_basic_emotion","final_courtroom_affect"]], size=10, h=5.8)

    a = row("DCBWoWhsTpA_turn06570")
    s = slide_frame(prs, "Example: Weak / Unresolved Candidate", "DCBWoWhsTpA_turn06570: qualified response and low basic-emotion confidence")
    add_table(s, ["Field", "Value"], [[k, a.get(k, "")] for k in [
        "utterance_id","utterance_text","phase1_basic_emotion","phase1_basic_emotion_confidence",
        "basic_emotion_review_candidate","basic_emotion_review_candidate_confidence","proposed_courtroom_affect",
        "proposed_courtroom_affect_confidence","negative_activation_candidate","distress_corroboration_present",
        "speaker_emotion_evidence_present","emotion_target_scope","critical_conflict","annotation_status",
        "annotation_tier","final_basic_emotion","final_courtroom_affect","acceptance_gate_reason"]], size=10, h=5.8)

    s = slide_frame(prs, "Decision and Next Action", "What this pilot proves and what must happen before scaling")
    add_bullets(s, [
        "The gate separates machine-accepted SILVER rows from unresolved rows without overwriting Phase 1 evidence.",
        "The EDA demonstrates domain-shift cases where reported sadness or distress is not necessarily the witness's own emotion.",
        "First inspect the 75 unresolved rows and 7 critical-conflict rows using the acceptance-gate inspection script.",
        "Human-review a stratified sample of SILVER rows to estimate precision before using them for weak supervision.",
        "Run DeBERTa target/temporal scope inference and compare its hypotheses with the existing scope fields.",
        "Only after review should final labels be merged into the annotation manifest and used for controlled training.",
    ], size=17)
    prs.save(OUTPUT_PPTX)


def build_docx(df):
    doc = Document(); sec = doc.sections[0]; sec.top_margin=Inches(.65); sec.bottom_margin=Inches(.65); sec.left_margin=Inches(.8); sec.right_margin=Inches(.8)
    normal=doc.styles["Normal"]; normal.font.name="Aptos"; normal.font.size=Pt(10.5)
    title=doc.add_paragraph(); title.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=title.add_run("LegalMemoCMT Phase 2\nSecond Review Call Speaking Document"); r.bold=True; r.font.size=Pt(21); r.font.color.rgb=RGBColor(20,48,87)
    p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; p.add_run("Student-level technical explanation of the current gated annotation stage").italic=True
    doc.add_paragraph("This is a speaking document, not a copy of the slide text. I use it to explain what each artifact means, how the code produced it, and what can and cannot be claimed from the current 200-row pilot. Clancy is the primary benchmark; the tribunal branch remains secondary bootstrap evidence, while Tupac/Keffe D and Indian-SIM are planned expansion/adaptation branches.")
    doc.add_heading("Project title, abstract, and goal", 1)
    doc.add_paragraph("Project title: LegalMemoCMT: An Explainable Multilingual Multimodal Emotional Cue Analysis Framework for Indian Courtroom Testimony Using Cross-Modal Transformers.")
    doc.add_paragraph("The Phase 2 goal is to construct an utterance- or turn-level courtroom corpus in which transcript text, audio, video, speaker role, witness identity status, and affect evidence remain traceable. The research novelty is not claiming that a pretrained model understands courtroom emotion automatically. It is the evidence-aware pipeline that separates what was said, who spoke, how the person presented, whose emotion was described, and what a model predicted.")
    doc.add_heading("Slide-by-slide speaking guidance", 1)
    # Explain inherited slides by their visible first title where possible.
    prs = Presentation(OUTPUT_PPTX)
    for i, slide in enumerate(prs.slides, 1):
        texts=[shape.text.strip() for shape in slide.shapes if hasattr(shape,"text") and shape.text.strip()]
        heading=texts[0] if texts else f"Slide {i}"
        doc.add_heading(f"Slide {i}: {heading[:120]}", 2)
        doc.add_paragraph("I explain this slide as part of the evidence chain rather than reading it aloud. The important question is which input manifest, model output, or validation rule supports the claim. A source URL is discovery evidence; a verified record, aligned turn, clip path, and preserved provenance are stronger evidence. Any candidate label remains separate from a human-validated label.")
    doc.add_heading("Current gated EDA: how to explain it", 1)
    doc.add_paragraph(f"The gated pilot contains {len(df)} rows. {int((df.annotation_status == 'AUTO_ADJUDICATED').sum())} rows are AUTO_ADJUDICATED/SILVER and {int((df.annotation_status == 'UNRESOLVED').sum())} are UNRESOLVED/WEAK. There are {int((df.critical_conflict == 'YES').sum())} critical-conflict rows. These counts describe the behavior of the acceptance policy, not the prevalence of emotions in the courtroom.")
    doc.add_paragraph("The final fields are intentionally conservative. Accepted rows receive final_basic_emotion and final_courtroom_affect from machine-assisted candidates. Weak rows receive final_basic_emotion=UNRESOLVED, although a supported affect candidate may be retained for review. This prevents a confident affect score from hiding a weak or conflicting basic-emotion decision.")
    doc.add_heading("Hugging Face models: student explanation", 1)
    models=[
        ("faster-whisper tiny.en", "converts speech to text and word timestamps for alignment; it can make recognition errors, so transcript text remains the source text and ASR is evidence."),
        ("pyannote/speaker-diarization-3.1", "finds speech regions and anonymous voice clusters. It answers who spoke when acoustically, not whether that person is a witness, lawyer, or judge."),
        ("google/vit-base-patch16-224-in21k", "converts sampled face crops into 768-dimensional video features. These features support a loader or visual analysis; they are not labels."),
        ("3loi/SER-Odyssey-Baseline-WavLM-Multi-Attributes", "provides continuous valence, excitement, and dominance-like acoustic evidence. Excitement is an activation signal, not a diagnosis of distress."),
        ("speechbrain/emotion-recognition-wav2vec2-IEMOCAP", "provides a small categorical audio cross-check. Its neutral/anger/joy/sadness-like outputs do not cover every MELD class."),
        ("MoritzLaurer/deberta-v3-large-zeroshot-v2.0", "compares transcript context with natural-language NLI hypotheses for target scope and temporal scope. It proposes labels; it does not adjudicate truth or credibility."),
    ]
    for name, explanation in models:
        doc.add_paragraph(f"{name}: {explanation}", style="List Bullet")
    doc.add_heading("Worked example: DCBWoWhsTpA_turn06801", 1)
    a=df[df.utterance_id=="DCBWoWhsTpA_turn06801"].iloc[0]
    doc.add_paragraph("I would explain this as a clean separation of dimensions. The original Phase 1 model predicted sadness at 0.544674, which is retained as evidence of domain shift. The transcript is QUOTED_SPEECH, and speaker_emotion_evidence_present is NO, so the quoted emotional content is not automatically assigned to the witness. Negative activation is YES, but distress corroboration is NO; therefore negative acoustic activation is not treated as DISTRESSED. The integrated review candidate is neutral at 0.75 and the behavioral candidate is CALM_COMPOSED at 0.60. The gate accepts this row as AUTO_ADJUDICATED/SILVER, not as human gold truth.")
    doc.add_paragraph("The final machine-assisted fields are:")
    for k in ["phase1_basic_emotion","phase1_basic_emotion_confidence","basic_emotion_review_candidate","basic_emotion_review_candidate_confidence","proposed_courtroom_affect","proposed_courtroom_affect_confidence","negative_activation_candidate","distress_corroboration_present","speaker_emotion_evidence_present","emotion_target_scope","annotation_status","annotation_tier","final_basic_emotion","final_courtroom_affect"]:
        doc.add_paragraph(f"{k}={a.get(k, '')}", style="List Bullet")
    doc.add_heading("Worked example: DCBWoWhsTpA_turn06570", 1)
    b=df[df.utterance_id=="DCBWoWhsTpA_turn06570"].iloc[0]
    doc.add_paragraph("This row demonstrates why the gate does not accept every confident-looking affect candidate. The witness says, 'I don't believe...' and the transcript contains a non-speech marker. The revised heuristic proposes HESITANT_UNCERTAIN, not generic TENSE, but the basic candidate remains low-confidence. Therefore final_basic_emotion is UNRESOLVED and the row is WEAK. The affect candidate is retained for a reviewer; it is not presented as a final human annotation.")
    for k in ["phase1_basic_emotion","phase1_basic_emotion_confidence","basic_emotion_review_candidate","basic_emotion_review_candidate_confidence","proposed_courtroom_affect","proposed_courtroom_affect_confidence","negative_activation_candidate","distress_corroboration_present","speaker_emotion_evidence_present","emotion_target_scope","critical_conflict","annotation_status","annotation_tier","final_basic_emotion","final_courtroom_affect","acceptance_gate_reason"]:
        doc.add_paragraph(f"{k}={b.get(k, '')}", style="List Bullet")
    doc.add_heading("What I will say about readiness", 1)
    doc.add_paragraph("The pilot is ready for controlled review and calibration, not for claiming fully human-validated courtroom emotion labels. The next action is to inspect unresolved and critical rows, manually validate a stratified sample of SILVER rows, run the DeBERTa scope pilot, and only then merge approved human labels into a training manifest. I will not infer deception, credibility, truthfulness, or reliability from any model output.")
    doc.save(OUTPUT_DOCX)


def main():
    df=pd.read_csv(GATED,dtype=str).fillna("")
    build_pptx(df); build_docx(df)
    print(json.dumps({"pptx":str(OUTPUT_PPTX),"docx":str(OUTPUT_DOCX),"rows_used":len(df)}, indent=2))


if __name__ == "__main__": main()
