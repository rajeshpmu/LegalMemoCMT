from __future__ import annotations

import json
import subprocess
from pathlib import Path

import matplotlib.pyplot as plt
from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_MELD_Facial_Cues_ViT_Implementation_Plan.docx"
PPTX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_MELD_Facial_Cues_ViT_Implementation_Plan.pptx"
FIG_DIR = ROOT / "implementation_docments" / "figures" / "meld_vit_facecue_ppt"
FIG_DIR.mkdir(parents=True, exist_ok=True)

VIT_PNG = FIG_DIR / "vit_step.png"
VIT_SVG = FIG_DIR / "vit_step.svg"
FACE_PNG = FIG_DIR / "face_crop_vs_full_frame.png"
FACE_SVG = FIG_DIR / "face_crop_vs_full_frame.svg"
WARM_PNG = FIG_DIR / "warm_start_flow.png"
WARM_SVG = FIG_DIR / "warm_start_flow.svg"
FLOW_PNG = FIG_DIR / "utterance_to_embedding_flow.png"
FLOW_SVG = FIG_DIR / "utterance_to_embedding_flow.svg"
HANDOFF_PNG = FIG_DIR / "phase1_to_phase2_handoff.png"
HANDOFF_SVG = FIG_DIR / "phase1_to_phase2_handoff.svg"
VALCURVE_PNG = FIG_DIR / "validation_curve.png"
VALCURVE_SVG = FIG_DIR / "validation_curve.svg"
COMPARE_PNG = FIG_DIR / "run_comparison.png"
COMPARE_SVG = FIG_DIR / "run_comparison.svg"
FCEC_PNG = FIG_DIR / "face_crop_next_step.png"
FCEC_SVG = FIG_DIR / "face_crop_next_step.svg"


def render_mermaid(code: str, svg_path: Path, png_path: Path) -> None:
    mmd_path = png_path.with_suffix(".mmd")
    mmd_path.write_text(code, encoding="utf-8")
    subprocess.run(
        ["npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(mmd_path), "-o", str(svg_path), "-b", "white"],
        check=True,
        cwd=str(ROOT),
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    subprocess.run(
        ["npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(mmd_path), "-o", str(png_path), "-b", "white"],
        check=True,
        cwd=str(ROOT),
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )


def add_bg(slide, prs: Presentation):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 48, 87)
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None, y: float = 0.35):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(12.2), Inches(0.65))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(20, 48, 87)
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.58), Inches(y + 0.5), Inches(12.1), Inches(0.5))
        p2 = tx2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(12.5)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(72, 72, 72)


def add_body(slide, bullets, x=0.75, y=1.45, w=5.9, h=5.55, font_size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.name = "Aptos"
        r.font.size = Pt(font_size)
        r.font.color.rgb = RGBColor(35, 35, 35)
    return box


def add_code_box(slide, text: str, left: float, top: float, width: float, height: float, font_size: int = 10):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
    shape.line.color.rgb = RGBColor(180, 190, 205)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Courier New"
    r.font.size = Pt(font_size)
    r.font.color.rgb = RGBColor(25, 25, 25)
    return shape


def add_table(slide, headers: list[str], rows: list[list[str]], left: float, top: float, width: float, height: float, col_widths: list[float], font_size: float = 10):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for i, header in enumerate(headers):
        table.columns[i].width = Inches(col_widths[i])
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(20, 48, 87)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Aptos"
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = RGBColor(255, 255, 255)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_idx < len(row) - 1 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Aptos"
                    r.font.size = Pt(font_size - 0.2)
                    r.font.color.rgb = RGBColor(35, 35, 35)


def add_picture(slide, path: Path, left: float, top: float, width: float):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def build_assets() -> None:
    render_mermaid(
        """flowchart LR
  A["Raw MELD clip"] --> B["Sample frames"]
  B --> C["Pretrained ViT"]
  C --> D["Facial embedding"]
  D --> E["Saved .npy"]
""",
        VIT_SVG,
        VIT_PNG,
    )
    render_mermaid(
        """flowchart LR
  A["Video clip"] --> B["Full sampled frames"]
  A --> C["Face crop before ViT"]
  B --> D["Robust baseline"]
  C --> E["Cleaner facial cue"]
  D --> F["Compare"]
  E --> F
""",
        FACE_SVG,
        FACE_PNG,
    )
    render_mermaid(
        """sequenceDiagram
  participant A as Weighted-CE MELD checkpoint
  participant B as Warm-start fine-tuning
  participant C as Add facial-cue branch
  participant D as Small learning rate
  participant E as Compare against baseline
  A->>B: reuse the learned boundary
  B->>C: attach facial cues
  C->>D: protect the backbone
  D->>E: evaluate the change
""",
        WARM_SVG,
        WARM_PNG,
    )
    render_mermaid(
        """sequenceDiagram
  participant A as Annotation row
  participant B as find_meld_clip(...)
  participant C as Raw clip
  participant D as sample_video_frames(...)
  participant E as ViT embeddings
  participant F as np.save(... .npy)
  participant G as meld_vit_facecue.csv
  A->>B: resolve the utterance clip
  B->>C: load the raw video
  C->>D: sample and resize frames
  D->>E: infer facial embeddings
  E->>F: write the .npy file
  F->>G: point the manifest row to it
""",
        FLOW_SVG,
        FLOW_PNG,
    )
    render_mermaid(
        """sequenceDiagram
  participant A as Phase 1 closure
  participant B as Keep best checkpoint
  participant C as Freeze baseline narrative
  participant D as Stop Phase 1 tuning
  participant E as Move to courtroom-testimony Phase 2
  A->>B: keep the strongest model
  B->>C: lock the baseline story
  C->>D: stop further tuning
  D->>E: move the remaining effort forward
""",
        HANDOFF_SVG,
        HANDOFF_PNG,
    )
    render_mermaid(
        """sequenceDiagram
  participant A as MELD utterance clip
  participant B as Face detector
  participant C as Crop + resize
  participant D as Pretrained ViT
  participant E as Face embedding .npy
  participant F as Compare vs full-frame baseline
  A->>B: read the utterance video
  B->>C: isolate and standardize the face
  C->>D: extract facial features
  D->>E: save the embedding file
  E->>F: test whether faces help
""",
        FCEC_SVG,
        FCEC_PNG,
    )
    epochs = [1, 2, 3, 4]
    train_loss = [1.4195, 1.5937, 1.3378, 1.1202]
    val_loss = [1.6160, 1.6063, 1.7225, 1.9397]
    train_acc = [0.5975, 0.5576, 0.6492, 0.7261]
    val_acc = [0.5414, 0.6010, 0.5893, 0.5399]

    fig, axes = plt.subplots(2, 1, figsize=(12, 6.8), sharex=True)
    loss_ax, acc_ax = axes

    loss_ax.plot(epochs, train_loss, marker="o", linewidth=2.2, label="Train loss", color="#143057")
    loss_ax.plot(epochs, val_loss, marker="o", linewidth=2.2, label="Val loss", color="#c44e52")
    loss_ax.set_ylabel("Loss")
    loss_ax.grid(True, alpha=0.25)
    loss_ax.legend(loc="upper right", frameon=False)
    loss_ax.set_title("MELD Phase 1 Validation Curve", fontsize=14, weight="bold")

    acc_ax.plot(epochs, train_acc, marker="o", linewidth=2.2, label="Train acc", color="#2a7f62")
    acc_ax.plot(epochs, val_acc, marker="o", linewidth=2.2, label="Val acc", color="#f28e2b")
    acc_ax.set_xlabel("Epoch")
    acc_ax.set_ylabel("Accuracy")
    acc_ax.grid(True, alpha=0.25)
    acc_ax.legend(loc="lower right", frameon=False)
    acc_ax.set_xticks(epochs)
    acc_ax.set_ylim(0.45, 0.8)

    fig.suptitle(
        "Where Phase 1 Is Right Now: validation improves early, then starts to weaken",
        fontsize=12.5,
        y=0.99,
    )
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    fig.savefig(VALCURVE_PNG, dpi=220, bbox_inches="tight")
    fig.savefig(VALCURVE_SVG, bbox_inches="tight")
    plt.close(fig)

    runs = [
        ("Weighted CE", json.loads((ROOT / "results/paper_aligned_meld_cv/cmt_min/fold_2/metrics.json").read_text())),
        ("Focal scratch", json.loads((ROOT / "results/improvement/class_balanced_focal/meld_selected/cmt_min/fold_2/metrics.json").read_text())),
        ("Warm-start focal", json.loads((ROOT / "results/improvement/warmstart_focal/meld_selected/cmt_min/fold_2/metrics.json").read_text())),
        ("Warm-start + ViT", json.loads((ROOT / "results/improvement/warmresume_focal/meld_fold_2/metrics.json").read_text())),
    ]
    labels = [name for name, _ in runs]
    weighted_f1 = [item[1]["weighted_f1"] for item in runs]
    macro_f1 = [item[1]["macro_f1"] for item in runs]
    accuracy = [item[1]["accuracy"] for item in runs]

    x = list(range(len(labels)))
    bar_w = 0.24
    fig, ax = plt.subplots(figsize=(11.8, 5.8))
    ax.bar([i - bar_w for i in x], accuracy, width=bar_w, label="Accuracy", color="#143057")
    ax.bar(x, weighted_f1, width=bar_w, label="Weighted F1", color="#2a7f62")
    ax.bar([i + bar_w for i in x], macro_f1, width=bar_w, label="Macro F1", color="#f28e2b")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=0)
    ax.set_ylim(0.0, 0.75)
    ax.set_ylabel("Score")
    ax.set_title("MELD Fold 2 Run Comparison", fontsize=13, weight="bold")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False, ncol=3, loc="upper center", bbox_to_anchor=(0.5, 1.12))
    for idx, val in enumerate(accuracy):
        ax.text(idx - bar_w, val + 0.015, f"{val:.3f}", ha="center", va="bottom", fontsize=8)
    for idx, val in enumerate(weighted_f1):
        ax.text(idx, val + 0.015, f"{val:.3f}", ha="center", va="bottom", fontsize=8)
    for idx, val in enumerate(macro_f1):
        ax.text(idx + bar_w, val + 0.015, f"{val:.3f}", ha="center", va="bottom", fontsize=8)
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    fig.savefig(COMPARE_PNG, dpi=220, bbox_inches="tight")
    fig.savefig(COMPARE_SVG, bbox_inches="tight")
    plt.close(fig)


def build_pptx() -> None:
    build_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "LegalMemoCMT MELD Facial Cues ViT Implementation Plan", "How the facial-cue branch is built, why it is warm-started, and how the outputs are interpreted.")
    add_body(
        slide,
        [
            "This deck explains the next MELD extension: adding facial cues using ViT on top of the weighted-CE Phase 1 baseline.",
            "The main idea is to keep the existing text+audio backbone and test whether visual evidence from the speaker's face improves emotion recognition.",
            "The current pipeline uses full sampled frames first, then learns the ViT facial embeddings as reusable .npy features.",
            "A student should read the deck as a pipeline story: raw clip -> frames -> ViT -> embeddings -> warm-start training -> validation trend -> next step.",
        ],
        x=0.75,
        y=1.55,
        w=6.0,
        h=5.2,
        font_size=15,
    )
    add_picture(slide, VIT_PNG, 6.9, 1.55, 5.75)
    add_code_box(
        slide,
        "Target path: HuBERT + BERT backbone + ViT facial cues\nGoal: improve the current weighted-CE MELD baseline without changing the core text/audio story first.",
        6.95,
        5.75,
        5.75,
        0.8,
        font_size=10,
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "1. What the ViT Step Means in Practice", "The ViT step turns raw video clips into reusable facial feature vectors.")
    add_body(
        slide,
        [
            "For each MELD utterance, the script resolves the matching raw clip, samples fixed frames, and feeds those frames into a pretrained ViT.",
            "ViT does not read the whole video directly here; it reads sampled RGB frames and converts them into embeddings.",
            "Those embeddings are saved as .npy files so the training script can reuse them without running ViT again.",
            "The student idea is simple: raw clip in, learned face representation out.",
        ],
        x=0.75,
        y=1.5,
        w=5.7,
        h=5.4,
        font_size=15,
    )
    add_picture(slide, VIT_PNG, 6.75, 1.65, 5.9)
    add_code_box(
        slide,
        "find_meld_clip(...) -> sample_video_frames(...) -> AutoImageProcessor(...) -> AutoModel(...) -> np.save(... .npy)",
        6.8,
        5.9,
        5.85,
        0.6,
        font_size=10,
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "2. Full Frames vs Face Crop for Courtroom Testimony", "Face crop is the better facial-cue interpretation, but full frames are the safer first baseline.")
    add_body(
        slide,
        [
            "For courtroom testimony, the most informative visual signal is usually the speaker's face, not the background.",
            "Face cropping focuses ViT on the actual emotion-bearing region: eyes, brows, mouth, and micro-expressions.",
            "Full sampled frames are still useful as a baseline because they are simpler and more robust when face detection is unreliable.",
            "The practical plan is to compare both and keep the version that improves MELD most cleanly.",
        ],
        x=0.75,
        y=1.5,
        w=5.7,
        h=5.35,
        font_size=15,
    )
    add_picture(slide, FACE_PNG, 6.85, 1.55, 5.75)
    add_code_box(
        slide,
        "Primary claim: facial cues should ideally come from the face.\nFallback baseline: full frames keep the first experiment stable and reproducible.",
        6.9,
        5.85,
        5.7,
        0.75,
        font_size=10,
    )

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "3. Why Warm-Start from the Existing Checkpoint", "The baseline is reused because the text/audio model already learned a useful MELD boundary.")
    add_body(
        slide,
        [
            "The weighted-CE MELD Phase 1 checkpoint is the backbone because it is the strongest current conversational result.",
            "Warm-starting means we reuse those learned text/audio weights and then add the facial branch on top.",
            "This keeps the comparison clean: any gain can be attributed more directly to the new visual information.",
            "The small learning rate protects the existing solution while the facial branch learns to align.",
        ],
        x=0.75,
        y=1.5,
        w=5.6,
        h=5.4,
        font_size=15,
    )
    add_picture(slide, WARM_PNG, 6.9, 1.6, 5.8)
    add_code_box(
        slide,
        "Why this baseline? It is already competitive.\nWhy warm-start? It isolates the effect of facial cues instead of retraining the whole system from scratch.",
        6.95,
        5.8,
        5.75,
        0.78,
        font_size=10,
    )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "4. Utterance-to-Embedding Flow With Code Lines", "This is the exact code path from a MELD row to a saved facial embedding file.")
    add_picture(slide, FLOW_PNG, 0.75, 1.45, 5.2)
    add_table(
        slide,
        ["Stage", "Input", "Output", "Example"],
        [
            ["1. Locate clip", "Annotation row", "Raw video path", "data/MELD/raw/.../dia0_utt0.mp4"],
            ["2. Sample frames", "Raw clip", "Fixed frame tensor", "(16, 224, 224, 3)"],
            ["3. Run ViT", "RGB frames", "CLS embeddings", "(16, 768) float32"],
            ["4. Save feature", "Embedding array", ".npy file", "train_dia0_utt0.npy"],
        ],
        left=6.1,
        top=1.45,
        width=6.45,
        height=3.2,
        col_widths=[1.05, 1.35, 1.45, 2.6],
        font_size=9.0,
    )
    add_code_box(
        slide,
        "clip_path = find_meld_clip(...)\nframes = sample_video_frames(video_path, cfg)\ninputs = image_processor(images=list(batch), return_tensors=\"pt\")\noutputs = vit_model(**inputs)\nhidden = outputs.last_hidden_state[:, 0, :]\nnp.save(video_feat_path, vit_embeddings)",
        0.85,
        4.95,
        11.65,
        1.2,
        font_size=9,
    )
    add_code_box(
        slide,
        "Code locations:\nbuild_meld_vit_facecue_manifest.py:123-135  |  preprocessing.py:41-75  |  build_meld_vit_facecue_manifest.py:101-102",
        0.85,
        6.25,
        11.65,
        0.45,
        font_size=9,
    )

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "5. Epoch Trend So Far", "The warm-start facial-cue run peaks early, so the best checkpoint is likely an early epoch.")
    add_table(
        slide,
        ["Epoch", "Train Loss", "Train Acc", "Val Loss", "Val Acc", "Meaning"],
        [
            ["1", "1.4195", "0.5975", "1.6160", "0.5414", "Warm-start begins; stable but not best yet"],
            ["2", "1.5937", "0.5576", "1.6063", "0.6010", "Best validation so far"],
            ["3", "1.3378", "0.6492", "1.7225", "0.5893", "Training improves, validation slips"],
            ["4", "1.1202", "0.7261", "1.9397", "0.5399", "Clear overfitting"],
        ],
        left=0.75,
        top=1.55,
        width=12.0,
        height=2.75,
        col_widths=[0.8, 1.25, 1.25, 1.25, 1.15, 4.3],
        font_size=9.0,
    )
    add_code_box(
        slide,
        "Conclusion: epoch 2 is the likely best checkpoint because validation accuracy peaks there and later epochs overfit.\nNext step: use the best-validation checkpoint and continue the MELD facial-cue comparison against the weighted-CE baseline.",
        0.85,
        4.65,
        11.7,
        0.85,
        font_size=10,
    )
    add_body(
        slide,
        [
            "Student takeaway: do not judge the run by the last epoch alone; judge it by the best validation point.",
        ],
        x=0.95,
        y=5.75,
        w=11.25,
        h=0.8,
        font_size=15,
    )

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "Current Run Comparison: Where the project stands now", "This compares the four MELD training strategies on Fold 2.")
    add_picture(slide, COMPARE_PNG, 0.55, 1.45, 8.35)
    add_body(
        slide,
        [
            "Weighted CE is still the strongest overall Fold 2 baseline.",
            "Focal loss from scratch is clearly too unstable for this setup.",
            "Warm-start focal is better than focal from scratch, but still below the weighted-CE backbone.",
            "Warm-start + ViT is the strongest improvement run, so it is the most useful bridge into the next research step.",
        ],
        x=9.05,
        y=1.55,
        w=3.55,
        h=4.1,
        font_size=13.5,
    )
    add_code_box(
        slide,
        "Story to tell:\n1) keep the weighted-CE model as the reference\n2) use warm-start + ViT as the best improvement path\n3) treat the visual branch as the bridge to the next plan",
        9.05,
        5.25,
        3.55,
        1.05,
        font_size=10,
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "Current Status: What the experiments say", "A technical reading of the run comparison.")
    add_body(
        slide,
        [
            "The weighted-CE baseline remains the best stable result for MELD Fold 2.",
            "The focal-from-scratch run shows that focal loss alone is not enough in this setup.",
            "Warm-starting helps because it starts from a better learned boundary instead of relearning from zero.",
            "Adding ViT on top of the warm-start improves the improvement run, but it still does not beat the strongest baseline.",
        ],
        x=0.75,
        y=1.5,
        w=6.2,
        h=5.35,
        font_size=15,
    )
    add_picture(slide, HANDOFF_PNG, 7.05, 1.65, 5.55)
    add_code_box(
        slide,
        "Interpretation:\nThe project has a strong baseline already.\nThe new branches are informative, but the story is now about controlled extension rather than baseline replacement.",
        7.0,
        5.8,
        5.6,
        0.85,
        font_size=10,
    )

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "Next Plan: Where the remaining work should go", "Use the best current model as the base for the courtroom-testimony direction.")
    add_body(
        slide,
        [
            "Keep the weighted-CE MELD checkpoint as the reference baseline.",
            "Use the warm-start + ViT experiment as the most promising extension path.",
            "Shift the remaining implementation effort toward the courtroom-testimony novelty.",
            "Focus the next stage on whether facial cues add value in the legal-domain setting.",
        ],
        x=0.75,
        y=1.5,
        w=6.0,
        h=5.25,
        font_size=15,
    )
    add_table(
        slide,
        ["Next action", "Reason"],
        [
            ["Freeze baseline", "The weighted-CE model is already the strongest stable anchor"],
            ["Use ViT branch", "The facial-cue path is the best current improvement candidate"],
            ["Move to courtroom data", "The novelty is in the legal-domain adaptation"],
            ["Stop Phase 1 tuning", "Further loss-only optimization is now diminishing returns"],
        ],
        left=7.0,
        top=1.65,
        width=5.55,
        height=2.35,
        col_widths=[2.05, 3.5],
        font_size=9.5,
    )
    add_code_box(
        slide,
        "Technical conclusion:\nThe next useful work is not more baseline tuning.\nIt is the controlled transition into the courtroom-testimony novelty implementation.",
        7.0,
        4.3,
        5.55,
        1.0,
        font_size=10,
    )

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "Project Direction: Final recommendation for the current stage", "What the deck should tell the reviewer.")
    add_body(
        slide,
        [
            "The current project status is: Phase 1 has a strong baseline and a measured improvement path.",
            "The next plan is: carry the best current checkpoint into the courtroom-testimony novelty work.",
            "The technical story is now about controlled extension, not just benchmark chasing.",
            "That framing keeps the thesis novelty focused on adaptation to the legal domain.",
        ],
        x=0.75,
        y=1.5,
        w=6.0,
        h=5.25,
        font_size=15,
    )
    add_picture(slide, HANDOFF_PNG, 7.05, 1.65, 5.55)
    add_code_box(
        slide,
        "Current-stage message:\nkeep the best baseline\nshow the best improvement path\nmove the remaining work into the legal-domain novelty",
        7.0,
        5.8,
        5.55,
        0.9,
        font_size=10,
    )

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_title(slide, "Thesis Summary: What to do in the next two weeks", "Use one controlled face-crop ablation before moving fully into the courtroom-testimony implementation.")
    add_body(
        slide,
        [
            "Yes, you should still test MELD face crops before leaving Phase 1.",
            "That test is the cleanest way to answer whether the visual branch works better when it focuses on the speaker's face instead of the whole frame.",
            "Keep the same MELD splits, keep the current best checkpoint as the reference, and change only the visual input path.",
            "If face crops help, freeze that setting; if they do not, keep the full-frame path and move on without spending more time on visuals.",
        ],
        x=0.75,
        y=1.45,
        w=6.0,
        h=5.45,
        font_size=15,
    )
    add_picture(slide, FCEC_PNG, 6.85, 1.6, 5.9)
    add_code_box(
        slide,
        "Implementation logic:\n1) keep the same manifest and fold split\n2) insert face detection + crop before ViT\n3) resize faces to 224x224\n4) extract embeddings\n5) compare against full-frame results\n6) freeze the better setting",
        6.9,
        5.9,
        5.85,
        0.75,
        font_size=10,
    )

    prs.save(PPTX_PATH)


def main() -> None:
    build_pptx()
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
