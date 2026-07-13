from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import subprocess

import matplotlib.pyplot as plt
from matplotlib.table import Table
from docx import Document
from docx.shared import Pt
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph
from docx.table import _Row
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Presentation.pptx"
DOCX_PATH = ROOT / "implementation_docments" / "LegalMemoCMT_Phase1_ESA_Prep_Guide.docx"
ASSET_DIR = ROOT / "implementation_docments" / "phase1_esa_assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)
EDA_MMD = ASSET_DIR / "meld_eda_context.mmd"
EDA_SVG = ASSET_DIR / "meld_eda_context.svg"
EDA_PNG = ASSET_DIR / "meld_eda_context.png"
TABLE_PNG = ASSET_DIR / "meld_eda_table10_summary.png"

BODY = RGBColor(42, 42, 42)
MUTED = RGBColor(92, 92, 92)
PALE = RGBColor(249, 246, 222)
PALE2 = RGBColor(235, 231, 248)
ACCENT = RGBColor(112, 98, 184)

MERMAID = """%%{init: {'themeVariables': {'fontSize': '30px', 'fontFamily': 'Aptos'}}}%%
flowchart TB
  A[MELD is multi-party] --> B[Speaker context matters]
  B --> C[Emotion shifts are common]
  C --> D[Contextual distance matters]
  D --> E[Weighted metrics are needed]
  E --> F[Phase 1 evaluates context-aware behavior]
"""


def style_text(run, size=12, bold=False, color=BODY):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def insert_paragraph_before(paragraph, text="", style=None):
    new_p = OxmlElement("w:p")
    paragraph._p.addprevious(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if style is not None:
        new_para.style = style
    if text:
        new_para.add_run(text)
    return new_para


def remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def move_slide(prs, old_idx, new_idx):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[old_idx]
    sldIdLst.remove(sldId)
    sldIdLst.insert(new_idx, sldId)


def render_assets():
    EDA_MMD.write_text(MERMAID, encoding="utf-8")
    for out in [EDA_SVG, EDA_PNG]:
        subprocess.run([
            "npx", "-y", "@mermaid-js/mermaid-cli", "-i", str(EDA_MMD), "-o", str(out), "-b", "white", "-s", "3"
        ], cwd=str(ROOT), check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

    rows = [
        ["Joy / Happy", "648", "2308"],
        ["Anger", "1103", "1607"],
        ["Disgust", "2", "361"],
        ["Sadness", "1084", "1002"],
        ["Surprise", "107", "1636"],
        ["Neutral", "1708", "6436"],
        ["Avg utterance length", "15.8", "8.0"],
        ["Unique words", "3598", "10643"],
        ["Avg conversation length", "49.2", "9.6"],
    ]
    fig, ax = plt.subplots(figsize=(9.8, 4.8), dpi=220)
    ax.axis("off")
    ax.text(0.0, 1.04, "Table 10 summary: why MELD matters for Phase 1", fontsize=14, fontweight="bold", ha="left", va="bottom")
    ax.text(0.0, 0.985, "MELD vs IEMOCAP statistics from the paper", fontsize=10, color="#555555", ha="left", va="bottom")
    table = Table(ax, bbox=[0.0, 0.0, 1.0, 0.92])
    col_w = [0.46, 0.22, 0.22]
    headers = ["Metric", "IEMOCAP", "MELD"]
    for c, h in enumerate(headers):
        cell = table.add_cell(0, c, width=col_w[c], height=0.085, text=h, loc="center", facecolor="#d9e2f3")
        cell.set_edgecolor("#7a8ca8")
        cell.get_text().set_fontsize(10.5)
        cell.get_text().set_weight("bold")
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            face = "#ffffff" if r % 2 else "#f7f7fb"
            cell = table.add_cell(r, c, width=col_w[c], height=0.08, text=val, loc="center", facecolor=face)
            cell.set_edgecolor("#c1c5d0")
            cell.get_text().set_fontsize(10)
            if c == 0:
                cell._loc = 'left'
                cell.get_text().set_ha('left')
                cell.get_text().set_x(0.02)
    ax.add_table(table)
    fig.savefig(TABLE_PNG, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def update_pptx():
    render_assets()
    prs = Presentation(str(PPTX_PATH))

    # Add a new slide after current Slide 6.
    slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(slide_layout)
    # Move it directly after slide index 5 (current Slide 6).
    move_slide(prs, len(prs.slides._sldIdLst) - 1, 6)
    slide = prs.slides[6]

    # Build the slide using the existing visual language from Slide 6.
    # Background boxes.
    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.30), Inches(0.82), Inches(5.95), Inches(5.80)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = PALE
    slide.shapes[-1].line.color.rgb = RGBColor(214, 205, 159)

    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.15), Inches(0.90), Inches(5.95), Inches(1.75)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = PALE2
    slide.shapes[-1].line.color.rgb = ACCENT

    slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.15), Inches(2.78), Inches(5.95), Inches(3.65)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = PALE2
    slide.shapes[-1].line.color.rgb = ACCENT

    # Title and slide number.
    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(9.8), Inches(0.35))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "MELD Dataset EDA & Contextual Factors"
    style_text(r, 18, True, BODY)

    num = slide.shapes.add_textbox(Inches(11.35), Inches(0.12), Inches(0.8), Inches(0.35))
    p = num.text_frame.paragraphs[0]
    p.alignment = 2
    r = p.add_run()
    r.text = "07"
    style_text(r, 18, True, BODY)

    # Left bullet box.
    left = slide.shapes.add_textbox(Inches(0.55), Inches(1.08), Inches(5.45), Inches(4.95))
    tf = left.text_frame
    tf.word_wrap = True
    bullets = [
        "MELD is a multi-party, multimodal benchmark, so speaker context and turn history matter.",
        "Table 10 shows the dataset is shorter and much more neutral-heavy than IEMOCAP.",
        "Inter-speaker influence means one speaker’s utterance often depends on another speaker’s prior turn.",
        "Emotion shifts and contextual distance show that labels cannot always be read from one utterance alone.",
        "This is why Phase 1 uses weighted metrics and dialogue-aware analysis instead of accuracy only.",
    ]
    for i, t in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.level = 0
        p.bullet = True
        for run in p.runs:
            style_text(run, 12.2, False, BODY)

    # Top-right summary note.
    note = slide.shapes.add_textbox(Inches(6.35), Inches(1.08), Inches(5.35), Inches(1.35))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Table 10 is useful because it shows why MELD needs context-aware and imbalance-aware modeling: the dataset is dialogue-heavy, shorter, and skewed toward neutral."
    )
    for run in p.runs:
        style_text(run, 10.7, False, MUTED)

    # Table image on the bottom-right.
    pic = slide.shapes.add_picture(str(TABLE_PNG), Inches(6.28), Inches(2.95), width=Inches(5.55), height=Inches(3.15))
    pic.name = "MELDTable10Summary"

    # Speaker-note box at bottom-left / under bullets.
    explain = slide.shapes.add_textbox(Inches(0.62), Inches(5.58), Inches(5.20), Inches(0.68))
    tf = explain.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Use this slide to connect the MELD benchmark to the project’s evaluation strategy: the context is conversational, the labels are imbalanced, and the error pattern is shaped by nearby turns and speaker interactions."
    for run in p.runs:
        style_text(run, 10.0, False, MUTED)

    # Update slide number text boxes for all slides from the inserted slide onward.
    for idx in range(6, len(prs.slides)):
        s = prs.slides[idx]
        for shp in s.shapes:
            if hasattr(shp, "text") and shp.text.strip().isdigit():
                try:
                    shp.text = f"{idx + 1:02d}"
                except Exception:
                    pass
                break

    prs.save(str(PPTX_PATH))


def update_docx():
    doc = Document(str(DOCX_PATH))
    table = doc.tables[0]

    # Insert a new row after Slide 6.
    new_tr = deepcopy(table.rows[6]._tr)
    table.rows[6]._tr.addnext(new_tr)
    new_row = _Row(new_tr, table)
    new_row.cells[0].text = "Slide 7"
    new_row.cells[1].text = (
        "Explain why MELD is worth presenting with a compact EDA summary. Show that the dataset is multi-party, neutral-heavy, short-turn, and context-sensitive, and connect Table 10 plus the additional analysis to the evaluation choices in Phase 1."
    )
    new_row.cells[2].text = "Keep it technical, concise, and traceable."

    # Shift later row labels by +1.
    for idx in range(7, len(table.rows)):
        current = table.rows[idx]
        txt = current.cells[0].text.strip()
        if txt.startswith("Slide "):
            try:
                n = int(txt.split()[-1])
                current.cells[0].text = f"Slide {n + 1}"
            except Exception:
                pass

    # Insert new section before the literature review section.
    paras = list(doc.paragraphs)
    anchor = next((p for p in doc.paragraphs if p.text.strip().startswith("2.2 Literature Review Slides")), None)
    if anchor is None:
        anchor = doc.paragraphs[-1]
    h = insert_paragraph_before(anchor, style="Heading 2")
    h.add_run("2.2 Slide 7 MELD Dataset EDA & Contextual Factors")

    sections = [
        "EDA summary: MELD is a multi-party, multimodal conversational dataset, so the project cannot treat each utterance as an isolated sentence. The statistics from the MELD paper show a strong neutral skew, shorter utterances, and many dialogue turns, which is exactly why Phase 1 should discuss weighted metrics and context-aware evaluation rather than accuracy alone.",
        "Why Table 10 is useful: Table 10 compares MELD with IEMOCAP and shows that MELD has more emotion-labeled utterances, more unique words, and much shorter average utterances and conversations. That is useful for the ESA because it explains why conversational context is necessary and why the benchmark is harder than a simple single-utterance classification task.",
        "Inter-speaker influence: the additional analysis in the MELD paper shows that many correct predictions depend on utterances from different speakers, not only the target speaker. For this project, that means dialogue context and speaker-aware reasoning are not optional; they are part of the problem definition.",
        "Emotion shifts: the paper reports that detecting emotion shifts is possible, but detecting the shift and the correct emotion at the same time is much harder. This matters for LegalMemoCMT because emotional behavior can change quickly across turns, and a model that ignores shifts may look good on one utterance but still miss the surrounding conversational pattern.",
        "Contextual distance: the paper also shows that the model attends not only to nearby turns but also to farther historical and future utterances. That supports the idea that the project should keep the broader conversational timeline visible when explaining predictions and later courtroom adaptation.",
        "Project impact: these MELD observations justify the use of weighted F1, confusion-matrix analysis, and fold-safe dialogue grouping in Phase 1, and they also motivate the later Phase 2 idea of adding legal-domain metadata and explainable emotional timelines.",
    ]
    for t in sections:
        p = insert_paragraph_before(anchor)
        r = p.add_run(t)
        r.font.name = "Times New Roman"
        r.font.size = Pt(11.5)
        p.alignment = 3

    # Update existing heading labels to keep the guide consistent after slide insertion.
    replacements = {
        "2.2 Literature Review Slides (Slides 5-9)": "2.3 Literature Review Slides",
        "2.3 Slide 11 Extended Talking Points": "2.4 Slide 11 Extended Talking Points",
        "2.4 Slide 12 Extended Talking Points": "2.5 Slide 12 Extended Talking Points",
        "2.5 Slide 13 Extended Talking Points": "2.6 Slide 13 Extended Talking Points",
        "2.6 Slide 14 Design Properties & Implications": "2.7 Slide 14 Design Properties & Implications",
        "2.7 Slide 16 System Architecture": "2.8 Slide 16 System Architecture",
        "2.8 Slide 17 Model Architecture": "2.9 Slide 17 Model Architecture",
    }
    for p in doc.paragraphs:
        t = p.text.strip()
        if t in replacements:
            p.text = replacements[t]

    doc.save(str(DOCX_PATH))


def main():
    update_pptx()
    update_docx()
    print(f"Updated {PPTX_PATH}")
    print(f"Updated {DOCX_PATH}")


if __name__ == "__main__":
    main()
