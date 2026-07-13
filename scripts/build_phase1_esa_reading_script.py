from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor
from pptx import Presentation


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Presentation.pptx"
OUT_PATH = ROOT / "implementation_docments/LegalMemoCMT_Phase1_ESA_Reading_Script.docx"


def style_run(run, *, size=11, bold=False, color="000000", name="Aptos"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def add_para(doc, text, *, size=11, bold=False, color="000000", indent=0.0, italic=False):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.space_before = Pt(0)
    if indent:
        p.paragraph_format.left_indent = Pt(indent)
    r = p.add_run(text)
    style_run(r, size=size, bold=bold, color=color)
    r.italic = italic
    return p


def slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            txt = shape.text_frame.text.strip()
            if txt:
                return txt.splitlines()[0].strip()
    return "Untitled"


def build_slide_text():
    return {
        1: [
            "Start by stating the project title and the specific review stage. Then clarify that Phase 1 is a controlled implementation and validation effort, not the end state of the overall project.",
            "Explain that the real goal is to build a reproducible multimodal pipeline that can later support courtroom-testimony analysis. The examiner should immediately understand that this phase is about proving technical traceability, not making an overstated legal claim.",
            "Mention that the main evidence chain is raw data to manifests to cached features to checkpointed predictions. That chain matters because the project is only defensible if each output can be traced back to a specific input and code path.",
        ],
        2: [
            "Use the agenda to narrate the review in order: problem, literature, design, implementation, results, and future direction.",
            "The point of the agenda is not just organization. It gives you a speaking structure so the examiner can follow how each slide supports the overall claim of a reproducible Phase 1 benchmark.",
        ],
        3: [
            "Open the problem statement by saying that courtroom testimony is a multimodal communication problem. In practice, emotion is conveyed through wording, tone, and visible facial behavior rather than through text alone.",
            "Then explain why MELD is used in Phase 1: it gives a conversational benchmark where the model must deal with turn-level emotion, context, and class imbalance before any legal-domain adaptation is attempted.",
            "If asked why MemoCMT is relevant, say that it provides the closest paper-style multimodal reference for reproducing a stable text-audio baseline before extending the system with visual cues.",
        ],
        4: [
            "Treat the abstract as two linked ideas. The first is reproducibility of the paper-aligned baseline, and the second is the visual support path that turns raw video into reusable facial embeddings.",
            "Explain that the scope of Phase 1 is deliberately narrow: reproduce the benchmark, study the error pattern, and keep the implementation transparent enough that every result can be explained during review.",
        ],
        5: [
            "When presenting MemoCMT, explain that it is the closest published reference point for this project because it combines pretrained text and audio encoders with a cross-modal transformer fusion block rather than a plain concatenation pipeline.",
            "The paper uses BERT for text and HuBERT for audio, so the two branches first learn modality-specific representations before being fused through cross-attention. That matters for Phase 1 because the project keeps the same high-level design principle in the paper-aligned baseline.",
            "Make the fusion logic explicit: the CMT block lets the two modalities attend to each other, and then the paper compares aggregation strategies such as CLS, MEAN, MAX, and MIN. The Phase 1 baseline follows the same idea and uses MIN as the strongest paper-aligned choice for MELD.",
            "Tie the literature review directly to the implementation: the reason the baseline starts with text and audio is that this paper shows that multimodal attention fusion is more effective than treating the modalities as independent features. That is exactly why the Phase 1 work first reproduces the text-plus-audio baseline before adding the visual path.",
            "If asked why this paper is the anchor, say that it is not just a citation. It is the architectural template for the baseline, the metric reference for MELD, and the justification for using a paper-aligned training and evaluation story before moving to facial cues.",
        ],
        6: [
            "For MELD, say that the benchmark matters because it is dialogue-based and multi-party. This means the model has to deal with context, speaker interaction, and label ambiguity rather than isolated sentence emotion.",
            "Point out that the neutral-heavy distribution makes accuracy misleading by itself. That is why the project keeps weighted F1, macro F1, and confusion analysis in the evaluation story.",
        ],
        7: [
            "Use the EDA slide to show that the dataset statistics are not just descriptive. They explain why the pipeline must be dialogue-safe and why the evaluation must emphasize class imbalance.",
            "Explain that Table 10 gives evidence for the short-utterance profile, neutral skew, and conversational structure. Those properties justify the fold design and the use of weighted metrics.",
            "A good viva line is that MELD is a realistic conversational benchmark, which makes it useful for exposing exactly the kinds of weaknesses a simple classifier would hide.",
        ],
        8: [
            "This slide should be explained as the contextual layer of the benchmark. Inter-speaker influence means a nearby speaker’s turn can change how the current utterance should be interpreted.",
            "Emotion shifts matter because the same speaker can move across affective states inside a dialogue. That is why one isolated prediction is not enough to describe conversational behavior.",
            "Contextual distance is also important: nearby turns are usually strongest, but longer-range turns can still influence the model. In Phase 1, this is handled indirectly through the data pipeline and fusion behavior, while explicit speaker-memory modeling is future work.",
        ],
        9: [
            "Present BERT and HuBERT as pretrained branches that reduce the burden on Phase 1. BERT gives contextual text embeddings, while HuBERT provides speech representations learned from raw audio.",
            "Emphasize that both encoders are already strong before training starts. That is why the project can focus on adaptation and multimodal interaction instead of learning language and speech from scratch.",
        ],
        10: [
            "For ViT, explain that the video path does not consume an entire raw video stream at once. Instead, the system samples RGB frames, optionally crops the face region, and then converts those frames into embeddings.",
            "That design matters because it turns visual emotion into a compact representation that can be fused with the text and audio branches. It is a support path, not a replacement for the conversational baseline.",
        ],
        11: [
            "Weighted cross entropy and focal loss are the imbalance-aware training tools for this project. They matter because MELD is not balanced, so a model can appear strong while still ignoring minority classes.",
            "Explain that weighted CE increases the contribution of rare classes, while focal loss reduces the impact of easy examples so the model spends more learning pressure on difficult cases.",
            "The important conclusion is that loss functions change training emphasis; they do not automatically solve the class-imbalance problem. That is why evaluation still has to include macro F1 and the confusion matrix.",
        ],
        12: [
            "This slide is where you show that the review feedback was absorbed into the presentation structure. The strongest version of the story is to keep the baseline central, show the additional branches as controlled experiments, and use the outputs to justify the next step.",
            "If the examiner asks why this matters, say that the project is not trying to impress with a single number. It is trying to show an implementation path whose behavior can be defended technically.",
        ],
        13: [
            "Describe the design approach as a controlled experimental ladder. First validate the baseline, then add the visual branch, then test gated fusion and auxiliary loss as separate changes.",
            "This is the right approach for an MTech project because it preserves attribution. If a result changes, you can explain whether it came from the data path, the visual signal, the fusion rule, or the loss function.",
            "Mention that early fusion, late fusion, and video-only systems are valid alternatives, but they make causal explanation harder. The project favors interpretability over uncontrolled complexity.",
        ],
        14: [
            "This slide is where you defend the assumptions that make the design valid. The first assumption is that every sample_id must stay aligned across manifest, cache, and checkpoint.",
            "The second is that pretrained encoders and ffmpeg are available, because otherwise the data pipeline is incomplete. The third is that the class imbalance is real and must be addressed in the loss and evaluation.",
            "A useful line here is that these are not minor implementation details. They are the conditions under which the project can produce reproducible and reviewable results.",
        ],
        15: [
            "Use this slide to talk through the design details as a layered system. The project does not start from the model only; it starts from raw data, preprocessing, caching, and then moves into the model and analysis layers.",
            "Explain that the design is built to support later extension toward legal-domain adaptation. The idea is to keep the core multimodal backbone stable while allowing the task-specific layer to evolve later.",
        ],
        16: [
            "The properties slide is about engineering quality. Interoperability means the system can move a sample from video to audio to transcript to cached features without breaking the identity of the sample.",
            "Reliability means the outputs can be checked through fold metrics and confusion matrices, portability means it can run locally or on RunPod, and maintainability means the pipeline remains modular instead of monolithic.",
            "These properties matter because the project is not just a model; it is a workflow that must remain explainable when the reviewer asks how one artifact leads to another.",
        ],
        17: [
            "The methodology slide should sound like a plan of work rather than a static architecture slide. Walk the panel through training, checkpoint selection, evaluation, and then the decision about the next experiment.",
            "A strong technical point here is that the confusion matrix is what decides whether a change is meaningful. An increase in a single accuracy value is not enough if the minority-class confusion is still bad.",
        ],
        18: [
            "Use the system architecture slide to move from raw inputs to outputs in one clean story. The raw video, transcript, and audio are first preprocessed, then cached, then fed into the model, and finally analyzed through metrics and confusion tables.",
            "Explain that the cache layer is what makes the pipeline practical. Once the preprocessed artifacts exist, the same sample can be reused for training, evaluation, and demo without recomputing everything each time.",
        ],
        19: [
            "This is the internal model view. The text branch handles semantics, the audio branch handles delivery and prosody, and the video branch handles facial cues.",
            "Then explain that the fusion module decides how those branches interact. This is where the architecture becomes multimodal rather than three independent classifiers.",
        ],
        20: [
            "Keep this slide practical. It is there to show that the implementation rests on a standard and explainable stack: Python, PyTorch, Hugging Face transformers, numpy, pandas, sklearn, OpenCV, and ffmpeg.",
            "If asked about environments, say the code can run locally for development and on GPU platforms for heavier training and demo execution.",
        ],
        21: [
            "The progress slide should sound confident but not exaggerated. Say that the Phase 1 baseline is complete, the facial-cue branch is implemented, and the current behavior is understood through analysis rather than guessed.",
            "The most important message is that the project has reached a stable checkpoint. It is ready for review, and the next changes should be based on measured error patterns rather than intuition alone.",
        ],
        22: [
            "References can be handled briefly. Say that the citations cover the baseline paper, MELD, the encoder references, and the imbalance-loss literature.",
            "If the examiner asks for justification, mention that the references are chosen because they directly support the implementation choices rather than being a general reading list.",
        ],
        23: [
            "This slide is your error-analysis defense. Tell the panel that top-3 values are softmax probabilities, which means they show the model’s belief distribution, not a metric like accuracy or F1.",
            "When a prediction is wrong, describe whether it is a near miss or a confident wrong case. Then connect the error to label ambiguity, sampling limitations, or modality dominance instead of trying to hide it.",
            "The point of this slide is to show that the project can explain its failures, not just its successes.",
        ],
        24: [
            "The thank-you slide should close the review by summarizing the current state: the Phase 1 baseline is reproducible, the visual support branch is explainable, and the error analysis is honest.",
            "Leave the examiner with the sense that the work is stable enough for the current phase and structured enough for the next extension.",
        ],
        25: [
            "Open this slide as the formal comparison between the published MemoCMT MELD result, the paper-aligned baseline in this repository, and the gated+aux branch. The point is not to claim a leaderboard win; the point is to show whether the implementation is in the right performance range and whether the added visual branch behaves in a technically sensible way.",
            "Explain the metric columns carefully. Accuracy and weighted accuracy are the same sample-level correctness quantity in this codebase, weighted F1 reflects the class distribution, macro F1 treats all classes equally, and unweighted accuracy is the mean per-class accuracy. The paper reports the MELD benchmark in a slightly different framing, so the comparison is best interpreted directionally.",
            "State the actual numbers clearly: the MemoCMT paper reports 64.18% accuracy and 62.52% F1 on MELD CMT+MIN. The paper-aligned MELD cross-validation mean in this repository is 62.47% accuracy and 61.95% weighted F1. The best baseline fold 2 is 63.75% accuracy and 62.54% weighted F1, which is extremely close to the paper. The gated+aux branch reaches 60.54% accuracy and 60.22% weighted F1 on fold 2, and 59.92% accuracy and 60.56% weighted F1 on fold 4.",
            "The main interpretation is that the baseline is still the closest overall reproduction of the paper, while gated+aux gives selective gains on some clips but does not improve the full MELD benchmark universally. That is the right technical message for Phase 1 because it keeps the comparison honest and avoids overclaiming.",
            "If the examiner asks why gated+aux is not always better, explain that the visual branch and gating add more expressive capacity, but MELD is class-imbalanced and neutral-heavy. More capacity can sharpen the right boundary on some clips, such as neutral boundary cases, but it can also make wrong classes more confident on ambiguous clips. That is why fold-level confusion analysis still matters.",
            "A concise speaking line is: the paper-aligned baseline remains the strongest stable MELD result, and gated+aux is a selective improvement branch that helps some examples but does not yet surpass the baseline on the overall MELD cross-validation story.",
        ],
        26: [
            "Use this slide as the live evidence slide. Explain that each row uses the same raw mp4 clip for both checkpoints, so the only thing changing is the model checkpoint and its learned fusion behavior.",
            "The slide shows the transcript cue, ground truth, baseline prediction, baseline confidence, baseline top-3 probabilities, gated+aux prediction, gated+aux confidence, gated+aux top-3 probabilities, and a short reason for the outcome.",
            "The technical conclusion is selective improvement: the gated+aux checkpoint fixes the near-boundary neutral case, keeps one non-neutral success, but remains wrong on the other hard clips and can become more confident in the wrong answer.",
        ],
        27: [
            "Future work should be presented as the next technical layer, not as something already completed in Phase 1.",
            "Explain that explicit speaker modeling, emotion-shift modeling, long-range context, and courtroom-role metadata will help because the current errors still show contextual ambiguity and speaker dependence.",
            "If asked about timing, say that only a small context ablation might fit before Phase 1 closes, but the full redesign is better treated as the bridge into Phase 2.",
        ],
    }


def build_doc():
    ppt = Presentation(str(PPTX_PATH))
    doc = Document()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("LegalMemoCMT Phase 1 ESA Reading Script")
    style_run(r, size=18, bold=True, color="122F55")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("A slide-by-slide speaking script for the ESA presentation")
    style_run(r, size=11, color="5E5E5E")

    add_para(doc, "This document is written as a full speaking script. It is intentionally more detailed than the short speaker notes, so you can rehearse a natural explanation of each slide without sounding like you are reciting the text on the slide.")
    add_para(doc, "Use it to practise transitions, technical explanations, and answers to likely viva questions. The focus should stay on the project story: why the benchmark matters, how the pipeline works, what the results mean, and what the next step should be.")

    for i, slide in enumerate(ppt.slides, start=1):
        title = slide_title(slide)
        h = doc.add_paragraph()
        h.style = "Heading 1"
        rr = h.add_run(f"Slide {i}: {title}")
        style_run(rr, size=14, bold=True, color="122F55")

        for para in build_slide_text().get(i, [f"Explain the slide '{title}' by focusing on the technical meaning, the implementation choice behind it, and how it supports the Phase 1 story."]):
            add_para(doc, para, indent=10)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    build_doc()
