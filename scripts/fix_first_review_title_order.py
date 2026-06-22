from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.util import Pt


ROOT = Path("/Users/rajeshpmu/Desktop/LegalMemoCMT")
PPTX_PATH = ROOT / "implementation_docments" / "First_Review_LegalMemoCMT.pptx"


def move_slide_to_index(prs: Presentation, slide_index: int, target_index: int) -> None:
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[slide_index]
    sldIdLst.remove(sldId)
    sldIdLst.insert(target_index, sldId)


def title_shape(slide):
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if float(shape.top) > float(Inches(1.25)):
            continue
        candidates.append(shape)
    if not candidates:
        return None

    def score(shape):
        max_font = 0.0
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    max_font = max(max_font, float(run.font.size))
        return (max_font, -float(shape.top))

    return max(candidates, key=score)


def renumber_titles(prs: Presentation) -> None:
    seq = 1
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        ts = title_shape(slide)
        if ts is None:
            continue
        tf = ts.text_frame
        old = tf.text.strip()
        base = re.sub(r"^\d+\.\s*", "", old)
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{seq}. {base}"
        r.font.name = "Aptos Display"
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = RGBColor(20, 48, 87)
        seq += 1


def fix_title_slide(prs: Presentation) -> None:
    slide = prs.slides[0]
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        txt = shape.text_frame.text.strip()
        if txt == "1. First Review":
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "First Review"
            r.font.name = "Aptos Display"
            r.font.size = Pt(24)
            r.font.bold = True
            r.font.color.rgb = RGBColor(20, 48, 87)
            break


def main() -> None:
    prs = Presentation(str(PPTX_PATH))

    # Move the title slide back to the first position if a previous insertion shuffled it.
    # The title slide is the one that contains the project header text.
    title_idx = None
    for i, slide in enumerate(prs.slides):
        combined = "\n".join(
            sh.text_frame.text for sh in slide.shapes if hasattr(sh, "text_frame") and sh.text_frame
        )
        if "MTech Final Project" in combined and "LegalMemoCMT" in combined:
            title_idx = i
            break
    if title_idx is not None and title_idx != 0:
        ids = list(prs.slides._sldIdLst)
        title_id = ids[title_idx]
        new_order = [title_id] + ids[:title_idx] + ids[title_idx + 1 :]
        for sldId in ids:
            prs.slides._sldIdLst.remove(sldId)
        for sldId in new_order:
            prs.slides._sldIdLst.append(sldId)

    fix_title_slide(prs)
    renumber_titles(prs)
    prs.save(str(PPTX_PATH))
    print("Fixed First Review slide order and numbering.")


if __name__ == "__main__":
    main()
